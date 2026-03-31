"""
JAKALA Nordic — Outreach Campaign Package Builder
Generates for each of 25 Norwegian accounts:
  - account-overview.md  (copied from existing)
  - outreach-email.md    (tailored Norwegian email)
  - offering.md          (tailored offering one-pager)
  + one shared JAKALA company presentation PPTX
  → packaged as outreach-norge-2026-03.zip
"""

import json, os, shutil, zipfile, textwrap
from pathlib import Path
from datetime import date

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE   = Path(__file__).parent.parent
DATA   = BASE / "intelligence/outreach-accounts.json"
ACCTS  = BASE / "Accounts"
OUT    = BASE / "outreach-package"
PPTX_OUT = OUT / "_jakala-company-overview.pptx"

# ── Load account data ─────────────────────────────────────────────────────────
with open(DATA) as f:
    accounts = json.load(f)

# ── Email templates per GTM strategy ──────────────────────────────────────────
EMAIL_TEMPLATES = {
    "Commerce Optimization": """Emne: {name} — ett konkret grep for å løfte digital omsetning

Hei {buyer_name},

Jeg tar kontakt fordi vi nettopp fullførte et prosjekt for Maxbo i Norge — samme type utfordring jeg tror dere kjenner igjen: {core_problem}

Vi er JAKALA, en global digital transformation-partner med sterk nordisk tilstedeværelse. Vi jobber med ledende retailere på nettopp det skjæringspunktet mellom produktdata, søk og commerce-plattformer.

**Hva vi ser i {name}:**
{timing_signal}

**Hva vi foreslår som første steg:**
{entry_offer}

Dette er ikke en stor plattformprosjekt. Det er en avgrenset analyse som gir dere et klart bilde av hva som koster mest og hva som gir raskest effekt — leveransetid 6–8 uker, estimert investering {deal_value_range}.

Hadde det vært interessant med en 30-minutters samtale der jeg viser hva vi fant hos Maxbo og hvordan det oversettes til {name}?

Med vennlig hilsen
Jacob Skaue
JAKALA Nordic
jacob.skaue@jakala.com
""",

    "Data Revenue Unlock": """Emne: {name} — er dataene dere har undervurdert?

Hei {buyer_name},

Grunnen til at jeg skriver akkurat nå: {timing_signal}

Vi er JAKALA — global data- og digital transformation-partner. Vi hjelper selskaper som {name} med å forstå den kommersielle verdien i dataene de allerede sitter på.

**Problemstillingen vi ser:**
{core_problem}

**Hva vi kan gjøre:**
{entry_offer}

Vi kaller det en Data Revenue Diagnostic — en strukturert 3-dagers analyse som svarer på: hva er dataene deres faktisk verdt, for dere selv og for partnere? Investering: {deal_value_range}.

Kan jeg sende over en én-sides beskrivelse av metoden? Eller om du foretrekker det — 20 minutter på telefon denne uken?

Med vennlig hilsen
Jacob Skaue
JAKALA Nordic
jacob.skaue@jakala.com
""",

    "AI Readiness Accelerator": """Emne: {name} — AI-beredskap som faktisk gir forretningsverdi

Hei {buyer_name},

{timing_signal}

Jeg er Jacob Skaue fra JAKALA — vi er en global digital transformation-partner som de siste 18 månedene har hjulpet finansinstitusjoner i Europa med å gå fra AI-ambisjon til AI-leveranse.

**Det vi observerer fra utsiden:**
{core_problem}

**Hva vi tilbyr som inngang:**
{entry_offer}

Dette er ikke et stort AI-transformasjonsprosjekt. Det er en avgrenset sprint — 4–6 uker — som ender med en konkret AI-roadmap forankret i reell forretningsverdi. Investering: {deal_value_range}.

Er dette noe som treffer en utfordring dere jobber med nå? Gjerne en rask samtale.

Med vennlig hilsen
Jacob Skaue
JAKALA Nordic
jacob.skaue@jakala.com
""",

    "Experience Transformation": """Emne: {name} — digital studentopplevelse som et konkurransefortrinn

Hei {buyer_name},

{timing_signal}

Jeg kontakter deg fra JAKALA — vi er en global digital transformation-partner med erfaring fra høyere utdanning og kundeopplevelse i Europa.

**Det vi ser som mulighet for {name}:**
{core_problem}

**Vårt forslag til inngang:**
{entry_offer}

Vi starter alltid med en avgrenset diagnostisk fase — 4–6 uker — som gir en klar anbefaling om hva som gir størst effekt på rekruttering, retention og omdømme. Investering: {deal_value_range}.

Er dette noe du ønsker å høre mer om? Jeg kan sende over en kort presentasjon, eller vi tar en 20-minutters samtale.

Med vennlig hilsen
Jacob Skaue
JAKALA Nordic
jacob.skaue@jakala.com
""",
}

# ── Offering one-pager template ───────────────────────────────────────────────
OFFERING_TEMPLATE = """# JAKALA × {name}
## Tailored Offering — {gtm_strategy}
*Konfidensielt · Mars 2026*

---

## Situasjon

{name} er en ledende aktør innen **{sector}**. {core_problem}

**Timing:** {timing_signal}

**Tech stack:** {tech_stack_highlight}

---

## Hva vi tilbyr

### Entry-tilbud: {entry_offer_short}

{entry_offer}

**Scope:**
- Varighet: 6–8 uker
- Estimert investering: **{deal_value_range}**
- Leveranse: Diagnose · Prioritert roadmap · Proof-of-concept

---

## Hvorfor JAKALA

| Vi er | Det betyr for deg |
|-------|-------------------|
| Global transformation-partner (5.000+ ansatte, 40+ land) | Kapasitet og erfaring på enterprise-skala |
| Shopify Premier Agency (1 av 3 i Italia, 30+ sertifiseringer) | Best-in-class commerce implementering |
| Aktiv leveranse i Norge — Maxbo (Shopify Plus + Enterspeed) | Lokal referanse, ikke et konsept |
| Data- og AI-spesialist på tvers av retail, finans og utdanning | Forretningsverdi fra dag én |

---

## Ekspansjonssti

```
Entry: {entry_offer_short} ({deal_value_range})
  → Fase 2: Optimalisering og skalering
    → Fase 3: AI-drevet transformasjon
```

---

## Kjøper og kontakt

| | |
|--|--|
| **Primær kontakt** | {primary_buyer_name} — {primary_buyer_title} |
| **Inngangsvinkel** | {buyer_outreach_angle} |
| **JAKALA-kontakt** | Jacob Skaue · jacob.skaue@jakala.com |

---

*JAKALA Nordic · jakala.com*
"""

# ── Strategy description map ──────────────────────────────────────────────────
STRATEGY_DESC = {
    "Commerce Optimization": "Commerce Optimization — produktdata, søk og commerce-plattform",
    "Data Revenue Unlock": "Data Revenue Unlock — kommersiell verdi i eksisterende data",
    "AI Readiness Accelerator": "AI Readiness Accelerator — fra AI-ambisjon til AI-leveranse",
    "Experience Transformation": "Experience Transformation — digital kundeopplevelse som driver vekst",
}

# ── Generate per-account files ─────────────────────────────────────────────────
if OUT.exists():
    shutil.rmtree(OUT)
OUT.mkdir(parents=True)

for acc in accounts:
    folder = acc["folder"]
    name   = acc["name"]
    slug   = folder
    acct_dir = OUT / slug
    acct_dir.mkdir()

    # 1. Copy existing overview.md
    src_overview = ACCTS / folder / "overview.md"
    if src_overview.exists():
        shutil.copy(src_overview, acct_dir / "account-overview.md")

    # 2. Generate outreach email
    template = EMAIL_TEMPLATES.get(acc["gtm_strategy"], EMAIL_TEMPLATES["Commerce Optimization"])
    buyer_name = acc["primary_buyer_name"] if acc["primary_buyer_name"] != "TBD" else acc["primary_buyer_title"]
    email_content = template.format(
        name=name,
        buyer_name=buyer_name,
        core_problem=acc["core_problem"],
        timing_signal=acc["timing_signal"],
        entry_offer=acc["entry_offer"],
        deal_value_range=acc["deal_value_range"],
        gtm_strategy=STRATEGY_DESC.get(acc["gtm_strategy"], acc["gtm_strategy"]),
    )
    # Add header metadata
    email_header = f"""# Outreach Email — {name}
**Til:** {acc['primary_buyer_name']} · {acc['primary_buyer_title']}
**GTM-strategi:** {acc['gtm_strategy']}
**Timing-vinkel:** {acc['timing_signal']}

---

"""
    with open(acct_dir / "outreach-email.md", "w") as f:
        f.write(email_header + email_content)

    # 3. Generate offering one-pager
    # Short entry offer (first sentence only)
    entry_short = acc["entry_offer"].split("—")[0].strip() if "—" in acc["entry_offer"] else acc["entry_offer"][:60].strip()

    offering_content = OFFERING_TEMPLATE.format(
        name=name,
        gtm_strategy=STRATEGY_DESC.get(acc["gtm_strategy"], acc["gtm_strategy"]),
        sector=acc["sector"],
        core_problem=acc["core_problem"],
        timing_signal=acc["timing_signal"],
        tech_stack_highlight=acc["tech_stack_highlight"],
        entry_offer=acc["entry_offer"],
        entry_offer_short=entry_short,
        deal_value_range=acc["deal_value_range"],
        primary_buyer_name=acc["primary_buyer_name"],
        primary_buyer_title=acc["primary_buyer_title"],
        buyer_outreach_angle=acc["buyer_outreach_angle"],
    )
    with open(acct_dir / "offering.md", "w") as f:
        f.write(offering_content)

print(f"Generated content for {len(accounts)} accounts")

# ── Build JAKALA Company Overview PPTX ───────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE  = RGBColor(0x15, 0x3E, 0xED)
NAVY  = RGBColor(0x02, 0x02, 0x66)
RED   = RGBColor(0xF6, 0x57, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY  = RGBColor(0xAA, 0xAA, 0xCC)
BG    = RGBColor(0x08, 0x08, 0x18)
GREEN = RGBColor(0x22, 0xDD, 0x88)
AMBER = RGBColor(0xFF, 0xBB, 0x33)
FONT  = 'Raleway'
W = Inches(13.33); H = Inches(7.5)

prs = Presentation()
prs.slide_width = W; prs.slide_height = H
blank = prs.slide_layouts[6]

def ns(): return prs.slides.add_slide(blank)
def bg(s):
    r = s.shapes.add_shape(1,0,0,W,H); r.fill.solid(); r.fill.fore_color.rgb=BG; r.line.fill.background()
def rect(s,x,y,w,h,c):
    r=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h)); r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background(); return r
def ln(s,x,y,w,c=BLUE):
    r=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(0.02)); r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background()
def tb(s,text,x,y,w,h,size=14,color=WHITE,bold=False,align=PP_ALIGN.LEFT,italic=False):
    t=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=t.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text
    r.font.name=FONT; r.font.size=Pt(size); r.font.color.rgb=color; r.font.bold=bold; r.font.italic=italic; return t

# ── Slide 1: Cover ────────────────────────────────────────────────────────────
s=ns(); bg(s)
rect(s,0,0,0.5,7.5,BLUE)
rect(s,0.5,5.3,12.83,0.06,BLUE)
tb(s,"JAKALA",1.0,1.0,11.0,1.0,size=72,color=BLUE,bold=True)
tb(s,"Vi gjør data og teknologi\ntil forretningsvekst.",1.0,2.4,10.0,2.2,size=36,color=WHITE,bold=True)
tb(s,"Digital Transformation · Data · Commerce · AI",1.0,5.0,10.0,0.5,size=16,color=GREY)
tb(s,"jakala.com  ·  Nordic 2026",1.0,5.5,6.0,0.4,size=13,color=GREY)

# ── Slide 2: Who We Are ───────────────────────────────────────────────────────
s=ns(); bg(s)
rect(s,0,0,0.5,7.5,BLUE)
ln(s,0.8,1.3,11.5)
tb(s,"HVEM ER JAKALA",0.8,0.4,7.0,0.45,size=11,color=BLUE,bold=True)
tb(s,"Global skala. Nordisk tilstedeværelse.",0.8,0.7,11.0,0.65,size=28,color=WHITE,bold=True)
stats=[("5.000+","ansatte globalt"),("40+","land representert"),("€600M+","global omsetning"),("20+","år i drift"),("30+","Shopify-sertifiseringer"),("1 av 3","Shopify Premier Agency i Italia")]
for i,(num,lbl) in enumerate(stats):
    col=i%3; row=i//3
    x=0.8+col*4.2; y=1.65+row*1.4
    rect(s,x,y,3.9,1.2,NAVY)
    tb(s,num,x+0.15,y+0.1,3.6,0.7,size=32,color=BLUE,bold=True)
    tb(s,lbl,x+0.15,y+0.75,3.6,0.38,size=12,color=GREY)
tb(s,"JAKALA er en av Europas raskest voksende digital transformation-partnere — med sterk posisjon i Norden gjennom aktiv leveranse i Norge (Maxbo), Sverige, Danmark og Finland.",0.8,4.5,11.5,0.8,size=12,color=WHITE)

# ── Slide 3: What We Do ───────────────────────────────────────────────────────
s=ns(); bg(s)
rect(s,0,0,0.5,7.5,BLUE)
ln(s,0.8,1.3,11.5)
tb(s,"HVA VI GJØR",0.8,0.4,7.0,0.45,size=11,color=BLUE,bold=True)
tb(s,"Fire inngangsstrategier. En partner.",0.8,0.7,11.0,0.65,size=28,color=WHITE,bold=True)
strategies=[
    (BLUE,"Commerce\nOptimization","Produktdata, søk og commerce-plattform. Vi forbedrer digital omsetning fra eksisterende infrastruktur.","Retail · B2B · Omnichannel"),
    (GREEN,"Data Revenue\nUnlock","Vi gjør lojalitets-, kjøps- og kundedata om til ny kommersiell verdi. Retail media, clean rooms, ekstern monetisering.","Retail · Finans · Lojalitetsprogrammer"),
    (AMBER,"AI Readiness\nAccelerator","Fra AI-strategi til konkrete AI-løsninger med bevist ROI. Vi builder — ikke bare rådgiver.","Finans · B2B · Enterprise"),
    (RED,"Experience\nTransformation","Digital kundeopplevelse som driver rekruttering, retention og omdømme.","Utdanning · Helse · Offentlig"),
]
for i,(col,title,desc,sector) in enumerate(strategies):
    x=0.8+i*3.15
    rect(s,x,1.55,2.95,5.5,RGBColor(0x10,0x10,0x26))
    rect(s,x,1.55,2.95,0.08,col)
    tb(s,title,x+0.12,1.7,2.7,0.8,size=15,color=WHITE,bold=True)
    tb(s,desc,x+0.12,2.55,2.7,1.8,size=10,color=GREY)
    rect(s,x+0.12,4.5,2.71,0.32,col)
    tb(s,sector,x+0.15,4.55,2.65,0.28,size=9,color=WHITE,bold=True)

# ── Slide 4: Norwegian Presence ───────────────────────────────────────────────
s=ns(); bg(s)
rect(s,0,0,0.5,7.5,GREEN)
ln(s,0.8,1.3,11.5,GREEN)
tb(s,"NORSK TILSTEDEVÆRELSE",0.8,0.4,9.0,0.45,size=11,color=GREEN,bold=True)
tb(s,"Vi leverer allerede i Norge — ikke et løfte, en realitet.",0.8,0.7,11.0,0.65,size=24,color=WHITE,bold=True)
# Maxbo case
rect(s,0.8,1.55,7.5,3.2,NAVY)
rect(s,0.8,1.55,7.5,0.06,GREEN)
tb(s,"LIVE REFERANSE — MAXBO",1.0,1.67,7.0,0.35,size=10,color=GREEN,bold=True)
tb(s,"Maxbo — Norges ledende byggevare- og hagesenter",1.0,2.07,7.0,0.5,size=18,color=WHITE,bold=True)
details=[
    "✓  JAKALA composable frontend bygget på Shopify Plus",
    "✓  Enterspeed som headless content delivery layer",
    "✓  Speedtrain onboarding — raskeste implementeringsmetodikk",
    "✓  Live i produksjon — ikke et pilotprosjekt",
    "✓  Eneste aktive kombinasjon av Shopify Plus + Enterspeed i norsk retail",
]
for i,d in enumerate(details):
    tb(s,d,1.0,2.65+i*0.37,7.0,0.34,size=11,color=WHITE if i>0 else AMBER)
# Stats boxes
for i,(val,lbl) in enumerate([("Shopify Plus","Platform"),("Enterspeed","Content Layer"),("Speedtrain","Metode")]):
    x=8.55+i*0.01; y_start=1.55+i*1.1
    rect(s,8.5,y_start,4.65,0.95,RGBColor(0x05,0x14,0x08) if i==0 else (RGBColor(0x08,0x10,0x04) if i==1 else RGBColor(0x04,0x10,0x06)))
    rect(s,8.5,y_start,4.65,0.06,GREEN)
    tb(s,val,8.65,y_start+0.1,4.35,0.45,size=16,color=WHITE,bold=True)
    tb(s,lbl,8.65,y_start+0.57,4.35,0.3,size=11,color=GREY)
# Transferability
rect(s,0.8,4.92,11.5,1.75,RGBColor(0x08,0x14,0x08))
tb(s,"Hva betyr dette for deg?",1.0,5.05,11.0,0.38,size=12,color=GREEN,bold=True)
tb(s,"Vi kommer ikke med et konsept. Vi kommer med en løsning som allerede kjører i norsk retail — med lokal erfaring, lokale folk og et bevist resultat. Det betyr kortere oppstartstid, lavere risiko og en partner som kjenner det norske markedet.",1.0,5.45,11.2,0.95,size=11,color=WHITE)

# ── Slide 5: Our Approach ─────────────────────────────────────────────────────
s=ns(); bg(s)
rect(s,0,0,0.5,7.5,BLUE)
ln(s,0.8,1.3,11.5)
tb(s,"VÅR METODE",0.8,0.4,7.0,0.45,size=11,color=BLUE,bold=True)
tb(s,"Vi starter alltid lite. Vi beviser verdi raskt.",0.8,0.7,11.0,0.65,size=28,color=WHITE,bold=True)
phases=[
    ("1","Diagnostikk","2–4 uker","Vi kartlegger hva som faktisk koster og hva som gir mest effekt. Ingen antagelser — bare data.",BLUE),
    ("2","Speedtrain","4–8 uker","Raskt bevist konsept på avgrensede scope. Du ser resultater før vi skalerer.",GREEN),
    ("3","Optimalisering","Løpende","Vi skalerer det som virker. Kontinuerlig forbedring basert på data og resultater.",AMBER),
    ("4","Transformasjon","12M+","AI, plattform og organisasjonsutvikling for de som vil lede markedet.",RED),
]
for i,(num,title,dur,desc,col) in enumerate(phases):
    x=0.8+i*3.15
    rect(s,x,1.55,2.95,4.5,RGBColor(0x10,0x10,0x26))
    rect(s,x,1.55,2.95,0.06,col)
    tb(s,f"Fase {num}",x+0.12,1.68,2.7,0.3,size=9,color=col,bold=True)
    tb(s,title,x+0.12,2.0,2.7,0.55,size=18,color=WHITE,bold=True)
    tb(s,dur,x+0.12,2.58,2.7,0.3,size=10,color=AMBER,bold=True)
    tb(s,desc,x+0.12,2.95,2.7,1.4,size=10,color=GREY)
    if i<3:
        tb(s,"→",4.0+i*3.15-0.25,3.5,0.4,0.4,size=20,color=col,bold=True,align=PP_ALIGN.CENTER)
rect(s,0.8,6.22,11.5,0.9,RGBColor(0x05,0x05,0x25))
tb(s,"Alle engasjementer starter med Fase 1 — aldri et blankt ark. Du betaler for innsikt, ikke for tid.",1.0,6.35,11.2,0.55,size=12,color=WHITE,italic=True)

# ── Slide 6: Why JAKALA ───────────────────────────────────────────────────────
s=ns(); bg(s)
rect(s,0,0,0.5,7.5,BLUE)
ln(s,0.8,1.3,11.5)
tb(s,"HVORFOR JAKALA",0.8,0.4,7.0,0.45,size=11,color=BLUE,bold=True)
tb(s,"Seks grunner til at ledende selskaper velger oss.",0.8,0.7,11.0,0.65,size=26,color=WHITE,bold=True)
reasons=[
    (BLUE,"Vi leverer, vi rådgiver ikke bare","Alle engasjementer avsluttes med et bevist resultat — ikke en rapport."),
    (GREEN,"Norsk referanse, ikke et løfte","Maxbo er live på vår teknologi akkurat nå."),
    (AMBER,"Global kapasitet, lokal innsikt","5.000 ansatte globalt · team som kjenner norsk retail og finans."),
    (RED,"Shopify Premier Agency","1 av 3 i Europa · 30+ sertifiseringer · 20+ fullførte prosjekter."),
    (BLUE,"Data- og AI-spesialister","Vi bygger AI-løsninger som faktisk kjøres i produksjon."),
    (GREEN,"Fartsfokusert metode","Speedtrain: fra kontrakt til live resultat på 6–8 uker."),
]
for i,(col,title,desc) in enumerate(reasons):
    row=i//2; column=i%2
    x=0.8+column*6.2; y=1.55+row*1.6
    rect(s,x,y,5.9,1.45,RGBColor(0x10,0x10,0x26))
    rect(s,x,y,0.06,1.45,col)
    tb(s,title,x+0.2,y+0.12,5.5,0.42,size=13,color=WHITE,bold=True)
    tb(s,desc,x+0.2,y+0.58,5.5,0.65,size=10,color=GREY)

# ── Slide 7: Call to Action ───────────────────────────────────────────────────
s=ns(); bg(s)
rect(s,0,0,0.5,7.5,BLUE)
rect(s,0.5,0,12.83,7.5,NAVY)
tb(s,"Neste steg er enkelt.",1.2,1.5,10.0,1.2,size=44,color=WHITE,bold=True)
steps=[
    ("1","30-minutters samtale","Vi viser deg hva vi ser i ditt marked og hvordan vi kan hjelpe."),
    ("2","Gratis diagnostikk-brief","Vi sender en én-sides vurdering av din situasjon — ingen forpliktelser."),
    ("3","Pilot-forslag","Dersom det gir mening starter vi med en avgrenset fase 1."),
]
for i,(num,title,desc) in enumerate(steps):
    y=3.0+i*1.2
    rect(s,1.2,y,0.7,0.85,BLUE)
    tb(s,num,1.2,y+0.1,0.68,0.65,size=22,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
    tb(s,title,2.1,y+0.05,9.0,0.38,size=15,color=WHITE,bold=True)
    tb(s,desc,2.1,y+0.48,9.0,0.35,size=11,color=GREY)
ln(s,1.2,6.4,11.13,BLUE)
tb(s,"Jacob Skaue  ·  jacob.skaue@jakala.com  ·  JAKALA Nordic  ·  jakala.com",1.2,6.55,11.0,0.5,size=12,color=GREY)

prs.save(PPTX_OUT)
print(f"PPTX saved: {PPTX_OUT}")

# ── Create README for the package ────────────────────────────────────────────
readme = f"""# JAKALA Nordic — Outreach Package · Mars 2026

Generert: {date.today().isoformat()}
Accounts: {len(accounts)}

## Struktur

Hver account-mappe inneholder:
- `account-overview.md` — full account research og ICP/deal score
- `outreach-email.md`   — tilpasset norsk outreach-email klar til sending
- `offering.md`         — tailored offering one-pager for dette selskapet

Delt på tvers av alle accounts:
- `_jakala-company-overview.pptx` — JAKALA company presentation (7 slides)

## Accounts inkludert ({len(accounts)})

| # | Account | Sektor | Strategi | Deal-verdi | Primær kjøper |
|---|---------|--------|----------|-----------|---------------|
"""
for i, a in enumerate(accounts, 1):
    readme += f"| {i} | {a['name']} | {a['sector']} | {a['gtm_strategy']} | {a['deal_value_range']} | {a['primary_buyer_name']} |\n"

readme += f"""
## Bruksveiledning

1. Start med accounts der kjøper er navngitt og timing-signalet er sterkt
2. Tilpass subject-linja i outreach-email til din egen stil
3. Bruk offering.md som vedlegg eller som basis for en samtale
4. _jakala-company-overview.pptx brukes i møter for å introdusere JAKALA

## Prioritert rekkefølge (send først)

1. Sport Outlet — CTO + CDO vakante · CEO er beslutningsstaker NÅ
2. Trumf — Rikke Etholm-Idsøe ny Commercial Director · dag {(date.today() - date(2026,2,27)).days}
3. Vinmonopolet — Espen Terland ny CDO · honeymoon-fase
4. Helly Hansen — post-Kontoor oppkjøp · ny agenda åpen
5. Elkjøp Nordic — Morten Syversen identifisert · B2B-ekspansjon pågår
"""

with open(OUT / "_README.md", "w") as f:
    f.write(readme)

# ── Zip everything ─────────────────────────────────────────────────────────────
zip_path = BASE / "outreach-norge-2026-03.zip"
with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
    for path in sorted(OUT.rglob("*")):
        if path.is_file():
            arcname = path.relative_to(OUT)
            zf.write(path, arcname)

print(f"\nZip created: {zip_path}")
print(f"Contents: {len(list(OUT.rglob('*')))} files across {len(accounts)} account folders")

# Print summary
print(f"\n{'='*60}")
print(f"OUTREACH PACKAGE SUMMARY")
print(f"{'='*60}")
print(f"Accounts:  {len(accounts)}")
print(f"Files:     {len(accounts)*3 + 2} (3 per account + README + PPTX)")
print(f"Output:    outreach-norge-2026-03.zip")
print(f"{'='*60}")
