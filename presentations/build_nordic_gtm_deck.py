"""
JAKALA Nordic GTM — New Business Lines Deck
For: Paolo Pedersoli, Managing Director & Acting Global CMO
Date: March 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── Design Constants ────────────────────────────────────────────────────────
BLUE_BRIGHT = RGBColor(0x15, 0x3E, 0xED)  # #153EED
BLUE_DARK   = RGBColor(0x02, 0x02, 0x66)  # #020266
RED         = RGBColor(0xF6, 0x57, 0x4A)  # #F6574A
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY        = RGBColor(0xAA, 0xAA, 0xBB)
BG_COLOR    = RGBColor(0x08, 0x08, 0x18)  # #080818
FONT        = 'Raleway'
W           = Inches(13.33)  # widescreen 16:9
H           = Inches(7.5)

# ─── Helpers ─────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)

def set_bg(slide, color=BG_COLOR):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, x, y, w, h,
                size=16, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=Pt(4)):
    from pptx.util import Pt as PPt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def add_rect(slide, x, y, w, h, color=BLUE_DARK, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_line(slide, x, y, w, color=BLUE_BRIGHT, thickness=Pt(1.5)):
    line = slide.shapes.add_shape(1, x, y, w, Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def add_card(slide, x, y, w, h, title, body_lines,
             title_color=BLUE_BRIGHT, body_color=WHITE,
             bg=BLUE_DARK, title_size=14, body_size=11):
    add_rect(slide, x, y, w, h, bg)
    # title
    txBox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.12),
                                      w - Inches(0.3), Inches(0.38))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = title_color
    # body
    bxBox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.42),
                                      w - Inches(0.3), h - Inches(0.55))
    btf = bxBox.text_frame
    btf.word_wrap = True
    for i, line in enumerate(body_lines):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(body_size)
        run.font.color.rgb = body_color

# ─── Slides ──────────────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = blank_slide(prs)
    set_bg(slide)

    # Left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    # Top-left: JAKALA wordmark area
    add_textbox(slide, "JAKALA", Inches(0.25), Inches(0.3), Inches(3), Inches(0.5),
                size=13, bold=True, color=BLUE_BRIGHT)
    add_textbox(slide, "Nordic GTM", Inches(0.25), Inches(0.62), Inches(3), Inches(0.4),
                size=11, color=GREY)

    # Main headline — left aligned, large
    add_textbox(slide,
                "New Business Lines\n& GTM Integration",
                Inches(0.25), Inches(1.6), Inches(8), Inches(2.2),
                size=44, bold=True, color=WHITE)

    # Subtitle line
    add_line(slide, Inches(0.25), Inches(3.75), Inches(5), BLUE_BRIGHT)

    # Subtitle text
    add_textbox(slide,
                "How Shopify, Hello Growth and Communication & Media\nfit our Nordic go-to-market — and what the opportunity looks like.",
                Inches(0.25), Inches(3.9), Inches(8.5), Inches(1.0),
                size=14, color=GREY)

    # Bottom meta
    add_textbox(slide, "For: Paolo Pedersoli, Managing Director & Acting Global CMO",
                Inches(0.25), H - Inches(1.1), Inches(8), Inches(0.4),
                size=11, color=GREY)
    add_textbox(slide, "March 2026  ·  Internal",
                Inches(0.25), H - Inches(0.75), Inches(4), Inches(0.35),
                size=10, color=GREY)

    # Right side visual accent block
    add_rect(slide, Inches(9.5), Inches(0), Inches(3.83), H, BLUE_DARK)
    add_rect(slide, Inches(9.6), Inches(1.5), Inches(3.5), Inches(0.06), BLUE_BRIGHT)
    add_textbox(slide, "€6.8M\nNordic Pipeline",
                Inches(9.8), Inches(1.8), Inches(3.0), Inches(1.0),
                size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_textbox(slide, "45 active accounts\n3 integrated practices",
                Inches(9.8), Inches(2.85), Inches(3.0), Inches(0.8),
                size=13, color=GREY)
    add_textbox(slide, "DK · NO · SE",
                Inches(9.8), Inches(6.7), Inches(2.5), Inches(0.4),
                size=11, color=BLUE_BRIGHT, bold=True)


def slide_agenda(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)
    add_textbox(slide, "Agenda", Inches(0.25), Inches(0.3), Inches(12), Inches(0.55),
                size=28, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(0.88), Inches(5), BLUE_BRIGHT)

    items = [
        ("01", "Our GTM thesis — how we enter and expand"),
        ("02", "New business lines — what's been added"),
        ("03", "Shopify practice — credentials and GTM fit"),
        ("04", "The connected model — Commerce · Data & AI · Growth"),
        ("05", "Nordic market opportunity — pipeline snapshot"),
        ("06", "Revenue architecture — entry → expand → transform"),
        ("07", "Priority accounts — top 6 ranked"),
        ("08", "What we need to accelerate"),
    ]

    for i, (num, text) in enumerate(items):
        y = Inches(1.2) + i * Inches(0.72)
        add_rect(slide, Inches(0.25), y, Inches(0.55), Inches(0.5), BLUE_DARK)
        add_textbox(slide, num, Inches(0.25), y + Inches(0.04), Inches(0.55), Inches(0.45),
                    size=14, bold=True, color=BLUE_BRIGHT, align=PP_ALIGN.CENTER)
        add_textbox(slide, text, Inches(0.92), y + Inches(0.06), Inches(9), Inches(0.4),
                    size=14, color=WHITE)


def slide_gtm_thesis(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "01 · Our GTM Thesis", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.5), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "We don't sell transformation.\nWe earn the right to deliver it.",
                Inches(0.25), Inches(0.7), Inches(10), Inches(1.3),
                size=30, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(1.95), Inches(6), BLUE_BRIGHT)

    add_textbox(slide,
                "Enterprise clients don't start with large platform programmes.\nThey start with urgent, specific business problems.",
                Inches(0.25), Inches(2.1), Inches(9), Inches(0.8),
                size=14, color=GREY)

    # Three step arrows
    steps = [
        ("DIAGNOSE", "Small, paid diagnostic.\nIdentify the problem.\nBuild trust.", BLUE_DARK),
        ("PILOT", "Focused 6–12 week\ndelivery. Prove value.\nCreate internal champion.", BLUE_DARK),
        ("TRANSFORM", "Expand into full\nprogramme. Platform +\nData + AI.", RGBColor(0x05, 0x05, 0x40)),
    ]
    for i, (title, body, bg) in enumerate(steps):
        x = Inches(0.25) + i * Inches(4.3)
        add_rect(slide, x, Inches(3.1), Inches(4.0), Inches(3.2), bg)
        add_line(slide, x, Inches(3.1), Inches(4.0),
                 BLUE_BRIGHT if i < 2 else RED)
        add_textbox(slide, title, x + Inches(0.2), Inches(3.2),
                    Inches(3.6), Inches(0.5), size=16, bold=True,
                    color=BLUE_BRIGHT if i < 2 else RED)
        add_textbox(slide, body, x + Inches(0.2), Inches(3.75),
                    Inches(3.6), Inches(2.4), size=12, color=WHITE)
        if i < 2:
            add_textbox(slide, "→", x + Inches(3.85), Inches(4.35),
                        Inches(0.5), Inches(0.5), size=22, bold=True,
                        color=BLUE_BRIGHT)

    add_textbox(slide,
                "This motion works across all three practices: Commerce, Data & AI, and Growth.",
                Inches(0.25), Inches(6.5), Inches(11), Inches(0.5),
                size=11, color=GREY, italic=True)


def slide_new_business_lines(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "02 · New Business Lines", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.5), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "Three new practices added to the Nordic GTM in 2025–2026",
                Inches(0.25), Inches(0.7), Inches(11), Inches(0.7),
                size=26, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(1.38), Inches(7), BLUE_BRIGHT)

    cards = [
        {
            "title": "Shopify Commerce Practice",
            "tag": "NEW · Commerce",
            "lines": [
                "Full Shopify Plus implementation practice.",
                "Premier Agency status (1 of 3 in Italy).",
                "30+ certified specialists · 20+ projects delivered.",
                "Entry offer: Shopify Check-up · €15–30K",
                "Full build: €80K–€500K",
                "",
                "Clients: Kiton, La Prairie, Barilla,",
                "Pomellato, Fratelli Rossetti, Jacob Cohën",
            ],
            "tag_color": BLUE_BRIGHT,
        },
        {
            "title": "Hello Growth",
            "tag": "NEW · Growth",
            "lines": [
                "SaaS marketplace acceleration programme.",
                "Azure Marketplace activation + ABM + demand gen.",
                "4-phase programme · Fixed fee DKK 150–250K",
                "Targets Nordic SaaS scale-ups &",
                "Microsoft Solution Partners (SDCs, ISVs).",
                "",
                "Separate ICP track — high-velocity,",
                "lower ticket, faster cycle.",
            ],
            "tag_color": RED,
        },
        {
            "title": "Communication & Media",
            "tag": "NEW · Growth + Data",
            "lines": [
                "Data-driven media strategy & activation.",
                "Awareness · Performance · Influencer.",
                "Proprietary tech stack unifying brand",
                "and performance via first-party data.",
                "Entry: Media Activation Audit · €50K+",
                "Retainer/campaign: €50K–€500K+",
                "",
                "Natural upsell from any Data/Commerce engagement.",
            ],
            "tag_color": RGBColor(0xF6, 0xA5, 0x4A),
        },
    ]

    for i, c in enumerate(cards):
        x = Inches(0.25) + i * Inches(4.35)
        add_rect(slide, x, Inches(1.6), Inches(4.1), Inches(5.5), BLUE_DARK)
        add_rect(slide, x, Inches(1.6), Inches(4.1), Inches(0.08), c["tag_color"])
        add_textbox(slide, c["tag"], x + Inches(0.15), Inches(1.72),
                    Inches(3.8), Inches(0.3), size=9, bold=True, color=c["tag_color"])
        add_textbox(slide, c["title"], x + Inches(0.15), Inches(2.0),
                    Inches(3.8), Inches(0.55), size=15, bold=True, color=WHITE)
        add_line(slide, x + Inches(0.15), Inches(2.55), Inches(3.6), c["tag_color"])
        body_box = slide.shapes.add_textbox(x + Inches(0.15), Inches(2.65),
                                             Inches(3.8), Inches(4.2))
        btf = body_box.text_frame
        btf.word_wrap = True
        for j, line in enumerate(c["lines"]):
            if j == 0:
                p = btf.paragraphs[0]
            else:
                p = btf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.name = FONT
            run.font.size = Pt(11)
            run.font.color.rgb = GREY if not line.startswith("Entry") and not line.startswith("Full") else WHITE
            if line.startswith("Entry") or line.startswith("Full"):
                run.font.bold = True


def slide_shopify_practice(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "03 · Shopify Practice", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.5), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "Premier Shopify Agency — integrated into our Nordic GTM",
                Inches(0.25), Inches(0.7), Inches(11), Inches(0.65),
                size=26, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(1.35), Inches(8), BLUE_BRIGHT)

    # Left: credentials
    add_rect(slide, Inches(0.25), Inches(1.6), Inches(5.8), Inches(5.5), BLUE_DARK)
    add_textbox(slide, "Our credentials", Inches(0.45), Inches(1.72),
                Inches(5.4), Inches(0.4), size=13, bold=True, color=BLUE_BRIGHT)
    creds = [
        "✓  Shopify Premier Agency — 1 of 3 in Italy",
        "✓  7+ years delivering Shopify commerce",
        "✓  30+ Shopify certifications across the team",
        "✓  20+ projects delivered in last 2 years",
        "✓  Clients: Kiton · La Prairie · Barilla · Pomellato",
        "    Fratelli Rossetti · Alessandra Rich · Jacob Cohën",
        "    Molteni & C",
        "✓  Full capability: UX/UI · Platform · Integrations",
        "    Subscriptions (Recharge) · Loyalty (Growave)",
        "    Marketing automation (Klaviyo) · CX (Gorgias)",
    ]
    cred_box = slide.shapes.add_textbox(Inches(0.45), Inches(2.15),
                                         Inches(5.4), Inches(4.7))
    ctf = cred_box.text_frame
    ctf.word_wrap = True
    for j, line in enumerate(creds):
        p = ctf.paragraphs[0] if j == 0 else ctf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(11.5)
        run.font.color.rgb = WHITE if line.startswith("✓") else GREY

    # Right: GTM fit
    add_rect(slide, Inches(6.3), Inches(1.6), Inches(6.78), Inches(2.5), RGBColor(0x05, 0x05, 0x40))
    add_textbox(slide, "GTM fit — who we sell to in the Nordics",
                Inches(6.5), Inches(1.72), Inches(6.4), Inches(0.4),
                size=13, bold=True, color=BLUE_BRIGHT)
    fit_lines = [
        "· Retailers on Magento 2 / Adobe Commerce wanting to replatform",
        "· D2C brands needing owned ecommerce channel",
        "· Fashion, lifestyle, furniture, beauty brands (DK/NO/SE)",
        "· Mid-market accounts: €80K–€500K project range",
        "· Example: Sport Outlet (18 parallel webshops)",
        "· Example: Ferner Jacobsen (Magento 2 replatform signal)",
        "· Example: Follestad (WooCommerce upgrade candidate)",
    ]
    fbox = slide.shapes.add_textbox(Inches(6.5), Inches(2.15), Inches(6.4), Inches(1.75))
    ftf = fbox.text_frame
    ftf.word_wrap = True
    for j, line in enumerate(fit_lines):
        p = ftf.paragraphs[0] if j == 0 else ftf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE

    # Delivery model
    add_rect(slide, Inches(6.3), Inches(4.25), Inches(6.78), Inches(2.85), BLUE_DARK)
    add_textbox(slide, "6-phase delivery model",
                Inches(6.5), Inches(4.37), Inches(6.4), Inches(0.4),
                size=13, bold=True, color=BLUE_BRIGHT)
    phases = [
        "1 · Discovery / Check-up  (Entry: €15–30K)",
        "2 · Analysis & UX Design",
        "3 · POC & Platform Selection",
        "4 · Roadmap & Business Case",
        "5 · Build & Delivery  (€80K–€500K)",
        "6 · Run & Scale  (ongoing retainer)",
    ]
    pbox = slide.shapes.add_textbox(Inches(6.5), Inches(4.8), Inches(6.4), Inches(2.0))
    ptf = pbox.text_frame
    ptf.word_wrap = True
    for j, line in enumerate(phases):
        p = ptf.paragraphs[0] if j == 0 else ptf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(11)
        run.font.color.rgb = BLUE_BRIGHT if "Entry" in line or "€80K" in line else WHITE


def slide_connected_model(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "04 · The Connected GTM Model", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.5), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "Three practices. One integrated commercial motion.",
                Inches(0.25), Inches(0.72), Inches(11), Inches(0.65),
                size=28, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(1.38), Inches(8), BLUE_BRIGHT)

    add_textbox(slide,
                "Each practice opens independently — and every engagement creates an expansion path into the others.",
                Inches(0.25), Inches(1.55), Inches(11), Inches(0.5),
                size=13, color=GREY)

    practices = [
        {
            "name": "COMMERCE",
            "color": BLUE_BRIGHT,
            "offerings": [
                "Shopify Commerce Build",
                "Commerce Optimization Pilot",
                "DXP Transformation",
            ],
            "icp": "Head of Ecommerce · CMO · CTO",
            "entry": "€15–80K",
            "expand": "→ Data layer → Speedtrain → AI",
        },
        {
            "name": "DATA & AI",
            "color": RGBColor(0x60, 0x9A, 0xFF),
            "offerings": [
                "Data Revenue Diagnostic + Speedtrain",
                "AI Readiness Diagnostic",
                "CDP / CRM / Customer 360",
                "BI & Data Visualisation",
            ],
            "icp": "CDO · Head of Data · CTO",
            "entry": "€50–100K",
            "expand": "→ AI architecture → DXP",
        },
        {
            "name": "GROWTH",
            "color": RED,
            "offerings": [
                "Hello Growth (SaaS Marketplace)",
                "Communication & Media Activation",
                "Pricing & Revenue Management",
                "Strategy & Business Growth",
            ],
            "icp": "CMO · CRO · CEO · CFO",
            "entry": "€20–80K",
            "expand": "→ Commerce → Data → Full retainer",
        },
    ]

    for i, pr in enumerate(practices):
        x = Inches(0.25) + i * Inches(4.35)
        add_rect(slide, x, Inches(2.25), Inches(4.1), Inches(4.9), BLUE_DARK)
        add_rect(slide, x, Inches(2.25), Inches(4.1), Inches(0.1), pr["color"])

        add_textbox(slide, pr["name"], x + Inches(0.15), Inches(2.4),
                    Inches(3.8), Inches(0.45), size=16, bold=True, color=pr["color"])

        obox = slide.shapes.add_textbox(x + Inches(0.15), Inches(2.9),
                                         Inches(3.8), Inches(2.0))
        otf = obox.text_frame
        otf.word_wrap = True
        for j, off in enumerate(pr["offerings"]):
            p = otf.paragraphs[0] if j == 0 else otf.add_paragraph()
            run = p.add_run()
            run.text = f"· {off}"
            run.font.name = FONT
            run.font.size = Pt(11)
            run.font.color.rgb = WHITE

        add_line(slide, x + Inches(0.15), Inches(4.95), Inches(3.7), pr["color"])

        add_textbox(slide, f"Buyer: {pr['icp']}", x + Inches(0.15), Inches(5.05),
                    Inches(3.8), Inches(0.35), size=10, color=GREY)
        add_textbox(slide, f"Entry: {pr['entry']}", x + Inches(0.15), Inches(5.38),
                    Inches(3.8), Inches(0.35), size=10, bold=True, color=WHITE)
        add_textbox(slide, pr["expand"], x + Inches(0.15), Inches(5.72),
                    Inches(3.8), Inches(0.35), size=10, color=pr["color"], bold=True)

    # Bottom connector text
    add_rect(slide, Inches(0.25), Inches(7.0), Inches(12.83), Inches(0.35),
             RGBColor(0x05, 0x05, 0x35))
    add_textbox(slide,
                "Speedtrain (proprietary data orchestration) is the connective tissue — it sits between every commerce implementation and the AI/personalisation layer.",
                Inches(0.4), Inches(7.02), Inches(12.5), Inches(0.3),
                size=10, color=GREY)


def slide_nordic_opportunity(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "05 · Nordic Market Opportunity", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.5), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "Denmark · Norway · Sweden",
                Inches(0.25), Inches(0.72), Inches(11), Inches(0.65),
                size=28, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(1.38), Inches(5), BLUE_BRIGHT)

    # KPI boxes
    kpis = [
        ("€6.8M", "Total pipeline\n(unweighted)"),
        ("45", "Active\nopportunities"),
        ("18", "Named buyers\nconfirmed"),
        ("€420K", "Base case Q2 2026\nforecast"),
    ]
    for i, (val, label) in enumerate(kpis):
        x = Inches(0.25) + i * Inches(3.25)
        add_rect(slide, x, Inches(1.65), Inches(3.0), Inches(1.4), BLUE_DARK)
        add_textbox(slide, val, x + Inches(0.15), Inches(1.75),
                    Inches(2.7), Inches(0.65), size=28, bold=True,
                    color=BLUE_BRIGHT, align=PP_ALIGN.LEFT)
        add_textbox(slide, label, x + Inches(0.15), Inches(2.38),
                    Inches(2.7), Inches(0.55), size=10, color=GREY)

    # Sector breakdown
    add_textbox(slide, "Account mix by sector", Inches(0.25), Inches(3.25),
                Inches(6), Inches(0.4), size=13, bold=True, color=WHITE)
    sectors = [
        ("Retail & Ecommerce", "22 accounts", "Core ICP — Commerce + Data"),
        ("Fashion & Lifestyle", "8 accounts", "Shopify replatform + D2C"),
        ("Financial Services", "6 accounts", "AI Readiness + Data"),
        ("Education", "4 accounts", "AI Readiness — transformation window"),
        ("Grocery & FMCG", "5 accounts", "Data Revenue + Loyalty activation"),
    ]
    for i, (sector, count, note) in enumerate(sectors):
        y = Inches(3.75) + i * Inches(0.6)
        add_rect(slide, Inches(0.25), y, Inches(6.5), Inches(0.5), RGBColor(0x0D, 0x0D, 0x28))
        add_textbox(slide, sector, Inches(0.4), y + Inches(0.07),
                    Inches(2.5), Inches(0.35), size=11, bold=True, color=WHITE)
        add_textbox(slide, count, Inches(2.95), y + Inches(0.07),
                    Inches(1.0), Inches(0.35), size=11, color=BLUE_BRIGHT, bold=True)
        add_textbox(slide, note, Inches(4.0), y + Inches(0.07),
                    Inches(2.8), Inches(0.35), size=10, color=GREY)

    # Timing signals panel
    add_rect(slide, Inches(7.0), Inches(1.65), Inches(6.08), Inches(5.5), BLUE_DARK)
    add_textbox(slide, "HOT TIMING SIGNALS — ACT NOW",
                Inches(7.2), Inches(1.77), Inches(5.7), Inches(0.4),
                size=11, bold=True, color=RED)
    signals = [
        ("Sport Outlet", "CTO + CDO both vacant · Contact CEO directly"),
        ("Trumf", "New Commercial Director · First 90 days window"),
        ("Vinmonopolet", "New CDO (ex-XXL) · Honeymoon phase · Agenda not set"),
        ("Skeidar", "CIO confirmed · 'Best furniture portal' ambition declared"),
        ("BI Handelshøyskolen", "Rector leaving Aug 2026 · Transition window open"),
    ]
    for i, (account, signal) in enumerate(signals):
        y = Inches(2.25) + i * Inches(0.95)
        add_rect(slide, Inches(7.2), y, Inches(5.7), Inches(0.85), RGBColor(0x0A, 0x0A, 0x22))
        add_textbox(slide, account, Inches(7.35), y + Inches(0.05),
                    Inches(5.3), Inches(0.35), size=12, bold=True, color=RED)
        add_textbox(slide, signal, Inches(7.35), y + Inches(0.42),
                    Inches(5.3), Inches(0.35), size=10, color=WHITE)


def slide_revenue_architecture(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "06 · Revenue Architecture", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.5), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "Entry · Expand · Transform",
                Inches(0.25), Inches(0.72), Inches(11), Inches(0.65),
                size=28, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(1.38), Inches(5), BLUE_BRIGHT)

    add_textbox(slide,
                "Every engagement follows the same commercial logic — regardless of which practice opens the account.",
                Inches(0.25), Inches(1.55), Inches(10), Inches(0.5),
                size=13, color=GREY)

    # Revenue ladder
    levels = [
        {
            "phase": "ENTRY",
            "range": "€15K – €100K",
            "offers": "Shopify Check-up · Data Revenue Diagnostic · AI Readiness Diagnostic\nCommerce Optimization Pilot · Media Audit · Hello Growth",
            "w": Inches(4.5),
            "color": BLUE_BRIGHT,
        },
        {
            "phase": "EXPAND",
            "range": "€100K – €500K",
            "offers": "Shopify Commerce Build · Speedtrain Implementation\nAI Architecture Programme · CDP/CRM Implementation · Full Media Activation",
            "w": Inches(7.0),
            "color": RGBColor(0x60, 0x9A, 0xFF),
        },
        {
            "phase": "TRANSFORM",
            "range": "€500K – €3M+",
            "offers": "DXP Transformation Programme · Multi-market Commerce + Data + AI\nFull integrated platform — Commerce · Orchestration · Personalisation",
            "w": Inches(12.58),
            "color": RED,
        },
    ]

    for i, lvl in enumerate(levels):
        y = Inches(2.3) + i * Inches(1.55)
        add_rect(slide, Inches(0.25), y, lvl["w"], Inches(1.35),
                 RGBColor(0x05, 0x05, 0x30) if i < 2 else RGBColor(0x18, 0x02, 0x08))
        add_rect(slide, Inches(0.25), y, Inches(0.08), Inches(1.35), lvl["color"])
        add_textbox(slide, lvl["phase"], Inches(0.45), y + Inches(0.08),
                    Inches(1.8), Inches(0.4), size=14, bold=True, color=lvl["color"])
        add_textbox(slide, lvl["range"], Inches(2.3), y + Inches(0.08),
                    Inches(2.5), Inches(0.4), size=14, bold=True, color=WHITE)
        add_textbox(slide, lvl["offers"], Inches(0.45), y + Inches(0.55),
                    lvl["w"] - Inches(0.4), Inches(0.75), size=10, color=GREY)

    # Bottom insight
    add_rect(slide, Inches(0.25), Inches(7.0), Inches(12.83), Inches(0.35),
             RGBColor(0x05, 0x05, 0x35))
    add_textbox(slide,
                "Target revenue per account: Entry €50–80K → Year 1 expansion €200–500K → Full transformation €500K–3M+  ·  "
                "No account should stay at entry level beyond 12 months.",
                Inches(0.4), Inches(7.02), Inches(12.5), Inches(0.3),
                size=10, color=GREY)


def slide_top_accounts(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "07 · Priority Accounts", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.5), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "Top 6 ranked by value × strategic fit",
                Inches(0.25), Inches(0.72), Inches(11), Inches(0.65),
                size=28, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(1.38), Inches(6), BLUE_BRIGHT)

    accounts = [
        ("H&M Group", "SE", "€900K", "9", "Data Revenue Unlock", "Data Revenue Diagnostic → Speedtrain → DXP"),
        ("Matas Group", "DK", "€700K", "9", "AI Readiness", "AI Readiness Diagnostic → Speedtrain → CDP/CRM"),
        ("Elkjøp Nordic", "NO", "€700K", "8", "Commerce Optimization", "Commerce Optimization → B2B → Speedtrain"),
        ("Varner Group", "NO", "€450K", "9", "Data Revenue Unlock", "Data Revenue Diagnostic → Speedtrain (Sitoo gap)"),
        ("Trumf", "NO", "€450K", "9", "Data Revenue Unlock", "Data Revenue Diagnostic → AI Personalisation"),
        ("Clas Ohlson", "SE", "€350K", "7", "Commerce Optimization", "Commerce Optimization Pilot → Shopify / DXP"),
    ]

    headers = ["#", "Account", "Mkt", "Pipeline", "ICP", "GTM Strategy", "Expansion Path"]
    col_widths = [Inches(0.4), Inches(1.9), Inches(0.5), Inches(0.9), Inches(0.5), Inches(2.2), Inches(5.48)]
    col_x = [Inches(0.25)]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)

    # Header row
    add_rect(slide, Inches(0.25), Inches(1.6), Inches(12.83), Inches(0.42), BLUE_DARK)
    for j, (hdr, cx) in enumerate(zip(headers, col_x)):
        add_textbox(slide, hdr, cx + Inches(0.05), Inches(1.65),
                    col_widths[j] - Inches(0.05), Inches(0.35),
                    size=10, bold=True, color=BLUE_BRIGHT)

    for i, (acc, mkt, val, icp, gtm, exp) in enumerate(accounts):
        y = Inches(2.1) + i * Inches(0.78)
        row_bg = RGBColor(0x0D, 0x0D, 0x28) if i % 2 == 0 else RGBColor(0x0A, 0x0A, 0x20)
        add_rect(slide, Inches(0.25), y, Inches(12.83), Inches(0.72), row_bg)

        row_data = [str(i+1), acc, mkt, val, icp, gtm, exp]
        for j, (cell, cx) in enumerate(zip(row_data, col_x)):
            color = WHITE
            if j == 0:
                color = GREY
            elif j == 3:
                color = BLUE_BRIGHT
            elif j == 4:
                # ICP score coloring
                score = int(icp)
                color = BLUE_BRIGHT if score >= 9 else (WHITE if score >= 8 else GREY)
            add_textbox(slide, cell, cx + Inches(0.05), y + Inches(0.18),
                        col_widths[j] - Inches(0.05), Inches(0.42),
                        size=10, bold=(j == 1), color=color)


def slide_next_steps(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "08 · What We Need to Accelerate", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.5), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "From pipeline to closed — three asks",
                Inches(0.25), Inches(0.72), Inches(11), Inches(0.65),
                size=28, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(1.38), Inches(6), BLUE_BRIGHT)

    asks = [
        {
            "num": "01",
            "title": "Referenceable client cases — Nordic-relevant",
            "body": (
                "The Shopify practice has strong Italian references (Kiton, La Prairie, Barilla).\n"
                "We need 1–2 cases positioned for Nordic retail: furniture, fashion, outdoor/sport.\n"
                "Ask: Can we package Molteni or Fratelli Rossetti as a Nordic-facing case study?\n"
                "Impact: Unlocks Helly Hansen, Sport Outlet, Ferner Jacobsen, Follestad conversations."
            ),
        },
        {
            "num": "02",
            "title": "Shopify + Speedtrain joint narrative",
            "body": (
                "The most powerful differentiator we have is: Shopify implementation backed by\n"
                "proprietary data orchestration (Speedtrain). No competitor offers this combination.\n"
                "Ask: One-page joint narrative showing Shopify Build → Speedtrain integration path.\n"
                "Impact: Justifies premium pricing and creates a natural €300K+ expansion path."
            ),
        },
        {
            "num": "03",
            "title": "First meetings — conversion from 0%",
            "body": (
                "We have 45 active accounts. 0 first meetings booked. The pipeline is warm — not cold.\n"
                "5 accounts have critical leadership transitions happening NOW (Sport Outlet, Trumf,\n"
                "Vinmonopolet, Skeidar, BI Handelshøyskolen).\n"
                "Ask: Nordic CCO endorsement for outreach on the top 5 timing-critical accounts.\n"
                "Impact: First €350–500K in entry engagements booked by end of Q2 2026."
            ),
        },
    ]

    for i, ask in enumerate(asks):
        y = Inches(1.7) + i * Inches(1.8)
        add_rect(slide, Inches(0.25), y, Inches(12.83), Inches(1.68),
                 RGBColor(0x08, 0x08, 0x25))
        add_rect(slide, Inches(0.25), y, Inches(0.08), Inches(1.68), BLUE_BRIGHT)
        # Number badge
        add_rect(slide, Inches(0.4), y + Inches(0.12), Inches(0.5), Inches(0.5), BLUE_DARK)
        add_textbox(slide, ask["num"], Inches(0.4), y + Inches(0.15),
                    Inches(0.5), Inches(0.42), size=13, bold=True, color=BLUE_BRIGHT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, ask["title"], Inches(1.05), y + Inches(0.12),
                    Inches(11.7), Inches(0.42), size=14, bold=True, color=WHITE)
        add_textbox(slide, ask["body"], Inches(1.05), y + Inches(0.55),
                    Inches(11.7), Inches(1.05), size=11, color=GREY)

    # Footer
    add_rect(slide, Inches(0.25), Inches(7.1), Inches(12.83), Inches(0.3),
             RGBColor(0x05, 0x05, 0x35))
    add_textbox(slide,
                "Prepared by: Nordic GTM Team · March 2026  ·  For: Paolo Pedersoli, Managing Director & Acting Global CMO",
                Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.25),
                size=9, color=GREY)


# ─── Main ────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_cover(prs)
    slide_agenda(prs)
    slide_gtm_thesis(prs)
    slide_new_business_lines(prs)
    slide_shopify_practice(prs)
    slide_connected_model(prs)
    slide_nordic_opportunity(prs)
    slide_revenue_architecture(prs)
    slide_top_accounts(prs)
    slide_next_steps(prs)

    out = "/Users/jacobskaue/Desktop/jakala-commercial-os/presentations/nordic-gtm-new-business-lines-2026-03.pptx"
    prs.save(out)
    print(f"✓ Saved: {out}")
    print(f"  Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
