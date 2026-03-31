"""
JAKALA Nordic — Norges Topp 10 Leads 2026
Shopify + Enterspeed fokus · Norsk · Mars 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand constants ──────────────────────────────────────────────────────────
BLUE   = RGBColor(0x15, 0x3E, 0xED)   # #153EED
NAVY   = RGBColor(0x02, 0x02, 0x66)   # #020266
RED    = RGBColor(0xF6, 0x57, 0x4A)   # #F6574A
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0xAA, 0xAA, 0xCC)
BG     = RGBColor(0x08, 0x08, 0x18)   # #080818
GREEN  = RGBColor(0x22, 0xDD, 0x88)
AMBER  = RGBColor(0xFF, 0xBB, 0x33)

FONT = 'Raleway'
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(blank)

def bg(slide):
    rect = slide.shapes.add_shape(1, 0, 0, W, H)
    rect.fill.solid(); rect.fill.fore_color.rgb = BG
    rect.line.fill.background()

def rect(slide, x, y, w, h, color, alpha=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def line(slide, x, y, w, color=BLUE, thickness=Pt(1.5)):
    ln = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.02))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background()

def tb(slide, text, x, y, w, h, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = FONT
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.italic = italic
    return txb

def para(tf, text, size=13, color=WHITE, bold=False, align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False):
    from pptx.util import Pt as Pt2
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name  = FONT
    run.font.size  = Pt2(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.italic = italic
    return p

def tag(slide, label, x, y, color=BLUE):
    """Small pill label"""
    rect(slide, x, y, len(label)*0.095 + 0.2, 0.28, color)
    tb(slide, label, x+0.08, y+0.02, len(label)*0.12 + 0.2, 0.26, size=9, color=WHITE, bold=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(); bg(s)
rect(s, 0, 0, 0.5, 7.5, BLUE)
rect(s, 0.5, 5.5, 12.83, 0.06, BLUE)

tb(s, "NORGE", 1.0, 1.2, 11.0, 0.8, size=14, color=BLUE, bold=True)
tb(s, "Topp 10 prioriterte\naccounts 2026", 1.0, 1.9, 11.0, 2.0, size=54, color=WHITE, bold=True)
tb(s, "Shopify · Enterspeed · Headless Commerce", 1.0, 4.4, 11.0, 0.6, size=18, color=GREY)
tb(s, "JAKALA Nordic   ·   Mars 2026", 1.0, 5.7, 8.0, 0.5, size=13, color=GREY)
tb(s, "KONFIDENSIELT", 9.5, 5.7, 3.5, 0.5, size=11, color=RED, bold=True, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — NORSK MARKED: MULIGHETSVINDU
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(); bg(s)
rect(s, 0, 0, 0.5, 7.5, BLUE)
line(s, 0.8, 1.3, 11.5)
tb(s, "NORSK MARKED", 0.8, 0.4, 6.0, 0.5, size=11, color=BLUE, bold=True)
tb(s, "Mulighetsvindu åpner seg nå", 0.8, 0.7, 11.0, 0.7, size=30, color=WHITE, bold=True)

signals = [
    ("Sport Outlet",     "CTO + CDO begge vakante pr. mars 2026. Kontakt CEO direkte nå.",   RED,   "🔴 Kritisk timing"),
    ("Trumf",            "Rikke Etholm-Idsøe ny Commercial Director. Første 90 dager = vindu.", AMBER, "🟡 Ny rolle"),
    ("Vinmonopolet",     "Espen Terland ny CDO (ex-XXL 15 år). Agenda ikke satt ennå.",       AMBER, "🟡 Ny rolle"),
    ("Helly Hansen",     "55 Adobe Commerce-sider. Kontoor Brands-oppkjøp juni 2025 = ny agenda.", AMBER, "🟡 Etter oppkjøp"),
    ("Skeidar",          "CIO Sujit Nath identifisert + 'beste møbelportal'-ambisjon erklært.", GREEN, "🟢 Erklært ambisjon"),
]

for i, (acc, desc, col, badge) in enumerate(signals):
    y = 1.6 + i * 1.0
    rect(s, 0.8, y, 11.5, 0.85, RGBColor(0x10, 0x10, 0x28))
    rect(s, 0.8, y, 0.06, 0.85, col)
    tb(s, acc, 1.1, y+0.06, 3.5, 0.32, size=13, color=WHITE, bold=True)
    tb(s, badge, 8.5, y+0.06, 3.5, 0.32, size=10, color=col, bold=True, align=PP_ALIGN.RIGHT)
    tb(s, desc, 1.1, y+0.4, 10.5, 0.38, size=11, color=GREY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — VÅR NORSKE REFERANSE + JOINT TILBUD
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(); bg(s)
rect(s, 0, 0, 0.5, 7.5, BLUE)
line(s, 0.8, 1.3, 11.5)
tb(s, "POSISJON OG TILBUD", 0.8, 0.4, 6.0, 0.5, size=11, color=BLUE, bold=True)
tb(s, "Vår norske referanse og joint offer", 0.8, 0.7, 11.0, 0.7, size=28, color=WHITE, bold=True)

# Maxbo box
rect(s, 0.8, 1.55, 5.6, 2.5, NAVY)
rect(s, 0.8, 1.55, 5.6, 0.06, BLUE)
tb(s, "LIVE NORSK REFERANSE", 1.0, 1.65, 5.0, 0.35, size=10, color=BLUE, bold=True)
tb(s, "Maxbo", 1.0, 2.0, 5.0, 0.55, size=24, color=WHITE, bold=True)
txb = s.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(5.0), Inches(1.3))
tf = txb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "·  Jakala composable frontend på Shopify Plus"
r.font.name = FONT; r.font.size = Pt(11); r.font.color.rgb = WHITE
for line_text in [
    "·  Enterspeed som content delivery layer",
    "·  Speedtrain onboarding · aktiv leveranse",
    "·  Norges eneste live referanse i kombinasjonen",
]:
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = line_text
    r2.font.name = FONT; r2.font.size = Pt(11); r2.font.color.rgb = WHITE

# Joint offer box
rect(s, 6.8, 1.55, 5.7, 2.5, RGBColor(0x05, 0x05, 0x35))
rect(s, 6.8, 1.55, 5.7, 0.06, BLUE)
tb(s, "JOINT TILBUD", 7.0, 1.65, 5.0, 0.35, size=10, color=BLUE, bold=True)
tb(s, "Headless Commerce\nAccelerator", 7.0, 2.0, 5.2, 0.8, size=20, color=WHITE, bold=True)
txb2 = s.shapes.add_textbox(Inches(7.0), Inches(2.85), Inches(5.2), Inches(1.1))
tf2 = txb2.text_frame; tf2.word_wrap = True
p0 = tf2.paragraphs[0]; p0.alignment = PP_ALIGN.LEFT
r0 = p0.add_run(); r0.text = "Jakala frontend   ×   Shopify Plus   ×   Enterspeed"
r0.font.name = FONT; r0.font.size = Pt(10); r0.font.color.rgb = BLUE; r0.font.bold = True
for lt in [
    "→  Raskere time-to-market enn custom builds",
    "→  Lavere TCO enn enterprise platform",
    "→  Skalerbart til 50+ markeder",
]:
    pp = tf2.add_paragraph(); pp.alignment = PP_ALIGN.LEFT
    rr = pp.add_run(); rr.text = lt
    rr.font.name = FONT; rr.font.size = Pt(11); rr.font.color.rgb = WHITE

# Capability bar
caps = ["Jakala Premier Agency", "Shopify Plus", "Enterspeed CDP", "30+ sertifiseringer", "Maxbo referanse"]
cap_colors = [BLUE, GREEN, AMBER, BLUE, RED]
total_w = 11.5
cap_w = total_w / len(caps)
for i, (cap, col) in enumerate(zip(caps, cap_colors)):
    rect(s, 0.8 + i*cap_w, 4.3, cap_w-0.05, 0.55, col)
    tb(s, cap, 0.85 + i*cap_w, 4.37, cap_w-0.1, 0.4, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Why combine
rect(s, 0.8, 5.05, 11.5, 1.8, RGBColor(0x10, 0x10, 0x28))
tb(s, "Hvorfor dette kombinasjonstilbudet er unikt i det norske markedet", 1.0, 5.15, 11.0, 0.4, size=12, color=BLUE, bold=True)
whys = [
    ("Bevist",     "Maxbo er live på denne nøyaktige stacken i Norge — ikke et konsept"),
    ("Fart",       "Enterspeed reduserer content-latency med 40-60 % · Shopify eliminerer infrastrukturoppbygging"),
    ("Utvidbarhet","Fra én nettbutikk til omnichannel platform uten å bytte teknologi"),
]
for i, (head, desc) in enumerate(whys):
    xpos = 1.0 + i * 3.83
    tb(s, head, xpos, 5.55, 3.5, 0.35, size=11, color=AMBER, bold=True)
    tb(s, desc, xpos, 5.9, 3.6, 0.7, size=10, color=GREY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TOPP 10 OVERSIKT
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(); bg(s)
rect(s, 0, 0, 0.5, 7.5, BLUE)
line(s, 0.8, 1.3, 11.5)
tb(s, "PIPELINE OVERSIKT", 0.8, 0.4, 6.0, 0.5, size=11, color=BLUE, bold=True)
tb(s, "Topp 10 prioriterte norske accounts", 0.8, 0.7, 11.0, 0.7, size=28, color=WHITE, bold=True)

# Table header
rect(s, 0.8, 1.45, 11.5, 0.42, BLUE)
headers = ["#", "Account", "Sektor", "Tilbud", "ICP", "Deal", "Shopify", "Enterspeed", "Kontakt"]
col_x   = [0.85, 1.3, 3.2, 4.8, 6.55, 7.0, 7.55, 8.55, 9.7]
col_w   = [0.45, 1.9, 1.6, 1.7, 0.45, 0.45, 1.0, 1.0, 2.5]
for hdr, cx, cw in zip(headers, col_x, col_w):
    tb(s, hdr, cx, 1.49, cw, 0.34, size=9, color=WHITE, bold=True)

rows = [
    ("1",  "Elkjøp Nordic",   "Elektronikk", "Commerce Opt.",  "8", "9", "◦", "●●", "Morten Syversen · CDO"),
    ("2",  "Trumf",           "Lojalitet",   "Data Revenue",   "9", "9", "◦", "●",  "Rikke Etholm-Idsøe · Comm.Dir."),
    ("3",  "Sport Outlet",    "Sport",       "Shopify + ES",   "8", "9", "●●","●●", "CEO Tor-André Skeie"),
    ("4",  "Helly Hansen",    "Outdoor",     "Shopify + ES",   "8", "8", "●●","●●", "Joumana Lovstad · CMO"),
    ("5",  "Vinmonopolet",    "Drikke",      "Data Revenue",   "8", "8", "◦", "●",  "Espen Terland · CDO"),
    ("6",  "Skeidar",         "Møbler",      "Commerce Opt.",  "8", "8", "◦", "●●", "Sujit Nath · CIO"),
    ("7",  "Norrøna",         "Outdoor",     "Shopify + ES",   "8", "7", "●●","●",  "CMO · TBD"),
    ("8",  "GANT Norway",     "Fashion",     "Shopify + ES",   "7", "7", "●●","●",  "Fredrik Malm · CEO"),
    ("9",  "Jernia",          "Jernvare",    "Commerce Opt.",  "7", "8", "◦", "●●", "Ina Kristin Haugen · CMO"),
    ("10", "Strai Kjøkken",   "Kjøkken",     "Shopify + ES",   "7", "6", "●●","●",  "Monica Lohnås Roland · DL"),
]

row_colors = [RGBColor(0x10, 0x10, 0x28), RGBColor(0x0D, 0x0D, 0x22)]
shopify_cols = {"●●": GREEN, "●": AMBER, "◦": GREY}

for i, row in enumerate(rows):
    y = 1.87 + i * 0.50
    rect(s, 0.8, y, 11.5, 0.48, row_colors[i % 2])
    for j, (val, cx, cw) in enumerate(zip(row, col_x, col_w)):
        if j in (6, 7):  # Shopify / Enterspeed columns
            col = shopify_cols.get(val, WHITE)
            tb(s, val, cx, y+0.07, cw, 0.35, size=11, color=col, bold=(val=="●●"))
        elif j == 4 or j == 5:  # ICP / Deal score
            score = int(val)
            col = GREEN if score >= 9 else (AMBER if score >= 7 else WHITE)
            tb(s, val, cx, y+0.07, cw, 0.35, size=11, color=col, bold=True)
        elif j == 0:
            tb(s, val, cx, y+0.07, cw, 0.35, size=10, color=GREY, bold=False)
        else:
            tb(s, val, cx, y+0.07, cw, 0.35, size=10, color=WHITE)

# Legend
tb(s, "●● = Primær    ● = Sekundær    ◦ = Ikke aktuelt   |   ICP/Deal: score av 10", 0.8, 7.15, 11.5, 0.3, size=9, color=GREY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SHOPIFY-FOKUS (5 accounts)
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(); bg(s)
rect(s, 0, 0, 0.5, 7.5, GREEN)
line(s, 0.8, 1.3, 11.5, GREEN)
tb(s, "SHOPIFY PLUS", 0.8, 0.4, 5.0, 0.5, size=11, color=GREEN, bold=True)
tb(s, "5 accounts — primær Shopify-mulighet", 0.8, 0.7, 11.0, 0.7, size=28, color=WHITE, bold=True)

shopify_accounts = [
    {
        "name": "Sport Outlet",
        "sub":  "18 nettbutikker · Sport & outdoor",
        "why":  "CTO og CDO begge vakante. CEO Tor-André Skeie tar alle beslutninger. Maximalt timing-vindu.",
        "tech":  "Ukjent plattform · 18 sider · høy fragmentering",
        "entry": "Konsolidering av 18 butikker på Shopify Plus + Enterspeed",
        "buyer": "CEO Tor-André Skeie",
        "value": "€250K",
        "urgency": RED,
        "urgency_label": "KRITISK — KONTAKT NÅ",
    },
    {
        "name": "Helly Hansen",
        "sub":  "55 Adobe Commerce-sider · Outdoor",
        "why":  "Kontoor Brands-oppkjøp juni 2025. 55 Adobe-sider = teknisk gjeld. Ny agenda åpen.",
        "tech":  "55 × Adobe Commerce · 65.000 SKU pr. side",
        "entry": "Pilotmigrering av 3-5 markeder til Shopify Plus + Enterspeed",
        "buyer": "Joumana Lovstad (CMO)",
        "value": "€400K",
        "urgency": AMBER,
        "urgency_label": "AKTIVER DENNE UKEN",
    },
    {
        "name": "Norrøna",
        "sub":  "Premium outdoor · Direkte Shopify-kandidat",
        "why":  "JAKALA identifisert som primær Shopify-partner for Norrøna. DTC-ambisjon tydelig.",
        "tech":  "Ukjent — sannsynlig SaaS/Shopify-klar stack",
        "entry": "Shopify Plus DTC + Enterspeed for innholdshastighet",
        "buyer": "CMO (TBD — LinkedIn-kartlegging nødvendig)",
        "value": "€200K",
        "urgency": GREEN,
        "urgency_label": "PIPELINE — AKTIVER",
    },
    {
        "name": "GANT Norway",
        "sub":  "Fashion · IMPACT Commerce ny partner feb 2026",
        "why":  "Ny CEO Fredrik Malm. IMPACT Commerce inn som ny partner. Vindu åpent før det lukkes.",
        "tech":  "Salesforce Commerce Cloud + Dynamic Yield",
        "entry": "Shopify Plus-migrering eller parallel DTC-løp ved siden av SFCC",
        "buyer": "Fredrik Malm (ny CEO)",
        "value": "€180K",
        "urgency": AMBER,
        "urgency_label": "AKTIVER INNEN 2 UKER",
    },
    {
        "name": "Strai Kjøkken",
        "sub":  "Kjøkken · Odoo ERP · Rekordvekst",
        "why":  "Odoo nylig implementert. Bekreftet investeringsvilje. Ny kjøkkenkjede med digital ambisjon.",
        "tech":  "Odoo ERP · Enkel nettbutikk i dag",
        "entry": "Shopify Plus headless frontend + Enterspeed innholdslag",
        "buyer": "Monica Lohnås Roland (daglig leder)",
        "value": "€120K",
        "urgency": GREEN,
        "urgency_label": "PIPELINE",
    },
]

positions = [(0.8, 1.45), (4.65, 1.45), (8.5, 1.45), (0.8, 4.3), (4.65, 4.3)]
card_w, card_h = 3.65, 2.75

for acc, (px, py) in zip(shopify_accounts, positions):
    rect(s, px, py, card_w, card_h, RGBColor(0x0A, 0x12, 0x0A))
    rect(s, px, py, card_w, 0.05, acc["urgency"])
    # Urgency badge
    rect(s, px+card_w-2.1, py+0.08, 2.0, 0.28, acc["urgency"])
    tb(s, acc["urgency_label"], px+card_w-2.08, py+0.1, 1.95, 0.24, size=7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, acc["name"], px+0.12, py+0.1, card_w-2.3, 0.38, size=14, color=WHITE, bold=True)
    tb(s, acc["sub"], px+0.12, py+0.5, card_w-0.2, 0.3, size=9, color=GREY, italic=True)
    tb(s, acc["why"], px+0.12, py+0.82, card_w-0.2, 0.62, size=9, color=WHITE)
    tb(s, "Entry: " + acc["entry"], px+0.12, py+1.48, card_w-0.2, 0.42, size=9, color=BLUE)
    tb(s, f"Kjøper: {acc['buyer']}   |   Estimert verdi: {acc['value']}", px+0.12, py+1.95, card_w-0.2, 0.35, size=9, color=AMBER, bold=True)

# 5th card right column
rect(s, 8.5, 4.3, card_w, card_h, RGBColor(0x12, 0x0A, 0x04))
acc = shopify_accounts[4]; px, py = 8.5, 4.3
rect(s, px, py, card_w, 0.05, acc["urgency"])
rect(s, px+card_w-2.1, py+0.08, 2.0, 0.28, acc["urgency"])
tb(s, acc["urgency_label"], px+card_w-2.08, py+0.1, 1.95, 0.24, size=7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tb(s, acc["name"], px+0.12, py+0.1, card_w-2.3, 0.38, size=14, color=WHITE, bold=True)
tb(s, acc["sub"], px+0.12, py+0.5, card_w-0.2, 0.3, size=9, color=GREY, italic=True)
tb(s, acc["why"], px+0.12, py+0.82, card_w-0.2, 0.62, size=9, color=WHITE)
tb(s, "Entry: " + acc["entry"], px+0.12, py+1.48, card_w-0.2, 0.42, size=9, color=BLUE)
tb(s, f"Kjøper: {acc['buyer']}   |   Estimert verdi: {acc['value']}", px+0.12, py+1.95, card_w-0.2, 0.35, size=9, color=AMBER, bold=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ENTERSPEED + COMMERCE OPTIMIZATION (Elkjøp, Skeidar, Jernia)
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(); bg(s)
rect(s, 0, 0, 0.5, 7.5, AMBER)
line(s, 0.8, 1.3, 11.5, AMBER)
tb(s, "ENTERSPEED + COMMERCE OPTIMIZATION", 0.8, 0.4, 10.0, 0.5, size=11, color=AMBER, bold=True)
tb(s, "3 accounts — Enterprise Commerce-transformasjon", 0.8, 0.7, 11.0, 0.7, size=28, color=WHITE, bold=True)

commerce_accounts = [
    {
        "name": "Elkjøp Nordic",
        "tag": "#1 PRIORITET",
        "tag_col": GREEN,
        "sub": "Nordens største elektronikkjede",
        "value": "€700K",
        "icp": "8/10", "deal": "9/10",
        "situation": "Morten Syversen (Chief Brand & Digital Officer) identifisert som navngitt kjøper. Elkjøp er #3 i Nordic-pipeline. Commerce-transformasjon pågår.",
        "tech": "Hybris/SAP Commerce · DXP i transisjon",
        "entry": "Enterspeed som headless delivery layer for kampanjeinnhold — 40% fartsforbedring",
        "expansion": "Full Headless Commerce Accelerator → AI-drevet personalisering",
        "buyer": "Morten Syversen · Chief Brand & Digital Officer",
        "actions": ["LinkedIn-kontakt fra norsk partner", "Demo basert på Maxbo-referansen", "Enterspeed fartsstudie"],
    },
    {
        "name": "Skeidar",
        "tag": "ERKLÆRT AMBISJON",
        "tag_col": AMBER,
        "sub": "Norges ledende møbelkjede",
        "value": "€250K",
        "icp": "8/10", "deal": "8/10",
        "situation": "CIO Sujit Nath identifisert. Skeidar har erklært ambisjonen om 'Norges beste møbelportal'. Microsoft Dynamics 365 Commerce + SpectrumOne CDP allerede live.",
        "tech": "D365 Commerce (BE-terna) · SpectrumOne CDP",
        "entry": "Enterspeed som innholdslag på toppen av D365 — ingen stack-bytte nødvendig",
        "expansion": "Personalisert produktopplevelse + AI-drevet anbefalingsmotor",
        "buyer": "Sujit Nath (CIO) + Martin Andresen (CEO)",
        "actions": ["Kontakt Sujit Nath direkte", "Presentér Enterspeed + D365-integrasjon", "Referer Maxbo"],
    },
    {
        "name": "Jernia",
        "tag": "PIM NETTOPP LIVE",
        "tag_col": BLUE,
        "sub": "Norges ledende jernvare- og verktøykjede",
        "value": "€200K",
        "icp": "7/10", "deal": "8/10",
        "situation": "SAP Commerce Cloud + Bluestone PIM nylig live. PIM er i plass — neste fase er klar. CEO Espen Karlsen og CMO Ina Kristin Haugen begge identifisert.",
        "tech": "SAP Commerce Cloud · Bluestone PIM · Infor CloudSuite ERP",
        "entry": "Enterspeed delivery layer for fart + commerce personalisering på PIM-data",
        "expansion": "AI-produktanbefalinger basert på nytt PIM-fundament",
        "buyer": "Ina Kristin Haugen (CMO) — primær kontakt",
        "actions": ["Kontakt CMO om PIM + neste fase", "Enterspeed demo på SAP Commerce stack", "Vis ROI fra Maxbo"],
    },
]

for i, acc in enumerate(commerce_accounts):
    x0 = 0.8 + i * 4.1
    y0 = 1.55
    bw = 3.95; bh = 5.6

    rect(s, x0, y0, bw, bh, RGBColor(0x10, 0x10, 0x26))
    rect(s, x0, y0, bw, 0.06, acc["tag_col"])

    # Tag badge
    rect(s, x0+bw-2.3, y0+0.1, 2.2, 0.3, acc["tag_col"])
    tb(s, acc["tag"], x0+bw-2.28, y0+0.12, 2.15, 0.26, size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    tb(s, acc["name"], x0+0.15, y0+0.1, bw-2.5, 0.42, size=16, color=WHITE, bold=True)
    tb(s, acc["sub"], x0+0.15, y0+0.55, bw-0.2, 0.28, size=10, color=GREY, italic=True)

    # Scores
    tb(s, f"ICP {acc['icp']}  ·  Deal {acc['deal']}  ·  {acc['value']}", x0+0.15, y0+0.87, bw-0.2, 0.3, size=10, color=AMBER, bold=True)

    line(s, x0+0.1, y0+1.22, bw-0.2, GREY)

    tb(s, "Situasjon", x0+0.15, y0+1.32, bw-0.2, 0.28, size=9, color=BLUE, bold=True)
    tb(s, acc["situation"], x0+0.15, y0+1.6, bw-0.2, 0.9, size=9, color=WHITE)

    tb(s, "Tech stack", x0+0.15, y0+2.55, bw-0.2, 0.28, size=9, color=BLUE, bold=True)
    tb(s, acc["tech"], x0+0.15, y0+2.82, bw-0.2, 0.35, size=9, color=GREY)

    tb(s, "Entry-tilbud", x0+0.15, y0+3.22, bw-0.2, 0.28, size=9, color=BLUE, bold=True)
    tb(s, acc["entry"], x0+0.15, y0+3.5, bw-0.2, 0.42, size=9, color=WHITE)

    tb(s, "Kjøper", x0+0.15, y0+3.97, bw-0.2, 0.28, size=9, color=BLUE, bold=True)
    tb(s, acc["buyer"], x0+0.15, y0+4.25, bw-0.2, 0.3, size=10, color=WHITE, bold=True)

    tb(s, "Neste steg: " + acc["actions"][0], x0+0.15, y0+4.65, bw-0.2, 0.35, size=9, color=AMBER)
    tb(s, acc["actions"][1], x0+0.15, y0+4.98, bw-0.2, 0.35, size=9, color=AMBER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DATA REVENUE (Trumf, Vinmonopolet)
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(); bg(s)
rect(s, 0, 0, 0.5, 7.5, BLUE)
line(s, 0.8, 1.3, 11.5)
tb(s, "DATA REVENUE UNLOCK", 0.8, 0.4, 7.0, 0.5, size=11, color=BLUE, bold=True)
tb(s, "2 accounts — norsk data- og lojalitetsmulighet", 0.8, 0.7, 11.0, 0.7, size=28, color=WHITE, bold=True)

data_accounts = [
    {
        "name": "Trumf",
        "value": "€450K", "icp": "9/10", "deal": "9/10",
        "sub": "Norges største lojalitetsprogram · NorgesGruppen",
        "tag": "13 DAGER TIMING-VINDU",
        "tag_col": RED,
        "situation": "Rikke Etholm-Idsøe tiltrådte som ny Commercial Director i en nyopprettet stilling. Første 90 dager = ubeskrevet blad. Kristin Bigseth (MD) allerede identifisert.",
        "data_asset": "13 millioner norske lojalitetsmedlemmer · Unike handledata på tvers av 1.800+ butikker",
        "problem": "Lojalitetsdata monetisert sub-optimalt. Ingen ekstern data-revenue-modell. GDPR-compliant data clean room = åpenbar mulighet.",
        "entry": "Data Revenue Diagnostic — 3-dagers analyse: hva er dataene verdt for leverandører og mediahuset?",
        "expansion": "Retail Media-plattform → Data Clean Room → Ekstern datamonetisering",
        "buyers": ["Rikke Etholm-Idsøe (ny Commercial Director)", "Kristin Bigseth (MD)"],
        "actions": [
            "1. Send outreach til Rikke Etholm-Idsøe DENNE UKEN — ref. ny rolle",
            "2. Bruk Trumf-størrelse som oppsett: '13M nordmenn — hva er det verdt?'",
            "3. Tilby gratis Retail Media-benchmark som inngangsportalen",
        ],
    },
    {
        "name": "Vinmonopolet",
        "value": "€350K", "icp": "8/10", "deal": "8/10",
        "sub": "Norsk statlig drikkevareforhandler · 340 butikker",
        "tag": "NY CDO — AKTIVER NÅ",
        "tag_col": AMBER,
        "situation": "Espen Terland tiltrådte som ny CDO (ex-XXL 15 år). Agenda ikke satt. Lars Thorenfeldt også identifisert. Klassisk honeymoon-vindu for ny CDO.",
        "data_asset": "Unike kjøpsdata på alle alkoholkjøp i Norge · Demografisk bredde · Ingen privat konkurrent",
        "problem": "Data brukt internt men ikke optimalisert for personalisering, anbefalinger eller ekstern verdi. Ny CDO vil bevise seg raskt.",
        "entry": "Data Revenue Diagnostic — kartlegg personaliserings- og anbefalingsmuligheter for en CDO som vil levere i sitt første år",
        "expansion": "AI-drevet produktanbefaling → Personalisert opplevelse → Ekstern datainnsikt for leverandører",
        "buyers": ["Espen Terland (CDO — ny)", "Lars Thorenfeldt (identifisert)"],
        "actions": [
            "1. Kontakt Espen Terland med CDO-onboarding-vinkel: 'Hjelp oss forstå verdien i dataene'",
            "2. Posisjonér som CDO-enabler ikke vendor",
            "3. Referer til XXL-bakgrunnen hans — Jakala jobbet i retail-segmentet",
        ],
    },
]

for i, acc in enumerate(data_accounts):
    x0 = 0.8 + i * 6.1
    y0 = 1.55
    bw = 5.85; bh = 5.6

    rect(s, x0, y0, bw, bh, RGBColor(0x0A, 0x0A, 0x20))
    rect(s, x0, y0, bw, 0.06, acc["tag_col"])

    rect(s, x0+bw-3.0, y0+0.1, 2.85, 0.32, acc["tag_col"])
    tb(s, acc["tag"], x0+bw-2.98, y0+0.13, 2.8, 0.26, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    tb(s, acc["name"], x0+0.15, y0+0.12, bw-3.2, 0.45, size=22, color=WHITE, bold=True)
    tb(s, acc["sub"], x0+0.15, y0+0.62, bw-0.2, 0.3, size=10, color=GREY, italic=True)
    tb(s, f"ICP {acc['icp']}  ·  Deal {acc['deal']}  ·  Estimert verdi: {acc['value']}", x0+0.15, y0+0.98, bw-0.2, 0.3, size=11, color=AMBER, bold=True)
    line(s, x0+0.1, y0+1.35, bw-0.25, GREY)

    tb(s, "Situasjon", x0+0.15, y0+1.45, bw-0.2, 0.28, size=9, color=BLUE, bold=True)
    tb(s, acc["situation"], x0+0.15, y0+1.73, bw-0.2, 0.72, size=10, color=WHITE)

    tb(s, "Data-eiendel", x0+0.15, y0+2.5, bw-0.2, 0.28, size=9, color=BLUE, bold=True)
    tb(s, acc["data_asset"], x0+0.15, y0+2.78, bw-0.2, 0.45, size=10, color=WHITE)

    tb(s, "Forretningsproblem", x0+0.15, y0+3.28, bw-0.2, 0.28, size=9, color=BLUE, bold=True)
    tb(s, acc["problem"], x0+0.15, y0+3.56, bw-0.2, 0.5, size=10, color=WHITE)

    tb(s, "Entry-tilbud", x0+0.15, y0+4.12, bw-0.2, 0.28, size=9, color=BLUE, bold=True)
    tb(s, acc["entry"], x0+0.15, y0+4.4, bw-0.2, 0.4, size=10, color=BLUE)

    tb(s, "Steg denne uken:", x0+0.15, y0+4.88, bw-0.2, 0.28, size=9, color=RED, bold=True)
    tb(s, acc["actions"][0], x0+0.15, y0+5.16, bw-0.2, 0.3, size=9, color=AMBER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TIMING-KART (alle 10 accounts)
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(); bg(s)
rect(s, 0, 0, 0.5, 7.5, RED)
line(s, 0.8, 1.3, 11.5, RED)
tb(s, "TIMING-SIGNALER", 0.8, 0.4, 7.0, 0.5, size=11, color=RED, bold=True)
tb(s, "Nå-vinduet for alle 10 accounts", 0.8, 0.7, 11.0, 0.7, size=28, color=WHITE, bold=True)

timing_rows = [
    ("Sport Outlet",    "CTO + CDO vakante",               "Kontakt CEO NÅ",                    RED,   "KRITISK"),
    ("Trumf",           "Ny Commercial Director — dag 13",  "Send outreach i dag",               RED,   "KRITISK"),
    ("Vinmonopolet",    "Ny CDO — honeymoon-fase",          "CDO-onboarding outreach",           AMBER, "DENNE UKEN"),
    ("Helly Hansen",    "Post-Kontoor oppkjøp — ny agenda", "Pilot-migrering 3 markeder",        AMBER, "DENNE UKEN"),
    ("Elkjøp Nordic",   "Commerce-transformasjon pågår",    "Demo basert på Maxbo-referansen",   AMBER, "DENNE UKEN"),
    ("Skeidar",         "Erklært 'beste møbelportal'",      "Kontakt CIO Sujit Nath",            GREEN, "DENNE MND."),
    ("GANT Norway",     "IMPACT Commerce ny partner",       "Kontakt CEO Fredrik Malm",          GREEN, "DENNE MND."),
    ("Norrøna",         "DTC-ambisjon tydelig",             "Identifisér CMO-kontakt",           GREEN, "DENNE MND."),
    ("Jernia",          "PIM nylig live — klar for fase 2", "CMO-outreach om neste fase",        GREEN, "DENNE MND."),
    ("Strai Kjøkken",   "Odoo live + rekordvekst",          "Shopify Plus pilot-samtale",        GREEN, "Q2 2026"),
]

# Header
rect(s, 0.8, 1.48, 11.5, 0.38, RGBColor(0x15, 0x00, 0x00))
for lbl, cx in [("Account", 0.9), ("Signal", 3.3), ("Anbefalt handling", 6.5), ("Urgency", 10.8)]:
    tb(s, lbl, cx, 1.52, 3.0, 0.3, size=9, color=WHITE, bold=True)

for i, (name, signal, action, col, urgency) in enumerate(timing_rows):
    y = 1.88 + i * 0.5
    bg_c = RGBColor(0x18, 0x05, 0x05) if col == RED else (RGBColor(0x18, 0x12, 0x03) if col == AMBER else RGBColor(0x05, 0x14, 0x08))
    rect(s, 0.8, y, 11.5, 0.47, bg_c)
    rect(s, 0.8, y, 0.06, 0.47, col)
    tb(s, name, 0.9, y+0.07, 2.3, 0.34, size=11, color=WHITE, bold=True)
    tb(s, signal, 3.3, y+0.07, 3.1, 0.34, size=10, color=GREY)
    tb(s, action, 6.5, y+0.07, 4.1, 0.34, size=10, color=col if col != GREEN else WHITE)
    # Badge
    rect(s, 10.7, y+0.07, 1.5, 0.3, col)
    tb(s, urgency, 10.7, y+0.09, 1.5, 0.26, size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — NESTE STEG
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(); bg(s)
rect(s, 0, 0, 0.5, 7.5, BLUE)
line(s, 0.8, 1.3, 11.5)
tb(s, "HANDLINGSPLAN", 0.8, 0.4, 7.0, 0.5, size=11, color=BLUE, bold=True)
tb(s, "Neste steg — mars 2026", 0.8, 0.7, 11.0, 0.7, size=30, color=WHITE, bold=True)

# This week
rect(s, 0.8, 1.55, 11.5, 0.38, RED)
tb(s, "DENNE UKEN — Kritiske handlinger", 1.0, 1.61, 10.0, 0.3, size=12, color=WHITE, bold=True)

this_week = [
    ("1", "Sport Outlet",   "Kontakt CEO Tor-André Skeie — CTO/CDO vakant-vindu",               "Jacob Skaue"),
    ("2", "Trumf",          "Send outreach til Rikke Etholm-Idsøe — dag 13 av ny rolle",        "Jacob Skaue"),
    ("3", "Vinmonopolet",   "CDO onboarding-outreach til Espen Terland",                         "Jacob Skaue"),
    ("4", "Helly Hansen",   "LinkedIn-kontakt Joumana Lovstad — post-Kontoor agenda",            "Jacob Skaue"),
]

for i, (num, acc, action, owner) in enumerate(this_week):
    y = 1.97 + i * 0.52
    bg_r = RGBColor(0x18, 0x05, 0x05)
    rect(s, 0.8, y, 11.5, 0.49, bg_r)
    rect(s, 0.8, y, 0.35, 0.49, RED)
    tb(s, num, 0.82, y+0.1, 0.3, 0.3, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, acc, 1.22, y+0.06, 2.3, 0.34, size=11, color=WHITE, bold=True)
    tb(s, action, 3.6, y+0.06, 6.7, 0.34, size=11, color=GREY)
    tb(s, owner, 10.4, y+0.06, 1.8, 0.34, size=10, color=BLUE, bold=True)

# This month
rect(s, 0.8, 4.15, 11.5, 0.38, NAVY)
tb(s, "DENNE MÅNEDEN — Viktige handlinger", 1.0, 4.21, 10.0, 0.3, size=12, color=WHITE, bold=True)

this_month = [
    ("5", "Elkjøp Nordic",  "Demo: Enterspeed på Maxbo-referansen til Morten Syversen",          "Jacob Skaue"),
    ("6", "Skeidar",        "Kontakt CIO Sujit Nath — Enterspeed + D365 pitch",                  "Jacob Skaue"),
    ("7", "GANT Norway",    "Ny CEO Fredrik Malm — Shopify Plus vindu",                          "Jacob Skaue"),
    ("8", "Norrøna",        "Identifisér CMO-kontakt — LinkedIn kartlegging",                    "Jacob Skaue"),
]

for i, (num, acc, action, owner) in enumerate(this_month):
    y = 4.57 + i * 0.5
    bg_r = RGBColor(0x0A, 0x0A, 0x20)
    rect(s, 0.8, y, 11.5, 0.47, bg_r)
    rect(s, 0.8, y, 0.35, 0.47, BLUE)
    tb(s, num, 0.82, y+0.09, 0.3, 0.3, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, acc, 1.22, y+0.07, 2.3, 0.32, size=11, color=WHITE, bold=True)
    tb(s, action, 3.6, y+0.07, 6.7, 0.32, size=11, color=GREY)
    tb(s, owner, 10.4, y+0.07, 1.8, 0.32, size=10, color=BLUE, bold=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — SISTE SIDE / KONTAKT
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(); bg(s)
rect(s, 0, 0, 0.5, 7.5, BLUE)
rect(s, 0.5, 0, 12.83, 7.5, NAVY)

tb(s, "JAKALA Nordic", 1.2, 2.0, 10.0, 0.8, size=16, color=BLUE, bold=True)
tb(s, "La oss snakke om\nditt neste steg.", 1.2, 2.7, 10.0, 2.0, size=46, color=WHITE, bold=True)
tb(s, "jacob.skaue@jakala.com", 1.2, 5.2, 8.0, 0.55, size=18, color=GREY)
tb(s, "jakala.com/nordic", 1.2, 5.75, 8.0, 0.45, size=14, color=GREY, italic=True)

rect(s, 1.2, 6.4, 11.13, 0.05, BLUE)
tb(s, "Data-drevet vekst · Shopify Plus · Enterspeed · Headless Commerce · AI Readiness", 1.2, 6.5, 11.0, 0.4, size=10, color=GREY, italic=True)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "presentations/norge-topp10-leads-2026-03.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
