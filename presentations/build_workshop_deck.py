from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Constants ──────────────────────────────────────────────────────────────────
BLUE       = RGBColor(0x15, 0x3E, 0xED)
BLUE_DARK  = RGBColor(0x02, 0x02, 0x66)
RED        = RGBColor(0xF6, 0x57, 0x4A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
MUTED      = RGBColor(0x88, 0x88, 0xAA)
BG         = RGBColor(0x08, 0x08, 0x18)
CARD       = RGBColor(0x12, 0x12, 0x28)
CARD2      = RGBColor(0x0D, 0x0D, 0x22)
GREEN      = RGBColor(0x00, 0xD4, 0xA0)
FONT       = 'Raleway'
W          = Inches(13.33)
H          = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)

def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, text, x, y, w, h, size=14, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = FONT
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def tb_multi(slide, lines, x, y, w, h, size=13, color=WHITE,
             align=PP_ALIGN.LEFT, line_spacing=1.1):
    """lines = list of (text, bold, color, size_override)"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, col, sz = item, False, color, size
        elif len(item) == 2:
            text, bold = item; col = color; sz = size
        elif len(item) == 3:
            text, bold, col = item; sz = size
        else:
            text, bold, col, sz = item
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.name  = FONT
        run.font.size  = Pt(sz)
        run.font.bold  = bold
        run.font.color.rgb = col
    return txb

def rect(slide, x, y, w, h, color=BLUE_DARK, alpha=None):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def line_h(slide, x, y, w, color=BLUE, thick=2):
    from pptx.util import Pt as PPt
    ln = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.02))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    return ln

def dot(slide, x, y, r=0.08, color=BLUE):
    from pptx.util import Inches as I
    shape = slide.shapes.add_shape(
        9, I(x - r), I(y - r), I(r * 2), I(r * 2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs); bg(s)
rect(s, 0, 0, 13.33, 7.5, RGBColor(0x08, 0x08, 0x18))

# Left accent bar
rect(s, 0, 0, 0.06, 7.5, BLUE)

# Top label
tb(s, "JAKALA NORDIC  ·  INTERNAL WORKSHOP  ·  MARCH 2026",
   0.35, 0.3, 9, 0.4, size=9, color=MUTED, bold=False)

# Main headline
tb(s, "Norway:", 0.35, 1.4, 10, 1.2, size=60, bold=True, color=WHITE)
tb(s, "Building a Commercial Engine", 0.35, 2.55, 10, 1.1, size=38, bold=False, color=BLUE)

line_h(s, 0.35, 3.85, 5.5, BLUE)

# Sub
tb(s, "Market · Pipeline · GTM Approach · Acceleration Services",
   0.35, 4.05, 9, 0.5, size=15, color=MUTED)

# Author
tb(s, "Jacob Skaue  ·  Commercial Director, Norway",
   0.35, 5.2, 7, 0.45, size=13, color=WHITE, bold=True)

# Right side — big number
tb(s, "€3.5M", 8.5, 1.6, 4.5, 1.5, size=72, bold=True, color=BLUE,
   align=PP_ALIGN.RIGHT)
tb(s, "Norway Pipeline", 8.5, 3.1, 4.5, 0.4, size=14, color=MUTED,
   align=PP_ALIGN.RIGHT)
tb(s, "32 accounts  ·  12 named buyers", 8.5, 3.55, 4.5, 0.4, size=12, color=MUTED,
   align=PP_ALIGN.RIGHT)

# Bottom bar
rect(s, 0, 6.9, 13.33, 0.6, BLUE_DARK)
tb(s, "Confidential  ·  JAKALA Nordic  ·  2026", 0.3, 7.0, 8, 0.4,
   size=9, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE OPPORTUNITY: WHY NORWAY, WHY NOW
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs); bg(s)
rect(s, 0, 0, 0.06, 7.5, BLUE)
tb(s, "01  ·  THE OPPORTUNITY", 0.35, 0.3, 8, 0.35, size=9, color=BLUE, bold=True)
tb(s, "Why Norway. Why Now.", 0.35, 0.7, 10, 0.8, size=36, bold=True, color=WHITE)
line_h(s, 0.35, 1.6, 12.5, BLUE)

# 3 columns
cols = [
    ("MARKET SIZE", "Norway is JAKALA Nordic's largest pipeline market — €3.5M across 32 accounts. Retail, commerce and data-driven loyalty dominate. High digital maturity, strong buying power."),
    ("TIMING SIGNALS", "Multiple honeymoon windows open now: new CDOs at Vinmonopolet and Sport Outlet. New Commercial Director at Trumf. Elkjøp B2B commerce live. Every signal identified and mapped."),
    ("OUR POSITION", "We mapped 32 accounts with ICP scores, named buyers and entry offers before sending a single message. No other Nordic market has this depth of pre-sales intelligence."),
]
for i, (title, body) in enumerate(cols):
    cx = 0.35 + i * 4.35
    rect(s, cx, 1.8, 4.1, 4.5, CARD)
    line_h(s, cx, 1.8, 4.1, BLUE, thick=3)
    tb(s, title, cx + 0.2, 2.0, 3.8, 0.4, size=10, bold=True, color=BLUE)
    tb(s, body, cx + 0.2, 2.55, 3.7, 3.5, size=12.5, color=WHITE)

# Bottom stat strip
rect(s, 0.35, 6.45, 12.6, 0.75, BLUE_DARK)
stats = [("38%", "of Nordic total pipeline"), ("12", "named buyers confirmed"), ("4", "GTM entry strategies mapped"), ("1", "active delivery: Maxbo")]
sw = 3.15
for i, (val, lbl) in enumerate(stats):
    sx = 0.55 + i * sw
    tb(s, val, sx, 6.52, 2, 0.35, size=18, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
    tb(s, lbl, sx + 0.7, 6.6, 2.3, 0.3, size=9, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — GTM APPROACH
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs); bg(s)
rect(s, 0, 0, 0.06, 7.5, BLUE)
tb(s, "02  ·  GTM APPROACH", 0.35, 0.3, 8, 0.35, size=9, color=BLUE, bold=True)
tb(s, "Signal → Named Buyer → Entry Offer → Accelerate", 0.35, 0.7, 12, 0.8,
   size=30, bold=True, color=WHITE)
line_h(s, 0.35, 1.6, 12.5, BLUE)

# 4 stage pipeline
stages = [
    ("01", "SIGNAL", BLUE, "Identify timing triggers:\n· New leadership (CDO/CIO/CMO)\n· Platform migration\n· Public AI statement\n· Org restructure\n· Budget signal"),
    ("02", "NAMED BUYER", BLUE, "No outreach without a name.\nConfirm the exact person:\n· Role & decision authority\n· Budget ownership\n· Entry to conversation"),
    ("03", "ENTRY OFFER", BLUE, "Match to one of 4 GTM strategies.\nLead with a diagnostic:\n· Low-risk for buyer\n· €30–100K entry\n· Immediate commercial value"),
    ("04", "ACCELERATE", GREEN, "Layer in Acceleration Services:\n· Speedtrain (data foundation)\n· AI Readiness Diagnostic\n· Commerce Optimization Pilot\n→ Expands to full program"),
]
sw = 3.1
for i, (num, title, col, body) in enumerate(stages):
    sx = 0.35 + i * sw
    rect(s, sx, 1.8, 2.85, 4.6, CARD)
    line_h(s, sx, 1.8, 2.85, col, thick=4)
    tb(s, num, sx + 0.18, 2.0, 0.6, 0.5, size=28, bold=True, color=col)
    tb(s, title, sx + 0.18, 2.6, 2.5, 0.4, size=11, bold=True, color=WHITE)
    tb(s, body, sx + 0.18, 3.1, 2.6, 3.1, size=11, color=MUTED)

# Arrow connectors (simple right arrows)
for i in range(3):
    ax = 0.35 + (i + 1) * sw - 0.22
    tb(s, "→", ax, 3.7, 0.4, 0.4, size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Bottom insight
rect(s, 0.35, 6.55, 12.6, 0.65, BLUE_DARK)
tb(s, "The goal is not to sell transformation first. The goal is a repeatable entry motion that creates the conditions for transformation.",
   0.6, 6.62, 12, 0.45, size=11, color=WHITE, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW WE LAYER ACCELERATION SERVICES
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs); bg(s)
rect(s, 0, 0, 0.06, 7.5, BLUE)
tb(s, "03  ·  ACCELERATION SERVICES", 0.35, 0.3, 8, 0.35, size=9, color=BLUE, bold=True)
tb(s, "How We Layer Services Into Every Proposition", 0.35, 0.7, 12, 0.8,
   size=30, bold=True, color=WHITE)
line_h(s, 0.35, 1.6, 12.5, BLUE)

# Left: the model
rect(s, 0.35, 1.75, 5.8, 4.7, CARD)
tb(s, "THE LAYERING MODEL", 0.55, 1.9, 5.4, 0.4, size=10, bold=True, color=BLUE)

layers = [
    ("ENTRY", "Diagnostic / Pilot  ·  €30–100K", BLUE, "Low-risk opening. Quantified output. Immediate value for buyer."),
    ("ACCELERATION", "Speedtrain / AI Readiness  ·  €100–300K", GREEN, "Build the foundation. Data, architecture, product information layer."),
    ("OPTIMIZATION", "Commerce Optimization  ·  €200–500K", BLUE, "Revenue impact. Search, discovery, merchandising performance."),
    ("TRANSFORMATION", "DXP Program  ·  €500K+", RGBColor(0xF6, 0x57, 0x4A), "Full platform. Multi-market. Long-term partnership."),
]
for i, (stage, value, col, desc) in enumerate(layers):
    ly = 2.4 + i * 0.95
    rect(s, 0.5, ly, 5.5, 0.8, RGBColor(0x10, 0x10, 0x25))
    line_h(s, 0.5, ly, 0.25, col, thick=8)
    tb(s, stage, 0.85, ly + 0.05, 2, 0.3, size=9, bold=True, color=col)
    tb(s, value, 0.85, ly + 0.38, 3, 0.25, size=9, color=MUTED)
    tb(s, desc, 3.1, ly + 0.08, 2.8, 0.65, size=9, color=WHITE)

# Right: Norway examples
rect(s, 6.5, 1.75, 6.5, 4.7, CARD)
tb(s, "NORWAY EXAMPLES", 6.7, 1.9, 6.0, 0.4, size=10, bold=True, color=BLUE)

examples = [
    ("Maxbo", "ACTIVE DELIVERY", GREEN, "Entry: Data Revenue Diagnostic → Speedtrain\nActive: Product data foundation, 1M+ SKUs\nNext: Commerce optimization layer"),
    ("Elkjøp", "READY TO ACTIVATE", BLUE, "Entry: Commerce Optimization Pilot\nSignal: B2B commerce live — new buyer segment\nBuyer: Morten Syversen (Chief Brand & Digital)"),
    ("Trumf", "READY TO ACTIVATE", BLUE, "Entry: Data Revenue Diagnostic\nSignal: Retail media 'not big enough' — stated gap\nBuyer: Bigseth + Etholm-Idsøe (new CDO)"),
    ("Varner Group", "PIPELINE", MUTED, "Entry: Data Revenue Diagnostic\nSignal: 7 brands, no shared PIM — clear wedge\nBuyer: Research ongoing"),
]
for i, (name, status, col, desc) in enumerate(examples):
    ey = 2.4 + i * 0.95
    rect(s, 6.6, ey, 6.2, 0.8, RGBColor(0x10, 0x10, 0x25))
    tb(s, name, 6.8, ey + 0.05, 2.5, 0.3, size=11, bold=True, color=WHITE)
    tb(s, status, 9.5, ey + 0.08, 3, 0.25, size=8, bold=True, color=col, align=PP_ALIGN.RIGHT)
    tb(s, desc, 6.8, ey + 0.38, 6.0, 0.5, size=9, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — NORWAY PIPELINE
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs); bg(s)
rect(s, 0, 0, 0.06, 7.5, BLUE)
tb(s, "04  ·  NORWAY PIPELINE", 0.35, 0.3, 8, 0.35, size=9, color=BLUE, bold=True)
tb(s, "€3.5M · 32 Accounts · 12 Named Buyers", 0.35, 0.7, 12, 0.8,
   size=32, bold=True, color=WHITE)
line_h(s, 0.35, 1.6, 12.5, BLUE)

# Pipeline table — top 7
headers = ["Account", "GTM Strategy", "Entry Offer", "Named Buyer", "Win %", "Value"]
col_x   = [0.35, 2.6, 5.0, 7.3, 10.2, 11.5]
col_w   = [2.2, 2.3, 2.2, 2.8, 1.2, 1.7]

# Header row
rect(s, 0.35, 1.7, 12.6, 0.42, BLUE_DARK)
for j, h in enumerate(headers):
    tb(s, h, col_x[j] + 0.08, 1.76, col_w[j], 0.3, size=9, bold=True, color=BLUE)

rows = [
    ("Maxbo",         "Data Revenue",         "Speedtrain (live)",          "Active delivery",            "✓",  "€539M rev", GREEN),
    ("Elkjøp Nordic", "Commerce Optim.",       "Commerce Pilot",             "Morten Syversen ✓",          "65%", "€700K",   GREEN),
    ("Trumf",         "Data Revenue",          "Data Revenue Diagnostic",    "Bigseth + Etholm-Idsøe ✓",   "40%", "€450K",   BLUE),
    ("Varner Group",  "Data Revenue",          "Data Revenue Diagnostic",    "Research ongoing",            "~25%","€500K+",  MUTED),
    ("Vinmonopolet",  "Data Revenue",          "Data Revenue Diagnostic",    "Espen Terland ✓ (new CDO)",   "35%", "€200K",   BLUE),
    ("Skeidar",       "Commerce Optim.",       "Commerce Pilot",             "Sujit Nath ✓ (CIO)",          "40%", "€200K",   BLUE),
    ("Bulder Bank",   "AI Readiness",          "AI Readiness Diagnostic",    "Simen Eilertsen ✓",           "30%", "€200K",   MUTED),
]
for i, row in enumerate(rows):
    ry = 2.14 + i * 0.63
    rc = CARD if i % 2 == 0 else RGBColor(0x0F, 0x0F, 0x22)
    rect(s, 0.35, ry, 12.6, 0.61, rc)
    row_color = row[6]
    for j, val in enumerate(row[:6]):
        c = row_color if j == 4 else WHITE
        if j == 4 and val == "✓":
            c = GREEN
        tb(s, str(val), col_x[j] + 0.08, ry + 0.1, col_w[j], 0.4,
           size=10.5, color=c, bold=(j == 0))

# Bottom note
tb(s, "* 32 accounts total mapped. Table shows top 7 by strategic priority. All entries have ICP scores, named buyer research and entry offer defined.",
   0.35, 7.1, 12.5, 0.3, size=9, color=MUTED, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CASE: MAXBO (PROOF POINT)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs); bg(s)
rect(s, 0, 0, 0.06, 7.5, GREEN)
tb(s, "05  ·  CASE STUDY", 0.35, 0.3, 8, 0.35, size=9, color=GREEN, bold=True)
tb(s, "Maxbo — Speedtrain in Action", 0.35, 0.7, 10, 0.8, size=34, bold=True, color=WHITE)
tb(s, "Norway's largest home improvement retailer  ·  Active delivery  ·  Acceleration Services live",
   0.35, 1.45, 12, 0.4, size=13, color=MUTED)
line_h(s, 0.35, 1.9, 12.5, GREEN)

# Left column
rect(s, 0.35, 2.05, 4.1, 4.8, CARD)
tb(s, "THE SITUATION", 0.55, 2.2, 3.8, 0.35, size=10, bold=True, color=GREEN)
tb(s, "1,000,000+ products across a 4-layer data pipeline: Perfion → Azure → Pimcore → Magento.\n\nProduct data inconsistency at this scale silently kills search relevance, discovery and conversion.\n\nNo AI or personalization layer. Digital ambitions ahead of data infrastructure.",
   0.55, 2.65, 3.7, 3.8, size=11.5, color=WHITE)

# Middle column
rect(s, 4.65, 2.05, 4.1, 4.8, CARD)
tb(s, "OUR APPROACH", 4.85, 2.2, 3.8, 0.35, size=10, bold=True, color=GREEN)
tb_multi(s, [
    ("Entry wedge:", True, BLUE, 11),
    ("Data Revenue Diagnostic — quantify where broken product data costs revenue.", False, WHITE, 11),
    (" ", False, MUTED, 6),
    ("Acceleration layer:", True, BLUE, 11),
    ("Speedtrain onboarding — build the product data foundation at scale.", False, WHITE, 11),
    (" ", False, MUTED, 6),
    ("Framing used:", True, BLUE, 11),
    ('"With 1M products across a multi-system pipeline, data inconsistency is silently costing you revenue in search and conversion."', False, MUTED, 10),
], 4.85, 2.65, 3.7, 4.2)

# Right column
rect(s, 8.95, 2.05, 4.05, 4.8, CARD)
tb(s, "OUTCOME & EXPANSION", 9.1, 2.2, 3.7, 0.35, size=10, bold=True, color=GREEN)
tb_multi(s, [
    ("Status:", True, GREEN, 11),
    ("Active delivery — Speedtrain onboarding in progress", False, WHITE, 11),
    (" ", False, MUTED, 5),
    ("Proven:", True, BLUE, 11),
    ("Diagnostic-led entry works. Commercial framing (revenue, not IT) opened the door.", False, WHITE, 11),
    (" ", False, MUTED, 5),
    ("Expansion path:", True, BLUE, 11),
    ("Commerce Optimization → AI Readiness → DXP transformation (Magento succession)", False, WHITE, 11),
    (" ", False, MUTED, 5),
    ("Key learning:", True, RED, 11),
    ("Entry at diagnostic scope, not full program. Buyer risk was low. Acceleration followed naturally.", False, MUTED, 10),
], 9.1, 2.65, 3.7, 4.2)

# Bottom bar — the lesson
rect(s, 0.35, 6.95, 12.6, 0.4, BLUE_DARK)
tb(s, "LESSON APPLIED ACROSS ALL NORWAY ACCOUNTS: Lead with revenue problem → confirm the buyer → deliver the diagnostic → accelerate.",
   0.6, 7.0, 12, 0.3, size=10.5, color=WHITE, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TOP 3 NEXT ACTIVATIONS
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs); bg(s)
rect(s, 0, 0, 0.06, 7.5, BLUE)
tb(s, "06  ·  NEXT ACTIVATIONS", 0.35, 0.3, 8, 0.35, size=9, color=BLUE, bold=True)
tb(s, "Three Deals Ready to Move — Now", 0.35, 0.7, 10, 0.8,
   size=34, bold=True, color=WHITE)
line_h(s, 0.35, 1.6, 12.5, BLUE)

deals = [
    {
        "rank": "01",
        "company": "Elkjøp Nordic",
        "label": "HIGHEST VALUE",
        "color": GREEN,
        "signal": "B2B commerce expanded across all Nordic markets. New buyer segment demands better product discovery.",
        "buyer": "Morten Syversen\nChief Brand & Digital Officer",
        "offer": "Commerce Optimization Pilot",
        "value": "€700K → €455K weighted",
        "win": "65%",
        "next": "Send LinkedIn outreach this week. Message drafted.",
    },
    {
        "rank": "02",
        "company": "Trumf (NorgesGruppen)",
        "label": "STRONGEST SIGNAL",
        "color": BLUE,
        "signal": "Retail media effort internally acknowledged as 'not big enough.' New Commercial Director role created — Rikke Etholm-Idsøe in first 90 days.",
        "buyer": "Kristin Bigseth (MD)\n+ Rikke Etholm-Idsøe (Comm. Dir.)",
        "offer": "Data Revenue Diagnostic",
        "value": "€450K → €180K weighted",
        "win": "40%",
        "next": "Rikke's 90-day window = now. Outreach drafted. Send this week.",
    },
    {
        "rank": "03",
        "company": "Vinmonopolet",
        "label": "HONEYMOON WINDOW",
        "color": RGBColor(0xF6, 0x57, 0x4A),
        "signal": "Espen Terland appointed new CDO (ex-XXL, 15 years). First 90 days — agenda not set. Window is open.",
        "buyer": "Espen Terland\nChief Digital Officer (new)",
        "offer": "Data Revenue Diagnostic",
        "value": "€200K → €70K weighted",
        "win": "35%",
        "next": "CDO honeymoon window. Frame as: 'What does your data estate look like today?'",
    },
]

cw = 4.1
for i, d in enumerate(deals):
    cx = 0.35 + i * 4.35
    col = d["color"]
    rect(s, cx, 1.75, cw, 5.3, CARD)
    line_h(s, cx, 1.75, cw, col, thick=4)
    tb(s, d["rank"], cx + 0.18, 1.9, 0.8, 0.5, size=26, bold=True, color=col)
    tb(s, d["label"], cx + 0.18, 2.42, 3.7, 0.3, size=8, bold=True, color=col)
    tb(s, d["company"], cx + 0.18, 2.78, 3.7, 0.45, size=15, bold=True, color=WHITE)
    tb(s, "SIGNAL", cx + 0.18, 3.32, 3.7, 0.25, size=8, bold=True, color=MUTED)
    tb(s, d["signal"], cx + 0.18, 3.6, 3.7, 0.85, size=10.5, color=WHITE)
    tb(s, "BUYER", cx + 0.18, 4.52, 3.7, 0.25, size=8, bold=True, color=MUTED)
    tb(s, d["buyer"], cx + 0.18, 4.78, 3.7, 0.5, size=10.5, color=WHITE)
    tb(s, "ENTRY  ·  " + d["offer"], cx + 0.18, 5.38, 3.7, 0.3, size=9, color=BLUE, bold=True)
    rect(s, cx, 6.45, cw, 0.6, RGBColor(0x0A, 0x0A, 0x20))
    tb(s, "→ " + d["next"], cx + 0.15, 6.52, 3.85, 0.45, size=9.5, color=col, bold=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CCO LENS: WHAT THIS SHOWS
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs); bg(s)
rect(s, 0, 0, 0.06, 7.5, BLUE)
tb(s, "07  ·  COMMERCIAL LEADERSHIP", 0.35, 0.3, 8, 0.35, size=9, color=BLUE, bold=True)
tb(s, "A Scalable Commercial Model — Not Just a Pipeline", 0.35, 0.7, 12, 0.8,
   size=30, bold=True, color=WHITE)
line_h(s, 0.35, 1.6, 12.5, BLUE)

# 2 x 2 cards
cards = [
    (0.35, 1.75, "MARKET INTELLIGENCE",
     "Built 32-account Norway market map from zero. Each account scored on ICP fit, deal strength, timing signal and named buyer. No guesswork — every account has an entry thesis.",
     BLUE),
    (6.85, 1.75, "COMMERCIAL ARCHITECTURE",
     "Designed a 4-strategy GTM model (Data Revenue · Commerce Optimization · AI Readiness · Experience Transformation) applicable across all Nordic markets. Replicable. Scalable.",
     BLUE),
    (0.35, 4.55, "ACCELERATION INTEGRATION",
     "Every Norway proposition is structured to layer in Acceleration Services at the entry point. Diagnostic leads to Speedtrain. Pilot leads to optimization. No standalone deals.",
     GREEN),
    (6.85, 4.55, "PIPELINE DISCIPLINE",
     "Named buyer rule enforced: no outreach without a confirmed decision-maker. Win probability capped at 25% without a name. Forecast is probability-weighted, not optimistic.",
     GREEN),
]
for (cx, cy, title, body, col) in cards:
    rect(s, cx, cy, 6.2, 2.5, CARD)
    line_h(s, cx, cy, 6.2, col, thick=3)
    tb(s, title, cx + 0.22, cy + 0.18, 5.7, 0.35, size=10, bold=True, color=col)
    tb(s, body, cx + 0.22, cy + 0.62, 5.7, 1.7, size=12, color=WHITE)

# Bottom
rect(s, 0.35, 7.05, 12.6, 0.35, BLUE_DARK)
tb(s, "The Norway build is a proof of concept for a Nordic commercial operating model. The same system applies to Denmark and Sweden.",
   0.6, 7.1, 12, 0.25, size=10, color=WHITE, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CLOSING / NEXT STEPS
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs); bg(s)
rect(s, 0, 0, 0.06, 7.5, BLUE)
rect(s, 0, 0, 13.33, 7.5, RGBColor(0x08, 0x08, 0x18))
rect(s, 0, 0, 0.06, 7.5, BLUE)

tb(s, "08  ·  NEXT STEPS", 0.35, 0.3, 8, 0.35, size=9, color=BLUE, bold=True)
tb(s, "Three Actions. This Week.", 0.35, 0.7, 10, 0.8, size=36, bold=True, color=WHITE)
line_h(s, 0.35, 1.6, 8, BLUE)

actions = [
    ("01", "Send Elkjøp outreach", "Morten Syversen. Commerce Optimization Pilot. Message is drafted. Highest weighted deal in Norway at €455K.", BLUE),
    ("02", "Send Trumf outreach", "Rikke Etholm-Idsøe. New Commercial Director. First 90-day window = now. Data Revenue Diagnostic as entry.", GREEN),
    ("03", "Send Vinmonopolet outreach", "Espen Terland. New CDO. Honeymoon window open. Frame as a data estate conversation, not a sales pitch.", BLUE),
]

for i, (num, title, body, col) in enumerate(actions):
    ay = 1.9 + i * 1.55
    rect(s, 0.35, ay, 8.5, 1.35, CARD)
    line_h(s, 0.35, ay, 8.5, col, thick=3)
    tb(s, num, 0.55, ay + 0.12, 0.7, 0.55, size=24, bold=True, color=col)
    tb(s, title, 1.35, ay + 0.14, 7.2, 0.4, size=14, bold=True, color=WHITE)
    tb(s, body, 1.35, ay + 0.62, 7.1, 0.6, size=11, color=MUTED)

# Right side summary
rect(s, 9.2, 1.75, 3.8, 5.1, CARD)
tb(s, "NORWAY SCORECARD", 9.4, 1.9, 3.5, 0.35, size=9, bold=True, color=BLUE)
summary_rows = [
    ("Pipeline", "€3.5M"),
    ("Accounts", "32"),
    ("Named buyers", "12"),
    ("Active delivery", "Maxbo"),
    ("Outreach ready", "3 (drafted)"),
    ("Meetings booked", "0"),
    ("Status", "READY TO ACTIVATE"),
]
for i, (lbl, val) in enumerate(summary_rows):
    ry = 2.38 + i * 0.62
    is_last = i == len(summary_rows) - 1
    col = RED if is_last else WHITE
    tb(s, lbl, 9.4, ry, 2, 0.38, size=10, color=MUTED)
    tb(s, val, 11.1, ry, 1.7, 0.38, size=10, bold=True, color=col, align=PP_ALIGN.RIGHT)
    if not is_last:
        line_h(s, 9.4, ry + 0.48, 3.3, RGBColor(0x20, 0x20, 0x35))

# ── Save ───────────────────────────────────────────────────────────────────────
out_dir = os.path.join(os.path.dirname(__file__))
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "jacob-norway-workshop-2026-03.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
