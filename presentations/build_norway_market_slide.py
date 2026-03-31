from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BLUE      = RGBColor(0x15, 0x3E, 0xED)
BLUE_DARK = RGBColor(0x02, 0x02, 0x66)
RED       = RGBColor(0xF6, 0x57, 0x4A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x88, 0x88, 0xAA)
BG        = RGBColor(0x08, 0x08, 0x18)
CARD      = RGBColor(0x12, 0x12, 0x28)
GREEN     = RGBColor(0x00, 0xD4, 0xA0)
FONT      = 'Raleway'

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def tb(slide, text, x, y, w, h, size=12, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.name   = FONT
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return txb

def rect(slide, x, y, w, h, color=CARD):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def line_h(slide, x, y, w, color=BLUE, thick=0.025):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(thick))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

# ══════════════════════════════════════════════════════════════════════════════
# SINGLE SLIDE — NORWAY MARKET + ACTIVATION SERVICES
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs); bg(s)

# Left accent bar
rect(s, 0, 0, 0.055, 7.5, BLUE)

# ── HEADER ──────────────────────────────────────────────────────────────────
tb(s, "NORWAY", 0.28, 0.22, 5, 0.55, size=32, bold=True, color=WHITE)
tb(s, "Market Opportunity & Acceleration Services",
   0.28, 0.78, 9, 0.35, size=13, color=MUTED)
line_h(s, 0.28, 1.22, 12.8, BLUE)

# ── LEFT COLUMN: MARKET SNAPSHOT ────────────────────────────────────────────
rect(s, 0.28, 1.35, 3.75, 5.85, CARD)
line_h(s, 0.28, 1.35, 3.75, BLUE, thick=0.04)
tb(s, "MARKET SNAPSHOT", 0.48, 1.5, 3.4, 0.3, size=9, bold=True, color=BLUE)

market_rows = [
    ("Pipeline",        "€3.5M",          WHITE),
    ("Accounts mapped", "32",              WHITE),
    ("Named buyers",    "12 confirmed",    WHITE),
    ("Active delivery", "Maxbo (live)",    GREEN),
    ("Top sectors",     "Retail · Fashion · Grocery · Finance · DIY", MUTED),
]
for i, (lbl, val, col) in enumerate(market_rows):
    ry = 1.95 + i * 0.72
    rect(s, 0.38, ry, 3.55, 0.62, RGBColor(0x0E, 0x0E, 0x22))
    tb(s, lbl, 0.55, ry + 0.06, 1.5, 0.25, size=9, color=MUTED)
    tb(s, val,  0.55, ry + 0.3, 3.2, 0.28, size=11.5, bold=True, color=col)

# Timing note
rect(s, 0.38, 5.63, 3.55, 0.85, RGBColor(0x0A, 0x14, 0x30))
line_h(s, 0.38, 5.63, 3.55, GREEN, thick=0.03)
tb(s, "OPEN WINDOWS NOW",
   0.55, 5.7, 3.2, 0.25, size=8, bold=True, color=GREEN)
tb(s, "New CDO: Vinmonopolet · Sport Outlet\nNew Comm. Dir.: Trumf (first 90 days)\nB2B live: Elkjøp",
   0.55, 5.98, 3.2, 0.46, size=9, color=WHITE)

# ── RIGHT COLUMNS: 4 ACTIVATION SERVICES ───────────────────────────────────
services = [
    {
        "title":   "DATA REVENUE\nDIAGNOSTIC + SPEEDTRAIN",
        "color":   BLUE,
        "value":   "€50–100K entry · €200–700K expansion",
        "why":     "Norway has 14+ accounts with large product catalogs, fragmented PIM/ERP setups and loyalty data sitting unused. Speedtrain is the acceleration layer.",
        "accounts": "Varner · Trumf · Vinmonopolet · Maxbo · GANT · NAF · Kitch'n",
        "signal":  "PIM live but not activated · loyalty data unused · retail media gap",
    },
    {
        "title":   "COMMERCE\nOPTIMIZATION PILOT",
        "color":   BLUE,
        "value":   "€40–80K entry · ongoing retainer",
        "why":     "Norwegian retail is investing heavily in logistics (DC automation, B2B expansion) but digital product experience lags. Measurable uplift in 8 weeks.",
        "accounts": "Elkjøp · Skeidar · Sport Outlet · Bohus · Strai Kjøkken",
        "signal":  "B2B commerce live · DC investment · low online share vs. store",
    },
    {
        "title":   "AI READINESS\nDIAGNOSTIC",
        "color":   RGBColor(0xF6, 0x57, 0x4A),
        "value":   "€50–100K entry · €200–500K program",
        "why":     "Post-merger banks and large retailers have AI mandates without the data infrastructure to execute. Architecture gap is the entry wedge.",
        "accounts": "Bulder Bank · Sparebanken Norge · Equinor · DNB · NorgesGruppen",
        "signal":  "AI partnership announced · post-merger · Chief AI Officer hired",
    },
    {
        "title":   "SHOPIFY\nCOMMERCE BUILD",
        "color":   GREEN,
        "value":   "€15–30K check-up · €80–500K build",
        "why":     "Several Norwegian retail brands are on ageing Magento 2 or WooCommerce platforms with replatform signals. Shopify check-up is the low-friction entry.",
        "accounts": "Helly Hansen · Ferner Jacobsen · Follestad · Strai Kjøkken",
        "signal":  "Magento 2 / WooCommerce · D2C ambition · multi-market expansion",
    },
]

cw = 2.28
for i, svc in enumerate(services):
    cx = 4.22 + i * (cw + 0.17)
    col = svc["color"]

    rect(s, cx, 1.35, cw, 5.85, CARD)
    line_h(s, cx, 1.35, cw, col, thick=0.045)

    # Title
    tb(s, svc["title"],
       cx + 0.15, 1.45, cw - 0.2, 0.65,
       size=9.5, bold=True, color=WHITE)

    # Value band
    rect(s, cx, 2.13, cw, 0.38, RGBColor(0x0A, 0x0A, 0x20))
    tb(s, svc["value"],
       cx + 0.12, 2.19, cw - 0.18, 0.28,
       size=8.5, bold=True, color=col)

    # Why label
    tb(s, "WHY NORWAY",
       cx + 0.12, 2.62, cw - 0.15, 0.22,
       size=7.5, bold=True, color=MUTED)

    # Why body
    tb(s, svc["why"],
       cx + 0.12, 2.86, cw - 0.15, 1.1,
       size=9.5, color=WHITE)

    # Accounts label
    tb(s, "ACCOUNTS",
       cx + 0.12, 4.02, cw - 0.15, 0.22,
       size=7.5, bold=True, color=MUTED)

    # Accounts
    rect(s, cx + 0.08, 4.26, cw - 0.16, 0.72, RGBColor(0x0E, 0x0E, 0x22))
    tb(s, svc["accounts"],
       cx + 0.18, 4.3, cw - 0.28, 0.62,
       size=9, color=col, bold=True)

    # Signal label
    tb(s, "KEY SIGNAL",
       cx + 0.12, 5.06, cw - 0.15, 0.22,
       size=7.5, bold=True, color=MUTED)

    # Signal
    tb(s, svc["signal"],
       cx + 0.12, 5.3, cw - 0.15, 0.72,
       size=9, color=WHITE, italic=True)

# ── BOTTOM BAR ──────────────────────────────────────────────────────────────
rect(s, 0, 7.12, 13.33, 0.38, BLUE_DARK)
tb(s, "Entry model: Diagnostic / Pilot → Acceleration Services → Full program. Never lead with transformation.",
   0.28, 7.18, 10, 0.28, size=9.5, color=WHITE, italic=True)
tb(s, "JAKALA Nordic  ·  2026",
   11.0, 7.18, 2.2, 0.28, size=9, color=MUTED, align=PP_ALIGN.RIGHT)

# ── SAVE ────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "norway-market-activation-services.pptx")
prs.save(out)
print(f"Saved: {out}")
