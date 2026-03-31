"""
JAKALA × Kingfisher Group — AI Readiness Pitch Deck
Type: pitch
Language: English
Audience: Client (Dr. Simon Jury, CDO / Tim Ellison, Group Digital CoE Director)
Date: March 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Design Constants ────────────────────────────────────────────────────────
BLUE_BRIGHT = RGBColor(0x15, 0x3E, 0xED)
BLUE_DARK   = RGBColor(0x02, 0x02, 0x66)
RED         = RGBColor(0xF6, 0x57, 0x4A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY        = RGBColor(0xAA, 0xAA, 0xBB)
BG_COLOR    = RGBColor(0x08, 0x08, 0x18)
FONT        = 'Raleway'
W           = Inches(13.33)
H           = Inches(7.5)

# ─── Helpers ─────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color=BG_COLOR):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, x, y, w, h,
                size=16, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_rect(slide, x, y, w, h, color=BLUE_DARK):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_line(slide, x, y, w, color=BLUE_BRIGHT):
    line = slide.shapes.add_shape(1, x, y, w, Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def add_body(slide, lines, x, y, w, h, size=11, color=WHITE):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, clr) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = clr

def list_box(slide, items, x, y, w, h, size=11, bullet="·"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = WHITE

# ─── Slides ──────────────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    # JAKALA brand
    add_textbox(slide, "JAKALA", Inches(0.25), Inches(0.3), Inches(3), Inches(0.5),
                size=13, bold=True, color=BLUE_BRIGHT)
    add_textbox(slide, "Data & AI Practice", Inches(0.25), Inches(0.62),
                Inches(3), Inches(0.4), size=11, color=GREY)

    # Main headline
    add_textbox(slide, "Kingfisher Group",
                Inches(0.25), Inches(1.6), Inches(9), Inches(0.7),
                size=20, bold=False, color=GREY)
    add_textbox(slide, "AI at Scale Needs\na Data Foundation\nThat Can Keep Up.",
                Inches(0.25), Inches(2.1), Inches(9), Inches(2.8),
                size=42, bold=True, color=WHITE)

    add_line(slide, Inches(0.25), Inches(4.85), Inches(6), BLUE_BRIGHT)

    add_textbox(slide,
                "How JAKALA helps Kingfisher protect and extend the £100M+\nin AI-driven revenue your data layer is now carrying.",
                Inches(0.25), Inches(5.0), Inches(9), Inches(0.9),
                size=14, color=GREY)

    # Right panel
    add_rect(slide, Inches(9.6), Inches(0), Inches(3.73), H, BLUE_DARK)
    add_rect(slide, Inches(9.6), Inches(1.8), Inches(3.6), Inches(0.06), BLUE_BRIGHT)

    stats = [
        ("£100M+", "AI-driven incremental\nrevenue (2025/26)"),
        ("+62%", "Marketplace GMV YoY\nacross 5 markets"),
        ("30+", "AI initiatives running\nin parallel"),
    ]
    for i, (val, lbl) in enumerate(stats):
        y = Inches(2.0) + i * Inches(1.7)
        add_textbox(slide, val, Inches(9.8), y,
                    Inches(3.2), Inches(0.65),
                    size=28, bold=True, color=WHITE)
        add_textbox(slide, lbl, Inches(9.8), y + Inches(0.65),
                    Inches(3.2), Inches(0.6),
                    size=11, color=GREY)

    add_textbox(slide, "Prepared for: Dr. Simon Jury, CDO · Tim Ellison, Group Digital CoE",
                Inches(0.25), H - Inches(1.05), Inches(9), Inches(0.4),
                size=11, color=GREY)
    add_textbox(slide, "March 2026  ·  Confidential",
                Inches(0.25), H - Inches(0.68), Inches(4), Inches(0.35),
                size=10, color=GREY)


def slide_signal(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "01 · Market Signal — Why Now", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.45), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "You have built something rare.\nNow the data layer has to carry it.",
                Inches(0.25), Inches(0.72), Inches(10), Inches(1.3),
                size=30, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(2.0), Inches(7), BLUE_BRIGHT)

    signals = [
        {
            "icon": "£100M+",
            "title": "AI-driven revenue confirmed",
            "body": "Spirit and Perception are delivering measurable commercial results.\nAt this scale, data quality issues compound — they don't stay small.",
            "color": BLUE_BRIGHT,
        },
        {
            "icon": "+62%",
            "title": "Marketplace GMV YoY",
            "body": "Five markets. Thousands of third-party SKUs added annually.\nCatalog data complexity is growing faster than governance can handle.",
            "color": BLUE_BRIGHT,
        },
        {
            "icon": "Athena",
            "title": "AI orchestration layer launched Jan 2026",
            "body": "B&Q marketplace conversion doubled post-Athena.\nAthena's output quality depends entirely on the data quality below it.",
            "color": RED,
        },
        {
            "icon": "28",
            "title": "AI engineers building in parallel",
            "body": "Multiple platforms, multiple brands, multiple markets.\nIntegration complexity grows non-linearly. External architecture review is the pressure valve.",
            "color": BLUE_BRIGHT,
        },
    ]

    for i, s in enumerate(signals):
        col = i % 2
        row = i // 2
        x = Inches(0.25) + col * Inches(6.55)
        y = Inches(2.25) + row * Inches(2.3)
        add_rect(slide, x, y, Inches(6.3), Inches(2.1), BLUE_DARK)
        add_rect(slide, x, y, Inches(6.3), Inches(0.07), s["color"])
        add_textbox(slide, s["icon"], x + Inches(0.18), y + Inches(0.15),
                    Inches(1.5), Inches(0.48), size=18, bold=True, color=s["color"])
        add_textbox(slide, s["title"], x + Inches(0.18), y + Inches(0.6),
                    Inches(5.9), Inches(0.4), size=13, bold=True, color=WHITE)
        add_textbox(slide, s["body"], x + Inches(0.18), y + Inches(1.02),
                    Inches(5.9), Inches(0.9), size=11, color=GREY)


def slide_company(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "02 · Company Snapshot", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.45), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "Kingfisher Group",
                Inches(0.25), Inches(0.72), Inches(11), Inches(0.65),
                size=30, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(1.38), Inches(5), BLUE_BRIGHT)

    # Left: facts
    add_rect(slide, Inches(0.25), Inches(1.6), Inches(5.8), Inches(5.55), BLUE_DARK)
    add_textbox(slide, "At a Glance", Inches(0.45), Inches(1.72),
                Inches(5.4), Inches(0.38), size=12, bold=True, color=BLUE_BRIGHT)

    facts = [
        ("HQ", "London, UK — listed on LSE (KGF)"),
        ("Brands", "B&Q · Screwfix · Castorama · Brico Dépôt · NeedHelp"),
        ("Footprint", "1,900+ stores across Europe · 82,000+ employees"),
        ("Revenue", "~£12.3B (2024/25)"),
        ("Ecommerce", "19% of total sales · 93% fulfilled in-store (click & collect)"),
        ("Marketplace", "B&Q marketplace live in UK, France, Poland, Iberia, Turkey"),
        ("AI Team", "28 ML engineers, data scientists, and AI engineers in-house"),
        ("AI Revenue", "£100M+ incremental revenue from AI (Spirit + Perception)"),
    ]
    fbox = slide.shapes.add_textbox(Inches(0.45), Inches(2.15), Inches(5.4), Inches(4.7))
    ftf = fbox.text_frame
    ftf.word_wrap = True
    for i, (label, val) in enumerate(facts):
        p = ftf.paragraphs[0] if i == 0 else ftf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{label}:  "
        run.font.name = FONT
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = GREY
        run2 = p.add_run()
        run2.text = val
        run2.font.name = FONT
        run2.font.size = Pt(11)
        run2.font.color.rgb = WHITE

    # Right: AI platforms
    add_rect(slide, Inches(6.3), Inches(1.6), Inches(6.78), Inches(5.55),
             RGBColor(0x05, 0x05, 0x30))
    add_textbox(slide, "Proprietary AI Platform — 6 Products",
                Inches(6.5), Inches(1.72), Inches(6.4), Inches(0.38),
                size=12, bold=True, color=BLUE_BRIGHT)

    platforms = [
        ("Spirit", "Group", "Personalisation engine — product recommendations"),
        ("Perception", "Group", "Demand forecasting and supply chain AI"),
        ("Inspect", "Group", "Content moderation — keeps product data current"),
        ("Athena", "Group", "AI orchestration layer — launched Jan 2026"),
        ("Screwfix Lens", "Screwfix UK", "Visual search — 100Ks of uses since launch"),
        ("Hello B&Q / Hello Casto", "B&Q · Castorama", "GenAI DIY assistant · 10% conversion uplift (FR)"),
    ]
    for i, (name, brand, desc) in enumerate(platforms):
        y = Inches(2.18) + i * Inches(0.82)
        add_rect(slide, Inches(6.3), y, Inches(6.78), Inches(0.75),
                 RGBColor(0x0A, 0x0A, 0x22) if i % 2 == 0 else RGBColor(0x0D, 0x0D, 0x28))
        add_textbox(slide, name, Inches(6.45), y + Inches(0.07),
                    Inches(1.8), Inches(0.32), size=11, bold=True, color=WHITE)
        add_textbox(slide, brand, Inches(8.3), y + Inches(0.07),
                    Inches(1.5), Inches(0.32), size=10, color=BLUE_BRIGHT)
        add_textbox(slide, desc, Inches(6.45), y + Inches(0.38),
                    Inches(6.4), Inches(0.3), size=10, color=GREY)


def slide_problem(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "03 · The Problem", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.45), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "The faster you build AI, the more\nyour data layer becomes the constraint.",
                Inches(0.25), Inches(0.72), Inches(10.5), Inches(1.5),
                size=30, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(2.2), Inches(8), BLUE_BRIGHT)

    # Central tension
    add_rect(slide, Inches(0.25), Inches(2.4), Inches(12.83), Inches(0.9),
             RGBColor(0x12, 0x02, 0x04))
    add_rect(slide, Inches(0.25), Inches(2.4), Inches(0.06), Inches(0.9), RED)
    add_textbox(slide,
                "Kingfisher is running 30+ AI initiatives across 6 platforms, 5 brands, and 5 markets simultaneously. "
                "The AI outputs are only as good as the data going in. "
                "At this scale and velocity, data quality governance is the single biggest risk to sustaining the £100M+ in AI-driven revenue.",
                Inches(0.45), Inches(2.5), Inches(12.2), Inches(0.72),
                size=12, color=WHITE)

    # Three problem dimensions
    problems = [
        {
            "title": "Multi-brand catalog complexity",
            "detail": [
                "B&Q, Screwfix, Castorama, and Brico Dépôt each have independent product catalogs.",
                "Marketplace GMV +62% YoY = thousands of new third-party SKUs added across 5 markets.",
                "Inconsistent taxonomy, duplicate SKUs, and attribution gaps compound as each market scales.",
                "Spirit's recommendations are only as accurate as the product data feeding them.",
            ],
        },
        {
            "title": "AI architecture growing faster than governance",
            "detail": [
                "30+ AI initiatives run in parallel. Each adds a new dependency on data quality.",
                "Athena (launched Jan 2026) orchestrates multiple AI models — errors propagate across all of them.",
                "Internal team of 28 is building features, not governance. That's appropriate — but creates a gap.",
                "External architecture review is a complement, not a challenge, to your internal capability.",
            ],
        },
        {
            "title": "Platform complexity between TCS and AI layer",
            "detail": [
                "TCS covers IT ops and automation. The gap is at the data architecture layer above it.",
                "The layer between Fivetran/Nucleus and Spirit/Perception/Athena is where quality degrades.",
                "Multi-market, multi-brand data flowing into AI models without a unified orchestration layer.",
                "This is not a TCS problem. It is not an internal team problem. It is an architecture gap.",
            ],
        },
    ]

    for i, prob in enumerate(problems):
        x = Inches(0.25) + i * Inches(4.35)
        add_rect(slide, x, Inches(3.55), Inches(4.1), Inches(3.7), BLUE_DARK)
        add_rect(slide, x, Inches(3.55), Inches(4.1), Inches(0.06), RED)
        add_textbox(slide, prob["title"], x + Inches(0.15), Inches(3.65),
                    Inches(3.8), Inches(0.5), size=12, bold=True, color=WHITE)
        dbox = slide.shapes.add_textbox(x + Inches(0.15), Inches(4.2),
                                         Inches(3.8), Inches(2.9))
        dtf = dbox.text_frame
        dtf.word_wrap = True
        for j, line in enumerate(prob["detail"]):
            p = dtf.paragraphs[0] if j == 0 else dtf.add_paragraph()
            p.space_before = Pt(3)
            run = p.add_run()
            run.text = f"·  {line}"
            run.font.name = FONT
            run.font.size = Pt(10.5)
            run.font.color.rgb = GREY


def slide_approach(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "04 · Our Approach", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.45), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "An external stress-test for the data architecture\nyour AI revenue depends on.",
                Inches(0.25), Inches(0.72), Inches(10.5), Inches(1.3),
                size=28, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(2.0), Inches(7), BLUE_BRIGHT)

    # Entry offer highlight box
    add_rect(slide, Inches(0.25), Inches(2.2), Inches(12.83), Inches(1.0),
             RGBColor(0x02, 0x04, 0x22))
    add_rect(slide, Inches(0.25), Inches(2.2), Inches(0.06), Inches(1.0), BLUE_BRIGHT)
    add_textbox(slide, "Entry Offer:", Inches(0.45), Inches(2.32),
                Inches(1.8), Inches(0.38), size=11, bold=True, color=BLUE_BRIGHT)
    add_textbox(slide,
                "AI Readiness Diagnostic  ·  £65–85K  ·  6 weeks  ·  "
                "Output: AI readiness score + data architecture gap analysis + prioritised fix roadmap",
                Inches(2.25), Inches(2.32), Inches(10.6), Inches(0.38),
                size=12, bold=True, color=WHITE)
    add_textbox(slide,
                "We assess the data layer between your PIM/catalog and your AI models. "
                "We identify where data quality degrades, where architecture creates bottlenecks, "
                "and where the £100M+ in AI revenue is most exposed.",
                Inches(0.45), Inches(2.72), Inches(12.2), Inches(0.38),
                size=11, color=GREY)

    # Three scope areas
    scope = [
        {
            "num": "01",
            "title": "Product Data Quality Audit",
            "lines": [
                "Assess catalog data completeness across B&Q, Screwfix, Castorama",
                "Identify taxonomy inconsistencies and attribution gaps",
                "Map impact on Spirit personalisation and Perception forecasting",
                "Quantify revenue exposure from data quality failures",
            ],
        },
        {
            "num": "02",
            "title": "AI Architecture Review",
            "lines": [
                "Stress-test the data layer between Fivetran/Nucleus and AI models",
                "Map data flows into Spirit, Perception, Inspect, Athena",
                "Identify where multi-brand complexity creates integration risk",
                "Validate governance framework against current AI velocity",
            ],
        },
        {
            "num": "03",
            "title": "Readiness Score + Roadmap",
            "lines": [
                "Score each AI platform against data architecture readiness",
                "Prioritised fix roadmap: quick wins vs. structural changes",
                "Business case: investment required vs. revenue risk protected",
                "Expansion plan: what comes after the diagnostic",
            ],
        },
    ]

    for i, s in enumerate(scope):
        x = Inches(0.25) + i * Inches(4.35)
        add_rect(slide, x, Inches(3.45), Inches(4.1), Inches(3.8), BLUE_DARK)
        add_rect(slide, x, Inches(3.45), Inches(0.55), Inches(3.8),
                 RGBColor(0x02, 0x03, 0x35))
        add_textbox(slide, s["num"], x + Inches(0.15), Inches(3.5),
                    Inches(0.6), Inches(0.38), size=13, bold=True, color=BLUE_BRIGHT)
        add_textbox(slide, s["title"], x + Inches(0.18), Inches(4.0),
                    Inches(3.8), Inches(0.45), size=13, bold=True, color=WHITE)
        add_line(slide, x + Inches(0.15), Inches(4.45), Inches(3.7), BLUE_BRIGHT)
        lbox = slide.shapes.add_textbox(x + Inches(0.15), Inches(4.55),
                                         Inches(3.8), Inches(2.55))
        ltf = lbox.text_frame
        ltf.word_wrap = True
        for j, line in enumerate(s["lines"]):
            p = ltf.paragraphs[0] if j == 0 else ltf.add_paragraph()
            p.space_before = Pt(3)
            run = p.add_run()
            run.text = f"·  {line}"
            run.font.name = FONT
            run.font.size = Pt(10.5)
            run.font.color.rgb = GREY


def slide_expansion(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "05 · Expansion Path & Value", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.45), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "Entry diagnostic. Expand across brands.",
                Inches(0.25), Inches(0.72), Inches(10.5), Inches(0.65),
                size=30, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(1.38), Inches(6), BLUE_BRIGHT)

    add_textbox(slide,
                "Multi-brand expansion is JAKALA's key differentiator. A group-level engagement scales "
                "across every brand automatically — each brand deepens the data foundation.",
                Inches(0.25), Inches(1.55), Inches(11), Inches(0.5),
                size=13, color=GREY)

    # Expansion ladder
    stages = [
        {
            "phase": "PHASE 1",
            "title": "AI Readiness Diagnostic",
            "timeline": "6 weeks",
            "value": "£65–85K",
            "color": BLUE_BRIGHT,
            "w_pct": "25%",
            "desc": "Data quality audit · AI architecture stress-test · Gap analysis · "
                    "Readiness score · Prioritised fix roadmap · Business case",
        },
        {
            "phase": "PHASE 2",
            "title": "Architecture Programme",
            "timeline": "3–6 months",
            "value": "£200–400K",
            "color": RGBColor(0x60, 0x9A, 0xFF),
            "w_pct": "50%",
            "desc": "Fix the gaps identified in Phase 1 · Data pipeline governance · "
                    "Unified catalog orchestration layer · Multi-brand data taxonomy · API standardisation",
        },
        {
            "phase": "PHASE 3",
            "title": "Data Orchestration at Scale",
            "timeline": "6–12 months",
            "value": "£400–600K",
            "color": RGBColor(0x90, 0xBB, 0xFF),
            "w_pct": "75%",
            "desc": "Speedtrain-layer across B&Q + Screwfix + Castorama · "
                    "Decoupled intelligence layer feeding Spirit/Athena/Perception · "
                    "Real-time product data quality at scale",
        },
        {
            "phase": "PHASE 4",
            "title": "Experience Transformation",
            "timeline": "12–24 months",
            "value": "£900K+",
            "color": RED,
            "w_pct": "100%",
            "desc": "Full digital experience transformation · B&Q → Screwfix → Castorama → Brico Dépôt · "
                    "Composable architecture · AI personalisation layer · Multi-market rollout",
        },
    ]

    for i, s in enumerate(stages):
        y = Inches(2.3) + i * Inches(1.15)
        add_rect(slide, Inches(0.25), y, Inches(12.83), Inches(1.02),
                 RGBColor(0x06, 0x06, 0x20) if i % 2 == 0 else RGBColor(0x09, 0x09, 0x26))
        add_rect(slide, Inches(0.25), y, Inches(0.06), Inches(1.02), s["color"])
        add_textbox(slide, s["phase"], Inches(0.45), y + Inches(0.07),
                    Inches(1.1), Inches(0.35), size=10, bold=True, color=s["color"])
        add_textbox(slide, s["title"], Inches(1.65), y + Inches(0.07),
                    Inches(2.8), Inches(0.35), size=13, bold=True, color=WHITE)
        add_textbox(slide, s["timeline"], Inches(4.55), y + Inches(0.07),
                    Inches(1.5), Inches(0.35), size=11, color=GREY)
        add_textbox(slide, s["value"], Inches(6.1), y + Inches(0.07),
                    Inches(1.5), Inches(0.35), size=13, bold=True, color=s["color"])
        add_textbox(slide, s["desc"], Inches(1.65), y + Inches(0.5),
                    Inches(11.2), Inches(0.42), size=10, color=GREY)

    # Total value
    add_rect(slide, Inches(0.25), Inches(7.0), Inches(12.83), Inches(0.38),
             RGBColor(0x02, 0x02, 0x25))
    add_textbox(slide,
                "Total programme value: £1.6M–£2M+  ·  "
                "Revenue hypothesis: £285M base case recoverable revenue protected  ·  "
                "Entry risk: £65–85K diagnostic",
                Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.28),
                size=10, color=GREY)


def slide_buyer(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "06 · Named Buyers & Outreach Plan", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.45), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "Two confirmed entry points into the group.",
                Inches(0.25), Inches(0.72), Inches(11), Inches(0.65),
                size=30, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(1.38), Inches(6), BLUE_BRIGHT)

    # Buyer 1 — Simon Jury
    add_rect(slide, Inches(0.25), Inches(1.65), Inches(6.3), Inches(5.5), BLUE_DARK)
    add_rect(slide, Inches(0.25), Inches(1.65), Inches(6.3), Inches(0.07), BLUE_BRIGHT)
    add_textbox(slide, "PRIMARY — DATA & AI ENTRY",
                Inches(0.45), Inches(1.75), Inches(5.9), Inches(0.32),
                size=9, bold=True, color=BLUE_BRIGHT)
    add_textbox(slide, "Dr. Simon Jury",
                Inches(0.45), Inches(2.08), Inches(5.9), Inches(0.5),
                size=20, bold=True, color=WHITE)
    add_textbox(slide, "Chief Data Officer · Kingfisher plc",
                Inches(0.45), Inches(2.55), Inches(5.9), Inches(0.38),
                size=13, color=GREY)
    add_line(slide, Inches(0.45), Inches(2.92), Inches(5.7), BLUE_BRIGHT)

    jury_facts = [
        "In role since January 2025 — came from Asda CDO (2 years)",
        "Mandate: data, analytics, and AI across the full Kingfisher group",
        "Budget holder for any data architecture or AI readiness programme",
        "14 months in role — priorities set, execution plan forming. Window is now.",
        "Academic credentials (Dr.) — appreciates rigour and data-led framing",
        "LinkedIn: uk.linkedin.com/in/simonjury",
    ]
    jbox = slide.shapes.add_textbox(Inches(0.45), Inches(3.02), Inches(5.9), Inches(2.5))
    jtf = jbox.text_frame
    jtf.word_wrap = True
    for i, fact in enumerate(jury_facts):
        p = jtf.paragraphs[0] if i == 0 else jtf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"·  {fact}"
        run.font.name = FONT
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE if i < 4 else BLUE_BRIGHT

    add_textbox(slide, "Outreach angle:",
                Inches(0.45), Inches(5.55), Inches(1.5), Inches(0.35),
                size=10, bold=True, color=GREY)
    add_textbox(slide,
                "Lead with the £100M+ result. Acknowledge the internal team. "
                "Frame as external architecture stress-test — not a challenge to internal capability.",
                Inches(0.45), Inches(5.88), Inches(5.9), Inches(0.95),
                size=10, color=GREY)

    # Buyer 2 — Tim Ellison
    add_rect(slide, Inches(6.78), Inches(1.65), Inches(6.3), Inches(5.5), BLUE_DARK)
    add_rect(slide, Inches(6.78), Inches(1.65), Inches(6.3), Inches(0.07),
             RGBColor(0x60, 0x9A, 0xFF))
    add_textbox(slide, "SECONDARY — PLATFORM / OPERATIONAL ENTRY",
                Inches(6.98), Inches(1.75), Inches(5.9), Inches(0.32),
                size=9, bold=True, color=RGBColor(0x60, 0x9A, 0xFF))
    add_textbox(slide, "Tim Ellison",
                Inches(6.98), Inches(2.08), Inches(5.9), Inches(0.5),
                size=20, bold=True, color=WHITE)
    add_textbox(slide, "Group Digital Centres of Excellence Director · Kingfisher plc",
                Inches(6.98), Inches(2.55), Inches(5.9), Inches(0.38),
                size=13, color=GREY)
    add_line(slide, Inches(6.98), Inches(2.92), Inches(5.7), RGBColor(0x60, 0x9A, 0xFF))

    ellison_facts = [
        "In role since December 2023 — board-level retail background",
        "Owns Group Digital CoE — directly responsible for Spirit, Perception, Inspect, Screwfix Lens",
        "Prior: TEAL Retail Consulting, Interim Head of E-Commerce B&Q (Covid era), Boots UK (15+ yrs)",
        "Closest to platform teams — good secondary if Jury unresponsive",
        "Likely receives any architecture or data quality conversation from Jury",
        "LinkedIn: linkedin.com/in/tim-ellison-36833b30",
    ]
    ebox = slide.shapes.add_textbox(Inches(6.98), Inches(3.02), Inches(5.9), Inches(2.5))
    etf = ebox.text_frame
    etf.word_wrap = True
    for i, fact in enumerate(ellison_facts):
        p = etf.paragraphs[0] if i == 0 else etf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"·  {fact}"
        run.font.name = FONT
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE if i < 4 else RGBColor(0x60, 0x9A, 0xFF)

    add_textbox(slide, "Outreach angle:",
                Inches(6.98), Inches(5.55), Inches(1.5), Inches(0.35),
                size=10, bold=True, color=GREY)
    add_textbox(slide,
                "Same AI complexity angle — data architecture will constrain Spirit and Athena at scale. "
                "Use only if Jury is unresponsive after 1 week.",
                Inches(6.98), Inches(5.88), Inches(5.9), Inches(0.95),
                size=10, color=GREY)


def slide_next_steps(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "07 · Next Steps", Inches(0.25), Inches(0.3),
                Inches(12), Inches(0.45), size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "Three actions. One goal: first conversation with Simon Jury.",
                Inches(0.25), Inches(0.72), Inches(11), Inches(0.65),
                size=28, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(1.38), Inches(7), BLUE_BRIGHT)

    steps = [
        {
            "num": "01",
            "title": "Send LinkedIn outreach to Dr. Simon Jury — this week",
            "urgency": "BEFORE 2026-03-13",
            "body": (
                "Message is written and ready. Peer-to-peer tone — CDO to data strategist. "
                "Lead with the £100M+ result. Acknowledge the internal team. "
                "End with a hypothesis, not a pitch: "
                '"I\'d be curious where you see the data architecture constraints appearing as you scale."'
            ),
            "color": RED,
        },
        {
            "num": "02",
            "title": "Brief the JAKALA internal team on Kingfisher AI context",
            "urgency": "THIS WEEK",
            "body": (
                "Before any call: team must know Spirit, Perception, Inspect, Athena, Screwfix Lens. "
                "TCS relationship — never position against them. "
                "Internal build culture — frame JAKALA as complement, not replacement. "
                "The gap is the data architecture layer above TCS, below the AI models."
            ),
            "color": BLUE_BRIGHT,
        },
        {
            "num": "03",
            "title": "Prepare AI Readiness Diagnostic one-pager — multi-brand retail",
            "urgency": "WITHIN 2 WEEKS",
            "body": (
                "Tailored to multi-brand, multi-market retail with proprietary AI platforms. "
                "Angle: marketplace GMV growth + catalog data complexity + AI architecture governance. "
                "Include revenue hypothesis (£285M base case) in business case section. "
                "If Jury unresponsive after 7 days: send Version B to Tim Ellison."
            ),
            "color": BLUE_BRIGHT,
        },
    ]

    for i, s in enumerate(steps):
        y = Inches(1.8) + i * Inches(1.78)
        add_rect(slide, Inches(0.25), y, Inches(12.83), Inches(1.65),
                 RGBColor(0x08, 0x08, 0x22))
        add_rect(slide, Inches(0.25), y, Inches(0.07), Inches(1.65), s["color"])
        add_rect(slide, Inches(0.4), y + Inches(0.12), Inches(0.5), Inches(0.48), BLUE_DARK)
        add_textbox(slide, s["num"], Inches(0.4), y + Inches(0.16),
                    Inches(0.5), Inches(0.42), size=13, bold=True, color=s["color"],
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, s["title"], Inches(1.05), y + Inches(0.1),
                    Inches(9.5), Inches(0.4), size=13, bold=True, color=WHITE)
        add_textbox(slide, s["urgency"], Inches(10.7), y + Inches(0.12),
                    Inches(2.2), Inches(0.35), size=10, bold=True, color=s["color"],
                    align=PP_ALIGN.RIGHT)
        add_textbox(slide, s["body"], Inches(1.05), y + Inches(0.55),
                    Inches(11.7), Inches(1.0), size=11, color=GREY)

    # Footer
    add_rect(slide, Inches(0.25), Inches(7.1), Inches(12.83), Inches(0.3),
             RGBColor(0x05, 0x05, 0x20))
    add_textbox(slide,
                "JAKALA · Data & AI Practice · March 2026  ·  "
                "Prepared for: Dr. Simon Jury, CDO & Tim Ellison, Group Digital CoE Director · Kingfisher plc",
                Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.25),
                size=9, color=GREY)


def slide_outreach(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    add_textbox(slide, "Appendix · Outreach Message — Dr. Simon Jury",
                Inches(0.25), Inches(0.3), Inches(12), Inches(0.45),
                size=11, color=BLUE_BRIGHT, bold=True)
    add_textbox(slide, "Ready to send. LinkedIn — Version B (recommended).",
                Inches(0.25), Inches(0.72), Inches(11), Inches(0.55),
                size=22, bold=True, color=WHITE)
    add_line(slide, Inches(0.25), Inches(1.3), Inches(6), BLUE_BRIGHT)

    add_rect(slide, Inches(0.25), Inches(1.5), Inches(12.83), Inches(4.8),
             RGBColor(0x05, 0x05, 0x28))
    add_rect(slide, Inches(0.25), Inches(1.5), Inches(0.06), Inches(4.8), BLUE_BRIGHT)

    msg = (
        'Hi Simon —\n\n'
        'The £100M+ in AI-driven incremental revenue Kingfisher reported is a strong result, '
        'and what your team has built with Spirit, Perception, and now Athena across B&Q, Screwfix, '
        'and Castorama is genuinely impressive.\n\n'
        'I work with JAKALA. We specialise in data architecture for retail groups running AI at scale — '
        'specifically the layer between the product catalog and the AI models that depend on it. '
        'With marketplace GMV growing 62% year-on-year across five markets, that layer is about to '
        'carry a lot more weight.\n\n'
        'Would you be open to a short conversation? I\'d be curious where you see the data architecture '
        'constraints appearing as you scale the programme further.'
    )

    mbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(12.3), Inches(4.45))
    mtf = mbox.text_frame
    mtf.word_wrap = True
    p = mtf.paragraphs[0]
    run = p.add_run()
    run.text = msg
    run.font.name = FONT
    run.font.size = Pt(13)
    run.font.color.rgb = WHITE

    add_rect(slide, Inches(0.25), Inches(6.45), Inches(12.83), Inches(0.82),
             RGBColor(0x02, 0x02, 0x18))
    notes = [
        "Tone: peer-to-peer · CDO to data strategist · no sales language",
        "Lead with their result, not our offer · end with a question, not a pitch",
        "If no response in 7 days: send secondary message to Tim Ellison (Group Digital CoE)",
        "Do NOT mention TCS · Do NOT lead with DXP or platform replacement",
    ]
    nbox = slide.shapes.add_textbox(Inches(0.45), Inches(6.52), Inches(12.5), Inches(0.7))
    ntf = nbox.text_frame
    ntf.word_wrap = True
    for i, note in enumerate(notes):
        p = ntf.paragraphs[0] if i == 0 else ntf.add_paragraph()
        run = p.add_run()
        run.text = f"·  {note}"
        run.font.name = FONT
        run.font.size = Pt(10)
        run.font.color.rgb = GREY


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    slide_cover(prs)
    slide_signal(prs)
    slide_company(prs)
    slide_problem(prs)
    slide_approach(prs)
    slide_expansion(prs)
    slide_buyer(prs)
    slide_next_steps(prs)
    slide_outreach(prs)

    out = "/Users/jacobskaue/Desktop/jakala-commercial-os/Accounts/kingfisher/kingfisher-pitch-2026-03.pptx"
    prs.save(out)
    print(f"✓ Saved: {out}")
    print(f"  Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
