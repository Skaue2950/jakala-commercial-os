"""
JAKALA × H&M Group — Pitch Deck
Type: pitch
Language: English
Audience: Adam Ull (Global Product Information Lead) + CDO
Date: March 2026
Offerings: Data Revenue Diagnostic → Speedtrain → J-IGNITE
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Design Constants ────────────────────────────────────────────────────────
BLUE_BRIGHT = RGBColor(0x15, 0x3E, 0xED)
BLUE_DARK   = RGBColor(0x02, 0x02, 0x66)
RED         = RGBColor(0xF6, 0x57, 0x4A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY        = RGBColor(0xAA, 0xAA, 0xBB)
BG_COLOR    = RGBColor(0x08, 0x08, 0x18)
FONT        = 'Raleway'
W           = Inches(13.33)
H           = Inches(7.5)

# ─── Helpers ─────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color=BG_COLOR):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, text, x, y, w, h, size=14, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def rect(slide, x, y, w, h, color=BLUE_DARK):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def line(slide, x, y, w, color=BLUE_BRIGHT):
    s = slide.shapes.add_shape(1, x, y, w, Inches(0.02))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def multiline(slide, items, x, y, w, h, size=11, gap=Pt(4)):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = gap
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color

def label(slide, text, x, y, w=Inches(3), size=9, color=BLUE_BRIGHT):
    tb(slide, text, x, y, w, Inches(0.3), size=size, bold=True, color=color)

# ─── Slides ──────────────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    tb(slide, "JAKALA", Inches(0.25), Inches(0.3), Inches(3), Inches(0.45),
       size=13, bold=True, color=BLUE_BRIGHT)
    tb(slide, "Data & AI · Commerce · Innovation", Inches(0.25), Inches(0.65),
       Inches(5), Inches(0.35), size=11, color=GREY)

    tb(slide, "H&M Group", Inches(0.25), Inches(1.55), Inches(10), Inches(0.6),
       size=18, color=GREY)
    tb(slide, "Your product data is the foundation\neverything else is building on.",
       Inches(0.25), Inches(2.05), Inches(9.2), Inches(2.2),
       size=40, bold=True, color=WHITE)

    line(slide, Inches(0.25), Inches(4.2), Inches(6.5), BLUE_BRIGHT)

    tb(slide, "How JAKALA helps H&M Group turn product data quality and innovation throughput\ninto compounding competitive advantage — across all eight brands.",
       Inches(0.25), Inches(4.35), Inches(9), Inches(0.85),
       size=14, color=GREY)

    # Right panel
    rect(slide, Inches(9.7), Inches(0), Inches(3.63), H, BLUE_DARK)
    rect(slide, Inches(9.7), Inches(1.9), Inches(3.5), Inches(0.055), BLUE_BRIGHT)

    stats = [
        ("8", "Fashion brands\nacross 70+ markets"),
        ("30+", "Parallel AI initiatives\ncurrently active"),
        ("€900K", "Estimated JAKALA\nopportunity (unweighted)"),
    ]
    for i, (val, lbl) in enumerate(stats):
        y = Inches(2.1) + i * Inches(1.65)
        tb(slide, val, Inches(9.9), y, Inches(3.2), Inches(0.65),
           size=28, bold=True, color=WHITE)
        tb(slide, lbl, Inches(9.9), y + Inches(0.65), Inches(3.2), Inches(0.55),
           size=11, color=GREY)

    tb(slide, "Prepared for: Adam Ull, Global Product Information Lead",
       Inches(0.25), H - Inches(1.0), Inches(9), Inches(0.35), size=11, color=GREY)
    tb(slide, "March 2026  ·  Confidential",
       Inches(0.25), H - Inches(0.65), Inches(4), Inches(0.3), size=10, color=GREY)


def slide_why_now(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    label(slide, "01 · Why Now", Inches(0.25), Inches(0.3))
    tb(slide, "Four signals converging at once.",
       Inches(0.25), Inches(0.72), Inches(10), Inches(0.62),
       size=30, bold=True, color=WHITE)
    line(slide, Inches(0.25), Inches(1.35), Inches(7), BLUE_BRIGHT)

    signals = [
        {
            "tag": "PUBLIC SIGNAL",
            "tag_color": BLUE_BRIGHT,
            "title": "Adam Ull presenting at eCommerce Expo 2026",
            "body": "H&M's Global Product Information Lead is keynoting on 'The Future of Commerce at H&M' — specifically on product data distribution and AI integration. A senior leader presenting on a problem publicly is the clearest possible buy signal.",
            "cta": "→  Connect before the conference. Not after.",
        },
        {
            "tag": "TECH SIGNAL",
            "tag_color": BLUE_BRIGHT,
            "title": "Google Cloud AI partnership — live",
            "body": "Supply chain, demand forecasting, virtual fitting rooms. The AI infrastructure is being built. But AI output quality is directly proportional to product data quality — and 30+ parallel initiatives need governance to stay coherent.",
            "cta": "→  The data layer is the constraint.",
        },
        {
            "tag": "MACRO SIGNAL",
            "tag_color": RED,
            "title": "Sweden National AI Strategy — SEK 12.5B (Feb 2026)",
            "body": "Sweden's government has committed SEK 12.5B to national AI investment. For H&M Group, headquartered in Stockholm, this raises the internal bar for AI readiness and creates board-level urgency to demonstrate progress.",
            "cta": "→  AI readiness is now a governance question, not just a technology one.",
        },
        {
            "tag": "SCALE SIGNAL",
            "tag_color": RED,
            "title": "30+ AI initiatives across 8 brands and 70+ markets",
            "body": "This is not an innovation deficit. This is an innovation throughput problem. Evaluating, prioritising, and coordinating 30+ concurrent AI initiatives across multiple brands without a structured pipeline is where good ideas — and budget — disappear.",
            "cta": "→  The bottleneck is governance, not ambition.",
        },
    ]

    for i, s in enumerate(signals):
        col = i % 2
        row = i // 2
        x = Inches(0.25) + col * Inches(6.55)
        y = Inches(1.6) + row * Inches(2.75)
        rect(slide, x, y, Inches(6.3), Inches(2.6), BLUE_DARK)
        rect(slide, x, y, Inches(6.3), Inches(0.065), s["tag_color"])
        tb(slide, s["tag"], x + Inches(0.18), y + Inches(0.12),
           Inches(5.9), Inches(0.28), size=9, bold=True, color=s["tag_color"])
        tb(slide, s["title"], x + Inches(0.18), y + Inches(0.4),
           Inches(5.9), Inches(0.42), size=13, bold=True, color=WHITE)
        tb(slide, s["body"], x + Inches(0.18), y + Inches(0.85),
           Inches(5.9), Inches(1.0), size=10.5, color=GREY)
        tb(slide, s["cta"], x + Inches(0.18), y + Inches(2.2),
           Inches(5.9), Inches(0.28), size=10, bold=True, color=s["tag_color"])


def slide_snapshot(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    label(slide, "02 · Company Snapshot", Inches(0.25), Inches(0.3))
    tb(slide, "H&M Group", Inches(0.25), Inches(0.72), Inches(11), Inches(0.62),
       size=30, bold=True, color=WHITE)
    line(slide, Inches(0.25), Inches(1.35), Inches(5), BLUE_BRIGHT)

    # Left block — facts
    rect(slide, Inches(0.25), Inches(1.55), Inches(5.8), Inches(5.7), BLUE_DARK)
    tb(slide, "At a Glance", Inches(0.45), Inches(1.67),
       Inches(5.4), Inches(0.35), size=11, bold=True, color=BLUE_BRIGHT)

    facts = [
        ("HQ", "Stockholm, Sweden · Listed on Nasdaq Stockholm"),
        ("Brands", "H&M · COS · & Other Stories · Arket · Monki · Weekday · NYDEN · Afound"),
        ("Markets", "70+ countries · 4,000+ stores + ecommerce"),
        ("Revenue", "~SEK 236B (~€21B) · ecommerce growing share"),
        ("Employees", "100,000+"),
        ("AI", "Google Cloud AI partnership active · Supply chain · Demand forecasting · Visual fitting"),
        ("AI Strategy", "30+ parallel AI initiatives · Sweden AI Strategy SEK 12.5B backdrop"),
        ("Product data", "Adam Ull leads global team of 25 managing product information across all brands"),
    ]
    items = []
    for label_text, val in facts:
        items.append((f"{label_text}:  {val}", WHITE))
        items.append(("", GREY))
    multiline(slide, [(f"{l}:  {v}" if l else "", WHITE if l else GREY)
                      for l, v in facts],
              Inches(0.45), Inches(2.08), Inches(5.4), Inches(4.9),
              size=11, gap=Pt(5))

    # Right block — brands
    rect(slide, Inches(6.28), Inches(1.55), Inches(6.8), Inches(5.7),
         RGBColor(0x05, 0x05, 0x30))
    tb(slide, "The 8 Brands — One Data Problem",
       Inches(6.48), Inches(1.67), Inches(6.4), Inches(0.35),
       size=11, bold=True, color=BLUE_BRIGHT)

    brands = [
        ("H&M", "Mass market · largest brand · ecommerce flagship"),
        ("COS", "Premium positioning · agile team · best entry point"),
        ("& Other Stories", "Independent creative studios · fragmented product data"),
        ("Arket", "Lifestyle concept · smaller team · faster decisions"),
        ("Monki", "Youth / digital-first · social commerce focus"),
        ("Weekday", "Trend-led · collaborative collections"),
        ("NYDEN", "Co-creation model · AI personalisation ambition"),
        ("Afound", "Outlet / off-price · multi-brand SKU consolidation challenge"),
    ]
    for i, (brand, desc) in enumerate(brands):
        y = Inches(2.1) + i * Inches(0.66)
        bg = RGBColor(0x0A, 0x0A, 0x22) if i % 2 == 0 else RGBColor(0x0D, 0x0D, 0x28)
        rect(slide, Inches(6.28), y, Inches(6.8), Inches(0.6), bg)
        tb(slide, brand, Inches(6.45), y + Inches(0.1),
           Inches(1.5), Inches(0.35), size=11, bold=True, color=WHITE)
        tb(slide, desc, Inches(8.0), y + Inches(0.1),
           Inches(4.85), Inches(0.35), size=10.5, color=GREY)


def slide_problem(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    label(slide, "03 · The Problem", Inches(0.25), Inches(0.3))
    tb(slide, "Two compounding challenges.\nOne shared root cause.",
       Inches(0.25), Inches(0.72), Inches(10), Inches(1.3),
       size=30, bold=True, color=WHITE)
    line(slide, Inches(0.25), Inches(2.0), Inches(6), BLUE_BRIGHT)

    # Root cause bar
    rect(slide, Inches(0.25), Inches(2.2), Inches(12.83), Inches(0.78),
         RGBColor(0x10, 0x02, 0x04))
    rect(slide, Inches(0.25), Inches(2.2), Inches(0.06), Inches(0.78), RED)
    tb(slide, "Root cause:",
       Inches(0.45), Inches(2.32), Inches(1.3), Inches(0.3),
       size=11, bold=True, color=RED)
    tb(slide, "At H&M Group's scale — 8 brands, 70+ markets, millions of SKUs — the quality of your product data and the speed of your innovation process are the same problem. "
       "Both degrade when intelligence is fragmented across brands rather than shared.",
       Inches(1.8), Inches(2.32), Inches(11.1), Inches(0.55),
       size=11, color=WHITE)

    # Two problem columns
    for col, (title, tag, tag_color, points) in enumerate([
        (
            "Challenge 1 — Product Data Quality",
            "DATA REVENUE AT RISK",
            BLUE_BRIGHT,
            [
                "8 brands × 70+ markets = product data managed independently per brand",
                "No confirmed group-level PIM — each brand likely maintains its own catalog standards",
                "Google Cloud AI partnership outputs (demand forecasting, recommendations) are directly dependent on input data quality",
                "30+ AI initiatives = 30+ data quality dependencies. One weak foundation affects all of them.",
                "Adam Ull's eCommerce Expo keynote topic ('product data distribution + AI integration') signals internal awareness of the gap",
                "Estimated revenue impact: even 5% improvement in product discovery conversion on H&M's ecommerce = material at €21B revenue scale",
            ]
        ),
        (
            "Challenge 2 — Innovation Throughput",
            "J-IGNITE FIT",
            RED,
            [
                "30+ parallel AI initiatives across 8 independent brands = evaluation and prioritisation overload",
                "Brand autonomy (the business model) creates coordination overhead — good ideas die between BUs",
                "No shared innovation pipeline framework = duplicated effort, inconsistent evaluation criteria",
                "12,000 hours/year industry benchmark spent on manual innovation evaluation — at H&M Group scale, this is conservative",
                "Sweden AI Strategy (SEK 12.5B) raises governance expectations — board needs visibility on ROI across all AI investments",
                "Without a structured pipeline: 90% of initiatives stall before proper evaluation",
            ]
        )
    ]):
        x = Inches(0.25) + col * Inches(6.55)
        rect(slide, x, Inches(3.18), Inches(6.3), Inches(4.05), BLUE_DARK)
        rect(slide, x, Inches(3.18), Inches(6.3), Inches(0.065), tag_color)
        tb(slide, tag, x + Inches(0.18), Inches(3.28),
           Inches(5.9), Inches(0.28), size=9, bold=True, color=tag_color)
        tb(slide, title, x + Inches(0.18), Inches(3.58),
           Inches(5.9), Inches(0.42), size=13, bold=True, color=WHITE)
        line(slide, x + Inches(0.18), Inches(4.02), Inches(5.9), tag_color)
        pts = [(f"·  {p}", GREY) for p in points]
        multiline(slide, pts, x + Inches(0.18), Inches(4.12),
                  Inches(5.9), Inches(2.9), size=10.5, gap=Pt(3))


def slide_approach(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    label(slide, "04 · Our Approach", Inches(0.25), Inches(0.3))
    tb(slide, "Two offerings. One connected programme.",
       Inches(0.25), Inches(0.72), Inches(10), Inches(0.62),
       size=30, bold=True, color=WHITE)
    line(slide, Inches(0.25), Inches(1.35), Inches(6), BLUE_BRIGHT)

    tb(slide, "We enter on the product data problem — the highest-urgency, named-buyer conversation. "
       "J-IGNITE is the natural second conversation once we're in the room.",
       Inches(0.25), Inches(1.52), Inches(11), Inches(0.45),
       size=13, color=GREY)

    # Offer 1 — Data Revenue Diagnostic
    rect(slide, Inches(0.25), Inches(2.15), Inches(6.3), Inches(5.1), BLUE_DARK)
    rect(slide, Inches(0.25), Inches(2.15), Inches(6.3), Inches(0.065), BLUE_BRIGHT)
    tb(slide, "ENTRY OFFER", Inches(0.45), Inches(2.25),
       Inches(5.9), Inches(0.28), size=9, bold=True, color=BLUE_BRIGHT)
    tb(slide, "Data Revenue Diagnostic", Inches(0.45), Inches(2.55),
       Inches(5.9), Inches(0.5), size=18, bold=True, color=WHITE)
    tb(slide, "€75–100K  ·  6–8 weeks  ·  Entry via Adam Ull",
       Inches(0.45), Inches(3.05), Inches(5.9), Inches(0.35),
       size=12, bold=True, color=BLUE_BRIGHT)
    line(slide, Inches(0.45), Inches(3.42), Inches(5.7), BLUE_BRIGHT)

    d1_points = [
        "Audit product data completeness across H&M Group's brand catalog",
        "Identify taxonomy inconsistencies and attribution gaps affecting AI model inputs",
        "Assess product information quality for Google Cloud AI workloads",
        "Map where data fragmentation degrades Spirit/Perception-equivalent outputs",
        "Quantify: revenue impact of search relevance gaps + zero-result rates",
        "Output: AI data readiness score + revenue hypothesis + prioritised fix roadmap",
        "Pilot scope: 1–2 brands (COS or Arket recommended for faster decision cycle)",
    ]
    multiline(slide, [(f"·  {p}", GREY) for p in d1_points],
              Inches(0.45), Inches(3.52), Inches(5.8), Inches(3.45),
              size=10.5, gap=Pt(3))

    # Offer 2 — J-IGNITE
    rect(slide, Inches(6.78), Inches(2.15), Inches(6.3), Inches(5.1),
         RGBColor(0x08, 0x02, 0x12))
    rect(slide, Inches(6.78), Inches(2.15), Inches(6.3), Inches(0.065), RED)
    tb(slide, "EXPANSION OFFER", Inches(6.98), Inches(2.25),
       Inches(5.9), Inches(0.28), size=9, bold=True, color=RED)
    tb(slide, "J-IGNITE", Inches(6.98), Inches(2.55),
       Inches(5.9), Inches(0.5), size=18, bold=True, color=WHITE)
    tb(slide, "€40–60K Discovery  ·  6 weeks  ·  1 business unit pilot",
       Inches(6.98), Inches(3.05), Inches(5.9), Inches(0.35),
       size=12, bold=True, color=RED)
    line(slide, Inches(6.98), Inches(3.42), Inches(5.7), RED)

    d2_points = [
        "Map H&M Group's current innovation evaluation process across brands",
        "Identify where AI initiatives stall — evaluation lag, coordination gaps, duplication",
        "Deploy AI-powered evaluation layer across one brand's innovation pipeline",
        "Demonstrate: 6-month decision cycle → 6-week structured progress",
        "Output: innovation process assessment + platform business case + rollout roadmap",
        "Sweden AI Strategy compliance framing: board-level visibility on all AI investments",
        "Pilot scope: 1 brand (H&M SE or COS) · scale across all 8 brands post-pilot",
    ]
    multiline(slide, [(f"·  {p}", GREY) for p in d2_points],
              Inches(6.98), Inches(3.52), Inches(5.8), Inches(3.45),
              size=10.5, gap=Pt(3))


def slide_expansion(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    label(slide, "05 · Expansion Path & Value", Inches(0.25), Inches(0.3))
    tb(slide, "Start with one brand. Scale across all eight.",
       Inches(0.25), Inches(0.72), Inches(11), Inches(0.62),
       size=30, bold=True, color=WHITE)
    line(slide, Inches(0.25), Inches(1.35), Inches(7), BLUE_BRIGHT)

    tb(slide, "H&M Group's multi-brand structure is JAKALA's expansion multiplier. "
       "Every programme validated on one brand becomes the blueprint for the next seven.",
       Inches(0.25), Inches(1.52), Inches(11), Inches(0.45), size=13, color=GREY)

    phases = [
        {
            "phase": "PHASE 1",
            "title": "Data Revenue Diagnostic",
            "sub": "Entry — 6–8 weeks · COS or Arket",
            "value": "€75–100K",
            "color": BLUE_BRIGHT,
            "desc": "Product data audit across 1–2 brands · AI readiness score · "
                    "Revenue hypothesis · Fix roadmap · Adam Ull is the buyer",
        },
        {
            "phase": "PHASE 2",
            "title": "Speedtrain Implementation",
            "sub": "Data orchestration — 3–6 months · 1–2 brands",
            "value": "€250–500K",
            "color": RGBColor(0x60, 0x9A, 0xFF),
            "desc": "Proprietary data orchestration layer between PIM and AI models · "
                    "Decouples intelligence from frontend · Enables real-time personalisation at scale",
        },
        {
            "phase": "PHASE 3",
            "title": "J-IGNITE Innovation Pipeline",
            "sub": "AI-powered governance — 6 weeks pilot → multi-brand rollout",
            "value": "€150–400K",
            "color": RED,
            "desc": "Deploy innovation pipeline management across H&M Group · "
                    "AI evaluates, scores and routes all 30+ initiatives · "
                    "Board-level visibility on ROI across all AI investments",
        },
        {
            "phase": "PHASE 4",
            "title": "AI Personalisation + DXP Transformation",
            "sub": "Full transformation — 12–24 months · all 8 brands",
            "value": "€500K–2M+",
            "color": RGBColor(0xFF, 0xCC, 0x44),
            "desc": "Full digital experience transformation across all brands · "
                    "AI personalisation layer · Multi-market composable architecture · "
                    "Data + Commerce + Innovation converging into one programme",
        },
    ]

    for i, ph in enumerate(phases):
        y = Inches(2.2) + i * Inches(1.2)
        bg = RGBColor(0x06, 0x06, 0x20) if i % 2 == 0 else RGBColor(0x09, 0x09, 0x26)
        rect(slide, Inches(0.25), y, Inches(12.83), Inches(1.1), bg)
        rect(slide, Inches(0.25), y, Inches(0.065), Inches(1.1), ph["color"])
        tb(slide, ph["phase"], Inches(0.45), y + Inches(0.08),
           Inches(1.1), Inches(0.32), size=9, bold=True, color=ph["color"])
        tb(slide, ph["title"], Inches(1.65), y + Inches(0.08),
           Inches(3.5), Inches(0.35), size=13, bold=True, color=WHITE)
        tb(slide, ph["sub"], Inches(1.65), y + Inches(0.46),
           Inches(3.5), Inches(0.3), size=10, color=GREY)
        tb(slide, ph["value"], Inches(5.3), y + Inches(0.08),
           Inches(1.6), Inches(0.35), size=14, bold=True, color=ph["color"])
        tb(slide, ph["desc"], Inches(7.0), y + Inches(0.08),
           Inches(5.9), Inches(0.88), size=10.5, color=GREY)

    rect(slide, Inches(0.25), Inches(7.08), Inches(12.83), Inches(0.32),
         RGBColor(0x03, 0x03, 0x22))
    tb(slide, "Total programme value: €975K–€3M+  ·  "
       "Entry investment: €75–100K  ·  "
       "Revenue hypothesis: significant uplift on €21B H&M Group revenue",
       Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.25), size=10, color=GREY)


def slide_buyer(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    label(slide, "06 · Named Buyer", Inches(0.25), Inches(0.3))
    tb(slide, "The right conversation. The right person. Now.",
       Inches(0.25), Inches(0.72), Inches(11), Inches(0.62),
       size=30, bold=True, color=WHITE)
    line(slide, Inches(0.25), Inches(1.35), Inches(7), BLUE_BRIGHT)

    # Main buyer card — Adam Ull
    rect(slide, Inches(0.25), Inches(1.6), Inches(7.5), Inches(5.65), BLUE_DARK)
    rect(slide, Inches(0.25), Inches(1.6), Inches(7.5), Inches(0.065), BLUE_BRIGHT)
    tb(slide, "PRIMARY BUYER — DATA REVENUE ENTRY",
       Inches(0.45), Inches(1.7), Inches(7.1), Inches(0.28),
       size=9, bold=True, color=BLUE_BRIGHT)
    tb(slide, "Adam Ull", Inches(0.45), Inches(2.0),
       Inches(7.1), Inches(0.55), size=24, bold=True, color=WHITE)
    tb(slide, "Global Product Information Lead · H&M Group · Stockholm",
       Inches(0.45), Inches(2.52), Inches(7.1), Inches(0.35),
       size=13, color=GREY)
    line(slide, Inches(0.45), Inches(2.9), Inches(7.1), BLUE_BRIGHT)

    facts = [
        ("Role", "Leads global team of 25 managing product information for all H&M Group online stores and marketplaces"),
        ("Background", "Launched Amazon Sweden · Led Klarna's international expansion — commercial sophistication, not just technical"),
        ("Conference", "Keynoting at eCommerce Expo 2026 (Sep 23–24, London) — topic: 'product data distribution and AI integration'"),
        ("Signal", "Public keynote = product data quality is a known, strategic gap he is actively working to solve"),
        ("Why he buys", "He owns the pain. His team produces the data that feeds every AI model. A diagnostic that quantifies the revenue impact of data gaps is directly in his mandate."),
        ("LinkedIn", "linkedin.com/in/adamull/"),
    ]
    fitems = []
    for l, v in facts:
        fitems.append((f"{l}:  {v}", WHITE if l != "LinkedIn" else BLUE_BRIGHT))
    multiline(slide, fitems, Inches(0.45), Inches(3.0),
              Inches(7.1), Inches(3.95), size=11, gap=Pt(5))

    # Right — outreach + secondary
    rect(slide, Inches(7.98), Inches(1.6), Inches(5.1), Inches(3.4),
         RGBColor(0x05, 0x05, 0x28))
    tb(slide, "OUTREACH MESSAGE — LINKEDIN",
       Inches(8.18), Inches(1.72), Inches(4.7), Inches(0.28),
       size=9, bold=True, color=BLUE_BRIGHT)

    msg = (
        "Hi Adam —\n\n"
        "I noticed your upcoming keynote at eCommerce Expo 2026 on product data distribution and AI "
        "integration. The topic is exactly the intersection we work in at JAKALA.\n\n"
        "We specialise in product data architecture for multi-brand retailers — specifically quantifying "
        "the revenue impact of data quality gaps and building the orchestration layer that feeds AI models "
        "reliably at scale.\n\n"
        "Would a short conversation before the conference be useful? I'd be curious what you're "
        "seeing as the hardest part of the data-to-AI pipeline at H&M Group's scale."
    )
    tb(slide, msg, Inches(8.18), Inches(2.05), Inches(4.7), Inches(2.7),
       size=10, color=WHITE)

    rect(slide, Inches(7.98), Inches(5.18), Inches(5.1), Inches(2.07),
         RGBColor(0x03, 0x03, 0x1A))
    tb(slide, "SECONDARY BUYER — J-IGNITE + AI READINESS",
       Inches(8.18), Inches(5.28), Inches(4.7), Inches(0.28),
       size=9, bold=True, color=RED)
    tb(slide, "CDO / Head of Innovation", Inches(8.18), Inches(5.58),
       Inches(4.7), Inches(0.4), size=14, bold=True, color=WHITE)
    tb(slide, "TBD — not yet named. Research path: H&M Group leadership page + LinkedIn. "
       "Identify who owns the group-level AI and innovation programme above brand level. "
       "This buyer unlocks J-IGNITE and moves win probability from 25% → 65%.",
       Inches(8.18), Inches(5.98), Inches(4.7), Inches(1.1),
       size=10, color=GREY)


def slide_next_steps(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    label(slide, "07 · Next Steps", Inches(0.25), Inches(0.3))
    tb(slide, "Four actions to open a €900K opportunity.",
       Inches(0.25), Inches(0.72), Inches(11), Inches(0.62),
       size=28, bold=True, color=WHITE)
    line(slide, Inches(0.25), Inches(1.35), Inches(7), BLUE_BRIGHT)

    steps = [
        {
            "num": "01",
            "color": BLUE_BRIGHT,
            "urgency": "THIS WEEK",
            "title": "Connect with Adam Ull on LinkedIn",
            "body": "Outreach message is written and ready. Reference eCommerce Expo 2026 keynote as opener. "
                    "Peer-to-peer tone — data practitioner to strategist. "
                    "End with a question about the hardest part of the data-to-AI pipeline at scale. "
                    "Do not mention price or a specific offering in the first message.",
        },
        {
            "num": "02",
            "color": BLUE_BRIGHT,
            "urgency": "THIS WEEK",
            "title": "Identify H&M Group CDO / Head of Innovation",
            "body": "30-minute LinkedIn research session. Look at H&M Group's leadership page and LinkedIn. "
                    "Find who owns the group-level AI and innovation programme above brand level. "
                    "This one name moves H&M's win probability from 25% to 65% and unlocks the J-IGNITE conversation. "
                    "Estimated forecast unlock: +€225K weighted.",
        },
        {
            "num": "03",
            "color": RED,
            "urgency": "WEEK 12",
            "title": "Prepare Data Revenue Diagnostic one-pager — multi-brand fashion",
            "body": "Tailored to 8-brand, 70-market fashion group with Google Cloud AI dependency. "
                    "Frame: 'your AI output quality is a function of your product data quality — let us quantify the gap.' "
                    "Include revenue hypothesis: 5% discovery improvement on H&M Group ecommerce = €X. "
                    "Recommend COS or Arket as pilot brand (shorter decision cycle, more agile team).",
        },
        {
            "num": "04",
            "color": RED,
            "urgency": "WEEK 12",
            "title": "Prepare J-IGNITE one-pager — multi-brand innovation governance",
            "body": "Angle: 30+ parallel AI initiatives need a structured evaluation and prioritisation layer. "
                    "Sweden AI Strategy (SEK 12.5B) creates board-level urgency for AI governance. "
                    "Frame J-IGNITE as the infrastructure that makes H&M's existing AI investments more productive — "
                    "not a new initiative, but the missing coordination layer.",
        },
    ]

    for i, s in enumerate(steps):
        y = Inches(1.75) + i * Inches(1.35)
        rect(slide, Inches(0.25), y, Inches(12.83), Inches(1.25),
             RGBColor(0x08, 0x08, 0x22))
        rect(slide, Inches(0.25), y, Inches(0.065), Inches(1.25), s["color"])
        rect(slide, Inches(0.4), y + Inches(0.1), Inches(0.48), Inches(0.45), BLUE_DARK)
        tb(slide, s["num"], Inches(0.4), y + Inches(0.13),
           Inches(0.48), Inches(0.38), size=13, bold=True, color=s["color"],
           align=PP_ALIGN.CENTER)
        tb(slide, s["title"], Inches(1.02), y + Inches(0.08),
           Inches(9.8), Inches(0.37), size=13, bold=True, color=WHITE)
        tb(slide, s["urgency"], Inches(11.0), y + Inches(0.1),
           Inches(2.0), Inches(0.3), size=10, bold=True, color=s["color"],
           align=PP_ALIGN.RIGHT)
        tb(slide, s["body"], Inches(1.02), y + Inches(0.5),
           Inches(11.7), Inches(0.68), size=10.5, color=GREY)

    rect(slide, Inches(0.25), Inches(7.1), Inches(12.83), Inches(0.3),
         RGBColor(0x04, 0x04, 0x1C))
    tb(slide, "JAKALA · Data & AI + Commerce + Innovation · March 2026  ·  "
       "Prepared for: Adam Ull, Global Product Information Lead · H&M Group",
       Inches(0.4), Inches(7.14), Inches(12.5), Inches(0.22), size=9, color=GREY)


# ─── Main ─────────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()
    slide_cover(prs)
    slide_why_now(prs)
    slide_snapshot(prs)
    slide_problem(prs)
    slide_approach(prs)
    slide_expansion(prs)
    slide_buyer(prs)
    slide_next_steps(prs)

    out = "/Users/jacobskaue/Desktop/jakala-commercial-os/Accounts/hm/hm-pitch-2026-03.pptx"
    prs.save(out)
    print(f"✓ Saved: {out}")
    print(f"  Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
