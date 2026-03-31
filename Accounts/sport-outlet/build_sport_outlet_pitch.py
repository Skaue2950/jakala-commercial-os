"""
JAKALA × Sport Outlet — Pitch Deck
Type: pitch
Language: Norwegian
Audience: Tor-André Skeie (CEO) · Innkommende CDO / CTO
Date: March 2026
Offerings: Commerce Optimization Pilot → Data Revenue Diagnostic + Speedtrain
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Design Constants ────────────────────────────────────────────────────────
BLUE_BRIGHT = RGBColor(0x15, 0x3E, 0xED)
BLUE_DARK   = RGBColor(0x02, 0x02, 0x66)
RED         = RGBColor(0xF6, 0x57, 0x4A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY        = RGBColor(0xAA, 0xAA, 0xBB)
BG_COLOR    = RGBColor(0x08, 0x08, 0x18)
FONT        = 'Raleway'
W           = Inches(13.33)
H           = Inches(7.5)

# ─── Helpers ─────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color=BG_COLOR):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, text, x, y, w, h, size=14, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def rect(slide, x, y, w, h, color=BLUE_DARK):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def line(slide, x, y, w, color=BLUE_BRIGHT):
    s = slide.shapes.add_shape(1, x, y, w, Inches(0.02))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def multiline(slide, items, x, y, w, h, size=11, gap=Pt(4)):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = gap
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color

def label(slide, text, x, y, w=Inches(4), size=9, color=BLUE_BRIGHT):
    tb(slide, text, x, y, w, Inches(0.3), size=size, bold=True, color=color)

# ─── Slides ──────────────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    tb(slide, "JAKALA", Inches(0.25), Inches(0.3), Inches(3), Inches(0.45),
       size=13, bold=True, color=BLUE_BRIGHT)
    tb(slide, "Data & AI · Commerce · Innovation", Inches(0.25), Inches(0.65),
       Inches(5), Inches(0.35), size=11, color=GREY)

    tb(slide, "Sport Outlet / Invent Sport", Inches(0.25), Inches(1.55),
       Inches(10), Inches(0.6), size=18, color=GREY)
    tb(slide, "18 nettbutikker.\nÉn felles mulighet.",
       Inches(0.25), Inches(2.05), Inches(9.2), Inches(2.2),
       size=40, bold=True, color=WHITE)

    line(slide, Inches(0.25), Inches(4.2), Inches(6.5), BLUE_BRIGHT)

    tb(slide, "Hvordan JAKALA hjelper Invent Sport å bygge skalerbar produktdata-infrastruktur\nog commerce-ytelse på tvers av alle 18 brands — før neste vekstfase.",
       Inches(0.25), Inches(4.35), Inches(9), Inches(0.85),
       size=14, color=GREY)

    # Right panel
    rect(slide, Inches(9.7), Inches(0), Inches(3.63), H, BLUE_DARK)
    rect(slide, Inches(9.7), Inches(1.9), Inches(3.5), Inches(0.055), BLUE_BRIGHT)

    stats = [
        ("111", "butikker · mål 150–200"),
        ("18", "parallelle nettbutikker\nunder Invent Sport"),
        ("NOK 2 mrd", "omsetning 2024\n+12,6 % YoY"),
    ]
    for i, (val, lbl) in enumerate(stats):
        y = Inches(2.1) + i * Inches(1.65)
        tb(slide, val, Inches(9.9), y, Inches(3.2), Inches(0.65),
           size=24, bold=True, color=WHITE)
        tb(slide, lbl, Inches(9.9), y + Inches(0.65), Inches(3.2), Inches(0.6),
           size=11, color=GREY)

    tb(slide, "Forberedt for: Tor-André Skeie, CEO · Sport Outlet / Invent Sport",
       Inches(0.25), H - Inches(1.0), Inches(9), Inches(0.35), size=11, color=GREY)
    tb(slide, "Mars 2026  ·  Konfidensielt",
       Inches(0.25), H - Inches(0.65), Inches(4), Inches(0.3), size=10, color=GREY)


def slide_why_now(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    label(slide, "01 · Hvorfor nå", Inches(0.25), Inches(0.3))
    tb(slide, "Fire signaler som alle peker i samme retning.",
       Inches(0.25), Inches(0.72), Inches(10), Inches(0.62),
       size=30, bold=True, color=WHITE)
    line(slide, Inches(0.25), Inches(1.35), Inches(7), BLUE_BRIGHT)

    signals = [
        {
            "tag": "LEDERSKAP-SIGNAL",
            "tag_color": BLUE_BRIGHT,
            "title": "CTO-stillingen er lyst ut — ikke besatt",
            "body": "Sport Outlet rekrutterer aktivt en CTO med ansvar for teknologiplattform, IT-sikkerhet og utviklingsteam. En ny CTO uten etablerte leverandørrelasjoner søker kompetente partnere i onboarding-fasen. Det er nå du kommer inn — ikke etter at agendaen er satt.",
            "cta": "→  Kontakt CEO før ny CTO setter kursen.",
        },
        {
            "tag": "LEDERSKAP-SIGNAL",
            "tag_color": BLUE_BRIGHT,
            "title": "CDO-stillingen er lyst ut — ikke besatt",
            "body": "Invent Sport rekrutterer en Chief Digital Officer med ansvar for tech, data, ecommerce og digital markedsføring — rapporterer direkte til CEO. Det betyr at digital agenda nå eskaleres til toppnivå. Stillingen er ny, agenda er åpen, og budsjett er uallokert.",
            "cta": "→  Ny CDO vil forme sin agenda i de første 90 dagene.",
        },
        {
            "tag": "VEKST-SIGNAL",
            "tag_color": RED,
            "title": "Fra 80 til 111 til 150–200 butikker",
            "body": "Invent Sport vokser fra ~80 til 111 butikker og har mål om 150–200. Hver nye butikk legger press på produktdata, pris og søk på tvers av alle 18 nettbutikker. Digital infrastruktur som ikke skalerer med butikkveksten = tapt omsetning på nett.",
            "cta": "→  Skalering krever én felles produktdatainfrastruktur.",
        },
        {
            "tag": "OPERASJONELT SIGNAL",
            "tag_color": RED,
            "title": "18 parallelle nettbutikker — én kompleks utfordring",
            "body": "Dynamic Brands / Invent Sport opererer 18 separate nettbutikker med felles logistikk men trolig uten delt PIM-system. Fragmentert produktdata betyr lavere søkerelevans, høyere zero-results-rate og manuelt dobbeltarbeid på tvers av brands — alt løst lenge etter at problemet har kostet dem omsetning.",
            "cta": "→  Multi-brand PIM-diagnose er den tydeligste inngangen.",
        },
    ]

    for i, s in enumerate(signals):
        col = i % 2
        row = i // 2
        x = Inches(0.25) + col * Inches(6.55)
        y = Inches(1.6) + row * Inches(2.75)
        rect(slide, x, y, Inches(6.3), Inches(2.6), BLUE_DARK)
        rect(slide, x, y, Inches(6.3), Inches(0.065), s["tag_color"])
        tb(slide, s["tag"], x + Inches(0.18), y + Inches(0.12),
           Inches(5.9), Inches(0.28), size=9, bold=True, color=s["tag_color"])
        tb(slide, s["title"], x + Inches(0.18), y + Inches(0.4),
           Inches(5.9), Inches(0.42), size=13, bold=True, color=WHITE)
        tb(slide, s["body"], x + Inches(0.18), y + Inches(0.85),
           Inches(5.9), Inches(1.0), size=10.5, color=GREY)
        tb(slide, s["cta"], x + Inches(0.18), y + Inches(2.2),
           Inches(5.9), Inches(0.28), size=10, bold=True, color=s["tag_color"])


def slide_snapshot(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    label(slide, "02 · Selskapet", Inches(0.25), Inches(0.3))
    tb(slide, "Sport Outlet / Invent Sport", Inches(0.25), Inches(0.72),
       Inches(11), Inches(0.62), size=30, bold=True, color=WHITE)
    line(slide, Inches(0.25), Inches(1.35), Inches(5), BLUE_BRIGHT)

    # Left block — facts
    rect(slide, Inches(0.25), Inches(1.55), Inches(5.8), Inches(5.7), BLUE_DARK)
    tb(slide, "Nøkkeltall", Inches(0.45), Inches(1.67),
       Inches(5.4), Inches(0.35), size=11, bold=True, color=BLUE_BRIGHT)

    facts = [
        ("Morselskap", "Invent Sport AS (tidl. Dynamic Brands-paraply)"),
        ("HQ", "Hagavik, Os (nær Bergen), Norge"),
        ("CEO", "Tor-André Skeie · Deputy: Odd Arne Larsen"),
        ("Investor", "Icon Capital (PE) — inne siden 2017"),
        ("Omsetning", "NOK 2 mrd (2024) · +12,6 % YoY · mål passert"),
        ("EBITDA", "NOK 350M+ (~17,5 % margin) — eksepsjonelt for sportsretail"),
        ("Butikker", "~111 (2025) · ekspansjonsmål 150–200"),
        ("Netthandel", "sportoutlet.no + 18 nettbutikker under Invent Sport"),
        ("Markedsandel", "~11,9 % av norsk sportsmarked (nr. 3 etter XXL og Sport Holding)"),
    ]
    multiline(slide, [(f"{l}:  {v}", WHITE) for l, v in facts],
              Inches(0.45), Inches(2.08), Inches(5.4), Inches(4.9),
              size=11, gap=Pt(5))

    # Right block — brands / webshops
    rect(slide, Inches(6.28), Inches(1.55), Inches(6.8), Inches(5.7),
         RGBColor(0x05, 0x05, 0x30))
    tb(slide, "18 Nettbutikker — Én Felles Infrastrukturutfordring",
       Inches(6.48), Inches(1.67), Inches(6.4), Inches(0.35),
       size=11, bold=True, color=BLUE_BRIGHT)

    brands = [
        ("Sport Outlet", "Kjernemerke · 111 butikker · discount/outlet"),
        ("Sportland", "Multi-brand sportsretail · utvalgte markeder"),
        ("Active Brands portfolio", "Kari Traa · Bula · Sweet Protection · Dæhlie"),
        ("Norrøna (distribusjon)", "Premium outdoor · nordisk distribusjon"),
        ("Tilleggsbrands", "Øvrige merker under Invent Sport-paraplyen"),
        ("Felles logistikk", "Sentrallager + multi-brand fulfillment"),
        ("Felles data?", "Trolig ikke — PIM-status ukjent, CTO ikke besatt"),
        ("Neste steg", "Skalering til 150–200 butikker krever felles fundament"),
    ]
    for i, (brand, desc) in enumerate(brands):
        y = Inches(2.1) + i * Inches(0.66)
        bg = RGBColor(0x0A, 0x0A, 0x22) if i % 2 == 0 else RGBColor(0x0D, 0x0D, 0x28)
        rect(slide, Inches(6.28), y, Inches(6.8), Inches(0.6), bg)
        tb(slide, brand, Inches(6.45), y + Inches(0.1),
           Inches(2.0), Inches(0.35), size=11, bold=True, color=WHITE)
        tb(slide, desc, Inches(8.5), y + Inches(0.1),
           Inches(4.35), Inches(0.35), size=10.5, color=GREY)


def slide_problem(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    label(slide, "03 · Problemet", Inches(0.25), Inches(0.3))
    tb(slide, "Vekst uten felles produktdata-fundament\ner vekst som koster mer enn den gir.",
       Inches(0.25), Inches(0.72), Inches(10), Inches(1.3),
       size=28, bold=True, color=WHITE)
    line(slide, Inches(0.25), Inches(2.0), Inches(6), BLUE_BRIGHT)

    # Root cause bar
    rect(slide, Inches(0.25), Inches(2.2), Inches(12.83), Inches(0.78),
         RGBColor(0x10, 0x02, 0x04))
    rect(slide, Inches(0.25), Inches(2.2), Inches(0.06), Inches(0.78), RED)
    tb(slide, "Kjerneproblemet:",
       Inches(0.45), Inches(2.32), Inches(1.6), Inches(0.3),
       size=11, bold=True, color=RED)
    tb(slide, "Sport Outlet og Invent Sport vokser raskt i butikk og omsetning — men digital infrastruktur for produktdata, søk og commerce "
       "er sannsynligvis ikke bygget for 18 parallelle nettbutikker og 150+ butikker. "
       "Ny CTO og CDO vil oppdage dette tidlig. JAKALA kan definere løsningen før de gjør det.",
       Inches(2.15), Inches(2.32), Inches(10.75), Inches(0.55),
       size=11, color=WHITE)

    for col, (title, tag, tag_color, points) in enumerate([
        (
            "Utfordring 1 — Produktdata og søk",
            "COMMERCE OPTIMIZATION",
            BLUE_BRIGHT,
            [
                "18 parallelle nettbutikker med trolig separate produktkatalogs-strukturer",
                "Ingen bekreftet felles PIM — CTO-søk antyder at dette er under vurdering",
                "Multi-brand søk uten felles taksonomi = lavere konvertering og høy zero-results-rate",
                "Prisreferanse-kontrovers (Forbrukertilsynet) synliggjør utfordringer med pris- og produktdatastyring",
                "Vekst til 150–200 butikker krever automatisert produktdata-flyt — manuell prosess skalerer ikke",
                "Estimert tap: selv 3–5 % bedre søkerelevans på tvers av 18 nettbutikker = betydelig omsetningsløft",
            ]
        ),
        (
            "Utfordring 2 — Data og AI-beredskap",
            "DATA REVENUE UNLOCK",
            RED,
            [
                "Ingen kjente AI-initiativer per mars 2026 — CDO-søk antyder ambisjon",
                "Ny CDO vil forvente AI-klar datainfrastruktur fra dag én — den finnes trolig ikke",
                "Multi-brand med sentrallager = rik transaksjons- og lagerdata som ikke aktiveres",
                "Personalisering på tvers av 18 nettbutikker krever felles datalayer — ikke 18 separate",
                "Invent Sports ekspansjonshastighet gir data-multipliserende effekt: hvert nytt brand = ny datastrøm",
                "Uten felles datalag: vekst øker kompleksitet, ikke innsikt",
            ]
        )
    ]):
        x = Inches(0.25) + col * Inches(6.55)
        rect(slide, x, Inches(3.18), Inches(6.3), Inches(4.05), BLUE_DARK)
        rect(slide, x, Inches(3.18), Inches(6.3), Inches(0.065), tag_color)
        tb(slide, tag, x + Inches(0.18), Inches(3.28),
           Inches(5.9), Inches(0.28), size=9, bold=True, color=tag_color)
        tb(slide, title, x + Inches(0.18), Inches(3.58),
           Inches(5.9), Inches(0.42), size=13, bold=True, color=WHITE)
        line(slide, x + Inches(0.18), Inches(4.02), Inches(5.9), tag_color)
        pts = [(f"·  {p}", GREY) for p in points]
        multiline(slide, pts, x + Inches(0.18), Inches(4.12),
                  Inches(5.9), Inches(2.9), size=10.5, gap=Pt(3))


def slide_approach(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    label(slide, "04 · Vår tilnærming", Inches(0.25), Inches(0.3))
    tb(slide, "To tilbud. Ett sammenhengende program.",
       Inches(0.25), Inches(0.72), Inches(10), Inches(0.62),
       size=30, bold=True, color=WHITE)
    line(slide, Inches(0.25), Inches(1.35), Inches(6), BLUE_BRIGHT)

    tb(slide, "Vi starter med Commerce Optimization — den mest konkrete og umiddelbart relevante inngangen. "
       "Data Revenue er den naturlige neste samtalen når vi er inne i rommet.",
       Inches(0.25), Inches(1.52), Inches(11), Inches(0.45),
       size=13, color=GREY)

    # Offer 1 — Commerce Optimization Pilot
    rect(slide, Inches(0.25), Inches(2.15), Inches(6.3), Inches(5.1), BLUE_DARK)
    rect(slide, Inches(0.25), Inches(2.15), Inches(6.3), Inches(0.065), BLUE_BRIGHT)
    tb(slide, "INNGANGS-TILBUD", Inches(0.45), Inches(2.25),
       Inches(5.9), Inches(0.28), size=9, bold=True, color=BLUE_BRIGHT)
    tb(slide, "Commerce Optimization Pilot", Inches(0.45), Inches(2.55),
       Inches(5.9), Inches(0.65), size=17, bold=True, color=WHITE)
    tb(slide, "€75–120K  ·  6–8 uker  ·  Inn via CEO eller innkommende CDO",
       Inches(0.45), Inches(3.2), Inches(5.9), Inches(0.35),
       size=12, bold=True, color=BLUE_BRIGHT)
    line(slide, Inches(0.45), Inches(3.57), Inches(5.7), BLUE_BRIGHT)

    d1_points = [
        "Kartlegging av produktdatakvalitet på tvers av Invent Sports 18 nettbutikker",
        "Identifiser taksonomi-inkonsistenser og attributt-gap som svekker søk og konvertering",
        "Benchmark søkerelevans og zero-results-rate per brand mot Nordic best practice",
        "Avdekk dobbeltarbeid og manuelle prosesser i produktdata-håndtering",
        "POC på 1–2 nettbutikker: demonstrer målbar forbedring i søk og discovery",
        "Leveranse: diagnose + revenue-hypotese + prioritert roadmap + skaleringsplan",
        "Utgangspunkt: én rapport som CEO og ny CDO/CTO kan handle på umiddelbart",
    ]
    multiline(slide, [(f"·  {p}", GREY) for p in d1_points],
              Inches(0.45), Inches(3.67), Inches(5.8), Inches(3.3),
              size=10.5, gap=Pt(3))

    # Offer 2 — Data Revenue Diagnostic + Speedtrain
    rect(slide, Inches(6.78), Inches(2.15), Inches(6.3), Inches(5.1),
         RGBColor(0x08, 0x02, 0x12))
    rect(slide, Inches(6.78), Inches(2.15), Inches(6.3), Inches(0.065), RED)
    tb(slide, "EKSPANSJONS-TILBUD", Inches(6.98), Inches(2.25),
       Inches(5.9), Inches(0.28), size=9, bold=True, color=RED)
    tb(slide, "Data Revenue Diagnostic + Speedtrain", Inches(6.98), Inches(2.55),
       Inches(5.9), Inches(0.65), size=17, bold=True, color=WHITE)
    tb(slide, "€50–100K diagnostic  ·  €200–500K Speedtrain  ·  3–9 måneder",
       Inches(6.98), Inches(3.2), Inches(5.9), Inches(0.35),
       size=12, bold=True, color=RED)
    line(slide, Inches(6.98), Inches(3.57), Inches(5.7), RED)

    d2_points = [
        "Kartlegg fullstendigheten av produktdata på tvers av alle brands og kanaler",
        "Identifiser gaps mellom data dere har og revenue den bør generere",
        "Vurder AI-beredskap: er dataen ren nok til å drive personalisering og anbefalinger?",
        "Speedtrain: JAKALAs proprietære dataorkestreringsplattform",
        "Kobler PIM/ERP fra commerce-laget — aktiverer sanntidspersonalisering på tvers av 18 brands",
        "Muliggjør én felles intelligens-lag for alle Invent Sport-nettbutikker",
        "Skalerbar fra 18 nettbutikker til 150+ butikkers datavolum",
    ]
    multiline(slide, [(f"·  {p}", GREY) for p in d2_points],
              Inches(6.98), Inches(3.67), Inches(5.8), Inches(3.3),
              size=10.5, gap=Pt(3))


def slide_expansion(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    label(slide, "05 · Ekspansjonssti og verdi", Inches(0.25), Inches(0.3))
    tb(slide, "Start med én diagnose. Bygg infrastrukturen for 150+ butikker.",
       Inches(0.25), Inches(0.72), Inches(11), Inches(0.62),
       size=28, bold=True, color=WHITE)
    line(slide, Inches(0.25), Inches(1.35), Inches(7), BLUE_BRIGHT)

    tb(slide, "Invent Sports multi-brand struktur er JAKALAs ekspansjonsmultiplikator. "
       "Hvert tilbud som valideres på ett brand skalerer direkte til alle 18 — og til den butikkveksten som kommer.",
       Inches(0.25), Inches(1.52), Inches(11), Inches(0.45), size=13, color=GREY)

    phases = [
        {
            "phase": "FASE 1",
            "title": "Commerce Optimization Pilot",
            "sub": "Inngang — 6–8 uker · 1–2 nettbutikker",
            "value": "€75–120K",
            "color": BLUE_BRIGHT,
            "desc": "Produktdata-diagnose på tvers av Invent Sport · Søkekvalitet og konverteringsanalyse · "
                    "POC på 1–2 brands · Roadmap for alle 18 · CEO eller innkommende CDO er kjøper",
        },
        {
            "phase": "FASE 2",
            "title": "Data Revenue Diagnostic",
            "sub": "Data-beredskap — 6 uker · alle brands",
            "value": "€50–100K",
            "color": RGBColor(0x60, 0x9A, 0xFF),
            "desc": "Full kartlegging av produktdata-gap · AI-beredskapscore · Revenue-hypotese · "
                    "Identifiser hvor Invent Sports transaksjonsdata kan drive personalisering og retail media",
        },
        {
            "phase": "FASE 3",
            "title": "Speedtrain Implementering",
            "sub": "Dataorkestrering — 3–9 måneder · alle 18 brands",
            "value": "€200–500K",
            "color": RED,
            "desc": "JAKALAs proprietære dataorkestreringsplattform · Én felles intelligens-lag for alle Invent Sport-brands · "
                    "Sanntidspersonalisering og produktanbefalinger på tvers av nettbutikker",
        },
        {
            "phase": "FASE 4",
            "title": "AI Personalisering + Full Commerce Optimisering",
            "sub": "Full transformasjon — 12–24 måneder",
            "value": "€300K–1M+",
            "color": RGBColor(0xFF, 0xCC, 0x44),
            "desc": "AI-drevet personalisering på tvers av alle 18 brands · "
                    "Søk og discovery optimalisert for 150–200 butikkers sortiment · "
                    "Datalag som skalerer med hver ny butikk og hvert nytt brand",
        },
    ]

    for i, ph in enumerate(phases):
        y = Inches(2.2) + i * Inches(1.2)
        bg = RGBColor(0x06, 0x06, 0x20) if i % 2 == 0 else RGBColor(0x09, 0x09, 0x26)
        rect(slide, Inches(0.25), y, Inches(12.83), Inches(1.1), bg)
        rect(slide, Inches(0.25), y, Inches(0.065), Inches(1.1), ph["color"])
        tb(slide, ph["phase"], Inches(0.45), y + Inches(0.08),
           Inches(1.1), Inches(0.32), size=9, bold=True, color=ph["color"])
        tb(slide, ph["title"], Inches(1.65), y + Inches(0.08),
           Inches(3.5), Inches(0.35), size=13, bold=True, color=WHITE)
        tb(slide, ph["sub"], Inches(1.65), y + Inches(0.46),
           Inches(3.5), Inches(0.3), size=10, color=GREY)
        tb(slide, ph["value"], Inches(5.3), y + Inches(0.08),
           Inches(1.6), Inches(0.35), size=14, bold=True, color=ph["color"])
        tb(slide, ph["desc"], Inches(7.0), y + Inches(0.08),
           Inches(5.9), Inches(0.88), size=10.5, color=GREY)

    rect(slide, Inches(0.25), Inches(7.08), Inches(12.83), Inches(0.32),
         RGBColor(0x03, 0x03, 0x22))
    tb(slide, "Total programverdi: €625K–€1.7M+  ·  "
       "Inngangs-investering: €75–120K  ·  "
       "Revenue-hypotese: 3–5 % konverteringsløft på NOK 2 mrd nettomsetning = NOK 60–100M+",
       Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.25), size=10, color=GREY)


def slide_buyer(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    label(slide, "06 · Kjøper og inngang", Inches(0.25), Inches(0.3))
    tb(slide, "Riktig person. Riktig tidspunkt. Riktig inngang.",
       Inches(0.25), Inches(0.72), Inches(11), Inches(0.62),
       size=30, bold=True, color=WHITE)
    line(slide, Inches(0.25), Inches(1.35), Inches(7), BLUE_BRIGHT)

    # Main buyer card — Tor-André Skeie
    rect(slide, Inches(0.25), Inches(1.6), Inches(7.5), Inches(5.65), BLUE_DARK)
    rect(slide, Inches(0.25), Inches(1.6), Inches(7.5), Inches(0.065), BLUE_BRIGHT)
    tb(slide, "PRIMÆR KJØPER — CEO-TRACK (NÅ)",
       Inches(0.45), Inches(1.7), Inches(7.1), Inches(0.28),
       size=9, bold=True, color=BLUE_BRIGHT)
    tb(slide, "Tor-André Skeie", Inches(0.45), Inches(2.0),
       Inches(7.1), Inches(0.55), size=24, bold=True, color=WHITE)
    tb(slide, "CEO · Sport Outlet / Invent Sport · Hagavik, Os",
       Inches(0.45), Inches(2.52), Inches(7.1), Inches(0.35),
       size=13, color=GREY)
    line(slide, Inches(0.45), Inches(2.9), Inches(7.1), BLUE_BRIGHT)

    facts = [
        ("Rolle", "CEO og grunnlegger av Sport Outlet / Invent Sport — beslutningstaker i alle strategiske spørsmål"),
        ("Bakgrunn", "Bygget Sport Outlet fra oppstart til NOK 2 mrd og 111 butikker — kommersiell skarphet, ikke bare operativt fokus"),
        ("Signal", "Uttalt preferanse for reinvestering fremfor salg (til tross for oppkjøpsinteresse) — investerer i vekst"),
        ("Timing", "Begge CTO og CDO-roller er åpne — CEO er eneste beslutningstaker i rommet akkurat nå"),
        ("Inngang", "Kort exec-POV om multi-brand produktdata og digital infrastruktur for skalering til 150+ butikker"),
        ("Sekundær", "Innkommende CDO — følg rekrutteringen og kontakt i onboarding-fasen"),
    ]
    fitems = []
    for l, v in facts:
        fitems.append((f"{l}:  {v}", WHITE))
    multiline(slide, fitems, Inches(0.45), Inches(3.0),
              Inches(7.1), Inches(3.95), size=11, gap=Pt(5))

    # Right — outreach message
    rect(slide, Inches(7.98), Inches(1.6), Inches(5.1), Inches(3.7),
         RGBColor(0x05, 0x05, 0x28))
    tb(slide, "OUTREACH-MELDING — LINKEDIN",
       Inches(8.18), Inches(1.72), Inches(4.7), Inches(0.28),
       size=9, bold=True, color=BLUE_BRIGHT)

    msg = (
        "Hei Tor-André —\n\n"
        "Imponerende vekst dere har bygget — fra 80 til 111 butikker og passert 2 mrd i omsetning "
        "er ekstraordinært for norsk sportsretail.\n\n"
        "Jeg ser at dere rekrutterer både CTO og CDO akkurat nå. Det er et bra tidspunkt for en kort "
        "samtale om et tema vi jobber mye med: når du opererer 18 parallelle nettbutikker under én "
        "paraply, er produktdata-kvalitet og søkerelevans det som avgjør om nettomsetningen holder "
        "følge med butikkveksten.\n\n"
        "JAKALA hjelper multi-brand retailere med akkurat dette. Ville 20 minutter vært nyttig?"
    )
    tb(slide, msg, Inches(8.18), Inches(2.05), Inches(4.7), Inches(3.0),
       size=10, color=WHITE)

    rect(slide, Inches(7.98), Inches(5.48), Inches(5.1), Inches(1.77),
         RGBColor(0x03, 0x03, 0x1A))
    tb(slide, "SEKUNDÆR KJØPER — CDO-TRACK",
       Inches(8.18), Inches(5.58), Inches(4.7), Inches(0.28),
       size=9, bold=True, color=RED)
    tb(slide, "Innkommende CDO (TBD)", Inches(8.18), Inches(5.88),
       Inches(4.7), Inches(0.4), size=14, bold=True, color=WHITE)
    tb(slide, "Stillingen er ikke besatt per mars 2026. Overvåk LinkedIn og FINN.no. "
       "Ny CDO vil søke kompetente partnere i onboarding-fasen — JAKALA bør være "
       "posisjonert som foretrukket partner FØR CDO setter sin 90-dagers-agenda.",
       Inches(8.18), Inches(6.28), Inches(4.7), Inches(0.85),
       size=10, color=GREY)


def slide_next_steps(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.08), H, BLUE_BRIGHT)

    label(slide, "07 · Neste steg", Inches(0.25), Inches(0.3))
    tb(slide, "Fire handlinger for å åpne en €300K+ mulighet.",
       Inches(0.25), Inches(0.72), Inches(11), Inches(0.62),
       size=28, bold=True, color=WHITE)
    line(slide, Inches(0.25), Inches(1.35), Inches(7), BLUE_BRIGHT)

    steps = [
        {
            "num": "01",
            "color": BLUE_BRIGHT,
            "urgency": "DENNE UKEN",
            "title": "Send LinkedIn-melding til Tor-André Skeie",
            "body": "Utkast er klart. Kort, konkret, kommersielt relevant. "
                    "Led med NOK 2 mrd + 18 nettbutikker — ikke med produkt eller pris. "
                    "Spør om 20 minutter for å snakke om digital infrastruktur for skalering til 150+ butikker. "
                    "Bruk vinduet mens CTO og CDO-rollene er åpne — CEO er eneste beslutningstaker nå.",
        },
        {
            "num": "02",
            "color": BLUE_BRIGHT,
            "urgency": "DENNE UKEN",
            "title": "Overvåk LinkedIn og FINN.no for CTO og CDO-ansettelser",
            "body": "Opprett varsler på Sport Outlet og Invent Sport på LinkedIn og FINN.no. "
                    "Kontakt ny CDO/CTO i onboarding-fasen (dag 1–30) — det er da de er mest mottakelig for eksternt input. "
                    "Forbered en kort intro-melding som posisjonerer JAKALA som foretrukket data/commerce-partner.",
        },
        {
            "num": "03",
            "color": RED,
            "urgency": "UKE 12–13",
            "title": "Kartlegg Dynamic Brands' 18 nettbutikker",
            "body": "Map alle 18 nettbutikker under Invent Sport: plattform, produktdata-kvalitet, søkefunksjon, "
                    "og synlige gap i attributter og kategorisering. "
                    "Identifiser de 2–3 nettbutikkene med svakest produktdata — disse er POC-kandidater i Fase 1. "
                    "Output: én side med funn til bruk i CEO-møtet.",
        },
        {
            "num": "04",
            "color": RED,
            "urgency": "UKE 13",
            "title": "Forbered Commerce Optimization Pilot-pitchdokument",
            "body": "Tilpas til Invent Sport / Dynamic Brands-kontekst: multi-brand, 18 nettbutikker, ekspansjon til 150+ butikker. "
                    "Inkluder revenue-hypotese: 3–5 % konverteringsløft på tvers av 18 brands = NOK 60–100M+. "
                    "Forbered to scenarioer: (a) CEO-track nå, (b) CDO-track ved ansettelse.",
        },
    ]

    for i, s in enumerate(steps):
        y = Inches(1.75) + i * Inches(1.35)
        rect(slide, Inches(0.25), y, Inches(12.83), Inches(1.25),
             RGBColor(0x08, 0x08, 0x22))
        rect(slide, Inches(0.25), y, Inches(0.065), Inches(1.25), s["color"])
        rect(slide, Inches(0.4), y + Inches(0.1), Inches(0.48), Inches(0.45), BLUE_DARK)
        tb(slide, s["num"], Inches(0.4), y + Inches(0.13),
           Inches(0.48), Inches(0.38), size=13, bold=True, color=s["color"],
           align=PP_ALIGN.CENTER)
        tb(slide, s["title"], Inches(1.02), y + Inches(0.08),
           Inches(9.8), Inches(0.37), size=13, bold=True, color=WHITE)
        tb(slide, s["urgency"], Inches(11.0), y + Inches(0.1),
           Inches(2.0), Inches(0.3), size=10, bold=True, color=s["color"],
           align=PP_ALIGN.RIGHT)
        tb(slide, s["body"], Inches(1.02), y + Inches(0.5),
           Inches(11.7), Inches(0.68), size=10.5, color=GREY)

    rect(slide, Inches(0.25), Inches(7.1), Inches(12.83), Inches(0.3),
         RGBColor(0x04, 0x04, 0x1C))
    tb(slide, "JAKALA · Data & AI + Commerce · Mars 2026  ·  "
       "Forberedt for: Tor-André Skeie, CEO · Sport Outlet / Invent Sport",
       Inches(0.4), Inches(7.14), Inches(12.5), Inches(0.22), size=9, color=GREY)


# ─── Main ─────────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()
    slide_cover(prs)
    slide_why_now(prs)
    slide_snapshot(prs)
    slide_problem(prs)
    slide_approach(prs)
    slide_expansion(prs)
    slide_buyer(prs)
    slide_next_steps(prs)

    out = "/Users/jacobskaue/Desktop/jakala-commercial-os/Accounts/sport-outlet/sport-outlet-pitch-2026-03.pptx"
    prs.save(out)
    print(f"✓ Gemt: {out}")
    print(f"  Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
