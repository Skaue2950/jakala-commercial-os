"""
Jakala Slide Generator — Shared Library
Reusable helpers for all /slide-deck generated presentations.

Usage:
    from knowledge.slide_generator import *
    prs = new_prs()
    s = add_slide(prs)
    set_bg(s)
    add_textbox(s, 'Hello', Inches(0.5), Inches(1), Inches(8), Inches(1))
    prs.save('output.pptx')
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────────
# JAKALA DESIGN SYSTEM
# ─────────────────────────────────────────────

BLUE_BRIGHT = RGBColor(0x15, 0x3E, 0xED)   # #153EED — headers, accents, CTAs
BLUE_DARK   = RGBColor(0x02, 0x02, 0x66)   # #020266 — boxes, secondary elements
RED         = RGBColor(0xF6, 0x57, 0x4A)   # #F6574A — risks, urgency
GREEN       = RGBColor(0x92, 0xD0, 0x50)   # #92D050 — positive signals
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)   # body text
GREY        = RGBColor(0xBB, 0xBB, 0xDD)   # muted text
MUTED       = RGBColor(0x88, 0x88, 0xAA)   # very muted
BG_COLOR    = RGBColor(0x08, 0x08, 0x18)   # #080818 — slide background
BG_CARD     = RGBColor(0x10, 0x10, 0x30)   # card/box background

FONT = 'Raleway'
W = Inches(9.84)    # 25cm — Jakala slide width
H = Inches(7.48)    # 19cm — Jakala slide height


# ─────────────────────────────────────────────
# PRESENTATION
# ─────────────────────────────────────────────

def new_prs():
    """Create a new presentation with Jakala dimensions."""
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def add_slide(prs):
    """Add a blank slide to the presentation."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


# ─────────────────────────────────────────────
# BACKGROUND
# ─────────────────────────────────────────────

def set_bg(slide, color=None):
    """Apply Jakala dark background to slide."""
    if color is None:
        color = BG_COLOR
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ─────────────────────────────────────────────
# SHAPES
# ─────────────────────────────────────────────

def add_rect(slide, x, y, w, h, color=BLUE_BRIGHT):
    """Add a filled rectangle with no border."""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_line(slide, x, y, w, color=BLUE_BRIGHT, thickness=2):
    """Add a thin horizontal accent line."""
    shape = slide.shapes.add_shape(1, x, y, w, Pt(thickness))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ─────────────────────────────────────────────
# TEXT
# ─────────────────────────────────────────────

def add_textbox(slide, text, x, y, w, h,
                font_size=14, bold=False, color=None,
                align=PP_ALIGN.LEFT, font=FONT, wrap=True):
    """Add a styled textbox to the slide."""
    if color is None:
        color = WHITE
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=13, bold=False, color=None, space_before=6):
    """Add a paragraph to an existing text frame."""
    if color is None:
        color = WHITE
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


# ─────────────────────────────────────────────
# STANDARD SLIDE COMPONENTS
# ─────────────────────────────────────────────

def slide_header(slide, tag, title, accent=None):
    """Add standard slide header: tag line + title + accent line."""
    if accent is None:
        accent = BLUE_BRIGHT
    add_line(slide, Inches(0.5), Inches(0.55), Inches(8.84), accent)
    add_textbox(slide, tag, Inches(0.5), Inches(0.22), Inches(8), Inches(0.38),
                font_size=10, bold=True, color=accent)
    add_textbox(slide, title, Inches(0.5), Inches(0.68), Inches(8.5), Inches(0.7),
                font_size=26, bold=True, color=WHITE)


def slide_footer(slide, text):
    """Add standard footer bar at bottom of slide."""
    add_rect(slide, Inches(0), H - Inches(0.35), W, Inches(0.35), BLUE_BRIGHT)
    add_textbox(slide, text, Inches(0.3), H - Inches(0.33), Inches(9), Inches(0.3),
                font_size=10, color=WHITE)


def success_bar(slide, text):
    """Add a success criterion bar near the bottom of the slide."""
    add_rect(slide, Inches(0.5), Inches(6.28), Inches(8.84), Inches(0.78), BG_CARD)
    add_textbox(slide, f'Suksesskriterium: {text}',
                Inches(0.7), Inches(6.38), Inches(8.4), Inches(0.58),
                font_size=12, color=WHITE)


def stat_card(slide, x, y, value, label, color=None):
    """Add a stat card (value + label) at given position."""
    if color is None:
        color = BLUE_DARK
    add_rect(slide, x, y, Inches(4.0), Inches(1.25), color)
    add_textbox(slide, value, x + Inches(0.15), y + Pt(6),
                Inches(3.7), Inches(0.65), font_size=32, bold=True, color=BLUE_BRIGHT)
    add_textbox(slide, label, x + Inches(0.15), y + Inches(0.6),
                Inches(3.7), Inches(0.5), font_size=12, color=WHITE)


def bullet_column(slide, x, y, w, h, sections):
    """
    Add a column of bullet sections.
    sections = [('Header', ['bullet 1', 'bullet 2']), ...]
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for header, bullets in sections:
        add_para(tf, header, 13, True, BLUE_BRIGHT, 0 if first else 12)
        first = False
        for b in bullets:
            add_para(tf, f'• {b}', 12, False, WHITE, 3)
    return tb


def period_pill(slide, text, x=None, y=Inches(0.75)):
    """Add a period indicator pill (right-aligned by default)."""
    if x is None:
        x = Inches(7.0)
    add_rect(slide, x, y, Inches(2.34), Inches(0.45), BLUE_DARK)
    add_textbox(slide, text, x + Inches(0.08), y + Pt(4),
                Inches(2.18), Inches(0.38), font_size=11,
                color=BLUE_BRIGHT, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# FULL SLIDE TEMPLATES
# ─────────────────────────────────────────────

def cover_slide(prs, account, subtitle, tag='JAKALA', date='2026'):
    """Generate a standard cover slide."""
    s = add_slide(prs)
    set_bg(s)
    add_rect(s, Inches(0), H - Inches(0.5), W, Inches(0.5), BLUE_BRIGHT)
    add_textbox(s, tag, Inches(0.5), Inches(1.2), Inches(8), Inches(0.45),
                font_size=11, color=MUTED)
    add_textbox(s, account, Inches(0.5), Inches(1.75), Inches(8.5), Inches(2.5),
                font_size=52, bold=True, color=WHITE)
    add_textbox(s, subtitle, Inches(0.5), Inches(4.1), Inches(7), Inches(0.8),
                font_size=24, bold=True, color=BLUE_BRIGHT)
    add_textbox(s, date, Inches(0.5), Inches(5.8), Inches(4), Inches(0.4),
                font_size=12, color=MUTED)
    return s


def phase_slide(prs, tag, title, period, value, accent, left_sections, right_sections, success_text):
    """Generate a standard phase slide with two columns."""
    s = add_slide(prs)
    set_bg(s)
    slide_header(s, tag, title, accent)
    period_pill(s, period)
    add_textbox(s, value, Inches(0.5), Inches(1.45), Inches(6), Inches(0.55),
                font_size=18, bold=True, color=accent)
    bullet_column(s, Inches(0.5), Inches(2.1), Inches(4.2), Inches(4.0), left_sections)
    bullet_column(s, Inches(5.0), Inches(2.1), Inches(4.34), Inches(4.0), right_sections)
    success_bar(s, success_text)
    return s


def next_steps_slide(prs, title, actions, footer_text):
    """
    Generate a next steps slide.
    actions = [('number', 'title', 'description', color), ...]
    """
    s = add_slide(prs)
    set_bg(s)
    slide_header(s, 'NÆSTE SKRIDT', title)
    for i, (num, atitle, desc, color) in enumerate(actions):
        y = Inches(1.6) + i * Inches(1.22)
        add_rect(s, Inches(0.5), y, Inches(0.55), Inches(0.55), color)
        add_textbox(s, num, Inches(0.5), y, Inches(0.55), Inches(0.55),
                    font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, atitle, Inches(1.2), y, Inches(8.1), Inches(0.45),
                    font_size=15, bold=True, color=WHITE)
        add_textbox(s, desc, Inches(1.2), y + Inches(0.42), Inches(8.1), Inches(0.58),
                    font_size=12, color=GREY)
    slide_footer(s, footer_text)
    return s
