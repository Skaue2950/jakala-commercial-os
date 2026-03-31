import os
import re
import json
import datetime
import io
from pathlib import Path
from flask import Flask, request, jsonify, render_template_string, Response, stream_with_context, send_file, session, redirect, url_for
import anthropic
from dotenv import load_dotenv

try:
    import bcrypt
    from models import init_db, SessionLocal, User, Industry, Account, Service, Activation, Signal, Prediction, Action, Meeting, WeeklyCommit, PartnerValidation
    CC_DB_OK = True
except Exception as _cc_err:
    print(f"[CC] Import error: {_cc_err}")
    CC_DB_OK = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

load_dotenv()

app = Flask(__name__)
app.secret_key = os.getenv("SECRET_KEY", "jakala-gtm-os-secret")
APP_PASSWORD = os.getenv("APP_PASSWORD", "JakalaQ12026")
BASE_DIR = Path(__file__).parent.parent  # jakala-commercial-os root
client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
MODEL = "claude-sonnet-4-6"

# ── File helpers ─────────────────────────────────────────────────────────────

def read_file(rel_path):
    p = BASE_DIR / rel_path
    return p.read_text(encoding="utf-8", errors="replace") if p.exists() else None

def write_file(rel_path, content):
    p = BASE_DIR / rel_path
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content, encoding="utf-8")

def get_accounts():
    d = BASE_DIR / "Accounts"
    if not d.exists():
        return []
    return sorted(
        f.name for f in d.iterdir()
        if f.is_dir() and not f.name.startswith(".")
    )

def load_account_files(account_slug):
    files = ["overview.md", "strategy.md", "stakeholders.md", "next-actions.md", "meetings.md"]
    parts = []
    for f in files:
        content = read_file(f"Accounts/{account_slug}/{f}")
        if content:
            parts.append(f"--- {f} ---\n{content}")
    return "\n\n".join(parts) if parts else None

def build_system_prompt():
    today = datetime.date.today().isoformat()
    knowledge = []
    for f in [
        "knowledge/gtm-strategy.md",
        "knowledge/strategy-mapping.md",
        "knowledge/icp-scoring.md",
        "knowledge/deal-scoring.md",
        "knowledge/offerings.md",
        "knowledge/outreach-playbook.md",
        "intelligence/top-opportunities.md",
        "intelligence/pipeline-dashboard.md",
    ]:
        c = read_file(f)
        if c:
            knowledge.append(f"=== {f} ===\n{c}")

    return f"""You are the JAKALA GTM OS — a senior commercial strategy assistant for JAKALA Nordic.

TODAY: {today}
GEOGRAPHIC SCOPE: Denmark, Norway, Sweden ONLY. UK and France accounts are inactive/excluded.
ACTIVE DELIVERY: Maxbo (Norway) — Speedtrain onboarding in progress.

YOUR ROLE:
- Help JAKALA sellers structure commercial thinking, analyze accounts, and plan next actions
- Generate outreach messages (LinkedIn/email), meeting briefs, revenue simulations, pitch summaries
- Score and prioritize pipeline opportunities
- Process meeting notes and suggest file updates
- Answer questions about accounts, signals, strategies and the pipeline

KNOWLEDGE BASE:
{"=" * 60}
{chr(10).join(knowledge)}
{"=" * 60}

SKILLS YOU CAN PERFORM:
- Outreach generation (LinkedIn/email, languages: en/no/da/sv) — max 300 words, soft CTA, peer-to-peer
- Pre-meeting brief (90-second battle card)
- Revenue simulation (3-lever model: data completeness, search quality, AI model input)
- Competitor intelligence mapping
- Account setup and research
- Pipeline forecast (probability-weighted, 7-factor model)
- Signal-to-action conversion
- Morning CCO briefing
- Commercial war room (full situation assessment)
- Meeting note processing — summarize, update next-actions and meetings files

GTM STRATEGIES (always map accounts to one):
1. Data Revenue Unlock — loyalty data, retail media, first-party data monetisation
2. AI Readiness Accelerator — AI ambition outpacing data architecture
3. Commerce Optimization — live ecommerce underperforming (search, discovery, catalog)
4. Experience Transformation — multi-brand, composable architecture, DXP

TONE & OUTPUT:
- Concise, structured, senior consultant style
- Use markdown tables and headers for clarity
- For outreach: no jargon, no bullets in message body, one soft question at end
- Score conservatively — never inflate ICP or deal scores without evidence
- When account context is provided in the conversation, use it directly
- When asked about a specific account not in context, say so and ask if the user wants to load it

RULES:
- Never fabricate company data or signals — only use what is confirmed or clearly labelled as estimated
- If a deal has no named buyer, cap win probability at 25%
- Always identify the GTM strategy and entry offer when reviewing an account
- Prefer updating files over long responses when processing meeting notes"""


SYSTEM_PROMPT = build_system_prompt()


def detect_accounts_in_message(message):
    """Detect account names mentioned in a message and return their slugs."""
    accounts = get_accounts()
    found = []
    msg_lower = message.lower()
    for slug in accounts:
        display = slug.replace("-", " ")
        if display in msg_lower or slug in msg_lower:
            found.append(slug)
    return found


# ── Auth ─────────────────────────────────────────────────────────────────────

LOGIN_HTML = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>JAKALA GTM OS — Login</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap" rel="stylesheet">
<style>
  * { box-sizing: border-box; margin: 0; padding: 0; }
  body { background: #F4F6FC; display: flex; align-items: center; justify-content: center; min-height: 100vh; font-family: 'Inter', -apple-system, sans-serif; }
  .card { background: #fff; border: 1px solid #E4E7F0; border-radius: 16px; padding: 48px 40px; width: 100%; max-width: 380px; box-shadow: 0 4px 24px rgba(21,62,237,.06); }
  .logo { color: #153EED; font-size: 12px; font-weight: 700; letter-spacing: 0.14em; text-transform: uppercase; margin-bottom: 32px; }
  h1 { color: #0F0F1A; font-size: 22px; font-weight: 700; margin-bottom: 8px; letter-spacing: -.02em; }
  p { color: #6B7280; font-size: 13px; margin-bottom: 32px; }
  label { display: block; color: #6B7280; font-size: 11px; font-weight: 600; margin-bottom: 8px; letter-spacing: 0.06em; text-transform: uppercase; }
  input { width: 100%; background: #F8F9FE; border: 1px solid #E4E7F0; border-radius: 8px; color: #0F0F1A; font-size: 14px; padding: 12px 14px; outline: none; transition: border-color 0.15s; font-family: inherit; }
  input:focus { border-color: rgba(21,62,237,.5); }
  button { width: 100%; background: #153EED; border: none; border-radius: 8px; color: #fff; cursor: pointer; font-size: 14px; font-weight: 600; margin-top: 20px; padding: 13px; transition: opacity 0.15s; font-family: inherit; }
  button:hover { opacity: 0.88; }
  .error { background: #FEF2F2; border: 1px solid #FECACA; border-radius: 8px; color: #DC2626; font-size: 13px; margin-bottom: 20px; padding: 10px 14px; }
  .back-link { display: block; text-align: center; margin-top: 20px; font-size: 12px; color: #6B7280; text-decoration: none; }
  .back-link:hover { color: #153EED; }
</style>
</head>
<body>
<div class="card">
  <div class="logo">JAKALA Nordic</div>
  <h1>GTM OS</h1>
  <p>Enter your access password to continue.</p>
  {% if error %}<div class="error">Incorrect password. Try again.</div>{% endif %}
  <form method="POST">
    <label>Password</label>
    <input type="password" name="password" autofocus placeholder="••••••••••••">
    <button type="submit">Sign in →</button>
  </form>
  <a href="/" class="back-link">← Back to overview</a>
</div>
</body>
</html>"""


@app.before_request
def require_login():
    # Landing page + CC routes are always accessible
    if request.path == "/" or request.path.startswith("/cc") or request.path.startswith("/api/cc"):
        return
    if request.endpoint in ("login", "static"):
        return
    if not session.get("authenticated"):
        return redirect(url_for("login"))


@app.route("/login", methods=["GET", "POST"])
def login():
    error = False
    if request.method == "POST":
        if request.form.get("password") == APP_PASSWORD:
            session["authenticated"] = True
            return redirect(url_for("gtm_app"))
        error = True
    return render_template_string(LOGIN_HTML, error=error)


@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("landing"))


# ── API routes ───────────────────────────────────────────────────────────────

@app.route("/api/accounts")
def api_accounts():
    if CC_DB_OK:
        try:
            db = SessionLocal()
            accounts = db.query(Account).all()
            ind_map = {i.id: i for i in db.query(Industry).all()}
            result = []
            for a in accounts:
                acts = [{"service_id": ac.service_id, "stage": ac.stage} for ac in a.activations]
                ind = ind_map.get(a.industry_id)
                result.append({
                    "slug":           a.slug or "",
                    "name":           a.name,
                    "country":        a.country,
                    "icp":            str(int(a.icp_score)) if a.icp_score is not None else "—",
                    "deal":           str(int(a.deal_score)) if a.deal_score is not None else "—",
                    "icp_score":      a.icp_score,
                    "deal_score":     a.deal_score,
                    "pipeline_value": a.pipeline_value,
                    "named_buyer":    a.named_buyer or "",
                    "buyer_role":     a.buyer_role or "",
                    "industry":       ind.name if ind else "Other",
                    "activations":    acts,
                })
            db.close()
            return jsonify(result)
        except Exception as e:
            print(f"[GTM] /api/accounts DB error: {e}")
    # Fallback: file-based
    accounts = get_accounts()
    result = []
    for slug in accounts:
        overview = read_file(f"Accounts/{slug}/overview.md") or ""
        icp = "—"
        m = re.search(r"ICP Score[:\s]+(\d+)/10", overview)
        if m:
            icp = m.group(1)
        deal = "—"
        m = re.search(r"Deal Score[:\s]+(\d+)/10", overview)
        if m:
            deal = m.group(1)
        country = "—"
        for line in overview.splitlines():
            if "Norway" in line or "Norge" in line or "NO)" in line:
                country = "NO"
                break
            if "Denmark" in line or "Danmark" in line or "DK)" in line:
                country = "DK"
                break
            if "Sweden" in line or "Sverige" in line or "SE)" in line:
                country = "SE"
                break
        result.append({"slug": slug, "name": slug.replace("-", " ").title(), "icp": icp, "deal": deal, "country": country})
    return jsonify(result)


@app.route("/api/account/<slug>")
def api_account(slug):
    if CC_DB_OK:
        try:
            db = SessionLocal()
            a = db.query(Account).filter(Account.slug == slug).first()
            if a:
                # Build overview text from DB fields
                overview = (
                    f"# {a.name}\n\n"
                    f"Revenue: {a.revenue or '—'}\n"
                    f"Tech Stack: {a.tech_stack or '—'}\n"
                    f"ICP Score: {a.icp_score}/10\n"
                    f"Deal Score: {a.deal_score}/10\n"
                    f"Named Buyer: {a.named_buyer or 'TBD'}\n"
                    f"Buyer Role: {a.buyer_role or '—'}\n"
                    f"Notes: {a.notes or ''}"
                )
                strategy = a.strategy_text or ""
                stakeholders = a.stakeholders_text or f"Named Buyer: {a.named_buyer or 'TBD'}\nRole: {a.buyer_role or '—'}"
                # Format open/pending actions
                open_actions = [ac for ac in a.actions if ac.status in ("open", "pending")]
                if open_actions:
                    na_lines = [f"# {a.name} — Next Actions\n"]
                    for ac in sorted(open_actions, key=lambda x: (x.priority or "medium")):
                        due = ac.due_date.strftime("%Y-%m-%d") if ac.due_date else "—"
                        na_lines.append(f"- [{ac.priority or 'medium'}] {ac.title} (due: {due})")
                    next_actions = "\n".join(na_lines)
                else:
                    next_actions = ""
                # Format meetings
                if a.meetings:
                    mtg_lines = [f"# {a.name} — Meetings\n"]
                    for m in sorted(a.meetings, key=lambda x: x.date, reverse=True):
                        mtg_lines.append(f"## {m.date.strftime('%Y-%m-%d')} — {m.summary or ''}")
                        if m.participants:
                            mtg_lines.append(f"Participants: {m.participants}")
                        if m.outcome:
                            mtg_lines.append(f"Outcome: {m.outcome}")
                        if m.next_step:
                            mtg_lines.append(f"Next step: {m.next_step}")
                        mtg_lines.append("")
                    meetings = "\n".join(mtg_lines)
                else:
                    meetings = ""
                # Build combined content string (same format as load_account_files)
                parts = []
                if overview:
                    parts.append(f"--- overview.md ---\n{overview}")
                if strategy:
                    parts.append(f"--- strategy.md ---\n{strategy}")
                if stakeholders:
                    parts.append(f"--- stakeholders.md ---\n{stakeholders}")
                if next_actions:
                    parts.append(f"--- next-actions.md ---\n{next_actions}")
                if meetings:
                    parts.append(f"--- meetings.md ---\n{meetings}")
                content = "\n\n".join(parts)
                db.close()
                return jsonify({"slug": slug, "content": content})
            db.close()
        except Exception as e:
            print(f"[GTM] /api/account/{slug} DB error: {e}")
    # Fallback: file-based
    content = load_account_files(slug)
    if not content:
        return jsonify({"error": "Account not found"}), 404
    return jsonify({"slug": slug, "content": content})


@app.route("/api/account/<slug>/file/<filename>", methods=["GET"])
def api_get_file(slug, filename):
    allowed = ["overview.md", "strategy.md", "stakeholders.md", "next-actions.md", "meetings.md"]
    if filename not in allowed:
        return jsonify({"error": "File not allowed"}), 403
    if CC_DB_OK:
        try:
            db = SessionLocal()
            a = db.query(Account).filter(Account.slug == slug).first()
            if a:
                if filename == "overview.md":
                    content = (
                        f"# {a.name}\n\n"
                        f"Revenue: {a.revenue or '—'}\n"
                        f"Tech Stack: {a.tech_stack or '—'}\n"
                        f"ICP Score: {a.icp_score}/10\n"
                        f"Deal Score: {a.deal_score}/10\n"
                        f"Named Buyer: {a.named_buyer or 'TBD'}\n"
                        f"Buyer Role: {a.buyer_role or '—'}\n"
                        f"Notes: {a.notes or ''}"
                    )
                elif filename == "strategy.md":
                    content = a.strategy_text or ""
                elif filename == "stakeholders.md":
                    content = a.stakeholders_text or f"Named Buyer: {a.named_buyer or 'TBD'}\nRole: {a.buyer_role or '—'}"
                elif filename == "next-actions.md":
                    open_actions = [ac for ac in a.actions if ac.status in ("open", "pending")]
                    lines = [f"# {a.name} — Next Actions\n"]
                    for ac in sorted(open_actions, key=lambda x: (x.priority or "medium")):
                        due = ac.due_date.strftime("%Y-%m-%d") if ac.due_date else "—"
                        lines.append(f"- [{ac.priority or 'medium'}] {ac.title} (due: {due})")
                    content = "\n".join(lines)
                elif filename == "meetings.md":
                    lines = [f"# {a.name} — Meetings\n"]
                    for m in sorted(a.meetings, key=lambda x: x.date, reverse=True):
                        lines.append(f"## {m.date.strftime('%Y-%m-%d')} — {m.summary or ''}")
                        if m.participants:
                            lines.append(f"Participants: {m.participants}")
                        if m.outcome:
                            lines.append(f"Outcome: {m.outcome}")
                        if m.next_step:
                            lines.append(f"Next step: {m.next_step}")
                        lines.append("")
                    content = "\n".join(lines)
                else:
                    content = None
                db.close()
                if content is None:
                    return jsonify({"error": "File not found"}), 404
                return jsonify({"content": content})
            db.close()
        except Exception as e:
            print(f"[GTM] api_get_file {slug}/{filename} DB error: {e}")
    # Fallback: file-based
    content = read_file(f"Accounts/{slug}/{filename}")
    if content is None:
        return jsonify({"error": "File not found"}), 404
    return jsonify({"content": content})


@app.route("/api/account/<slug>/file/<filename>", methods=["POST"])
def api_save_file(slug, filename):
    allowed = ["overview.md", "strategy.md", "stakeholders.md", "next-actions.md", "meetings.md"]
    if filename not in allowed:
        return jsonify({"error": "File not allowed"}), 403
    data = request.get_json()
    content = data.get("content", "")
    if CC_DB_OK:
        try:
            db = SessionLocal()
            a = db.query(Account).filter(Account.slug == slug).first()
            if a:
                if filename == "overview.md":
                    a.notes = content
                elif filename == "strategy.md":
                    a.strategy_text = content
                elif filename == "stakeholders.md":
                    a.stakeholders_text = content
                elif filename == "next-actions.md":
                    # Replace all open actions with a simple bulk note
                    db.query(Action).filter(
                        Action.account_id == a.id,
                        Action.status == "open"
                    ).delete(synchronize_session=False)
                    if content.strip():
                        new_action = Action(
                            account_id=a.id,
                            title="Next actions (imported)",
                            description=content,
                            status="open",
                            priority="medium",
                            action_type="follow-up",
                        )
                        db.add(new_action)
                elif filename == "meetings.md":
                    # Append-only: create a new meeting record from the content
                    new_meeting = Meeting(
                        account_id=a.id,
                        country=a.country,
                        date=datetime.datetime.utcnow(),
                        summary=content[:500],
                        outcome="neutral",
                    )
                    db.add(new_meeting)
                db.commit()
                db.close()
                return jsonify({"ok": True})
            db.close()
        except Exception as e:
            print(f"[GTM] api_save_file {slug}/{filename} DB error: {e}")
    # Fallback: file-based
    write_file(f"Accounts/{slug}/{filename}", content)
    return jsonify({"ok": True})


@app.route("/api/chat", methods=["POST"])
def api_chat():
    data = request.get_json()
    messages = data.get("messages", [])
    account_slug = data.get("account")

    injected_context = ""

    if account_slug:
        account_content = load_account_files(account_slug)
        if account_content:
            injected_context += f"\n\nACCOUNT CONTEXT LOADED — {account_slug.replace('-', ' ').title()}:\n{account_content}"

    if messages:
        last_msg = messages[-1].get("content", "")
        detected = detect_accounts_in_message(last_msg)
        for slug in detected:
            if slug != account_slug:
                content = load_account_files(slug)
                if content:
                    injected_context += f"\n\nACCOUNT CONTEXT — {slug.replace('-', ' ').title()}:\n{content}"

    api_messages = []
    for i, msg in enumerate(messages):
        role = msg.get("role")
        content = msg.get("content", "")
        if i == len(messages) - 1 and role == "user" and injected_context:
            content = f"{injected_context}\n\n---\n\nUser question: {content}"
        api_messages.append({"role": role, "content": content})

    def generate():
        try:
            with client.messages.stream(
                model=MODEL,
                max_tokens=4096,
                system=SYSTEM_PROMPT,
                messages=api_messages,
            ) as stream:
                for text in stream.text_stream:
                    yield f"data: {json.dumps({'text': text})}\n\n"
        except anthropic.AuthenticationError:
            yield f"data: {json.dumps({'text': '[ERROR] API key error — check ANTHROPIC_API_KEY in Railway environment variables.'})}\n\n"
        except anthropic.RateLimitError:
            yield f"data: {json.dumps({'text': '[ERROR] Rate limit reached — please wait a moment and try again.'})}\n\n"
        except anthropic.APIStatusError as e:
            yield f"data: {json.dumps({'text': f'[ERROR] API error {e.status_code}: {e.message}'})}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'text': f'[ERROR] Unexpected error: {str(e)}'})}\n\n"
        finally:
            yield "data: [DONE]\n\n"

    return Response(stream_with_context(generate()), mimetype="text/event-stream")


@app.route("/api/process-notes", methods=["POST"])
def api_process_notes():
    data = request.get_json()
    account_slug = data.get("account")
    notes = data.get("notes", "")

    if not account_slug or not notes:
        return jsonify({"error": "account and notes required"}), 400

    account_content = load_account_files(account_slug)
    today = datetime.date.today().isoformat()

    prompt = f"""You are processing meeting notes for JAKALA account: {account_slug.replace('-', ' ').title()}

EXISTING ACCOUNT FILES:
{account_content or '(no existing files)'}

MEETING NOTES TO PROCESS:
{notes}

TODAY: {today}

Your task — return a JSON object with EXACTLY these keys:
{{
  "summary": "2-3 sentence summary of the meeting",
  "meeting_entry": "Full markdown entry to append to meetings.md (include date, attendees if mentioned, key points, agreed actions)",
  "next_actions_updated": "Complete updated content for next-actions.md — incorporate new actions from this meeting, mark completed items if applicable",
  "key_insight": "One sentence: the most important commercial insight from this meeting"
}}

Rules:
- next_actions_updated should be the FULL new content of the file (not just the additions)
- Include today's date ({today}) in the meeting entry header
- Keep the same markdown format as existing files
- Prioritise actions by commercial impact
- Return ONLY valid JSON — no markdown fences, no explanation"""

    response = client.messages.create(
        model=MODEL,
        max_tokens=3000,
        messages=[{"role": "user", "content": prompt}],
    )

    raw = response.content[0].text.strip()
    raw = re.sub(r"^```json\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)

    try:
        result = json.loads(raw)
    except json.JSONDecodeError:
        return jsonify({"error": "Failed to parse AI response", "raw": raw}), 500

    return jsonify(result)


@app.route("/api/save-notes", methods=["POST"])
def api_save_notes():
    data = request.get_json()
    slug = data.get("account")
    meeting_entry = data.get("meeting_entry", "")
    next_actions = data.get("next_actions_updated", "")
    key_insight = data.get("key_insight", "")

    if not slug:
        return jsonify({"error": "account required"}), 400

    today = datetime.date.today().isoformat()

    if CC_DB_OK:
        try:
            db = SessionLocal()
            a = db.query(Account).filter(Account.slug == slug).first()
            if a:
                # Insert new meeting record
                if meeting_entry:
                    summary_text = meeting_entry[:500] if meeting_entry else key_insight
                    new_meeting = Meeting(
                        account_id=a.id,
                        country=a.country,
                        date=datetime.datetime.utcnow(),
                        summary=summary_text,
                        outcome="neutral",
                        next_step=key_insight or None,
                    )
                    db.add(new_meeting)
                    a.last_activity = datetime.datetime.utcnow()
                # Replace open actions with new set
                if next_actions:
                    db.query(Action).filter(
                        Action.account_id == a.id,
                        Action.status == "open"
                    ).delete(synchronize_session=False)
                    for line in next_actions.splitlines():
                        line = line.strip()
                        if line and (line.startswith("-") or line.startswith("*") or re.match(r"^\d+\.", line)):
                            title = re.sub(r"^[-*\d.]+\s*", "", line).strip()
                            if title:
                                new_action = Action(
                                    account_id=a.id,
                                    title=title[:300],
                                    status="open",
                                    priority="medium",
                                    action_type="follow-up",
                                )
                                db.add(new_action)
                db.commit()
                db.close()
                return jsonify({"ok": True})
            db.close()
        except Exception as e:
            print(f"[GTM] api_save_notes DB error: {e}")

    # Fallback: file-based
    existing_meetings = read_file(f"Accounts/{slug}/meetings.md") or f"# {slug.replace('-', ' ').title()} — Meetings\n\nLast updated: {today}\n\n---\n"
    existing_meetings = re.sub(r"Last updated: \d{4}-\d{2}-\d{2}", f"Last updated: {today}", existing_meetings)
    existing_meetings = existing_meetings + f"\n\n---\n\n{meeting_entry}"
    write_file(f"Accounts/{slug}/meetings.md", existing_meetings)

    if next_actions:
        next_actions = re.sub(r"Last updated: \d{4}-\d{2}-\d{2}", f"Last updated: {today}", next_actions)
        write_file(f"Accounts/{slug}/next-actions.md", next_actions)

    return jsonify({"ok": True})


# ── PPTX generation ──────────────────────────────────────────────────────────

BLUE  = RGBColor(0x15,0x3E,0xED) if PPTX_OK else None
NAVY  = RGBColor(0x02,0x02,0x66) if PPTX_OK else None
RED   = RGBColor(0xF6,0x57,0x4A) if PPTX_OK else None
GREEN = RGBColor(0x00,0xD4,0xA0) if PPTX_OK else None
WHITE = RGBColor(0xFF,0xFF,0xFF) if PPTX_OK else None
GREY  = RGBColor(0xBB,0xBB,0xDD) if PPTX_OK else None
MUTED = RGBColor(0x88,0x88,0xAA) if PPTX_OK else None
BG    = RGBColor(0x04,0x04,0x0F) if PPTX_OK else None
CARD  = RGBColor(0x0A,0x0A,0x22) if PPTX_OK else None
W = Inches(9.84) if PPTX_OK else None
H = Inches(7.48) if PPTX_OK else None
FONT  = "Calibri"


def _prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def _rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _txt(slide, text, x, y, w, h, size=14, bold=False, color=None, align=None):
    if color is None: color = WHITE
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    if align: p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.name = FONT; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    return tb

def _add_para(tf, text, size=13, bold=False, color=None, space=6):
    if color is None: color = WHITE
    p = tf.add_paragraph(); p.space_before = Pt(space)
    r = p.add_run(); r.text = str(text)
    r.font.name = FONT; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color

def _header(slide, tag, title):
    _rect(slide, Inches(0), Inches(0), W, Inches(0.08), BLUE)
    _txt(slide, tag,   Inches(0.5), Inches(0.14), Inches(8), Inches(0.35),
         size=9, bold=True, color=BLUE)
    _txt(slide, title, Inches(0.5), Inches(0.54), Inches(8.5), Inches(0.75),
         size=24, bold=True)

def _footer(slide, text):
    _rect(slide, Inches(0), H - Inches(0.32), W, Inches(0.32), BLUE)
    _txt(slide, text, Inches(0.3), H - Inches(0.30), Inches(9), Inches(0.28),
         size=9, color=WHITE)

def _bullet_col(slide, x, y, w, h, sections):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for header, bullets in sections:
        _add_para(tf, header, 12, True, BLUE, 0 if first else 10)
        first = False
        for b in bullets:
            _add_para(tf, f"• {b}", 11, False, WHITE, 3)

def build_account_deck(account_name, data):
    """Build a 5-slide discovery deck from Claude-generated JSON data."""
    prs = _prs()

    # Slide 1 — Cover
    s = _slide(prs); _bg(s)
    _rect(s, Inches(0), H - Inches(0.5), W, Inches(0.5), BLUE)
    _rect(s, Inches(0), Inches(0), Inches(0.08), H, BLUE)
    _txt(s, "JAKALA COMMERCIAL", Inches(0.3), Inches(1.2), Inches(8), Inches(0.4),
         size=10, bold=True, color=MUTED)
    _txt(s, account_name, Inches(0.3), Inches(1.75), Inches(8.5), Inches(2.2),
         size=48, bold=True)
    _txt(s, data.get("subtitle","Commercial Discovery"), Inches(0.3), Inches(3.9),
         Inches(7), Inches(0.65), size=22, bold=True, color=BLUE)
    _txt(s, data.get("date", datetime.date.today().isoformat()),
         Inches(0.3), Inches(5.6), Inches(4), Inches(0.35), size=11, color=MUTED)

    # Slide 2 — Why now
    s = _slide(prs); _bg(s)
    _header(s, "COMMERCIAL CONTEXT", data.get("context_title","Why Now"))
    _bullet_col(s, Inches(0.5), Inches(1.5), Inches(4.0), Inches(5.2),
                [("Signals & Timing", data.get("context_points",[]))])
    _bullet_col(s, Inches(5.0), Inches(1.5), Inches(4.34), Inches(5.2),
                [("Business Pressure", data.get("pressure_points",[]))])
    _footer(s, f"JAKALA — {account_name} — Confidential")

    # Slide 3 — GTM Strategy
    s = _slide(prs); _bg(s)
    _header(s, "GTM STRATEGY", data.get("gtm_title","Our Entry Approach"))
    _rect(s, Inches(0.5), Inches(1.5), Inches(3.8), Inches(0.55), NAVY)
    _txt(s, data.get("gtm_strategy","—"), Inches(0.6), Inches(1.52), Inches(3.6),
         Inches(0.5), size=14, bold=True, color=BLUE)
    _bullet_col(s, Inches(0.5), Inches(2.2), Inches(4.0), Inches(4.5),
                [("Entry Offer", data.get("entry_points",[])),
                 ("Expansion Path", data.get("expansion_points",[]))])
    _bullet_col(s, Inches(5.0), Inches(1.5), Inches(4.34), Inches(5.2),
                [("Likely Buyer", data.get("buyer_points",[])),
                 ("Why JAKALA Wins", data.get("why_jakala",[]))])
    _footer(s, f"JAKALA — {account_name} — Confidential")

    # Slide 4 — Business case
    s = _slide(prs); _bg(s)
    _header(s, "THE BUSINESS CASE", data.get("value_title","Value & Impact"))
    # Stat cards
    stats = data.get("stats", [])
    for i, st in enumerate(stats[:3]):
        x = Inches(0.5) + i * Inches(3.1)
        _rect(s, x, Inches(1.5), Inches(2.8), Inches(1.2), NAVY)
        _txt(s, st.get("value","—"), x + Inches(0.15), Inches(1.55),
             Inches(2.5), Inches(0.65), size=30, bold=True, color=BLUE)
        _txt(s, st.get("label",""), x + Inches(0.15), Inches(2.1),
             Inches(2.5), Inches(0.45), size=11, color=WHITE)
    _bullet_col(s, Inches(0.5), Inches(2.85), Inches(4.0), Inches(4.0),
                [("Impact Framing", data.get("value_points",[]))])
    _bullet_col(s, Inches(5.0), Inches(2.85), Inches(4.34), Inches(4.0),
                [("Risk of Inaction", data.get("risk_points",[]))])
    _footer(s, f"JAKALA — {account_name} — Confidential")

    # Slide 5 — Next steps
    s = _slide(prs); _bg(s)
    _header(s, "NEXT STEPS", data.get("next_title","Proposed Actions"))
    next_steps = data.get("next_steps",[])
    colors = [BLUE, GREEN, RED]
    for i, step in enumerate(next_steps[:4]):
        y = Inches(1.55) + i * Inches(1.2)
        col = colors[i % len(colors)]
        _rect(s, Inches(0.5), y, Inches(0.52), Inches(0.52), col)
        _txt(s, str(i+1), Inches(0.5), y, Inches(0.52), Inches(0.52),
             size=20, bold=True, align=PP_ALIGN.CENTER)
        _txt(s, step.get("title",""), Inches(1.15), y, Inches(8.1), Inches(0.42),
             size=14, bold=True)
        _txt(s, step.get("desc",""), Inches(1.15), y + Inches(0.4), Inches(8.1),
             Inches(0.55), size=11, color=GREY)
    _footer(s, f"JAKALA — {account_name} — Confidential")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


@app.route("/api/generate-deck/<slug>", methods=["POST"])
def api_generate_deck(slug):
    import traceback
    if not PPTX_OK:
        return jsonify({"error": "python-pptx not installed"}), 500
    try:
        return _do_generate_deck(slug)
    except Exception as e:
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500

def _do_generate_deck(slug):
    account_content = load_account_files(slug)
    account_name = slug.replace("-", " ").title()
    today = datetime.date.today().isoformat()

    prompt = f"""You are building a 5-slide commercial discovery deck for JAKALA about: {account_name}

ACCOUNT DATA:
{account_content or '(no files)'}

TODAY: {today}

Return ONLY valid JSON (no markdown fences) with this exact structure:
{{
  "subtitle": "Commercial Discovery — [GTM Strategy Name]",
  "date": "{today}",
  "context_title": "Why {account_name}, Why Now",
  "context_points": ["signal 1", "signal 2", "signal 3"],
  "pressure_points": ["business pressure 1", "pressure 2", "pressure 3"],
  "gtm_title": "Our Entry Approach",
  "gtm_strategy": "[one of: Data Revenue Unlock / AI Readiness Accelerator / Commerce Optimization / Experience Transformation]",
  "entry_points": ["entry offer detail 1", "detail 2"],
  "expansion_points": ["expansion 1", "expansion 2"],
  "buyer_points": ["Name — Title", "why they care"],
  "why_jakala": ["differentiator 1", "differentiator 2"],
  "value_title": "The Business Case",
  "stats": [
    {{"value": "€Xm", "label": "Estimated revenue impact"}},
    {{"value": "X/10", "label": "Deal score"}},
    {{"value": "Xwks", "label": "Time to first value"}}
  ],
  "value_points": ["value framing 1", "value 2", "value 3"],
  "risk_points": ["risk of inaction 1", "risk 2"],
  "next_title": "Proposed Next Steps",
  "next_steps": [
    {{"title": "action 1", "desc": "description"}},
    {{"title": "action 2", "desc": "description"}},
    {{"title": "action 3", "desc": "description"}}
  ]
}}

Rules:
- Be specific — use real names, real signals from account data
- Keep each bullet under 10 words
- GTM strategy must match the account's best fit
- Return ONLY the JSON object"""

    response = client.messages.create(
        model=MODEL,
        max_tokens=2000,
        messages=[{"role": "user", "content": prompt}],
    )

    raw = response.content[0].text.strip()
    raw = re.sub(r"^```json\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)

    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        return jsonify({"error": "Failed to parse AI response", "raw": raw}), 500

    buf = build_account_deck(account_name, data)
    filename = f"JAKALA-{slug}-discovery-{today}.pptx"
    return send_file(
        buf,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=filename
    )


# ── Partnership API ──────────────────────────────────────────────────────────

@app.route("/api/validate-partner", methods=["POST"])
def api_validate_partner():
    data      = request.get_json()
    partner   = data.get("partner", "").strip()
    market    = data.get("market", "Nordic")
    context   = data.get("context", "")
    if not partner:
        return jsonify({"error": "Partner name required"}), 400

    gtm      = read_file("knowledge/gtm-strategy.md") or ""
    services = read_file("knowledge/jakala-services.md") or ""
    offerings= read_file("knowledge/offerings.md") or ""
    today    = datetime.date.today().isoformat()

    prompt = f"""You are a senior JAKALA partnership strategist evaluating a potential partner.

TODAY: {today}
JAKALA CONTEXT:
- Nordic GTM focus: Denmark, Norway, Sweden
- Core services: Commerce Experience, Product Experience / Speedtrain PIM, DXP, Generative AI, CDP/CRM, BI, Data Architecture, Pricing, Hello Growth (SaaS)
- GTM strategies: Data Revenue Unlock · AI Readiness Accelerator · Commerce Optimization · Experience Transformation
- Buyer personas: CTO, CDO, CMO, Head of Ecommerce, Head of Data
- Active delivery: Maxbo (Norway) — Speedtrain onboarding

GTM STRATEGY SUMMARY:
{gtm[:1000]}

SERVICES:
{services[:1200]}

PARTNER TO EVALUATE: {partner}
MARKET FOCUS: {market}
ADDITIONAL CONTEXT: {context or "None provided"}

Evaluate this partnership across 5 dimensions and return a structured assessment.

Return ONLY valid JSON (no markdown fences):
{{
  "partner": "{partner}",
  "market": "{market}",
  "verdict": "STRONG FIT" | "POTENTIAL FIT" | "WEAK FIT" | "NOT RECOMMENDED",
  "verdict_reason": "One sentence summary of the verdict",
  "overall_score": 7,
  "dimensions": [
    {{
      "name": "GTM Fit",
      "score": 8,
      "max": 10,
      "rationale": "Does this partner serve the same buyer personas and complement JAKALA's GTM strategies?",
      "finding": "2-3 sentence specific assessment"
    }},
    {{
      "name": "Revenue Potential",
      "score": 7,
      "max": 10,
      "rationale": "Joint pipeline, referrals, co-delivery, or new market access potential",
      "finding": "2-3 sentence specific assessment"
    }},
    {{
      "name": "Market Positioning",
      "score": 8,
      "max": 10,
      "rationale": "Does the partnership enhance JAKALA credibility and reach in {market}?",
      "finding": "2-3 sentence specific assessment"
    }},
    {{
      "name": "Channel Conflict Risk",
      "score": 6,
      "max": 10,
      "rationale": "Risk of competing with JAKALA for the same deals or services (higher score = lower risk)",
      "finding": "2-3 sentence specific assessment"
    }},
    {{
      "name": "Activation Speed",
      "score": 7,
      "max": 10,
      "rationale": "How quickly can this partnership generate tangible pipeline or revenue?",
      "finding": "2-3 sentence specific assessment"
    }}
  ],
  "gtm_match": ["Data Revenue Unlock", "Commerce Optimization"],
  "buyer_overlap": ["CTO", "Head of Ecommerce"],
  "joint_offer": "One specific joint offer or entry point that combines both companies",
  "target_accounts": ["Account 1", "Account 2", "Account 3"],
  "risks": ["Risk 1", "Risk 2"],
  "first_step": "Concrete first action to activate this partnership",
  "partner_type": "Technology Vendor" | "System Integrator" | "Consulting Firm" | "Platform Vendor" | "Agency" | "Data Provider" | "Other"
}}

Rules:
- Be specific — if you know what this company does, use that knowledge
- Score honestly — a weak fit should score 3-4, not 6-7
- Target accounts should be from the Nordic pipeline (NO/SE/DK) if possible
- Return ONLY the JSON"""

    response = client.messages.create(
        model=MODEL,
        max_tokens=2000,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = re.sub(r"^```json\s*", "", response.content[0].text.strip())
    raw = re.sub(r"\s*```$", "", raw)
    try:
        result = json.loads(raw)
    except json.JSONDecodeError:
        return jsonify({"error": "AI parse failed", "raw": raw}), 500

    # Save result — DB preferred, fallback to file
    save_partner_result(partner, market, context, result, today)
    return jsonify(result)


def save_partner_result(partner, market, context, data, today):
    if CC_DB_OK:
        try:
            db = SessionLocal()
            pv = PartnerValidation(
                company_name=partner,
                market=market,
                context=context or "",
                verdict=data.get("verdict", ""),
                score=data.get("overall_score"),
                findings=json.dumps(data),
            )
            db.add(pv)
            db.commit()
            db.close()
            return
        except Exception as e:
            print(f"[GTM] save_partner_result DB error: {e}")
    # Fallback: write to file
    p = BASE_DIR / "intelligence" / "partnerships"
    p.mkdir(parents=True, exist_ok=True)
    slug = re.sub(r"[^a-z0-9]+", "-", partner.lower()).strip("-")
    fname = p / f"{slug}-{today}.md"
    lines = [
        f"# Partnership Assessment — {partner}",
        f"Date: {today} | Market: {market} | Verdict: {data.get('verdict','—')}",
        "",
        f"**Overall Score:** {data.get('overall_score','—')}/10",
        f"**Type:** {data.get('partner_type','—')}",
        f"**Verdict:** {data.get('verdict_reason','—')}",
        "",
        "## Dimension Scores",
    ]
    for d in data.get("dimensions", []):
        lines.append(f"- **{d['name']}:** {d['score']}/{d['max']} — {d['finding']}")
    lines += [
        "",
        f"## Joint Offer\n{data.get('joint_offer','—')}",
        "",
        f"## GTM Match\n{', '.join(data.get('gtm_match',[]))}",
        "",
        f"## Target Accounts\n{', '.join(data.get('target_accounts',[]))}",
        "",
        f"## First Step\n{data.get('first_step','—')}",
        "",
        "## Risks\n" + "\n".join(f"- {r}" for r in data.get("risks", [])),
    ]
    fname.write_text("\n".join(lines), encoding="utf-8")


@app.route("/api/partner-history")
def api_partner_history():
    if CC_DB_OK:
        try:
            db = SessionLocal()
            rows = db.query(PartnerValidation).order_by(PartnerValidation.created_at.desc()).limit(20).all()
            db.close()
            if rows:
                return jsonify([
                    {
                        "file":    f"{re.sub(r'[^a-z0-9]+', '-', r.company_name.lower()).strip('-')}-{r.created_at.strftime('%Y-%m-%d') if r.created_at else 'unknown'}",
                        "name":    r.company_name,
                        "verdict": r.verdict or "—",
                        "score":   str(r.score) if r.score is not None else "—",
                        "market":  r.market or "—",
                    }
                    for r in rows
                ])
            db.close()
        except Exception as e:
            print(f"[GTM] /api/partner-history DB error: {e}")
    # Fallback: file-based
    p = BASE_DIR / "intelligence" / "partnerships"
    if not p.exists():
        return jsonify([])
    results = []
    for f in sorted(p.glob("*.md"), reverse=True)[:20]:
        content = f.read_text(encoding="utf-8")
        verdict_m = re.search(r"Verdict: (.+)", content)
        score_m   = re.search(r"Overall Score:\*\* (\d+)/10", content)
        market_m  = re.search(r"Market: (\w+)", content)
        name_m    = re.search(r"# Partnership Assessment — (.+)", content)
        results.append({
            "file":    f.stem,
            "name":    name_m.group(1).strip()  if name_m    else f.stem,
            "verdict": verdict_m.group(1).strip() if verdict_m else "—",
            "score":   score_m.group(1)          if score_m   else "—",
            "market":  market_m.group(1)         if market_m  else "—",
        })
    return jsonify(results)


@app.route("/api/monthly-partnerships")
def api_monthly_partnerships():
    p = BASE_DIR / "intelligence" / "partnerships"
    if not p.exists():
        return jsonify({"content": None})
    files = sorted(p.glob("monthly-*.md"), reverse=True)
    if not files:
        return jsonify({"content": None})
    return jsonify({"content": files[0].read_text(encoding="utf-8")})


# ── Signals API ──────────────────────────────────────────────────────────────

HARDCODED_SIGNALS = [
    {"company":"Sport Outlet","text":"CTO + CDO both vacant March 2026. Entry via CEO Tor-André Skeie. Vacancy = budget already approved.","tag":"URGENT","tagColor":"red","icon":'<span style="width:8px;height:8px;background:#DC2626;border-radius:50%;display:inline-block;"></span>',"slug":"sport-outlet","type":"leadership","urgency":"critical"},
    {"company":"Trumf (NorgesGruppen)","text":"Rikke Etholm-Idsøe — new Commercial Director in newly created role. 90-day honeymoon window open now.","tag":"90-DAY WINDOW","tagColor":"red","icon":'<i data-lucide="zap" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i>',"slug":"trumf","type":"leadership","urgency":"critical"},
    {"company":"Vinmonopolet","text":"Espen Terland new CDO (ex-XXL 15 years). Agenda not set — honeymoon phase. Ideal discovery entry.","tag":"NEW EXEC","tagColor":"amber","icon":'<i data-lucide="plus-circle" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i>',"slug":"vinmonopolet","type":"leadership","urgency":"warning"},
    {"company":"Skeidar","text":"\"Best furniture portal in the Nordics\" declared publicly by CEO. CIO Sujit Nath confirmed buyer.","tag":"NAMED BUYER","tagColor":"amber","icon":'<i data-lucide="building-2" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i>',"slug":"skeidar","type":"market","urgency":"warning"},
    {"company":"BI Handelshøyskolen","text":"Rector Karen Spens leaving August 2026. Institution in transition — new leadership will reset priorities.","tag":"TRANSITION","tagColor":"amber","icon":'<i data-lucide="graduation-cap" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i>',"slug":"bi-handelshoyskolen","type":"leadership","urgency":"warning"},
    {"company":"GANT Norway","text":"New CEO Fredrik Malm + IMPACT Commerce new ecom partner (Feb 2026). Integration phase = JAKALA entry.","tag":"NEW CEO","tagColor":"amber","icon":'<i data-lucide="user" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i>',"slug":"gant-norway","type":"leadership","urgency":"warning"},
    {"company":"H&M Sweden","text":"ICP 9/10 · Deal 9/10 · €900K unweighted. No named buyer confirmed yet. Largest untouched opportunity.","tag":"TOP PRIORITY","tagColor":"blue","icon":'<i data-lucide="gem" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i>',"slug":"hm","type":"market","urgency":"info"},
    {"company":"Matas","text":"ICP 9/10 · Deal 9/10 · €700K. AI Readiness entry. Loyalty data + personalisation play — Matas More programme.","tag":"HIGH VALUE","tagColor":"blue","icon":'<i data-lucide="lightbulb" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i>',"slug":"matas","type":"market","urgency":"info"},
]


@app.route("/api/signals")
def api_signals():
    if CC_DB_OK:
        try:
            db = SessionLocal()
            db_signals = db.query(Signal).filter(Signal.is_active == True).order_by(Signal.severity, Signal.date.desc()).limit(20).all()
            db.close()
            if db_signals:
                severity_tag = {"critical": "URGENT", "warning": "ACTION", "info": "RADAR"}
                severity_color = {"critical": "red", "warning": "amber", "info": "blue"}
                result = []
                for s in db_signals:
                    result.append({
                        "company":  s.title,
                        "context":  s.description or "",
                        "text":     s.description or "",
                        "tag":      severity_tag.get(s.severity, "RADAR"),
                        "tagColor": severity_color.get(s.severity, "blue"),
                        "icon":     '<i data-lucide="zap" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i>',
                        "slug":     "",
                        "type":     s.signal_type or "market",
                        "urgency":  s.severity or "info",
                        "tags":     [s.signal_type or "market", s.vertical or "general"],
                    })
                return jsonify({"signals": result})
        except Exception as e:
            print(f"[GTM] /api/signals DB error: {e}")
    # Fallback: hardcoded + daily-leads files
    hot = HARDCODED_SIGNALS[:]
    extra = []
    intel = BASE_DIR / "intelligence" / "daily-leads"
    if intel.exists():
        for f in sorted(intel.iterdir(), reverse=True)[:1]:
            if f.suffix == ".md":
                for line in f.read_text(encoding="utf-8").splitlines():
                    if line.startswith("## ") and len(line) > 4:
                        extra.append({"company": line[3:].strip(), "text": f"From daily radar {f.stem}", "tag":"RADAR","tagColor":"blue","icon":'<i data-lucide="radio" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i>',"slug":"","type":"market","urgency":"info"})
    return jsonify({"signals": hot + extra[:4]})


# ── Live Dashboard API ────────────────────────────────────────────────────────

@app.route("/api/dashboard-live")
def api_dashboard_live():
    if CC_DB_OK:
        try:
            db = SessionLocal()
            accounts = db.query(Account).all()

            total_pipeline = sum(a.pipeline_value or 0 for a in accounts)
            named_buyers   = sum(1 for a in accounts if a.named_buyer and a.named_buyer not in ("TBD", ""))
            active_accounts = len(accounts)
            forecast_val   = sum((a.pipeline_value or 0) * (a.win_probability or 0) / 100.0 for a in accounts)

            # Format pipeline total
            if total_pipeline >= 1_000_000:
                pipeline_str = f"€{total_pipeline/1_000_000:.1f}M"
            else:
                pipeline_str = f"€{int(total_pipeline/1000)}K"

            if forecast_val >= 1_000_000:
                forecast_str = f"€{forecast_val/1_000_000:.1f}M"
            else:
                forecast_str = f"€{int(forecast_val/1000)}K"

            # Top 5 deals by pipeline_value
            top_accts = sorted([a for a in accounts if a.pipeline_value], key=lambda x: x.pipeline_value, reverse=True)[:6]
            top_deals = []
            for i, a in enumerate(top_accts, 1):
                win_pct = f"{int(a.win_probability or 0)}%"
                weighted_val = (a.pipeline_value or 0) * (a.win_probability or 0) / 100.0
                if weighted_val >= 1_000_000:
                    weighted_str = f"€{weighted_val/1_000_000:.1f}M"
                else:
                    weighted_str = f"€{int(weighted_val/1000)}K"
                entry_val = a.pipeline_value or 0
                if entry_val >= 1_000_000:
                    entry_str = f"€{entry_val/1_000_000:.1f}M"
                else:
                    entry_str = f"€{int(entry_val/1000)}K"
                top_deals.append({
                    "rank":      str(i),
                    "name":      a.name,
                    "country":   a.country,
                    "offering":  a.tech_stack[:30] if a.tech_stack else "—",
                    "icp":       str(int(a.icp_score)) if a.icp_score else "—",
                    "win_pct":   win_pct,
                    "entry_val": entry_str,
                    "weighted":  weighted_str,
                    "buyer":     a.named_buyer or "TBD",
                    "status":    a.deal_stage or "identified",
                    "slug":      a.slug or "",
                    "days_stale": None,
                })

            # Country split
            country_totals = {}
            country_counts = {}
            for a in accounts:
                c = a.country or "Unknown"
                country_totals[c] = country_totals.get(c, 0) + (a.pipeline_value or 0)
                country_counts[c] = country_counts.get(c, 0) + 1
            country_name_map = {"NO": "Norway", "SE": "Sweden", "DK": "Denmark"}
            country_split = []
            for code, total in sorted(country_totals.items(), key=lambda x: -x[1]):
                if total >= 1_000_000:
                    pip_str = f"€{total/1_000_000:.1f}M"
                else:
                    pip_str = f"€{int(total/1000)}K"
                country_split.append({
                    "country":  country_name_map.get(code, code),
                    "pipeline": pip_str,
                    "accounts": country_counts.get(code, 0),
                })

            # Top priority account with named buyer
            priority = None
            for d in top_deals:
                if d["buyer"] and d["buyer"] != "TBD" and "TBD" not in d["buyer"]:
                    buyer_name = re.split(r'\(', d["buyer"])[0].strip()
                    priority = {
                        "name":     buyer_name,
                        "company":  d["name"],
                        "country":  d["country"],
                        "win_pct":  d["win_pct"],
                        "weighted": d["weighted"],
                        "offering": d["offering"],
                        "status":   d["status"],
                        "slug":     d["slug"],
                    }
                    break

            # Top 3 active signals
            db_sigs = db.query(Signal).filter(Signal.is_active == True).order_by(
                Signal.severity, Signal.date.desc()
            ).limit(3).all()
            live_signals = [
                {"company": s.title, "text": (s.description or "")[:180], "date": s.date.strftime("%Y-%m-%d") if s.date else ""}
                for s in db_sigs
            ]

            # Status determination
            if forecast_val < 200_000:
                status = "RED"
            elif forecast_val < 500_000:
                status = "AMBER"
            else:
                status = "GREEN"

            db.close()
            return jsonify({
                "pipeline_total":    pipeline_str,
                "pipeline_weighted": forecast_str,
                "named_buyers":      named_buyers,
                "forecast_base":     forecast_str,
                "discovery_count":   sum(1 for a in accounts if (a.deal_stage or "identified") == "identified"),
                "account_count":     active_accounts,
                "status":            status,
                "last_updated":      datetime.datetime.utcnow().strftime("%Y-%m-%d"),
                "top_deals":         top_deals,
                "live_signals":      live_signals,
                "priority":          priority,
                "country_split":     country_split,
                "timestamp":         datetime.datetime.utcnow().isoformat() + "Z",
            })
        except Exception as e:
            print(f"[GTM] /api/dashboard-live DB error: {e}")

    # Fallback: file-based parsing
    dash = read_file("intelligence/pipeline-dashboard.md") or ""

    def _find(pattern, text, default=""):
        m = re.search(pattern, text)
        return m.group(1).strip() if m else default

    pipeline_val  = _find(r'Pipeline Value \(total\)\s*\|\s*([^\n|]+)', dash, "€10.5M")
    buyers_raw    = _find(r'Named buyers confirmed\s*\|\s*(\d+)', dash, "28")
    named_buyers  = int(buyers_raw) if buyers_raw.isdigit() else 28
    forecast_raw  = _find(r'Base case Q2 forecast\s*\|\s*\*\*(€[\d,K]+)', dash, "€600K")
    discovery_raw = _find(r'Accounts in Discovery\s*\|\s*(\d+)', dash, "14")
    discovery_n   = int(discovery_raw) if discovery_raw.isdigit() else 14
    last_updated  = _find(r'Last updated:\s*(.+)', dash, "unknown")
    weighted_raw  = _find(r'Probability-weighted forecast.*?\|\s*\*\*(€[\d.,M]+)', dash, "€1.8M")

    status = "AMBER"
    if "**Status: RED" in dash:   status = "RED"
    elif "**Status: GREEN" in dash: status = "GREEN"

    deals = []
    in_table = False
    for line in dash.splitlines():
        if "| Rank |" in line and "Account" in line:
            in_table = True
            continue
        if in_table and line.startswith("|"):
            if re.match(r'^\|\s*[-:]+', line):
                continue
            cells = [c.strip() for c in line.split("|")]
            cells = [c for c in cells if c != ""]
            if len(cells) >= 8 and cells[0].isdigit():
                deals.append({
                    "rank":      cells[0],
                    "name":      cells[1],
                    "country":   cells[2],
                    "offering":  cells[3],
                    "icp":       cells[4],
                    "win_pct":   cells[5],
                    "entry_val": cells[6],
                    "weighted":  cells[7],
                    "buyer":     cells[8] if len(cells) > 8 else "TBD",
                    "status":    cells[9] if len(cells) > 9 else "Active",
                })
        elif in_table and not line.startswith("|"):
            in_table = False

    now_ts = datetime.datetime.now().timestamp()
    for d in deals:
        raw_slug = d["name"].lower()
        for ch, rep in [(" nordic",""),(" group",""),("ø","o"),("å","a"),("æ","ae"),(" ","-"),("'","")]:
            raw_slug = raw_slug.replace(ch, rep)
        candidates = [raw_slug, raw_slug.split("-")[0]]
        d["days_stale"] = None
        d["slug"] = ""
        for cand in candidates:
            na = BASE_DIR / "Accounts" / cand / "next-actions.md"
            if na.exists():
                d["days_stale"] = int((now_ts - na.stat().st_mtime) / 86400)
                d["slug"] = cand
                break

    priority = None
    for d in deals:
        buyer = d.get("buyer", "TBD")
        if buyer and buyer != "TBD" and "TBD" not in buyer:
            buyer_name = re.split(r'\(', buyer)[0].strip()
            priority = {
                "name":     buyer_name,
                "company":  d["name"],
                "country":  d["country"],
                "win_pct":  d["win_pct"],
                "weighted": d["weighted"],
                "offering": d["offering"],
                "status":   d.get("status", ""),
                "slug":     d.get("slug", ""),
            }
            break

    live_signals = []
    intel_dir = BASE_DIR / "intelligence" / "daily-leads"
    if intel_dir.exists():
        for f in sorted(intel_dir.iterdir(), reverse=True)[:3]:
            if f.suffix != ".md":
                continue
            content = f.read_text(encoding="utf-8", errors="replace")
            for m in re.finditer(r'^##\s+(.+)$\s+([\s\S]+?)(?=\n##|\Z)', content, re.MULTILINE):
                company = m.group(1).strip()
                body    = m.group(2).strip()[:180]
                live_signals.append({"company": company, "text": body, "date": f.stem})
                if len(live_signals) >= 6:
                    break
            if len(live_signals) >= 6:
                break

    country_split = []
    for m in re.finditer(r'\|\s*(Norway|Sweden|Denmark)\s*\|\s*(€[\d.M]+)[^|]*\|\s*(\d+)\s*\|', dash):
        country_split.append({
            "country":  m.group(1),
            "pipeline": m.group(2),
            "accounts": int(m.group(3)),
        })

    return jsonify({
        "pipeline_total":   pipeline_val.strip(),
        "pipeline_weighted": weighted_raw.strip(),
        "named_buyers":     named_buyers,
        "forecast_base":    forecast_raw.strip(),
        "discovery_count":  discovery_n,
        "account_count":    len(get_accounts()),
        "status":           status,
        "last_updated":     last_updated,
        "top_deals":        deals[:6],
        "live_signals":     live_signals,
        "priority":         priority,
        "country_split":    country_split,
        "timestamp":        datetime.datetime.utcnow().isoformat() + "Z",
    })


# ── Pitch Simulator API ───────────────────────────────────────────────────────

@app.route("/api/pitch", methods=["POST"])
def api_pitch():
    data = request.get_json()
    messages = data.get("messages", [])
    account_slug = data.get("account", "")
    scoring = data.get("scoring", False)

    account_content = load_account_files(account_slug) or ""
    account_name = account_slug.replace("-", " ").title()

    # Extract buyer from stakeholders
    stakeholders = read_file(f"Accounts/{account_slug}/stakeholders.md") or ""
    buyer_match = re.search(r"###\s+(.+)", stakeholders)
    title_match = re.search(r"\*\*Title:\*\*\s*(.+)", stakeholders)
    buyer_name  = buyer_match.group(1).strip() if buyer_match else "The Decision Maker"
    buyer_title = title_match.group(1).strip() if title_match else "CDO / CMO"

    if scoring:
        system = f"""You are a senior sales coach. The user just completed a pitch practice session trying to sell JAKALA's services to {account_name}.

Review the conversation and provide a concise scorecard:

**Pitch Scorecard — {account_name}**
Score each dimension 1–10:
- Opening hook: X/10
- Value proposition clarity: X/10
- Objection handling: X/10
- Buyer fit: X/10
- Call to action: X/10
- **Overall: X/10**

Then give 1 key strength and 1 specific improvement to make the pitch 20% more effective."""
    else:
        system = f"""You are {buyer_name}, {buyer_title} at {account_name}.

ACCOUNT CONTEXT:
{account_content[:2000]}

YOUR PERSONA:
- You are a senior executive who is busy, slightly sceptical, and protective of budget
- You care deeply about business outcomes, not technology for its own sake
- You ask tough but fair questions
- You are open to the right partner if they demonstrate clear ROI and understand your specific challenges
- You do NOT know about JAKALA or their services — treat them as a cold approach
- Stay in character throughout. Respond as this buyer would in a real meeting.
- Keep responses concise (2-5 sentences) — you're in a meeting, not writing an essay
- Push back if the pitch is vague or generic
- If the pitch is compelling, show genuine interest

Start the conversation by briefly introducing yourself and asking what brings the salesperson to this meeting."""

    def generate():
        try:
            with client.messages.stream(
                model=MODEL,
                max_tokens=1024,
                system=system,
                messages=messages,
            ) as stream:
                for text in stream.text_stream:
                    yield f"data: {json.dumps({'text': text})}\n\n"
        except anthropic.AuthenticationError:
            yield f"data: {json.dumps({'text': '[ERROR] API key error — check ANTHROPIC_API_KEY in Railway environment variables.'})}\n\n"
        except anthropic.APIStatusError as e:
            yield f"data: {json.dumps({'text': f'[ERROR] API error {e.status_code}: {e.message}'})}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'text': f'[ERROR] Error: {str(e)}'})}\n\n"
        finally:
            yield "data: [DONE]\n\n"

    return Response(stream_with_context(generate()), mimetype="text/event-stream")


# ── Board Report API ──────────────────────────────────────────────────────────

@app.route("/api/board-report", methods=["POST"])
def api_board_report():
    import traceback
    if not PPTX_OK:
        return jsonify({"error": "python-pptx not installed"}), 500
    try:
        return _do_board_report()
    except Exception as e:
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500


def _do_board_report():
    today = datetime.date.today().isoformat()
    accounts = get_accounts()
    pipeline_lines = []
    for slug in accounts[:60]:
        ov = read_file(f"Accounts/{slug}/overview.md") or ""
        st = read_file(f"Accounts/{slug}/strategy.md") or ""
        icp_m  = re.search(r"ICP Score[:\s]+(\d+)/10", ov)
        deal_m = re.search(r"Deal Score[:\s]+(\d+)/10", ov)
        icp  = int(icp_m.group(1))  if icp_m  else 0
        deal = int(deal_m.group(1)) if deal_m else 0
        if icp >= 7 or deal >= 7:
            pipeline_lines.append(f"- {slug.replace('-',' ').title()}: ICP {icp}, Deal {deal}")

    top_ctx = read_file("intelligence/top-opportunities.md") or ""

    prompt = f"""You are building a board-level commercial review for JAKALA Nordic.

TODAY: {today}
PIPELINE (top accounts):
{chr(10).join(pipeline_lines[:20])}

TOP OPPORTUNITIES:
{top_ctx[:1500]}

Return ONLY valid JSON (no markdown fences):
{{
  "exec_summary": "2-sentence pipeline status for the board",
  "pipeline_value": "€6.8M",
  "pipeline_status": "AMBER",
  "top_deals": [
    {{"name":"H&M","value":"€900K","stage":"Prospecting","gtm":"Data Revenue Unlock","next":"Identify buyer — no named contact yet"}},
    {{"name":"Matas","value":"€700K","stage":"Prospecting","gtm":"AI Readiness","next":"Outreach to loyalty/data team"}},
    {{"name":"Elkjøp","value":"€700K","stage":"Prospecting","gtm":"Commerce Optimization","next":"Contact Morten Syversen"}},
    {{"name":"Varner Group","value":"€450K","stage":"Prospecting","gtm":"Data Revenue Unlock","next":"Confirm buyer — Elise Laupstad"}},
    {{"name":"Trumf","value":"€450K","stage":"Prospecting","gtm":"Data Revenue Unlock","next":"Reach Rikke Etholm-Idsøe (new role)"}}
  ],
  "q2_forecast": "€420K",
  "q2_confidence": "Base case",
  "q3_forecast": "€850K",
  "q3_confidence": "Upside",
  "risks": ["Zero first meetings booked","No named buyer on top 3 accounts","Pipeline age increasing"],
  "opportunities": ["5 hot timing signals active this week","45 ICP-scored accounts ready to contact","Trumf + Vinmonopolet honeymoon windows open"],
  "this_week": [
    {{"action":"Contact Sport Outlet CEO Tor-André Skeie","why":"CTO + CDO both vacant — entry window closing fast"}},
    {{"action":"Reach Trumf Commercial Director Rikke Etholm-Idsøe","why":"New role, 90-day honeymoon window"}},
    {{"action":"Vinmonopolet CDO Espen Terland outreach","why":"New CDO, agenda not set — perfect timing"}},
    {{"action":"H&M buyer identification","why":"€900K — largest opportunity, buyer TBD"}},
    {{"action":"Matas loyalty team outreach","why":"€700K, AI Readiness — data infrastructure opportunity"}}
  ],
  "gtm_split": {{"dru":12,"ai":11,"co":13,"xt":9}}
}}

Return ONLY the JSON — no explanation."""

    response = client.messages.create(
        model=MODEL, max_tokens=2500,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = re.sub(r"^```json\s*", "", response.content[0].text.strip())
    raw = re.sub(r"\s*```$", "", raw)

    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        return jsonify({"error": "AI parse failed", "raw": raw}), 500

    buf = _build_board_report(data, today)
    return send_file(buf,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=f"JAKALA-Board-Report-{today}.pptx")


def _build_board_report(data, today):
    prs = _prs()

    # Slide 1 — Cover
    s = _slide(prs); _bg(s)
    _rect(s, Inches(0), Inches(0), Inches(0.08), H, BLUE)
    _rect(s, Inches(0), H - Inches(0.5), W, Inches(0.5), BLUE)
    _txt(s, "JAKALA NORDIC", Inches(0.3), Inches(1.2), Inches(8), Inches(0.4), size=10, bold=True, color=MUTED)
    _txt(s, "Commercial\nBoard Review", Inches(0.3), Inches(1.75), Inches(8.5), Inches(2.5), size=44, bold=True)
    _txt(s, today, Inches(0.3), Inches(4.2), Inches(4), Inches(0.35), size=12, color=MUTED)
    status = data.get("pipeline_status", "AMBER")
    scol = GREEN if status == "GREEN" else (RED if status == "RED" else RGBColor(0xF5,0xA6,0x23))
    _rect(s, Inches(0.3), Inches(4.75), Inches(2.2), Inches(0.45), scol)
    _txt(s, f"STATUS: {status}", Inches(0.35), Inches(4.77), Inches(2.1), Inches(0.4), size=11, bold=True, color=WHITE)

    # Slide 2 — KPIs + Exec Summary
    s = _slide(prs); _bg(s)
    _header(s, "PIPELINE OVERVIEW", "Nordic Commercial Pipeline — Q1–Q2 2026")
    kpis = [
        (data.get("pipeline_value","€6.8M"), "Total Pipeline", BLUE),
        ("45", "Active Accounts", NAVY),
        (data.get("q2_forecast","€420K"), "Q2 Base Case", GREEN),
        ("18", "Named Buyers", RGBColor(0xF5,0xA6,0x23)),
    ]
    for i, (val, lbl, col) in enumerate(kpis):
        x = Inches(0.4) + i * Inches(2.35)
        _rect(s, x, Inches(1.5), Inches(2.1), Inches(1.3), col)
        _txt(s, val, x+Inches(0.12), Inches(1.56), Inches(1.86), Inches(0.72), size=26, bold=True)
        _txt(s, lbl, x+Inches(0.12), Inches(2.24), Inches(1.86), Inches(0.45), size=10, color=WHITE)
    _txt(s, data.get("exec_summary",""), Inches(0.5), Inches(3.1), Inches(9.0), Inches(1.0), size=13, color=GREY)
    _bullet_col(s, Inches(0.5), Inches(4.3), Inches(4.3), Inches(2.8), [("Key Risks", data.get("risks",[]))])
    _bullet_col(s, Inches(5.0), Inches(4.3), Inches(4.3), Inches(2.8), [("Key Opportunities", data.get("opportunities",[]))])
    _footer(s, f"JAKALA Nordic Board Review — {today} — Confidential")

    # Slide 3 — Top 5 Deals
    s = _slide(prs); _bg(s)
    _header(s, "TOP OPPORTUNITIES", "Highest-Value Active Deals")
    cols_cycle = [GREEN, BLUE, BLUE, NAVY, NAVY]
    for i, deal in enumerate(data.get("top_deals",[])[:5]):
        y = Inches(1.55) + i * Inches(1.04)
        col = cols_cycle[i]
        _rect(s, Inches(0.4), y, Inches(0.08), Inches(0.84), col)
        _txt(s, deal.get("name",""), Inches(0.6), y+Inches(0.02), Inches(3.0), Inches(0.44), size=14, bold=True)
        _txt(s, deal.get("gtm",""), Inches(0.6), y+Inches(0.44), Inches(3.0), Inches(0.35), size=10, color=MUTED)
        _txt(s, deal.get("value",""), Inches(3.8), y+Inches(0.02), Inches(1.5), Inches(0.44), size=18, bold=True, color=col)
        _txt(s, deal.get("stage",""), Inches(3.8), y+Inches(0.44), Inches(1.5), Inches(0.35), size=10, color=MUTED)
        _txt(s, f"\u2192 {deal.get('next','')}", Inches(5.5), y+Inches(0.2), Inches(3.8), Inches(0.45), size=11, color=GREY)
    _footer(s, f"JAKALA Nordic Board Review — {today} — Confidential")

    # Slide 4 — Forecast
    s = _slide(prs); _bg(s)
    _header(s, "REVENUE FORECAST", "Q2–Q3 2026 Probability-Weighted")
    for i, (period, val, conf) in enumerate([
        ("Q2 2026 Base", data.get("q2_forecast","€420K"), data.get("q2_confidence","Base case")),
        ("Q3 2026 Upside", data.get("q3_forecast","€850K"), data.get("q3_confidence","Upside")),
        ("FY 2026 Target", "€1.8M", "Stretch"),
    ]):
        x = Inches(0.5) + i * Inches(3.1)
        _rect(s, x, Inches(1.55), Inches(2.8), Inches(1.8), NAVY)
        _txt(s, period, x+Inches(0.15), Inches(1.62), Inches(2.5), Inches(0.44), size=11, bold=True, color=MUTED)
        _txt(s, val, x+Inches(0.15), Inches(2.1), Inches(2.5), Inches(0.72), size=32, bold=True, color=BLUE)
        _txt(s, conf, x+Inches(0.15), Inches(2.8), Inches(2.5), Inches(0.4), size=11, color=GREY)
    sp = data.get("gtm_split", {"dru":12,"ai":11,"co":13,"xt":9})
    _bullet_col(s, Inches(0.5), Inches(3.8), Inches(4.0), Inches(3.3), [("GTM Strategy Mix", [
        f"Commerce Optimization: {sp.get('co',0)} accounts",
        f"Data Revenue Unlock: {sp.get('dru',0)} accounts",
        f"AI Readiness: {sp.get('ai',0)} accounts",
        f"Experience Transformation: {sp.get('xt',0)} accounts",
    ])])
    _footer(s, f"JAKALA Nordic Board Review — {today} — Confidential")

    # Slide 5 — Priority Actions
    s = _slide(prs); _bg(s)
    _header(s, "PRIORITY ACTIONS", "This Week's Commercial Focus")
    colors_c = [BLUE, GREEN, RED, RGBColor(0xF5,0xA6,0x23), BLUE]
    for i, item in enumerate(data.get("this_week",[])[:5]):
        y = Inches(1.6) + i * Inches(1.1)
        col = colors_c[i % len(colors_c)]
        _rect(s, Inches(0.5), y, Inches(0.52), Inches(0.52), col)
        _txt(s, str(i+1), Inches(0.5), y, Inches(0.52), Inches(0.52), size=20, bold=True, align=PP_ALIGN.CENTER)
        _txt(s, item.get("action",""), Inches(1.15), y, Inches(8.1), Inches(0.42), size=14, bold=True)
        _txt(s, item.get("why",""), Inches(1.15), y+Inches(0.42), Inches(8.1), Inches(0.5), size=11, color=GREY)
    _footer(s, f"JAKALA Nordic Board Review — {today} — Confidential")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── Frontend ─────────────────────────────────────────────────────────────────

HTML = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>JAKALA GTM OS</title>
<style>
*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }

:root {
  --blue:       #153EED;
  --blue-mid:   #3558F0;
  --blue-light: #2952D8;
  --blue-glow:  rgba(21,62,237,0.2);
  --blue-dim:   rgba(21,62,237,0.08);
  --blue-dim2:  rgba(21,62,237,0.12);
  --red:        #DC2626;
  --red-dim:    rgba(220,38,38,0.1);
  --green:      #00A67E;
  --green-dim:  rgba(0,166,126,0.1);
  --amber:      #D97706;
  --amber-dim:  rgba(217,119,6,0.1);
  --purple:     #7C3AED;
  --bg:         #F4F6FC;
  --bg2:        #FFFFFF;
  --bg3:        #F0F2F9;
  --bg4:        #EAECF4;
  --border:     #E4E7F0;
  --border-hi:  #CDD1E4;
  --text:       #1F2937;
  --muted:      #6B7280;
  --muted2:     #9CA3AF;
  --white:      #0F0F1A;
  --sidebar:    272px;
  --font:       'Inter', -apple-system, system-ui, sans-serif;
}

body {
  font-family: var(--font);
  background: var(--bg);
  color: var(--text);
  height: 100vh;
  display: flex;
  overflow: hidden;
  font-size: 13px;
  -webkit-font-smoothing: antialiased;
}

/* ── Ambient background ── */
body::before {
  content: '';
  position: fixed;
  inset: 0;
  background:
    radial-gradient(ellipse 80% 60% at 15% 10%, rgba(21,62,237,0.05) 0%, transparent 60%),
    radial-gradient(ellipse 60% 50% at 85% 90%, rgba(124,58,237,0.03) 0%, transparent 60%);
  pointer-events: none;
  z-index: 0;
}

/* ── Grid overlay ── */
body::after {
  content: '';
  position: fixed;
  inset: 0;
  background-image:
    linear-gradient(rgba(21,62,237,0.035) 1px, transparent 1px),
    linear-gradient(90deg, rgba(21,62,237,0.035) 1px, transparent 1px);
  background-size: 48px 48px;
  pointer-events: none;
  z-index: 0;
}

/* ══════════════════════════════════════════
   SIDEBAR
══════════════════════════════════════════ */
#sidebar {
  width: var(--sidebar);
  min-width: var(--sidebar);
  background: var(--bg2);
  border-right: 1px solid var(--border);
  display: flex;
  flex-direction: column;
  height: 100vh;
  overflow: hidden;
  position: relative;
  z-index: 10;
  box-shadow: 2px 0 8px rgba(21,62,237,.04);
}

/* Logo */
#logo {
  padding: 20px 18px 16px;
  border-bottom: 1px solid var(--border);
  display: flex;
  align-items: center;
  gap: 11px;
  flex-shrink: 0;
}
.logo-mark {
  width: 32px; height: 32px;
  background: linear-gradient(135deg, var(--blue) 0%, var(--purple) 100%);
  border-radius: 8px;
  display: flex; align-items: center; justify-content: center;
  font-size: 15px; font-weight: 900; color: #fff; letter-spacing: -1px;
  box-shadow: 0 0 18px var(--blue-glow);
  flex-shrink: 0;
}
.logo-text { flex: 1; }
.logo-text .brand { font-size: 14px; font-weight: 800; letter-spacing: 2.5px; color: var(--white); }
.logo-text .sub   { font-size: 9px; color: var(--muted); letter-spacing: 1.5px; text-transform: uppercase; margin-top: 1px; }
.live-pill {
  display: flex; align-items: center; gap: 5px;
  padding: 3px 8px; border-radius: 12px;
  background: var(--green-dim);
  border: 1px solid rgba(0,212,160,0.2);
  font-size: 9px; font-weight: 700; color: var(--green);
  letter-spacing: 1px; text-transform: uppercase;
}
.live-dot {
  width: 5px; height: 5px; border-radius: 50%;
  background: var(--green);
  box-shadow: 0 0 6px var(--green);
  animation: livePulse 2s infinite;
}
@keyframes livePulse { 0%,100% { opacity:1; } 50% { opacity:0.25; } }

/* Nav */
.sidebar-label {
  padding: 14px 18px 7px;
  font-size: 9.5px; font-weight: 700;
  letter-spacing: 2px; text-transform: uppercase;
  color: var(--muted);
  flex-shrink: 0;
}
.nav-btn {
  display: flex; align-items: center; gap: 10px;
  width: calc(100% - 14px); margin: 1px 7px;
  padding: 9px 11px; border-radius: 7px;
  background: transparent; border: 1px solid transparent;
  color: var(--muted2); cursor: pointer;
  font-size: 12.5px; font-weight: 500; text-align: left;
  transition: all 0.18s;
  flex-shrink: 0;
}
.nav-btn:hover { background: rgba(21,62,237,0.06); color: var(--text); }
.nav-btn.active {
  background: var(--blue-dim2); color: var(--blue);
  border-color: rgba(21,62,237,0.22);
  box-shadow: inset 3px 0 0 var(--blue);
  font-weight: 600;
}
.nav-btn .icon { font-size: 14px; width: 20px; text-align: center; }

/* Skill buttons */
.skill-btn {
  display: flex; align-items: center; gap: 8px;
  width: calc(100% - 14px); margin: 1px 7px;
  padding: 7px 11px; border-radius: 6px;
  background: transparent; border: 1px solid transparent;
  color: var(--muted); cursor: pointer;
  font-size: 11.5px; font-weight: 500; text-align: left;
  transition: all 0.15s;
  flex-shrink: 0;
}
.skill-btn:hover {
  border-color: rgba(21,62,237,0.3);
  color: var(--blue); background: var(--blue-dim);
}

/* Account list */
#account-list { max-height: 480px; overflow-y: auto; padding-bottom: 8px; }
#account-list::-webkit-scrollbar { width: 3px; }
#account-list::-webkit-scrollbar-thumb { background: var(--border-hi); border-radius: 2px; }
.sidebar-accounts-toggle {
  display: flex; align-items: center; justify-content: space-between;
  width: calc(100% - 14px); margin: 4px 7px 2px;
  padding: 8px 11px; border-radius: 7px;
  background: transparent; border: 1px solid transparent;
  color: var(--muted2); cursor: pointer;
  font-size: 12px; font-weight: 600; text-align: left;
  letter-spacing: 0.06em; text-transform: uppercase;
  transition: all 0.18s;
}
.sidebar-accounts-toggle:hover { background: rgba(21,62,237,0.06); color: var(--text); }
.sidebar-accounts-toggle.open { color: var(--blue); }

#account-search {
  margin: 8px 10px;
  padding: 8px 12px;
  border-radius: 6px;
  background: #F8F9FE;
  border: 1px solid var(--border);
  color: var(--text); font-size: 12px;
  width: calc(100% - 20px); outline: none;
  transition: border-color 0.18s;
  flex-shrink: 0;
}
#account-search::placeholder { color: var(--muted); }
#account-search:focus { border-color: rgba(21,62,237,0.45); }

.account-item {
  display: flex; align-items: center; justify-content: space-between;
  padding: 9px 13px; margin: 2px 7px; border-radius: 6px;
  cursor: pointer; transition: all 0.12s;
}
.account-item:hover { background: rgba(21,62,237,0.05); }
.account-item.selected { background: var(--blue-dim2); border: 1px solid rgba(21,62,237,0.2); }
.account-item .aname { flex: 1; font-size: 13px; color: var(--text); white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }
.acc-badges { display: flex; gap: 4px; margin-left: 6px; align-items: center; }
.spill {
  font-size: 9.5px; font-weight: 700; padding: 1px 5px; border-radius: 3px;
  background: var(--blue-dim); color: var(--blue);
}
.spill.g { background: var(--green-dim); color: var(--green); }
.spill.c { font-size: 9px; color: var(--muted); background: transparent; }

/* ══════════════════════════════════════════
   MAIN
══════════════════════════════════════════ */
#main {
  flex: 1; display: flex; flex-direction: column;
  overflow: hidden; position: relative; z-index: 1;
}

.tab-pane { display: none; flex: 1; flex-direction: column; overflow: hidden; }
.tab-pane.active { display: flex; }

/* ══════════════════════════════════════════
   DASHBOARD
══════════════════════════════════════════ */
#tab-dashboard { overflow-y: auto; }
#dash-wrap { padding: 32px 36px 48px; min-height: 100%; }

.dash-top {
  display: flex; justify-content: space-between; align-items: flex-start;
  margin-bottom: 28px;
}
.dash-headline { font-size: 26px; font-weight: 800; color: var(--white); letter-spacing: -0.6px; line-height: 1.1; }
.dash-tagline { font-size: 13px; color: var(--muted2); margin-top: 5px; }
.dash-meta {
  text-align: right; padding: 12px 18px;
  background: var(--bg2);
  border: 1px solid var(--border); border-radius: 10px;
  min-width: 180px;
}
.dash-meta .dm-date { font-size: 13px; font-weight: 600; color: var(--text); }
.dash-meta .dm-time { font-size: 22px; font-weight: 800; color: var(--white); letter-spacing: -0.5px; margin-top: 2px; }
.dash-meta .dm-label { font-size: 9.5px; color: var(--muted); text-transform: uppercase; letter-spacing: 1.5px; margin-top: 3px; }

/* KPI row */
.kpi-row { display: grid; grid-template-columns: repeat(4, 1fr); gap: 14px; margin-bottom: 22px; }

.kpi {
  background: var(--bg2);
  border: 1px solid var(--border); border-radius: 12px;
  padding: 20px 20px 16px; position: relative; overflow: hidden;
  transition: border-color 0.2s, transform 0.2s, box-shadow 0.2s;
}
.kpi:hover { border-color: rgba(21,62,237,.25); transform: translateY(-2px); box-shadow: 0 4px 20px rgba(21,62,237,.07); }
.kpi::after {
  content: '';
  position: absolute; top: 0; left: 12px; right: 12px; height: 1px;
  background: linear-gradient(90deg, transparent, rgba(21,62,237,0.5), transparent);
}
.kpi-label { font-size: 10px; font-weight: 700; color: var(--muted); text-transform: uppercase; letter-spacing: 1.8px; margin-bottom: 10px; }
.kpi-num {
  font-size: 38px; font-weight: 900; color: var(--white);
  letter-spacing: -2px; line-height: 1;
  font-variant-numeric: tabular-nums;
}
.kpi-unit { font-size: 20px; font-weight: 700; color: var(--muted2); vertical-align: super; font-size: 16px; }
.kpi-suffix { font-size: 22px; font-weight: 800; color: var(--muted2); }
.kpi-sub { font-size: 11px; color: var(--muted2); margin-top: 8px; }
.kpi-badge {
  display: inline-flex; align-items: center; gap: 4px;
  padding: 3px 9px; border-radius: 5px;
  font-size: 10px; font-weight: 800; text-transform: uppercase; letter-spacing: 0.8px;
  margin-top: 10px;
}
.kpi-badge.amber { background: var(--amber-dim); color: var(--amber); border: 1px solid rgba(245,166,35,0.22); }
.kpi-badge.green { background: var(--green-dim); color: var(--green); border: 1px solid rgba(0,212,160,0.22); }
.kpi-badge.red   { background: var(--red-dim);   color: var(--red);   border: 1px solid rgba(246,87,74,0.22); }
.kpi-badge.blue  { background: var(--blue-dim2); color: var(--blue-light); border: 1px solid rgba(21,62,237,0.22); }

/* KPI accent bar */
.kpi-bar { height: 3px; border-radius: 2px; margin-top: 14px; background: rgba(0,0,0,0.07); overflow: hidden; }
.kpi-bar-fill { height: 100%; border-radius: 2px; transition: width 1.4s cubic-bezier(0.4,0,0.2,1); }

/* Dashboard two-col */
.dash-cols { display: grid; grid-template-columns: 1.5fr 1fr; gap: 18px; margin-bottom: 18px; }
.dash-card {
  background: var(--bg2);
  border: 1px solid var(--border); border-radius: 12px; padding: 22px;
}
.dash-card-head {
  display: flex; align-items: center; justify-content: space-between;
  margin-bottom: 16px; padding-bottom: 12px;
  border-bottom: 1px solid var(--border);
}
.dash-card-title { font-size: 10px; font-weight: 800; color: var(--muted2); text-transform: uppercase; letter-spacing: 2px; }
.dash-card-tag { font-size: 10px; color: var(--muted); }

/* Opportunity rows */
.opp {
  display: flex; align-items: center; gap: 12px;
  padding: 9px 8px; border-radius: 8px; cursor: pointer;
  transition: background 0.15s; margin: 0 -8px;
  border-bottom: 1px solid var(--border);
}
.opp:last-child { border-bottom: none; }
.opp:hover { background: rgba(21,62,237,0.04); border-radius: 6px; }
.opp-num { font-size: 10px; font-weight: 800; color: var(--muted); width: 18px; text-align: right; flex-shrink: 0; }
.opp-ring { position: relative; width: 38px; height: 38px; flex-shrink: 0; }
.opp-ring svg { position: absolute; top: 0; left: 0; }
.opp-ring-val {
  position: absolute; inset: 0;
  display: flex; align-items: center; justify-content: center;
  font-size: 11px; font-weight: 900; color: var(--white);
}
.opp-body { flex: 1; min-width: 0; }
.opp-name { font-size: 13.5px; font-weight: 700; color: var(--white); white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }
.opp-meta { font-size: 11px; color: var(--muted2); margin-top: 2px; }
.opp-right { text-align: right; flex-shrink: 0; }
.opp-val { font-size: 14px; font-weight: 800; color: var(--white); }
.opp-val-sub { font-size: 10px; color: var(--muted); margin-top: 1px; }

/* Signal rows */
.signal {
  display: flex; align-items: flex-start; gap: 11px;
  padding: 9px 0; border-bottom: 1px solid var(--border);
}
.signal:last-child { border-bottom: none; }
.signal-icon {
  width: 28px; height: 28px; border-radius: 6px;
  background: var(--red-dim); border: 1px solid rgba(246,87,74,0.2);
  display: flex; align-items: center; justify-content: center;
  font-size: 13px; flex-shrink: 0;
}
.signal-co { font-size: 13px; font-weight: 700; color: var(--white); }
.signal-txt { font-size: 11px; color: var(--muted2); margin-top: 2px; line-height: 1.45; }
.signal-tag {
  font-size: 9.5px; font-weight: 800; padding: 2px 7px; border-radius: 4px;
  background: var(--red-dim); color: var(--red); white-space: nowrap;
  align-self: flex-start; flex-shrink: 0; margin-top: 1px;
}
.signal-tag.amber { background: var(--amber-dim); color: var(--amber); }

/* Strategy split */
.strat-row { display: grid; grid-template-columns: repeat(4,1fr); gap: 12px; margin-top: 0; }
.strat-card {
  background: var(--bg2);
  border: 1px solid var(--border); border-radius: 10px;
  padding: 16px; text-align: center;
}
.strat-icon { font-size: 22px; margin-bottom: 8px; }
.strat-name { font-size: 10px; font-weight: 700; color: var(--muted2); text-transform: uppercase; letter-spacing: 1.2px; margin-bottom: 10px; line-height: 1.4; }
.strat-num { font-size: 28px; font-weight: 900; color: var(--white); letter-spacing: -1px; }
.strat-sub { font-size: 10px; color: var(--muted); margin-top: 4px; }
.strat-bar-wrap { height: 3px; background: rgba(0,0,0,0.07); border-radius: 2px; margin-top: 12px; overflow: hidden; }
.strat-bar { height: 100%; border-radius: 2px; transition: width 1.6s cubic-bezier(0.4,0,0.2,1); }

/* ══════════════════════════════════════════
   CHAT
══════════════════════════════════════════ */
#chat-header {
  padding: 14px 24px;
  border-bottom: 1px solid var(--border);
  display: flex; align-items: center; gap: 12px;
  background: var(--bg2);
  flex-shrink: 0;
}
#chat-header h2 { font-size: 14px; font-weight: 700; color: var(--white); }
#selected-account-tag {
  font-size: 11px; padding: 4px 10px; border-radius: 16px;
  background: var(--blue-dim2); color: var(--blue-light);
  border: 1px solid rgba(21,62,237,0.3);
  display: none; align-items: center; gap: 6px;
}
#clear-account { background: none; border: none; color: var(--muted); cursor: pointer; font-size: 12px; padding: 0; }
#clear-account:hover { color: var(--red); }
#deck-btn {
  padding: 5px 12px; border-radius: 6px; font-size: 11px; font-weight: 700;
  background: transparent; border: 1px solid rgba(0,212,160,0.35);
  color: var(--green); cursor: pointer; transition: all 0.15s; display: none; white-space: nowrap;
}
#deck-btn:hover { background: var(--green-dim); }
#deck-btn.loading { opacity: 0.5; cursor: not-allowed; }

#messages { flex: 1; overflow-y: auto; padding: 24px; display: flex; flex-direction: column; gap: 18px; }
#messages::-webkit-scrollbar { width: 3px; }
#messages::-webkit-scrollbar-thumb { background: var(--border-hi); border-radius: 2px; }

.msg { max-width: 80%; display: flex; flex-direction: column; gap: 5px; }
.msg.user { align-self: flex-end; align-items: flex-end; }
.msg.assistant { align-self: flex-start; align-items: flex-start; }
.msg-role { font-size: 9.5px; color: var(--muted); text-transform: uppercase; letter-spacing: 1.5px; padding: 0 4px; }
.msg-bubble {
  padding: 13px 17px; border-radius: 14px;
  font-size: 13.5px; line-height: 1.65; word-break: break-word;
}
.msg.user .msg-bubble {
  background: linear-gradient(135deg, var(--blue), var(--blue-light));
  color: #fff; border-bottom-right-radius: 4px;
  box-shadow: 0 4px 20px rgba(21,62,237,0.28);
}
.msg.assistant .msg-bubble {
  background: var(--bg2); color: var(--text);
  border: 1px solid var(--border); border-bottom-left-radius: 4px;
  box-shadow: 0 2px 8px rgba(0,0,0,.04);
}
.msg-bubble table { border-collapse: collapse; width: 100%; margin: 10px 0; font-size: 12.5px; }
.msg-bubble th { background: rgba(21,62,237,0.08); color: var(--blue); padding: 7px 11px; text-align: left; font-size: 10.5px; text-transform: uppercase; letter-spacing: 0.8px; }
.msg-bubble td { padding: 6px 11px; border-bottom: 1px solid var(--border); }
.msg-bubble tr:last-child td { border-bottom: none; }
.msg-bubble h1, .msg-bubble h2 { color: var(--blue); margin: 14px 0 7px; font-size: 15px; font-weight: 700; }
.msg-bubble h3 { color: var(--white); margin: 11px 0 5px; font-size: 13.5px; font-weight: 600; }
.msg-bubble ul, .msg-bubble ol { padding-left: 20px; margin: 5px 0; }
.msg-bubble li { margin: 3px 0; }
.msg-bubble code { background: rgba(21,62,237,0.09); padding: 2px 6px; border-radius: 4px; font-size: 12px; font-family: 'SF Mono','Fira Code',monospace; color: var(--blue); }
.msg-bubble pre { background: #F0F2F9; padding: 12px; border-radius: 8px; overflow-x: auto; border: 1px solid var(--border); margin: 8px 0; }
.msg-bubble strong { color: var(--white); }
.msg-bubble em { color: var(--muted); }
.msg-bubble blockquote { border-left: 3px solid var(--blue); padding-left: 13px; color: var(--muted2); margin: 8px 0; font-style: italic; }
.msg-bubble hr { border: none; border-top: 1px solid var(--border); margin: 12px 0; }

#typing-row { padding: 0 24px 10px; flex-shrink: 0; }
#typing-indicator {
  display: none; align-items: center; gap: 5px;
  padding: 11px 16px; background: var(--bg2);
  border: 1px solid var(--border); border-radius: 14px; border-bottom-left-radius: 4px;
  max-width: 90px;
}
#typing-indicator.visible { display: flex; }
.dot { width: 5px; height: 5px; border-radius: 50%; background: var(--blue); animation: dotBounce 1.2s infinite; }
.dot:nth-child(2) { animation-delay: 0.18s; }
.dot:nth-child(3) { animation-delay: 0.36s; }
@keyframes dotBounce { 0%,60%,100% { transform:translateY(0); opacity:0.35; } 30% { transform:translateY(-6px); opacity:1; } }

#chat-input-area {
  padding: 14px 24px 18px;
  border-top: 1px solid var(--border);
  background: var(--bg2);
  flex-shrink: 0;
}
#input-row { display: flex; gap: 10px; align-items: flex-end; }
#chat-input {
  flex: 1; padding: 11px 15px; border-radius: 10px;
  background: #F8F9FE; border: 1px solid var(--border);
  color: var(--text); font-size: 13.5px; font-family: var(--font);
  resize: none; outline: none; min-height: 44px; max-height: 140px; line-height: 1.5;
  transition: border-color 0.18s, box-shadow 0.18s;
}
#chat-input::placeholder { color: var(--muted); }
#chat-input:focus {
  border-color: rgba(21,62,237,0.5);
  box-shadow: 0 0 0 3px rgba(21,62,237,0.08);
}
#send-btn {
  padding: 11px 20px;
  background: linear-gradient(135deg, var(--blue), var(--blue-light));
  color: #fff; border: none; border-radius: 10px;
  cursor: pointer; font-size: 13.5px; font-weight: 700;
  transition: all 0.18s;
  box-shadow: 0 4px 14px rgba(21,62,237,0.3);
}
#send-btn:hover { transform: translateY(-1px); box-shadow: 0 6px 20px rgba(21,62,237,0.45); }
#send-btn:disabled { opacity: 0.38; cursor: not-allowed; transform: none; box-shadow: none; }
#input-hint { font-size: 10.5px; color: var(--muted); margin-top: 8px; }

/* Welcome screen */
#welcome {
  display: flex; flex-direction: column; align-items: center; justify-content: center;
  height: 100%; gap: 10px; text-align: center; padding: 40px;
}
.welcome-icon { font-size: 44px; margin-bottom: 6px; opacity: 0.7; }
#welcome .big { font-size: 24px; font-weight: 800; color: var(--white); letter-spacing: -0.4px; }
#welcome .sub { font-size: 14px; color: var(--muted2); max-width: 380px; line-height: 1.65; margin-top: 2px; }
.quick-chips { display: flex; flex-wrap: wrap; gap: 7px; justify-content: center; margin-top: 18px; }
.chip {
  padding: 7px 14px; border-radius: 20px;
  background: var(--bg2); border: 1px solid var(--border);
  color: var(--muted); font-size: 12px; cursor: pointer;
  transition: all 0.15s;
}
.chip:hover { border-color: rgba(21,62,237,0.4); color: var(--text); background: var(--blue-dim); }

/* ══════════════════════════════════════════
   NOTES
══════════════════════════════════════════ */
#tab-notes { overflow-y: auto; }
#notes-pane { padding: 32px; max-width: 760px; }
#notes-pane h2 { font-size: 20px; font-weight: 800; color: var(--white); letter-spacing: -0.3px; }
#notes-pane .desc { font-size: 13px; color: var(--muted2); margin: 6px 0 24px; line-height: 1.6; }
.form-group { margin-bottom: 18px; }
.form-group label { display: block; font-size: 10px; font-weight: 800; color: var(--muted2); text-transform: uppercase; letter-spacing: 1.8px; margin-bottom: 8px; }
.form-group select, .form-group textarea {
  width: 100%; padding: 10px 13px; border-radius: 8px;
  background: #F8F9FE; border: 1px solid var(--border);
  color: var(--text); font-size: 13.5px; font-family: var(--font); outline: none;
  transition: border-color 0.18s;
}
.form-group select:focus, .form-group textarea:focus { border-color: rgba(21,62,237,0.5); }
.form-group textarea { resize: vertical; min-height: 160px; line-height: 1.6; }
.form-group select option { background: var(--bg3); }
.btn-primary {
  padding: 11px 22px;
  background: linear-gradient(135deg, var(--blue), var(--blue-light));
  color: #fff; border: none; border-radius: 8px;
  cursor: pointer; font-size: 13px; font-weight: 700;
  transition: all 0.18s; box-shadow: 0 4px 12px rgba(21,62,237,0.25);
}
.btn-primary:hover { transform: translateY(-1px); box-shadow: 0 6px 20px rgba(21,62,237,0.4); }
.btn-primary:disabled { opacity: 0.38; cursor: not-allowed; transform: none; box-shadow: none; }
.btn-secondary {
  padding: 11px 22px; background: transparent; color: var(--text);
  border: 1px solid var(--border); border-radius: 8px;
  cursor: pointer; font-size: 13px; font-weight: 500;
  transition: border-color 0.15s; margin-left: 10px;
}
.btn-secondary:hover { border-color: var(--border-hi); }
#notes-result {
  margin-top: 26px; padding: 22px;
  background: var(--bg2); border: 1px solid var(--border);
  border-radius: 12px; display: none;
}
.result-label { font-size: 9.5px; font-weight: 800; color: var(--blue); text-transform: uppercase; letter-spacing: 2px; margin-bottom: 9px; }
.insight-box { padding: 13px 16px; background: var(--blue-dim2); border-left: 3px solid var(--blue); border-radius: 0 8px 8px 0; font-size: 13.5px; color: var(--text); line-height: 1.6; }
.summary-text { font-size: 13.5px; color: var(--muted2); line-height: 1.7; }
.file-preview {
  background: #F0F2F9; border: 1px solid var(--border); border-radius: 8px;
  padding: 13px; font-size: 11.5px; font-family: 'SF Mono','Fira Code',monospace;
  color: var(--muted); max-height: 180px; overflow-y: auto; white-space: pre-wrap; line-height: 1.6;
}
.result-section { margin-bottom: 18px; }
.save-actions { display: flex; gap: 10px; margin-top: 18px; }

/* ══════════════════════════════════════════
   ACCOUNTS
══════════════════════════════════════════ */
#tab-accounts { overflow-y: auto; }
#accounts-pane { padding: 32px; }
#accounts-pane h2 { font-size: 20px; font-weight: 800; color: var(--white); letter-spacing: -0.3px; }
#accounts-pane .desc { font-size: 13px; color: var(--muted2); margin: 6px 0 22px; }
.filter-row { display: flex; gap: 7px; margin-bottom: 20px; flex-wrap: wrap; }
.filter-btn {
  padding: 5px 13px; border-radius: 16px; font-size: 11.5px; font-weight: 600;
  background: transparent; border: 1px solid var(--border); color: var(--muted2); cursor: pointer; transition: all 0.15s;
}
.filter-btn.active, .filter-btn:hover { background: var(--blue-dim); border-color: rgba(21,62,237,0.3); color: var(--blue); }
.accounts-grid { display: grid; grid-template-columns: repeat(auto-fill, minmax(270px, 1fr)); gap: 11px; }
.account-card {
  background: var(--bg2); border: 1px solid var(--border); border-radius: 10px;
  padding: 15px 17px; cursor: pointer; transition: all 0.18s; position: relative; overflow: hidden;
}
.account-card:hover {
  border-color: rgba(21,62,237,0.3); background: rgba(21,62,237,0.03);
  transform: translateY(-2px); box-shadow: 0 6px 20px rgba(21,62,237,0.08);
}
.card-top { display: flex; justify-content: space-between; align-items: center; margin-bottom: 10px; }
.card-name { font-size: 13.5px; font-weight: 700; color: var(--white); }
.card-badges { display: flex; gap: 5px; }
.badge {
  font-size: 9.5px; font-weight: 800; padding: 2px 7px; border-radius: 4px;
  background: var(--blue-dim); color: var(--blue); border: 1px solid rgba(21,62,237,0.2);
}
.badge.c { background: var(--red-dim); color: var(--red); border-color: rgba(220,38,38,0.2); }
.badge.g { background: var(--green-dim); color: var(--green); border-color: rgba(0,166,126,0.2); }
.score-bars { display: flex; gap: 10px; }
.sb-item { flex: 1; }
.sb-label { font-size: 9.5px; color: var(--muted); margin-bottom: 4px; }
.sb-track { height: 3px; background: rgba(0,0,0,0.07); border-radius: 2px; overflow: hidden; }
.sb-fill { height: 100%; border-radius: 2px; }
.sb-fill.icp  { background: linear-gradient(90deg, var(--blue), var(--purple)); }
.sb-fill.deal { background: linear-gradient(90deg, var(--green), #00A880); }

/* ══════════════════════════════════════════
   TOAST
══════════════════════════════════════════ */
#toast {
  position: fixed; bottom: 24px; right: 24px;
  padding: 12px 20px; border-radius: 10px;
  font-size: 13px; font-weight: 700;
  background: linear-gradient(135deg, var(--blue), var(--blue-light));
  color: #fff; transform: translateY(80px); opacity: 0;
  transition: all 0.3s cubic-bezier(0.34,1.56,0.64,1);
  z-index: 1000; box-shadow: 0 8px 24px rgba(21,62,237,0.35);
}
#toast.show { transform: translateY(0); opacity: 1; }
#toast.error {
  background: linear-gradient(135deg, var(--red), #D94035);
  box-shadow: 0 8px 24px rgba(246,87,74,0.35);
}

::-webkit-scrollbar { width: 3px; height: 3px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: var(--border-hi); border-radius: 2px; }

/* ══════════════════════════════════════════
   BOARD REPORT BUTTON
══════════════════════════════════════════ */
.btn-board {
  padding: 8px 16px; border-radius: 8px;
  background: linear-gradient(135deg, var(--purple), #9B6BFF);
  color: #fff; border: none; cursor: pointer;
  font-size: 12px; font-weight: 700;
  transition: all 0.18s; box-shadow: 0 4px 14px rgba(123,92,245,0.3);
  white-space: nowrap; align-self: flex-start;
}
.btn-board:hover { transform: translateY(-1px); box-shadow: 0 6px 20px rgba(123,92,245,0.45); }
.btn-board:disabled { opacity: 0.4; cursor: not-allowed; transform: none; }

/* ══════════════════════════════════════════
   PIPELINE RADAR
══════════════════════════════════════════ */
#tab-radar { overflow-y: auto; }
#radar-wrap { padding: 28px 36px; }
.radar-title { font-size: 22px; font-weight: 800; color: var(--white); letter-spacing: -0.4px; }
.radar-sub { font-size: 12px; color: var(--muted2); margin-top: 4px; margin-bottom: 22px; }
.radar-layout { display: flex; gap: 22px; align-items: flex-start; }
.radar-svg-container {
  flex: 1; max-width: 580px;
  background: var(--bg2); border: 1px solid var(--border);
  border-radius: 16px; overflow: hidden; padding: 10px;
}
#radar-svg { width: 100%; height: auto; display: block; }
@keyframes radarSweep { from { transform: rotate(0deg); } to { transform: rotate(360deg); } }
#radar-sweep { animation: radarSweep 4s linear infinite; transform-origin: 300px 300px; }
.radar-sidebar { width: 230px; flex-shrink: 0; display: flex; flex-direction: column; gap: 14px; }
.radar-panel {
  background: var(--bg2); border: 1px solid var(--border);
  border-radius: 10px; padding: 16px;
}
.radar-panel-title { font-size: 9.5px; font-weight: 800; color: var(--muted2); text-transform: uppercase; letter-spacing: 2px; margin-bottom: 12px; }
.legend-item { display: flex; align-items: center; gap: 9px; padding: 3px 0; font-size: 11.5px; color: var(--muted2); }
.legend-dot { width: 9px; height: 9px; border-radius: 50%; flex-shrink: 0; }
.rhi { display: flex; align-items: center; gap: 9px; padding: 6px 0; border-bottom: 1px solid var(--border); cursor: pointer; transition: opacity 0.15s; }
.rhi:last-child { border-bottom: none; }
.rhi:hover { opacity: 0.75; }
.rhi-dot { width: 7px; height: 7px; border-radius: 50%; flex-shrink: 0; }
.rhi-name { flex: 1; font-size: 11.5px; color: var(--text); white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }
.rhi-score { font-size: 10px; font-weight: 800; color: var(--muted); }

/* ══════════════════════════════════════════
   PITCH SIMULATOR
══════════════════════════════════════════ */
#tab-simulator { overflow: hidden; display: none; flex-direction: column; }
#tab-simulator.active { display: flex; }
#sim-setup {
  padding: 22px 32px 18px; flex-shrink: 0;
  border-bottom: 1px solid var(--border);
  background: var(--bg2);
}
.sim-title { font-size: 18px; font-weight: 800; color: var(--white); letter-spacing: -0.3px; }
.sim-sub { font-size: 12px; color: var(--muted2); margin-top: 3px; margin-bottom: 14px; }
.sim-controls { display: flex; gap: 12px; align-items: flex-end; }
.sim-select-wrap { flex: 1; }
.sim-select-wrap label { display: block; font-size: 9.5px; font-weight: 800; color: var(--muted2); text-transform: uppercase; letter-spacing: 1.8px; margin-bottom: 7px; }
.sim-select-wrap select {
  width: 100%; padding: 9px 12px; border-radius: 8px;
  background: #F8F9FE; border: 1px solid var(--border);
  color: var(--text); font-size: 13px; font-family: var(--font); outline: none;
}
.sim-select-wrap select:focus { border-color: rgba(246,87,74,0.5); }
.btn-sim-start {
  padding: 9px 18px; border-radius: 8px;
  background: linear-gradient(135deg, #A52A2A, var(--red));
  color: #fff; border: none; cursor: pointer;
  font-size: 13px; font-weight: 700;
  transition: all 0.18s; box-shadow: 0 4px 14px rgba(246,87,74,0.3);
  white-space: nowrap;
}
.btn-sim-start:hover { transform: translateY(-1px); box-shadow: 0 6px 20px rgba(246,87,74,0.45); }
#sim-persona-bar {
  padding: 8px 32px; flex-shrink: 0;
  display: none; align-items: center; gap: 14px;
  border-bottom: 1px solid var(--border);
  background: rgba(246,87,74,0.04);
}
#sim-persona-bar.active { display: flex; }
.sim-persona {
  display: flex; align-items: center; gap: 10px;
  padding: 6px 12px; border-radius: 7px;
  background: var(--red-dim); border: 1px solid rgba(246,87,74,0.22);
}
.sim-persona-icon { font-size: 18px; }
.sim-persona-name { font-size: 12px; font-weight: 700; color: var(--red); }
.sim-persona-sub { font-size: 10.5px; color: var(--muted2); }
.btn-score {
  padding: 6px 13px; border-radius: 6px;
  background: transparent; border: 1px solid rgba(0,212,160,0.3);
  color: var(--green); font-size: 11px; font-weight: 700; cursor: pointer;
  transition: all 0.15s; white-space: nowrap; margin-left: auto;
}
.btn-score:hover { background: var(--green-dim); }
#sim-messages {
  flex: 1; overflow-y: auto; padding: 20px 32px;
  display: flex; flex-direction: column; gap: 14px;
}
#sim-messages::-webkit-scrollbar { width: 3px; }
#sim-messages::-webkit-scrollbar-thumb { background: var(--border-hi); border-radius: 2px; }
.sim-note { text-align: center; font-size: 10.5px; color: var(--muted); padding: 2px 0; }
.sim-opp .msg-bubble { background: rgba(246,87,74,0.07) !important; border: 1px solid rgba(246,87,74,0.18) !important; }
.sim-opp .msg-role { color: var(--red) !important; }
#sim-input-area {
  padding: 12px 32px 16px; border-top: 1px solid var(--border);
  background: var(--bg2); flex-shrink: 0;
}
#sim-input-row { display: flex; gap: 10px; align-items: flex-end; }
#sim-input {
  flex: 1; padding: 10px 14px; border-radius: 10px;
  background: #F8F9FE; border: 1px solid var(--border);
  color: var(--text); font-size: 13px; font-family: var(--font);
  resize: none; outline: none; min-height: 42px; max-height: 120px; line-height: 1.5;
  transition: border-color 0.18s;
}
#sim-input:focus { border-color: rgba(246,87,74,0.4); }
#sim-send-btn {
  padding: 10px 18px; border-radius: 10px;
  background: linear-gradient(135deg, #A52A2A, var(--red));
  color: #fff; border: none; cursor: pointer;
  font-size: 13px; font-weight: 700; transition: all 0.18s;
  box-shadow: 0 4px 12px rgba(246,87,74,0.3);
}
#sim-send-btn:hover { transform: translateY(-1px); }
#sim-send-btn:disabled { opacity: 0.38; cursor: not-allowed; transform: none; }

/* ══════════════════════════════════════════
   SIGNAL FEED
══════════════════════════════════════════ */
#tab-signals { overflow-y: auto; }
#signals-wrap { padding: 28px 36px; max-width: 860px; }
.signals-title { font-size: 22px; font-weight: 800; color: var(--white); letter-spacing: -0.4px; }
.signals-sub { font-size: 12px; color: var(--muted2); margin-top: 4px; margin-bottom: 22px; }
.signal-card {
  background: var(--bg2); border: 1px solid var(--border);
  border-radius: 10px; padding: 15px 18px; margin-bottom: 10px;
  display: flex; align-items: flex-start; gap: 14px;
  transition: border-color 0.18s, transform 0.15s;
}
.signal-card:hover { transform: translateX(4px); border-color: var(--border-hi); }
.signal-card.urg { border-color: rgba(246,87,74,0.3); background: rgba(246,87,74,0.03); }
.signal-card.amb { border-color: rgba(245,166,35,0.2); }
.sc-ico {
  width: 36px; height: 36px; border-radius: 8px; flex-shrink: 0;
  display: flex; align-items: center; justify-content: center; font-size: 16px;
}
.sc-ico.red   { background: var(--red-dim);   border: 1px solid rgba(246,87,74,0.2); }
.sc-ico.amber { background: var(--amber-dim); border: 1px solid rgba(245,166,35,0.2); }
.sc-ico.blue  { background: var(--blue-dim);  border: 1px solid rgba(21,62,237,0.2); }
.sc-body { flex: 1; min-width: 0; }
.sc-co { font-size: 14px; font-weight: 700; color: var(--white); }
.sc-txt { font-size: 12px; color: var(--muted2); margin-top: 3px; line-height: 1.5; }
.sc-right { display: flex; flex-direction: column; align-items: flex-end; gap: 7px; flex-shrink: 0; }
.sc-tag { font-size: 9.5px; font-weight: 800; padding: 3px 8px; border-radius: 4px; white-space: nowrap; }
.sc-tag.red   { background: var(--red-dim);   color: var(--red); }
.sc-tag.amber { background: var(--amber-dim); color: var(--amber); }
.sc-tag.blue  { background: var(--blue-dim2); color: var(--blue); }
.sc-act {
  padding: 4px 11px; border-radius: 5px; font-size: 10.5px; font-weight: 700;
  background: transparent; border: 1px solid var(--border);
  color: var(--muted2); cursor: pointer; transition: all 0.15s; white-space: nowrap;
}
.sc-act:hover { border-color: rgba(21,62,237,0.4); color: var(--blue); background: var(--blue-dim); }

/* ══════════════════════════════════════════
   PARTNERSHIP INTELLIGENCE
══════════════════════════════════════════ */
#tab-partners { overflow-y: auto; }
#partners-wrap { padding: 28px 36px; max-width: 900px; }
.partners-title { font-size: 22px; font-weight: 800; color: var(--white); letter-spacing: -0.4px; }
.partners-sub { font-size: 12px; color: var(--muted2); margin-top: 4px; margin-bottom: 26px; }

.partner-form-card {
  background: var(--bg2); border: 1px solid var(--border);
  border-radius: 12px; padding: 22px; margin-bottom: 22px;
}
.partner-form-title { font-size: 11px; font-weight: 800; color: var(--muted2); text-transform: uppercase; letter-spacing: 2px; margin-bottom: 16px; }
.partner-form-row { display: flex; gap: 12px; align-items: flex-end; flex-wrap: wrap; }
.pf-group { display: flex; flex-direction: column; gap: 6px; }
.pf-group label { font-size: 9.5px; font-weight: 800; color: var(--muted); text-transform: uppercase; letter-spacing: 1.5px; }
.pf-group input, .pf-group select, .pf-group textarea {
  padding: 9px 13px; border-radius: 7px;
  background: #F8F9FE; border: 1px solid var(--border);
  color: var(--text); font-size: 13px; font-family: var(--font); outline: none;
  transition: border-color 0.18s;
}
.pf-group input:focus, .pf-group select:focus, .pf-group textarea:focus { border-color: rgba(123,92,245,0.5); }
.pf-group.grow { flex: 1; min-width: 200px; }
.pf-group select option { background: var(--bg3); }
.btn-validate {
  padding: 10px 20px; border-radius: 8px;
  background: linear-gradient(135deg, var(--purple), #9B6BFF);
  color: #fff; border: none; cursor: pointer;
  font-size: 13px; font-weight: 700;
  transition: all 0.18s; box-shadow: 0 4px 14px rgba(123,92,245,0.3);
  white-space: nowrap; align-self: flex-end;
}
.btn-validate:hover { transform: translateY(-1px); box-shadow: 0 6px 20px rgba(123,92,245,0.45); }
.btn-validate:disabled { opacity: 0.4; cursor: not-allowed; transform: none; }

/* Validation result card */
#partner-result { display: none; }
.verdict-banner {
  padding: 14px 20px; border-radius: 10px; margin-bottom: 20px;
  display: flex; align-items: center; gap: 14px;
}
.verdict-banner.strong  { background: rgba(0,212,160,0.1);  border: 1px solid rgba(0,212,160,0.3); }
.verdict-banner.potential { background: rgba(21,62,237,0.1); border: 1px solid rgba(21,62,237,0.3); }
.verdict-banner.weak    { background: rgba(245,166,35,0.1);  border: 1px solid rgba(245,166,35,0.3); }
.verdict-banner.not-rec { background: rgba(246,87,74,0.1);   border: 1px solid rgba(246,87,74,0.3); }
.verdict-icon { font-size: 28px; }
.verdict-label { font-size: 16px; font-weight: 800; }
.verdict-label.strong  { color: var(--green); }
.verdict-label.potential { color: var(--blue); }
.verdict-label.weak    { color: var(--amber); }
.verdict-label.not-rec { color: var(--red); }
.verdict-reason { font-size: 12.5px; color: var(--muted2); margin-top: 3px; }
.verdict-score { margin-left: auto; font-size: 32px; font-weight: 900; color: var(--white); }
.verdict-score span { font-size: 14px; font-weight: 500; color: var(--muted); }

.dim-grid { display: grid; grid-template-columns: repeat(auto-fill, minmax(260px, 1fr)); gap: 12px; margin-bottom: 18px; }
.dim-card {
  background: var(--bg2); border: 1px solid var(--border);
  border-radius: 9px; padding: 14px;
}
.dim-name { font-size: 10px; font-weight: 800; color: var(--muted2); text-transform: uppercase; letter-spacing: 1.5px; margin-bottom: 8px; }
.dim-score-row { display: flex; align-items: center; gap: 10px; margin-bottom: 8px; }
.dim-score-num { font-size: 22px; font-weight: 900; color: var(--white); }
.dim-score-max { font-size: 12px; color: var(--muted); }
.dim-bar-track { flex: 1; height: 4px; background: rgba(0,0,0,0.07); border-radius: 2px; overflow: hidden; }
.dim-bar-fill { height: 100%; border-radius: 2px; transition: width 1s cubic-bezier(0.4,0,0.2,1); }
.dim-finding { font-size: 11.5px; color: var(--muted2); line-height: 1.5; }

.partner-meta-grid { display: grid; grid-template-columns: 1fr 1fr; gap: 14px; margin-bottom: 18px; }
.pmeta-card {
  background: var(--bg2); border: 1px solid var(--border);
  border-radius: 9px; padding: 14px;
}
.pmeta-label { font-size: 9.5px; font-weight: 800; color: var(--muted2); text-transform: uppercase; letter-spacing: 1.5px; margin-bottom: 8px; }
.pmeta-value { font-size: 12.5px; color: var(--text); line-height: 1.5; }
.pmeta-tag {
  display: inline-block; padding: 3px 9px; border-radius: 4px; font-size: 10px; font-weight: 700;
  background: var(--purple); color: #fff; margin: 2px 3px 2px 0;
}

.partner-history-title { font-size: 10px; font-weight: 800; color: var(--muted2); text-transform: uppercase; letter-spacing: 2px; margin: 24px 0 12px; }
.ph-row {
  display: flex; align-items: center; gap: 12px; padding: 9px 12px;
  background: var(--bg2); border: 1px solid var(--border);
  border-radius: 8px; margin-bottom: 7px; cursor: pointer; transition: border-color 0.15s;
}
.ph-row:hover { border-color: var(--border-hi); }
.ph-name { flex: 1; font-size: 13px; font-weight: 600; color: var(--text); }
.ph-market { font-size: 10px; color: var(--muted); }
.ph-verdict {
  font-size: 9.5px; font-weight: 800; padding: 2px 8px; border-radius: 4px;
}
.ph-verdict.strong   { background: rgba(0,212,160,0.12); color: var(--green); }
.ph-verdict.potential { background: var(--blue-dim2); color: var(--blue); }
.ph-verdict.weak     { background: rgba(245,166,35,0.12); color: var(--amber); }
.ph-score { font-size: 13px; font-weight: 800; color: var(--white); }

/* ── Deal Velocity Cards ─────────────────────────────────────────── */
.dv-card {
  background: var(--bg2); border: 1px solid var(--border);
  border-radius: 10px; padding: 14px 16px; cursor: pointer; transition: border-color 0.15s, transform 0.15s, box-shadow 0.15s;
  position: relative; overflow: hidden;
}
.dv-card:hover { border-color: var(--border-hi); transform: translateY(-1px); }
.dv-card.stale-hot { border-color: rgba(246,87,74,0.4); }
.dv-card.stale-warm { border-color: rgba(245,166,35,0.3); }
.dv-stale-bar {
  position: absolute; top: 0; left: 0; height: 3px; border-radius: 10px 10px 0 0;
}
.dv-header { display: flex; align-items: flex-start; justify-content: space-between; margin-bottom: 10px; }
.dv-name { font-size: 14px; font-weight: 700; color: var(--white); }
.dv-country { font-size: 10px; color: var(--muted); margin-top: 2px; }
.dv-win { font-size: 22px; font-weight: 900; }
.dv-win-label { font-size: 9px; color: var(--muted); text-transform: uppercase; letter-spacing: 1px; }
.dv-metrics { display: flex; gap: 10px; margin-bottom: 10px; }
.dv-metric { flex: 1; background: #F4F6FC; border-radius: 6px; padding: 7px 10px; text-align: center; }
.dv-metric-val { font-size: 13px; font-weight: 800; color: var(--white); }
.dv-metric-label { font-size: 9px; color: var(--muted); text-transform: uppercase; letter-spacing: 0.8px; margin-top: 1px; }
.dv-offering { font-size: 10.5px; color: var(--muted2); margin-bottom: 6px; }
.dv-stale-tag {
  display: inline-flex; align-items: center; gap: 4px;
  font-size: 9.5px; font-weight: 700; padding: 2px 7px; border-radius: 4px;
}
.dv-stale-tag.fresh { background: rgba(0,212,160,0.1); color: var(--green); }
.dv-stale-tag.warm  { background: rgba(245,166,35,0.12); color: var(--amber); }
.dv-stale-tag.hot   { background: rgba(246,87,74,0.12); color: var(--red); }
.dv-buyer { font-size: 10px; color: var(--muted); margin-top: 6px; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }

@keyframes pulse {
  0%,100% { opacity: 1; } 50% { opacity: 0.4; }
}

/* ── Command Palette ─────────────────────────────────────────────── */
#cmd-overlay {
  display: none; position: fixed; inset: 0; background: rgba(15,23,42,0.4);
  backdrop-filter: blur(6px); z-index: 9000; align-items: flex-start; justify-content: center;
  padding-top: 12vh;
}
#cmd-overlay.open { display: flex; }
#cmd-box {
  width: 100%; max-width: 580px; background: var(--bg2); border: 1px solid var(--border);
  border-radius: 14px; overflow: hidden; box-shadow: 0 8px 40px rgba(21,62,237,0.12), 0 2px 8px rgba(0,0,0,0.08);
}
#cmd-input-wrap {
  display: flex; align-items: center; gap: 10px; padding: 14px 18px;
  border-bottom: 1px solid var(--border);
}
#cmd-input-wrap svg { flex-shrink: 0; color: var(--muted); }
#cmd-input {
  flex: 1; background: none; border: none; outline: none; font-size: 16px;
  color: var(--text); font-family: inherit; caret-color: var(--blue);
}
#cmd-input::placeholder { color: var(--muted); }
#cmd-kb { font-size: 10px; color: var(--muted); background: var(--bg3); border: 1px solid var(--border); padding: 2px 6px; border-radius: 4px; flex-shrink:0; }
#cmd-results { max-height: 380px; overflow-y: auto; padding: 6px; }
.cmd-section { font-size: 9px; font-weight: 800; color: var(--muted2); text-transform: uppercase; letter-spacing: 2px; padding: 8px 10px 4px; }
.cmd-item {
  display: flex; align-items: center; gap: 12px; padding: 10px 12px;
  border-radius: 8px; cursor: pointer; transition: background 0.1s;
}
.cmd-item:hover, .cmd-item.selected { background: var(--blue-dim); }
.cmd-item.selected { background: var(--blue-dim2); }
.cmd-icon { width: 28px; height: 28px; display: flex; align-items: center; justify-content: center; background: var(--bg3); border: 1px solid var(--border); border-radius: 7px; flex-shrink: 0; color: var(--blue); }
.cmd-icon i { width: 14px; height: 14px; }
.cmd-label { flex: 1; font-size: 13.5px; color: var(--text); font-weight: 500; }
.cmd-sub { font-size: 10.5px; color: var(--muted); }
.cmd-arrow { font-size: 11px; color: var(--muted2); }
#cmd-empty { padding: 24px; text-align: center; font-size: 13px; color: var(--muted); }
#cmd-footer {
  display: flex; gap: 16px; padding: 8px 16px;
  border-top: 1px solid var(--border); background: var(--bg3); font-size: 10px; color: var(--muted2);
}
.cmd-hint { display: flex; align-items: center; gap: 5px; }
.cmd-hint kbd { background: var(--bg2); border: 1px solid var(--border); border-radius: 3px; padding: 1px 5px; font-family: inherit; color: var(--text); }

/* ── Outreach Modal ──────────────────────────────────────────────── */
#outreach-overlay {
  display: none; position: fixed; inset: 0; background: rgba(0,0,0,0.7);
  backdrop-filter: blur(4px); z-index: 9100; align-items: center; justify-content: center;
}
#outreach-overlay.open { display: flex; }
#outreach-box {
  width: 100%; max-width: 620px; background: #0e0e1a; border: 1px solid rgba(255,255,255,0.12);
  border-radius: 14px; overflow: hidden; box-shadow: 0 24px 80px rgba(0,0,0,0.7);
  display: flex; flex-direction: column; max-height: 80vh;
}
#outreach-header {
  padding: 18px 20px 14px; border-bottom: 1px solid rgba(255,255,255,0.07);
  display: flex; align-items: flex-start; justify-content: space-between;
}
#outreach-title { font-size: 16px; font-weight: 800; color: var(--white); }
#outreach-subtitle { font-size: 11px; color: var(--muted); margin-top: 3px; }
#outreach-close { background: none; border: none; color: var(--muted); cursor: pointer; font-size: 18px; padding: 0; line-height: 1; }
#outreach-body { flex: 1; overflow-y: auto; padding: 18px 20px; }
#outreach-text {
  font-size: 13.5px; color: var(--text); line-height: 1.7; white-space: pre-wrap;
  min-height: 80px;
}
#outreach-cursor { display: inline-block; width: 2px; height: 14px; background: #153EED; margin-left: 1px; animation: blink 1s infinite; vertical-align: middle; }
@keyframes blink { 0%,100% { opacity:1; } 50% { opacity:0; } }
#outreach-loading { display: flex; align-items: center; gap: 8px; color: var(--muted); font-size: 12px; }
#outreach-loading .dot { width: 6px; height: 6px; border-radius: 50%; background: #153EED; animation: bounce 1.2s infinite; }
#outreach-loading .dot:nth-child(2) { animation-delay: 0.2s; }
#outreach-loading .dot:nth-child(3) { animation-delay: 0.4s; }
@keyframes bounce { 0%,80%,100% { transform: scale(0.6); } 40% { transform: scale(1); } }
#outreach-actions {
  padding: 14px 20px; border-top: 1px solid rgba(255,255,255,0.07);
  display: flex; gap: 10px; align-items: center;
}
#outreach-copy-btn {
  background: #153EED; color: #fff; border: none; border-radius: 8px;
  padding: 10px 20px; font-size: 13px; font-weight: 700; cursor: pointer; font-family: inherit;
  transition: background 0.15s;
}
#outreach-copy-btn:hover { background: #1a4af0; }
#outreach-chat-btn {
  background: rgba(255,255,255,0.06); color: var(--text); border: 1px solid var(--border);
  border-radius: 8px; padding: 10px 16px; font-size: 13px; cursor: pointer; font-family: inherit;
}
#outreach-regen-btn {
  background: none; color: var(--muted); border: none; font-size: 12px; cursor: pointer; font-family: inherit; margin-left: auto;
}
</style>
<script src="https://unpkg.com/lucide@latest/dist/umd/lucide.min.js"></script>
</head>
<body>

<!-- ══════════════════ COMMAND PALETTE ══════════════════ -->
<div id="cmd-overlay" onclick="if(event.target===this)closePalette()">
  <div id="cmd-box">
    <div id="cmd-input-wrap">
      <svg width="16" height="16" viewBox="0 0 16 16" fill="none" stroke="currentColor" stroke-width="1.8"><circle cx="6.5" cy="6.5" r="4.5"/><line x1="10.5" y1="10.5" x2="14" y2="14"/></svg>
      <input id="cmd-input" placeholder="Search command or account…" oninput="cmdFilter()" onkeydown="cmdKey(event)" autocomplete="off" spellcheck="false">
      <span id="cmd-kb">ESC</span>
    </div>
    <div id="cmd-results"></div>
    <div id="cmd-footer">
      <span class="cmd-hint"><kbd>↑↓</kbd> navigate</span>
      <span class="cmd-hint"><kbd>↵</kbd> execute</span>
      <span class="cmd-hint"><kbd>ESC</kbd> close</span>
    </div>
  </div>
</div>

<!-- ══════════════════ OUTREACH MODAL ══════════════════ -->
<div id="outreach-overlay" onclick="if(event.target===this)closeOutreachModal()">
  <div id="outreach-box">
    <div id="outreach-header">
      <div>
        <div id="outreach-title">Generating outreach…</div>
        <div id="outreach-subtitle">LinkedIn · ready to send</div>
      </div>
      <button id="outreach-close" onclick="closeOutreachModal()">✕</button>
    </div>
    <div id="outreach-body">
      <div id="outreach-loading"><div class="dot"></div><div class="dot"></div><div class="dot"></div><span style="margin-left:4px">Writing message…</span></div>
      <div id="outreach-text" style="display:none"></div>
    </div>
    <div id="outreach-actions">
      <button id="outreach-copy-btn" onclick="copyOutreach()" style="display:none"><i data-lucide="clipboard" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i> Copy message</button>
      <button id="outreach-chat-btn" onclick="openOutreachInChat()" style="display:none">Open in chat →</button>
      <button id="outreach-regen-btn" onclick="regenOutreach()" style="display:none">↻ Regenerate</button>
    </div>
  </div>
</div>

<!-- ══════════════════ SIDEBAR ══════════════════ -->
<div id="sidebar">
  <div id="logo">
    <div class="logo-mark">J</div>
    <div class="logo-text">
      <div class="brand">JAKALA</div>
      <div class="sub">GTM OS — Close deals · Build pipeline</div>
    </div>
    <div class="live-pill"><div class="live-dot"></div>LIVE</div>
  </div>

  <div class="sidebar-label">Execute</div>
  <button class="nav-btn" onclick="showTab('dashboard')" id="nav-dashboard">
    <span class="icon"><svg viewBox="0 0 16 16" width="14" height="14" fill="none" stroke="currentColor" stroke-width="1.5"><rect x="1" y="1" width="6" height="6" rx="1"/><rect x="9" y="1" width="6" height="6" rx="1"/><rect x="1" y="9" width="6" height="6" rx="1"/><rect x="9" y="9" width="6" height="6" rx="1"/></svg></span> Deal Board
  </button>
  <div class="sidebar-label" style="margin-top:4px">Find &amp; Target</div>
  <button class="nav-btn" onclick="showTab('radar')" id="nav-radar">
    <span class="icon"><svg viewBox="0 0 16 16" width="14" height="14" fill="none" stroke="currentColor" stroke-width="1.5"><circle cx="8" cy="8" r="6"/><circle cx="8" cy="8" r="3"/><circle cx="8" cy="8" r="1" fill="currentColor"/></svg></span> Lead Radar
  </button>
  <button class="nav-btn" onclick="showTab('signals')" id="nav-signals">
    <span class="icon"><svg viewBox="0 0 16 16" width="14" height="14" fill="none" stroke="currentColor" stroke-width="1.5"><polyline points="1,9 4,4 7,11 10,5 13,8 15,6"/></svg></span> Buying Signals
  </button>
  <button class="nav-btn" onclick="showTab('accounts')" id="nav-accounts">
    <span class="icon"><svg viewBox="0 0 16 16" width="14" height="14" fill="none" stroke="currentColor" stroke-width="1.5"><rect x="2" y="4" width="12" height="11" rx="1"/><path d="M6 15V9h4v6"/><path d="M2 7h12"/><path d="M5 1h6v3H5z"/></svg></span> Accounts
  </button>

  <div class="sidebar-label" style="margin-top:8px">Pitch &amp; Close</div>
  <button class="nav-btn active" onclick="showTab('chat')" id="nav-chat">
    <span class="icon"><svg viewBox="0 0 16 16" width="14" height="14" fill="none" stroke="currentColor" stroke-width="1.5"><path d="M14 10a2 2 0 01-2 2H5l-3 2V4a2 2 0 012-2h8a2 2 0 012 2z"/></svg></span> GTM Assistant
  </button>
  <button class="nav-btn" onclick="showTab('notes')" id="nav-notes">
    <span class="icon"><svg viewBox="0 0 16 16" width="14" height="14" fill="none" stroke="currentColor" stroke-width="1.5"><rect x="2" y="1" width="12" height="14" rx="1"/><line x1="5" y1="5" x2="11" y2="5"/><line x1="5" y1="8" x2="11" y2="8"/><line x1="5" y1="11" x2="8" y2="11"/></svg></span> Meeting Notes
  </button>
  <button class="nav-btn" onclick="showTab('simulator')" id="nav-simulator">
    <span class="icon"><svg viewBox="0 0 16 16" width="14" height="14" fill="none" stroke="currentColor" stroke-width="1.5"><rect x="1" y="2" width="14" height="10" rx="1"/><line x1="8" y1="12" x2="8" y2="15"/><line x1="5" y1="15" x2="11" y2="15"/></svg></span> Pitch Builder
  </button>
  <button class="nav-btn" onclick="showTab('partners')" id="nav-partners">
    <span class="icon"><svg viewBox="0 0 16 16" width="14" height="14" fill="none" stroke="currentColor" stroke-width="1.5"><path d="M9.5 6.5l1-1a3 3 0 014.24 4.24l-2 2a3 3 0 01-4.24-4.24"/><path d="M6.5 9.5l-1 1a3 3 0 01-4.24-4.24l2-2a3 3 0 014.24 4.24"/></svg></span> Partnerships
  </button>
  <button class="nav-btn" onclick="insertSkill('websiteopt');showTab('chat')" id="nav-websiteopt">
    <span class="icon"><svg viewBox="0 0 16 16" width="14" height="14" fill="none" stroke="currentColor" stroke-width="1.5"><circle cx="7" cy="7" r="5"/><line x1="10.5" y1="10.5" x2="14" y2="14"/></svg></span> Website Optimizer
  </button>

  <button class="sidebar-accounts-toggle" id="skills-toggle" onclick="toggleSkillsPanel()">
    <span style="display:flex;align-items:center;gap:8px"><svg viewBox="0 0 16 16" width="14" height="14" fill="none" stroke="currentColor" stroke-width="1.5"><circle cx="8" cy="8" r="6"/><line x1="8" y1="5" x2="8" y2="8"/><line x1="8" y1="11" x2="8.01" y2="11"/></svg> Quick Skills</span>
    <svg id="skills-chevron" viewBox="0 0 16 16" width="12" height="12" fill="none" stroke="currentColor" stroke-width="2" style="transition:transform 0.2s;flex-shrink:0"><polyline points="4,6 8,10 12,6"/></svg>
  </button>
  <div id="skills-panel" style="max-height:0;overflow:hidden;transition:max-height 0.3s ease">
    <button class="skill-btn" onclick="insertSkill('contact')" style="background:#111;color:#fff;font-weight:700;"><i data-lucide="zap" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i> Who To Contact Today</button>
    <button class="skill-btn" onclick="insertSkill('blueprint')" style="background:linear-gradient(135deg,rgba(21,62,237,0.18),rgba(123,92,245,0.12));border-color:rgba(21,62,237,0.4);color:#6B8EF7;font-weight:700;"><i data-lucide="building-2" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i> Deal Blueprint</button>
    <button class="skill-btn" onclick="insertSkill('outreach')"><i data-lucide="mail" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i> Outreach Generator</button>
    <button class="skill-btn" onclick="insertSkill('brief')"><i data-lucide="clipboard" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i> Pre-Meeting Brief</button>
    <button class="skill-btn" onclick="insertSkill('revenue')"><i data-lucide="dollar-sign" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i> Revenue Simulation</button>
    <button class="skill-btn" onclick="insertSkill('pitch')"><i data-lucide="target" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i> Pitch Partner</button>
    <button class="skill-btn" onclick="insertSkill('signal')"><i data-lucide="zap" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i> Signal to Action</button>
    <button class="skill-btn" onclick="insertSkill('prospect')"><i data-lucide="search" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i> Prospect Hunt</button>
    <button class="skill-btn" onclick="insertSkill('strategic')"><i data-lucide="star" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i> Strategic Framework</button>
    <button class="skill-btn" onclick="insertSkill('morning')">Morning Briefing</button>
    <button class="skill-btn" onclick="insertSkill('websiteopt')"><i data-lucide="search" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i> Website Optimizer</button>
  </div>

  <button class="sidebar-accounts-toggle" id="accounts-toggle" onclick="toggleAccountsPanel()">
    <span style="display:flex;align-items:center;gap:8px"><svg viewBox="0 0 16 16" width="14" height="14" fill="none" stroke="currentColor" stroke-width="1.5"><rect x="2" y="4" width="12" height="11" rx="1"/><path d="M6 15V9h4v6"/><path d="M2 7h12"/><path d="M5 1h6v3H5z"/></svg> Accounts</span>
    <svg id="accounts-chevron" viewBox="0 0 16 16" width="12" height="12" fill="none" stroke="currentColor" stroke-width="2" style="transition:transform 0.2s;flex-shrink:0"><polyline points="4,6 8,10 12,6"/></svg>
  </button>
  <div id="accounts-panel" style="max-height:0;overflow:hidden;transition:max-height 0.3s ease">
    <input type="text" id="account-search" placeholder="Search accounts…" oninput="filterAccounts()" style="margin:4px 7px 6px;width:calc(100% - 14px)">
    <div id="account-list"></div>
  </div>
</div>

<!-- ══════════════════ MAIN ══════════════════ -->
<div id="main">

  <!-- ── DASHBOARD ── -->
  <div class="tab-pane" id="tab-dashboard">
    <div id="dash-wrap">
      <div class="dash-top">
        <div>
          <div class="dash-headline">Nordic Commercial Pipeline</div>
          <div class="dash-tagline">JAKALA · DK / NO / SE · Q1–Q2 2026</div>
        </div>
        <div style="display:flex;flex-direction:column;align-items:flex-end;gap:10px">
          <div style="display:flex;gap:8px;align-items:center">
            <button onclick="openPalette()" style="display:flex;align-items:center;gap:8px;background:rgba(255,255,255,0.05);border:1px solid rgba(255,255,255,0.1);color:var(--muted);border-radius:8px;padding:7px 14px;font-size:12px;cursor:pointer;font-family:inherit;transition:border-color 0.15s" onmouseenter="this.style.borderColor='rgba(255,255,255,0.2)'" onmouseleave="this.style.borderColor='rgba(255,255,255,0.1)'">
              <svg width="12" height="12" viewBox="0 0 16 16" fill="none" stroke="currentColor" stroke-width="2"><circle cx="6.5" cy="6.5" r="4.5"/><line x1="10.5" y1="10.5" x2="14" y2="14"/></svg>
              Search
              <kbd style="background:rgba(255,255,255,0.08);border-radius:3px;padding:1px 5px;font-size:10px;font-family:inherit">⌘K</kbd>
            </button>
            <button class="btn-board" id="board-report-btn" onclick="generateBoardReport()"><i data-lucide="bar-chart-2" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i> Board Report</button>
          </div>
          <div class="dash-meta">
            <div class="dm-label">Current time</div>
            <div class="dm-time" id="dash-time">--:--</div>
            <div class="dm-date" id="dash-date">Loading…</div>
          </div>
        </div>
      </div>

      <!-- Live status bar -->
      <div id="live-status-bar" style="display:flex;align-items:center;gap:12px;margin-bottom:14px;font-size:11px;color:var(--muted)">
        <div style="display:flex;align-items:center;gap:6px">
          <div id="live-dot" style="width:7px;height:7px;border-radius:50%;background:#00D4A0;box-shadow:0 0 6px #00D4A0;animation:pulse 2s infinite"></div>
          <span id="live-label">LIVE</span>
        </div>
        <span>·</span>
        <span>Last updated: <span id="live-last-updated">—</span></span>
        <span>·</span>
        <span>Next update in <span id="live-countdown" style="color:var(--text);font-weight:600">5:00</span></span>
        <span style="margin-left:auto">
          <button onclick="loadLiveDashboard(true)" style="background:none;border:1px solid var(--border);color:var(--muted);border-radius:4px;padding:2px 8px;font-size:10px;cursor:pointer;font-family:inherit">↻ Refresh now</button>
        </span>
      </div>

      <!-- ── DAILY ACTION PLAN ── -->
      <div id="daily-plan-section" style="margin-bottom:24px">
        <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:12px">
          <div>
            <div style="font-size:10px;font-weight:800;color:#153EED;text-transform:uppercase;letter-spacing:2px;margin-bottom:4px"><i data-lucide="zap" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i> Today's Action Plan</div>
            <div style="font-size:13px;color:var(--muted);font-weight:400" id="daily-plan-date">Generating plan…</div>
          </div>
          <button onclick="loadDailyPlan(true)" id="daily-plan-refresh" style="background:none;border:1px solid var(--border);color:var(--muted);border-radius:6px;padding:5px 12px;font-size:11px;cursor:pointer;font-family:inherit;transition:all .15s" onmouseenter="this.style.borderColor='#153EED';this.style.color='#153EED'" onmouseleave="this.style.borderColor='var(--border)';this.style.color='var(--muted)'">↻ Refresh</button>
        </div>
        <div id="daily-plan-cards" style="display:grid;grid-template-columns:repeat(3,1fr);gap:14px">
          <!-- Loading state -->
          <div style="grid-column:1/-1;background:rgba(21,62,237,0.04);border:1px dashed rgba(21,62,237,0.2);border-radius:12px;padding:32px;text-align:center;color:var(--muted);font-size:13px">
            <div style="font-size:20px;margin-bottom:8px"><i data-lucide="zap" style="width:20px;height:20px;display:inline-block;vertical-align:middle;"></i></div>
            Generating your action plan for today…
          </div>
        </div>
      </div>

      <!-- KPI row -->
      <div class="kpi-row">
        <div class="kpi">
          <div class="kpi-label">Total Pipeline</div>
          <div class="kpi-num"><span class="kpi-unit">€</span><span id="kpi-pipeline">0.0</span><span class="kpi-suffix">M</span></div>
          <div class="kpi-sub">Nordic only · unweighted</div>
          <div class="kpi-badge amber">● AMBER</div>
          <div class="kpi-bar"><div class="kpi-bar-fill" style="background:linear-gradient(90deg,#153EED,#7B5CF5);width:0" data-w="68%"></div></div>
        </div>
        <div class="kpi">
          <div class="kpi-label">Active Accounts</div>
          <div class="kpi-num"><span id="kpi-accounts">0</span></div>
          <div class="kpi-sub">ICP-scored opportunities</div>
          <div class="kpi-badge blue">NORDIC SCOPE</div>
          <div class="kpi-bar"><div class="kpi-bar-fill" style="background:linear-gradient(90deg,#4B6EF7,#7B5CF5);width:0" data-w="100%"></div></div>
        </div>
        <div class="kpi">
          <div class="kpi-label">Named Buyers</div>
          <div class="kpi-num"><span id="kpi-buyers">0</span></div>
          <div class="kpi-sub">Confirmed stakeholders</div>
          <div class="kpi-badge green">CONFIRMED</div>
          <div class="kpi-bar"><div class="kpi-bar-fill" style="background:linear-gradient(90deg,#00D4A0,#00A880);width:0" data-w="40%"></div></div>
        </div>
        <div class="kpi">
          <div class="kpi-label">Q2 Base Forecast</div>
          <div class="kpi-num"><span class="kpi-unit">€</span><span id="kpi-forecast">0</span><span class="kpi-suffix">K</span></div>
          <div class="kpi-sub">Probability-weighted · 0% closed</div>
          <div class="kpi-badge red">0 MEETINGS</div>
          <div class="kpi-bar"><div class="kpi-bar-fill" style="background:linear-gradient(90deg,#F6574A,#E04030);width:0" data-w="6%"></div></div>
        </div>
      </div>

      <!-- Two columns -->
      <div class="dash-cols">
        <!-- Top opportunities -->
        <div class="dash-card">
          <div class="dash-card-head">
            <div class="dash-card-title">Top Opportunities</div>
            <div class="dash-card-tag">By deal score · click to load</div>
          </div>
          <div id="top-opps"></div>
        </div>

        <!-- Hot signals -->
        <div class="dash-card">
          <div class="dash-card-head">
            <div class="dash-card-title">Hot Signals</div>
            <div class="dash-card-tag" id="signals-updated">Live fra intelligence</div>
          </div>
          <div id="signals-list"><div style="color:var(--muted);font-size:12px;padding:12px 0">Indlæser signaler…</div></div>
        </div>
      </div>

      <!-- Deal Velocity Cards -->
      <div class="dash-card" style="margin-top:18px">
        <div class="dash-card-head">
          <div class="dash-card-title">Top 6 Deals — Velocity</div>
          <div class="dash-card-tag">Win % · Weighted value · Days without action</div>
        </div>
        <div id="deal-velocity-cards" style="display:grid;grid-template-columns:repeat(auto-fill,minmax(280px,1fr));gap:12px;padding-top:4px">
          <div style="color:var(--muted);font-size:12px">Indlæser deals…</div>
        </div>
      </div>

      <!-- GTM Strategy split -->
      <div class="dash-card" style="margin-top:18px">
        <div class="dash-card-head">
          <div class="dash-card-title">GTM Strategy Distribution</div>
          <div class="dash-card-tag">Accounts by entry strategy</div>
        </div>
        <div class="strat-row">
          <div class="strat-card">
            <div class="strat-icon"><i data-lucide="lightbulb" style="width:20px;height:20px;display:inline-block;vertical-align:middle;"></i></div>
            <div class="strat-name">Data Revenue Unlock</div>
            <div class="strat-num" id="strat-dru">—</div>
            <div class="strat-sub">accounts</div>
            <div class="strat-bar-wrap"><div class="strat-bar" style="background:linear-gradient(90deg,#153EED,#7B5CF5);width:0" id="sb-dru"></div></div>
          </div>
          <div class="strat-card">
            <div class="strat-icon"><i data-lucide="bot" style="width:20px;height:20px;display:inline-block;vertical-align:middle;"></i></div>
            <div class="strat-name">AI Readiness Accelerator</div>
            <div class="strat-num" id="strat-ai">—</div>
            <div class="strat-sub">accounts</div>
            <div class="strat-bar-wrap"><div class="strat-bar" style="background:linear-gradient(90deg,#7B5CF5,#A080FF);width:0" id="sb-ai"></div></div>
          </div>
          <div class="strat-card">
            <div class="strat-icon"><i data-lucide="shopping-cart" style="width:20px;height:20px;display:inline-block;vertical-align:middle;"></i></div>
            <div class="strat-name">Commerce Optimization</div>
            <div class="strat-num" id="strat-co">—</div>
            <div class="strat-sub">accounts</div>
            <div class="strat-bar-wrap"><div class="strat-bar" style="background:linear-gradient(90deg,#00D4A0,#00A880);width:0" id="sb-co"></div></div>
          </div>
          <div class="strat-card">
            <div class="strat-icon"><i data-lucide="sparkles" style="width:20px;height:20px;display:inline-block;vertical-align:middle;"></i></div>
            <div class="strat-name">Experience Transformation</div>
            <div class="strat-num" id="strat-xt">—</div>
            <div class="strat-sub">accounts</div>
            <div class="strat-bar-wrap"><div class="strat-bar" style="background:linear-gradient(90deg,#F5A623,#E08800);width:0" id="sb-xt"></div></div>
          </div>
        </div>
      </div>

    </div>
  </div>

  <!-- ── CHAT ── -->
  <div class="tab-pane active" id="tab-chat">
    <div id="chat-header">
      <h2>GTM Assistant</h2>
      <span id="selected-account-tag">
        <span id="selected-account-name"></span>
        <button id="clear-account" onclick="clearAccount()">✕</button>
      </span>
      <button id="deck-btn" onclick="generateDeck()">⬇ Generate Deck</button>
    </div>
    <div id="messages">
      <div id="welcome">
        <div class="welcome-icon"><i data-lucide="zap" style="width:28px;height:28px;display:inline-block;vertical-align:middle;"></i></div>
        <div class="big">How can I help you win?</div>
        <div class="sub">Ask anything about the pipeline, accounts, or get outreach, briefs and commercial analysis. Select an account from the sidebar to pre-load context.</div>
        <div class="quick-chips">
          <div class="chip" onclick="insertSkill('morning')">Morning briefing</div>
          <div class="chip" onclick="insertSkill('warroom')">War room assessment</div>
          <div class="chip" onclick="insertSkill('forecast')">Q2 forecast</div>
          <div class="chip" onclick="insertSkill('outreach')">Write outreach</div>
          <div class="chip" onclick="insertSkill('signal')">Analyse signal</div>
        </div>
      </div>
    </div>
    <div id="typing-row" style="display:none">
      <div id="typing-indicator"><div class="dot"></div><div class="dot"></div><div class="dot"></div></div>
    </div>
    <div id="chat-input-area">
      <div id="input-row">
        <textarea id="chat-input" placeholder="Ask anything — or pick a Quick Skill from the sidebar…" rows="1"
          onkeydown="handleKey(event)" oninput="autoResize(this)"></textarea>
        <button id="send-btn" onclick="sendMessage()">Send ↑</button>
      </div>
      <div id="input-hint">Enter to send &nbsp;·&nbsp; Shift+Enter for new line</div>
    </div>
  </div>

  <!-- ── NOTES ── -->
  <div class="tab-pane" id="tab-notes">
    <div id="notes-pane">
      <h2><i data-lucide="file-text" style="width:20px;height:20px;display:inline-block;vertical-align:middle;"></i> Meeting Notes</h2>
      <p class="desc">Paste raw meeting notes below. GTM OS will summarise the meeting, extract next actions, and update account files automatically.</p>
      <div class="form-group">
        <label>Account</label>
        <select id="notes-account"><option value="">— Select account —</option></select>
      </div>
      <div class="form-group">
        <label>Meeting Notes</label>
        <textarea id="notes-text" placeholder="Paste your raw meeting notes here — attendees, discussion points, decisions, agreed actions, follow-ups…"></textarea>
      </div>
      <button class="btn-primary" onclick="processNotes()" id="process-btn">Process Notes</button>
      <div id="notes-result">
        <div class="result-section">
          <div class="result-label">Key Insight</div>
          <div id="result-insight" class="insight-box"></div>
        </div>
        <div class="result-section">
          <div class="result-label">Meeting Summary</div>
          <div id="result-summary" class="summary-text"></div>
        </div>
        <div class="result-section">
          <div class="result-label">Meeting Entry (appended to meetings.md)</div>
          <div id="result-meeting" class="file-preview"></div>
        </div>
        <div class="result-section">
          <div class="result-label">Updated Next Actions (replaces next-actions.md)</div>
          <div id="result-actions" class="file-preview"></div>
        </div>
        <div class="save-actions">
          <button class="btn-primary" onclick="saveNotes()" id="save-btn">Save to Account Files</button>
          <button class="btn-secondary" onclick="cancelNotes()">Cancel</button>
        </div>
      </div>
    </div>
  </div>

  <!-- ── ACCOUNTS ── -->
  <div class="tab-pane" id="tab-accounts">
    <div id="accounts-pane">
      <h2><i data-lucide="building" style="width:20px;height:20px;display:inline-block;vertical-align:middle;"></i> Accounts</h2>
      <p class="desc">All ICP-scored accounts in the Nordic GTM OS. Click any account to load it in the assistant.</p>
      <div class="filter-row">
        <button class="filter-btn active" onclick="filterGrid(this,'all')">All</button>
        <button class="filter-btn" onclick="filterGrid(this,'NO')">Norway</button>
        <button class="filter-btn" onclick="filterGrid(this,'DK')">Denmark</button>
        <button class="filter-btn" onclick="filterGrid(this,'SE')">Sweden</button>
        <button class="filter-btn" onclick="filterGrid(this,'scored')">ICP scored</button>
      </div>
      <div class="accounts-grid" id="accounts-grid"></div>
    </div>
  </div>


  <!-- ── PIPELINE RADAR ── -->
  <div class="tab-pane" id="tab-radar">
    <div id="radar-wrap">
      <div class="radar-title">Pipeline Radar</div>
      <div class="radar-sub">Each blip = one account &nbsp;·&nbsp; Distance from centre = deal score &nbsp;·&nbsp; Quadrant = GTM strategy &nbsp;·&nbsp; Colour = urgency &nbsp;·&nbsp; Click any blip to open account</div>
      <div class="radar-layout">
        <div class="radar-svg-container">
          <svg id="radar-svg" viewBox="0 0 600 600" xmlns="http://www.w3.org/2000/svg">
            <defs>
              <radialGradient id="sweepGrad" cx="0%" cy="0%" r="100%">
                <stop offset="0%" stop-color="#153EED" stop-opacity="0.55"/>
                <stop offset="100%" stop-color="#153EED" stop-opacity="0"/>
              </radialGradient>
            </defs>
            <!-- Rings -->
            <circle cx="300" cy="300" r="240" fill="none" stroke="rgba(255,255,255,0.04)" stroke-width="1"/>
            <circle cx="300" cy="300" r="155" fill="none" stroke="rgba(255,255,255,0.05)" stroke-width="1"/>
            <circle cx="300" cy="300" r="78"  fill="none" stroke="rgba(21,62,237,0.15)"   stroke-width="1"/>
            <!-- Axis lines -->
            <line x1="300" y1="55"  x2="300" y2="545" stroke="rgba(255,255,255,0.04)" stroke-width="1"/>
            <line x1="55"  y1="300" x2="545" y2="300" stroke="rgba(255,255,255,0.04)" stroke-width="1"/>
            <!-- Ring labels -->
            <text x="305" y="225" fill="rgba(0,212,160,0.4)" font-size="8" font-family="Inter,sans-serif">HOT</text>
            <text x="305" y="148" fill="rgba(75,110,247,0.4)" font-size="8" font-family="Inter,sans-serif">WARM</text>
            <text x="305" y="68"  fill="rgba(255,255,255,0.2)" font-size="8" font-family="Inter,sans-serif">COLD</text>
            <!-- Quadrant labels -->
            <text x="330" y="78"  fill="rgba(21,62,237,0.55)"   font-size="10" font-family="Inter,sans-serif" font-weight="600">Data Revenue Unlock</text>
            <text x="98"  y="78"  fill="rgba(0,212,160,0.55)"   font-size="10" font-family="Inter,sans-serif" font-weight="600">AI Readiness</text>
            <text x="330" y="528" fill="rgba(123,92,245,0.55)"  font-size="10" font-family="Inter,sans-serif" font-weight="600">Commerce Optimization</text>
            <text x="70"  y="528" fill="rgba(245,166,35,0.55)"  font-size="10" font-family="Inter,sans-serif" font-weight="600">Experience Transform.</text>
            <!-- Rotating sweep -->
            <g id="radar-sweep">
              <path d="M 300 300 L 300 60 A 240 240 0 0 1 470 130 Z" fill="url(#sweepGrad)"/>
            </g>
            <!-- Blips (populated by JS) -->
            <g id="radar-blips"></g>
            <!-- Tooltip (populated by JS) -->
            <g id="radar-tt" display="none">
              <rect id="tt-bg" rx="5" fill="rgba(4,4,15,0.95)" stroke="rgba(21,62,237,0.5)" stroke-width="1"/>
              <text id="tt-name" fill="white" font-size="11.5" font-family="Inter,sans-serif" font-weight="700"/>
              <text id="tt-deal" fill="#8080B0" font-size="10" font-family="Inter,sans-serif"/>
              <text id="tt-val"  fill="#4B6EF7" font-size="10.5" font-family="Inter,sans-serif" font-weight="700"/>
            </g>
          </svg>
        </div>
        <div class="radar-sidebar">
          <div class="radar-panel">
            <div class="radar-panel-title">Legend</div>
            <div class="legend-item"><span class="legend-dot" style="background:#00D4A0;box-shadow:0 0 5px #00D4A0"></span>Deal 8–10 — HOT</div>
            <div class="legend-item"><span class="legend-dot" style="background:#4B6EF7;box-shadow:0 0 5px #4B6EF7"></span>Deal 6–7 — WARM</div>
            <div class="legend-item"><span class="legend-dot" style="background:#F5A623;box-shadow:0 0 5px #F5A623"></span>Deal 4–5 — COOL</div>
            <div class="legend-item"><span class="legend-dot" style="background:#F6574A;box-shadow:0 0 5px #F6574A"></span>Deal &lt;4 — COLD</div>
          </div>
          <div class="radar-panel">
            <div class="radar-panel-title">Hottest Deals</div>
            <div id="radar-hot-list"></div>
          </div>
        </div>
      </div>
    </div>
  </div>

  <!-- ── PITCH SIMULATOR ── -->
  <div class="tab-pane" id="tab-simulator">
    <div id="sim-setup">
      <div class="sim-title">🎭 Pitch Simulator</div>
      <div class="sim-sub">Claude becomes the buyer. Practice your pitch before the real meeting — then get scored.</div>
      <div class="sim-controls">
        <div class="sim-select-wrap">
          <label>Target Account</label>
          <select id="sim-account"><option value="">— Select account to practice pitch on —</option></select>
        </div>
        <button class="btn-sim-start" onclick="startPitchSession()">▶ Start Session</button>
      </div>
    </div>
    <div id="sim-persona-bar">
      <div class="sim-persona">
        <div class="sim-persona-icon"><i data-lucide="target" style="width:20px;height:20px;display:inline-block;vertical-align:middle;"></i></div>
        <div>
          <div class="sim-persona-name" id="sim-persona-name">Buyer</div>
          <div class="sim-persona-sub" id="sim-persona-sub">Playing the decision maker</div>
        </div>
      </div>
      <button class="btn-score" onclick="scorePitch()"><i data-lucide="bar-chart-2" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i> Score My Pitch</button>
    </div>
    <div id="sim-messages">
      <div style="text-align:center;padding:60px 20px;color:var(--muted)">
        <div style="font-size:40px;margin-bottom:12px;opacity:0.4">🎭</div>
        <div style="font-size:15px;font-weight:700;color:var(--white)">No session active</div>
        <div style="font-size:13px;margin-top:6px">Select an account above and hit Start Session</div>
      </div>
    </div>
    <div id="sim-input-area" style="display:none">
      <div id="sim-input-row">
        <textarea id="sim-input" placeholder="Type your pitch…" rows="1"
          onkeydown="handleSimKey(event)" oninput="autoResize(this)"></textarea>
        <button id="sim-send-btn" onclick="sendSimMessage()">Send ↑</button>
      </div>
    </div>
  </div>

  <!-- ── PARTNERSHIPS ── -->
  <div class="tab-pane" id="tab-partners">
    <div id="partners-wrap">
      <div class="partners-title"><i data-lucide="handshake" style="width:20px;height:20px;display:inline-block;vertical-align:middle;"></i> Partnership Intelligence</div>
      <div class="partners-sub">Validate a potential partner against JAKALA's GTM strategy, positioning and revenue potential — or view monthly market analyses.</div>

      <!-- Validator form -->
      <div class="partner-form-card">
        <div class="partner-form-title">Validate a partner</div>
        <div class="partner-form-row">
          <div class="pf-group grow">
            <label>Company / Partner</label>
            <input type="text" id="partner-name" placeholder="e.g. Akeneo, Contentful, Algolia, Sitoo…" />
          </div>
          <div class="pf-group">
            <label>Market</label>
            <select id="partner-market">
              <option value="Nordic">Nordic (NO/SE/DK)</option>
              <option value="Norway">Norway</option>
              <option value="Sweden">Sweden</option>
              <option value="Denmark">Denmark</option>
            </select>
          </div>
          <div class="pf-group grow">
            <label>Context (optional)</label>
            <input type="text" id="partner-context" placeholder="e.g. 'PIM-vendor', 'meeting CEO next week', 'they work with Elkjøp'…" />
          </div>
          <button class="btn-validate" id="validate-btn" onclick="validatePartner()">Validate ↗</button>
        </div>
      </div>

      <!-- Result card -->
      <div id="partner-result">
        <div id="verdict-banner" class="verdict-banner">
          <div class="verdict-icon" id="verdict-icon">—</div>
          <div>
            <div class="verdict-label" id="verdict-label">—</div>
            <div class="verdict-reason" id="verdict-reason">—</div>
          </div>
          <div class="verdict-score"><span id="verdict-score">—</span><span>/10</span></div>
        </div>

        <div class="dim-grid" id="dim-grid"></div>

        <div class="partner-meta-grid">
          <div class="pmeta-card">
            <div class="pmeta-label">GTM Match</div>
            <div class="pmeta-value" id="pm-gtm"></div>
          </div>
          <div class="pmeta-card">
            <div class="pmeta-label">Buyer Overlap</div>
            <div class="pmeta-value" id="pm-buyers"></div>
          </div>
          <div class="pmeta-card">
            <div class="pmeta-label">Fælles Entry Offer</div>
            <div class="pmeta-value" id="pm-offer"></div>
          </div>
          <div class="pmeta-card">
            <div class="pmeta-label">Target Accounts</div>
            <div class="pmeta-value" id="pm-accounts"></div>
          </div>
          <div class="pmeta-card">
            <div class="pmeta-label">Første skridt</div>
            <div class="pmeta-value" id="pm-step"></div>
          </div>
          <div class="pmeta-card">
            <div class="pmeta-label">Risici</div>
            <div class="pmeta-value" id="pm-risks"></div>
          </div>
        </div>
      </div>

      <!-- History -->
      <div class="partner-history-title">Tidligere valideringer</div>
      <div id="partner-history"></div>
    </div>
  </div>

  <!-- ── SIGNAL FEED ── -->
  <div class="tab-pane" id="tab-signals">
    <div id="signals-wrap">
      <div class="signals-title"><i data-lucide="zap" style="width:18px;height:18px;display:inline-block;vertical-align:middle;"></i> Signal Feed</div>
      <div class="signals-sub">Live commercial signals from the Nordic market. Click "Write Outreach" to generate a message and load the account.</div>
      <div id="signals-content"></div>
    </div>
  </div>

</div>

<div id="toast"></div>

<script>
// ── State ─────────────────────────────────────────────────────────────────────
let messages = [];
let selectedAccount = null;
let allAccounts = [];
let processingResult = null;
let gridFilter = 'all';

// ── Boot ──────────────────────────────────────────────────────────────────────
async function boot() {
  updateClock();
  setInterval(updateClock, 1000);
  showTab('chat');
  try {
    const res = await fetch('/api/accounts');
    allAccounts = await res.json();
    renderAccountList(allAccounts);
    renderAccountsGrid(allAccounts);
    populateNotesSelect(allAccounts);
    populateSimSelect();
    renderDashboard(allAccounts);
    setTimeout(animateDashboard, 400);
    loadLiveDashboard();
    renderSignalFeed();
  } catch(e) {
    console.error('Boot error:', e);
  }
}
boot();

// ── Daily Action Plan ─────────────────────────────────────────────────────────
let _dailyPlanLoaded = false;

async function loadDailyPlan(force) {
  if (_dailyPlanLoaded && !force) return;
  const cards = document.getElementById('daily-plan-cards');
  const dateEl = document.getElementById('daily-plan-date');
  if (!cards) return;

  const urgencyColor  = { HOT: '#E53E3E', WARM: '#D97706', NOW: '#7B5CF5' };
  const urgencyBg     = { HOT: 'rgba(229,62,62,.08)', WARM: 'rgba(217,119,6,.08)', NOW: 'rgba(123,92,245,.08)' };
  const urgencyBorder = { HOT: 'rgba(229,62,62,.35)', WARM: 'rgba(217,119,6,.35)', NOW: 'rgba(123,92,245,.35)' };

  cards.innerHTML = '<div style="grid-column:1/-1;background:rgba(21,62,237,0.04);border:1px dashed rgba(21,62,237,0.2);border-radius:12px;padding:32px;text-align:center;color:var(--muted);font-size:13px"><div style="font-size:20px;margin-bottom:8px"><i data-lucide="zap" style="width:20px;height:20px;display:inline-block;vertical-align:middle;"></i></div>Analysing pipeline and signals…</div>';

  try {
    const r = await fetch('/api/gtm/daily-plan');
    const d = await r.json();
    if (d.error || !d.plan || !d.plan.length) {
      cards.innerHTML = '<div style="grid-column:1/-1;padding:24px;text-align:center;color:var(--muted);font-size:13px">No plan available — add accounts to the pipeline.</div>';
      return;
    }

    if (dateEl) dateEl.textContent = d.date || '';

    cards.innerHTML = d.plan.map(function(item) {
      const urg = item.urgency || 'WARM';
      const col  = urgencyColor[urg]  || '#153EED';
      const bg   = urgencyBg[urg]     || 'rgba(21,62,237,.06)';
      const bdr  = urgencyBorder[urg] || 'rgba(21,62,237,.2)';
      const slugGuess = (item.account || '').toLowerCase().replace(/[^a-z0-9]+/g,'-').replace(/-+/g,'-').replace(/^-|-$/g,'');
      const acc = allAccounts.find(function(a){ return a.name.toLowerCase() === (item.account||'').toLowerCase(); }) || allAccounts.find(function(a){ return a.slug === slugGuess; });
      const whyEsc = (item.why_today||'').replace(/'/g,"&#39;");
      const buyerEsc = (item.buyer||'').replace(/'/g,"&#39;");
      const accName = (item.account||'').replace(/'/g,"&#39;");
      const gtmEsc = (item.gtm_strategy||'').replace(/'/g,"&#39;");
      return '<div style="background:'+bg+';border:1px solid '+bdr+';border-radius:12px;padding:18px 20px;display:flex;flex-direction:column;gap:10px;position:relative;overflow:hidden">'
        + '<div style="position:absolute;top:0;right:0;width:80px;height:100%;background:linear-gradient(90deg,transparent,'+bg+');pointer-events:none"></div>'
        + '<div style="display:flex;align-items:flex-start;justify-content:space-between;gap:8px">'
        +   '<div>'
        +     '<div style="font-size:16px;font-weight:800;color:var(--white);letter-spacing:-.02em;margin-bottom:2px">'+(item.account||'')+'</div>'
        +     '<div style="font-size:11px;color:var(--muted)">'+(item.buyer||'—')+'</div>'
        +   '</div>'
        +   '<div style="flex-shrink:0;background:'+col+';color:#fff;border-radius:5px;padding:2px 8px;font-size:9.5px;font-weight:800;letter-spacing:.06em;text-transform:uppercase">'+urg+'</div>'
        + '</div>'
        + '<div style="background:rgba(0,0,0,.25);border-radius:7px;padding:10px 12px">'
        +   '<div style="font-size:9px;font-weight:700;color:'+col+';text-transform:uppercase;letter-spacing:.08em;margin-bottom:4px">Why today</div>'
        +   '<div style="font-size:12px;color:var(--muted);line-height:1.5">'+(item.why_today||'—')+'</div>'
        + '</div>'
        + '<div style="background:rgba(0,0,0,.2);border-radius:7px;padding:10px 12px">'
        +   '<div style="font-size:9px;font-weight:700;color:rgba(255,255,255,.4);text-transform:uppercase;letter-spacing:.08em;margin-bottom:4px">Opening line</div>'
        +   '<div style="font-size:12px;color:var(--text-light,#CBD5E1);line-height:1.6;font-style:italic">&ldquo;'+(item.opening_line||'—')+'&rdquo;</div>'
        + '</div>'
        + '<div style="display:flex;align-items:center;justify-content:space-between">'
        +   '<div style="font-size:11px;font-weight:700;color:'+col+'">'+(item.pipeline_value||'')+' · '+(item.gtm_strategy||'')+'</div>'
        +   '<div style="display:flex;gap:6px">'
        +     (acc ? '<button data-slug="'+acc.slug+'" data-name="'+acc.name+'" onclick="selectAccount(this.dataset.slug,this.dataset.name);showTab(\\'chat\\')" style="background:rgba(255,255,255,.08);border:1px solid rgba(255,255,255,.1);color:rgba(255,255,255,.7);border-radius:6px;padding:5px 10px;font-size:11px;font-weight:600;cursor:pointer;font-family:inherit">Open →</button>' : '')
        +     '<button onclick="insertCustomSkill(\\'Write a personalised LinkedIn outreach message for '+accName+'. Buyer: '+buyerEsc+'. GTM strategy: '+gtmEsc+'. Opening angle: '+whyEsc+'. Language: Norwegian.\\');showTab(\\'chat\\')" style="background:#153EED;border:none;color:#fff;border-radius:6px;padding:5px 10px;font-size:11px;font-weight:700;cursor:pointer;font-family:inherit">Write message →</button>'
        +   '</div>'
        + '</div>'
        + '</div>';
    }).join('');
    _dailyPlanLoaded = true;
  } catch(e) {
    cards.innerHTML = '<div style="grid-column:1/-1;padding:24px;text-align:center;color:var(--muted);font-size:13px">Kunne ikke generere plan: '+e.message+'</div>';
  }
}

function insertCustomSkill(text) {
  const inp = document.getElementById('chat-input');
  if (inp) { inp.value = text; inp.style.height = 'auto'; inp.style.height = inp.scrollHeight + 'px'; }
}

// ── Clock ─────────────────────────────────────────────────────────────────────
function updateClock() {
  const now = new Date();
  const t = document.getElementById('dash-time');
  const d = document.getElementById('dash-date');
  if (t) t.textContent = now.toLocaleTimeString('en-GB', { hour:'2-digit', minute:'2-digit' });
  if (d) d.textContent = now.toLocaleDateString('en-GB', { weekday:'short', day:'numeric', month:'long', year:'numeric' });
}

// ── Tab switching ─────────────────────────────────────────────────────────────
function showTab(name) {
  document.querySelectorAll('.tab-pane').forEach(p => p.classList.remove('active'));
  document.querySelectorAll('.nav-btn').forEach(b => b.classList.remove('active'));
  const tabEl = document.getElementById('tab-' + name);
  if (tabEl) tabEl.classList.add('active');
  const nb = document.getElementById('nav-' + name);
  if (nb) nb.classList.add('active');
  if (name === 'radar')     setTimeout(renderRadar, 50);
  if (name === 'partners')  setTimeout(loadPartnerHistory, 50);
  if (name === 'dashboard') setTimeout(loadDailyPlan, 100);
}

// ── Dashboard ─────────────────────────────────────────────────────────────────
const VALUE_MAP = {
  'hm':           { val:'€900K', strat:'Data Revenue Unlock', icp:9, deal:9 },
  'matas':        { val:'€700K', strat:'AI Readiness',        icp:9, deal:9 },
  'elkjop':       { val:'€700K', strat:'Commerce Opt.',       icp:8, deal:8 },
  'varner-group': { val:'€450K', strat:'Data Revenue Unlock', icp:9, deal:9 },
  'trumf':        { val:'€450K', strat:'Data Revenue Unlock', icp:9, deal:9 },
  'clas-ohlson':  { val:'€350K', strat:'Commerce Opt.',       icp:7, deal:7 },
  'boozt':        { val:'€300K', strat:'AI Readiness',        icp:8, deal:8 },
  'jysk':         { val:'€280K', strat:'Commerce Opt.',       icp:7, deal:7 },
};

function renderDashboard(accounts) {
  // Top 6 by deal score
  const top = accounts
    .filter(a => a.deal !== '\u2014')
    .sort((a,b) => parseInt(b.deal) - parseInt(a.deal))
    .slice(0, 6);

  const circ = 87.96;
  const container = document.getElementById('top-opps');
  if (container) {
    container.innerHTML = top.map((a, i) => {
      const deal = parseInt(a.deal) || 0;
      const icp  = parseInt(a.icp)  || 0;
      const vm   = VALUE_MAP[a.slug] || { val:'\u2014', strat:'GTM' };
      const dash = ((deal / 10) * circ).toFixed(1);
      const col  = deal >= 8 ? '#00D4A0' : deal >= 6 ? '#4B6EF7' : '#F5A623';
      return '<div class="opp" data-slug="' + a.slug + '" data-name="' + a.name + '" onclick="selectAccount(this.dataset.slug,this.dataset.name)">' +
        '<div class="opp-num">#' + (i+1) + '</div>' +
        '<div class="opp-ring">' +
          '<svg width="38" height="38" viewBox="0 0 38 38">' +
            '<circle cx="19" cy="19" r="15" fill="none" stroke="rgba(255,255,255,0.06)" stroke-width="2.5"/>' +
            '<circle cx="19" cy="19" r="15" fill="none" stroke="' + col + '" stroke-width="2.5" ' +
              'stroke-dasharray="' + dash + ' ' + circ + '" stroke-linecap="round" ' +
              'style="transform:rotate(-90deg);transform-origin:50% 50%"/>' +
          '</svg>' +
          '<div class="opp-ring-val">' + deal + '</div>' +
        '</div>' +
        '<div class="opp-body">' +
          '<div class="opp-name">' + a.name + '</div>' +
          '<div class="opp-meta">' + vm.strat + ' &nbsp;·&nbsp; ICP ' + icp + '/10 &nbsp;·&nbsp; ' + (a.country !== '\u2014' ? a.country : '') + '</div>' +
        '</div>' +
        '<div class="opp-right">' +
          '<div class="opp-val">' + vm.val + '</div>' +
          '<div class="opp-val-sub">unweighted</div>' +
        '</div>' +
      '</div>';
    }).join('');
  }

  // Active count
  const active = accounts.filter(a => a.deal !== '\u2014').length;
  const el = document.getElementById('kpi-accounts');
  if (el) el.setAttribute('data-target', active);

  // Strategy split (approximate categorization by known accounts)
  const dru = ['hm','varner-group','trumf','xxl-fraser-group','norgesgruppen','salling-group','coop-norge','naf','trumf','oda','europris','nille','vinmonopolet'];
  const ai  = ['matas','dustin-group','bestseller','komplett','dnb','saxo-bank','lyko','apotea','lindex','imerco','pandora','la-redoute','fnac-darty','gymgrossisten'];
  const co  = ['elkjop','clas-ohlson','jysk','boozt','xxl-sport','bohus','skeidar','halfords','currys','ao-world','ocado','dunelm','webhallen','sport-outlet'];
  const xt  = ['helly-hansen','loccitane','plantasjen','plantagen-sverige','norrona','kapphahl','polarn-o-pyret','gant-norway','follestad'];
  const allSlugs = accounts.map(a => a.slug);
  const druN = allSlugs.filter(s => dru.includes(s)).length;
  const aiN  = allSlugs.filter(s => ai.includes(s)).length;
  const coN  = allSlugs.filter(s => co.includes(s)).length;
  const xtN  = allSlugs.filter(s => xt.includes(s)).length;
  const maxN = Math.max(druN, aiN, coN, xtN, 1);
  ['dru','ai','co','xt'].forEach((k,i) => {
    const n = [druN,aiN,coN,xtN][i];
    const el = document.getElementById('strat-' + k);
    if (el) el.setAttribute('data-target', n);
    const bar = document.getElementById('sb-' + k);
    if (bar) bar.setAttribute('data-w', Math.round((n/maxN)*100) + '%');
  });
}

// ── Live Dashboard ────────────────────────────────────────────────────────────
let _liveRefreshTimer = null;
let _liveCountdown = 300; // 5 min in seconds

async function loadLiveDashboard(manual = false) {
  // Reset countdown
  _liveCountdown = 300;
  clearInterval(_liveRefreshTimer);
  _liveRefreshTimer = setInterval(() => {
    _liveCountdown--;
    const m = Math.floor(_liveCountdown / 60);
    const s = String(_liveCountdown % 60).padStart(2, '0');
    const el = document.getElementById('live-countdown');
    if (el) el.textContent = m + ':' + s;
    if (_liveCountdown <= 0) loadLiveDashboard();
  }, 1000);

  try {
    const res = await fetch('/api/dashboard-live');
    const d = await res.json();

    // --- KPIs ---
    // Pipeline: parse number from e.g. "€10.5M Commerce/Data + DKK 1.2M Hello Growth"
    const pipM = (d.pipeline_total || '').match(/([\\d.]+)M/);
    const pipNum = pipM ? parseFloat(pipM[1]) : 6.8;
    countUp('kpi-pipeline', pipNum, 1.4, 1);
    countUp('kpi-buyers',   d.named_buyers || 18, 1.1, 0);

    // Forecast: parse number from e.g. "€600K"
    const fcM = (d.forecast_base || '').match(/([\\d,]+)K/);
    const fcNum = fcM ? parseInt(fcM[1].replace(',','')) : 600;
    countUp('kpi-forecast', fcNum, 1.3, 0);

    // Status badge on pipeline KPI
    const statusBadge = document.querySelector('.kpi-badge.amber, .kpi-badge.red, .kpi-badge.green');
    if (statusBadge && d.status) {
      statusBadge.className = 'kpi-badge ' + d.status.toLowerCase();
      statusBadge.textContent = '● ' + d.status;
    }

    // Last updated
    const lu = document.getElementById('live-last-updated');
    if (lu) lu.textContent = d.last_updated || 'ukuendt';

    // --- Today's Priority Hero Card ---
    if (d.priority) {
      const p = d.priority;
      const card = document.getElementById('today-priority-card');
      if (card) card.style.display = 'block';
      const nameEl = document.getElementById('tp-name');
      if (nameEl) nameEl.textContent = p.name || '—';
      const metaEl = document.getElementById('tp-meta');
      if (metaEl) metaEl.textContent = (p.company || '') + ' · ' + (p.country || '') + ' · ' + (p.offering || '');
      const reasonEl = document.getElementById('tp-reason');
      if (reasonEl) reasonEl.textContent = p.status || 'Klar til outreach';
      const winEl = document.getElementById('tp-win');
      if (winEl) winEl.textContent = p.win_pct || '—';
      const wEl = document.getElementById('tp-weighted');
      if (wEl) wEl.textContent = p.weighted || '—';
    }

    // --- Deal Velocity Cards ---
    const dvc = document.getElementById('deal-velocity-cards');
    if (dvc && d.top_deals && d.top_deals.length) {
      const cards = d.top_deals.map(function(deal) {
        const stale = deal.days_stale;
        const staleClass = stale === null ? '' : stale > 14 ? 'stale-hot' : stale > 7 ? 'stale-warm' : '';
        const staleBarColor = stale > 14 ? '#F6574A' : stale > 7 ? '#F5A623' : '#00D4A0';
        const staleBarW = stale === null ? 0 : Math.min(100, (stale / 21) * 100);
        const staleTagClass = stale === null ? 'fresh' : stale > 14 ? 'hot' : stale > 7 ? 'warm' : 'fresh';
        const staleLabel = stale === null ? 'Active' : stale === 0 ? 'Today' : stale + ' days ago';
        const buyerShort = (deal.buyer || 'TBD').split('(')[0].trim().slice(0, 30);
        const hasBuyer = buyerShort && buyerShort !== 'TBD';
        const winColor = parseInt(deal.win_pct) >= 60 ? 'var(--green)' : parseInt(deal.win_pct) >= 40 ? 'var(--amber)' : 'var(--red)';
        const div = document.createElement('div');
        div.className = 'dv-card ' + staleClass;
        div.dataset.slug = deal.slug || '';
        div.dataset.name = deal.name || '';
        div.onclick = function() { selectAccount(this.dataset.slug, this.dataset.name); };
        div.innerHTML =
          '<div class="dv-stale-bar" style="width:' + staleBarW + '%;background:' + staleBarColor + '"></div>' +
          '<div class="dv-header">' +
            '<div><div class="dv-name">' + (deal.name || '') + '</div>' +
            '<div class="dv-country">' + (deal.country || '') + ' · ' + (deal.offering || '') + '</div></div>' +
            '<div style="text-align:right"><div class="dv-win" style="color:' + winColor + '">' + (deal.win_pct || '—') + '</div>' +
            '<div class="dv-win-label">win %</div></div>' +
          '</div>' +
          '<div class="dv-metrics">' +
            '<div class="dv-metric"><div class="dv-metric-val">' + (deal.entry_val || '—') + '</div><div class="dv-metric-label">Unweighted</div></div>' +
            '<div class="dv-metric"><div class="dv-metric-val" style="color:#7B5CF5">' + (deal.weighted || '—') + '</div><div class="dv-metric-label">Weighted</div></div>' +
            '<div class="dv-metric"><div class="dv-metric-val">' + (deal.icp || '—') + '</div><div class="dv-metric-label">ICP</div></div>' +
          '</div>' +
          '<div style="display:flex;align-items:center;justify-content:space-between;margin-top:6px">' +
            '<span class="dv-stale-tag ' + staleTagClass + '">● ' + staleLabel + '</span>' +
            (hasBuyer ? '<span class="dv-buyer"><i data-lucide="user" style="width:12px;height:12px;display:inline-block;vertical-align:middle;"></i> ' + buyerShort + '</span>' : '<span class="dv-buyer" style="color:var(--red)"><i data-lucide="alert-triangle" style="width:12px;height:12px;display:inline-block;vertical-align:middle;"></i> Buyer TBD</span>') +
          '</div>' +
          '<div style="margin-top:10px;padding-top:10px;border-top:1px solid rgba(255,255,255,0.06)">' +
            '<button class="dv-outreach-btn" style="width:100%;background:rgba(21,62,237,0.12);border:1px solid rgba(21,62,237,0.3);color:#6B8EF7;border-radius:6px;padding:7px 0;font-size:11.5px;font-weight:700;cursor:pointer;font-family:inherit;transition:background 0.15s"><i data-lucide="mail" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i> Generate outreach</button>' +
          '</div>';
        // Wire up outreach button without bubbling to card click
        const btn = div.querySelector('.dv-outreach-btn');
        const capturedSlug  = deal.slug  || '';
        const capturedName  = deal.name  || '';
        const capturedBuyer = buyerShort || '';
        btn.addEventListener('click', function(e) {
          e.stopPropagation();
          openOutreachModal(capturedSlug, capturedName, capturedBuyer);
        });
        btn.addEventListener('mouseenter', function(){ this.style.background = 'rgba(21,62,237,0.22)'; });
        btn.addEventListener('mouseleave', function(){ this.style.background = 'rgba(21,62,237,0.12)'; });
        return div;
      });
      dvc.innerHTML = '';
      cards.forEach(function(c) { dvc.appendChild(c); });
      if (typeof lucide !== 'undefined') lucide.createIcons();
    }

    // --- Live Signals ---
    const sl = document.getElementById('signals-list');
    if (sl) {
      // Always show hot hardcoded signals first, append live signals below
      const hardcoded = [
        {icon:'<span style="width:8px;height:8px;background:#DC2626;border-radius:50%;display:inline-block;"></span>', co:'Sport Outlet', txt:'CTO + CDO both vacant March 2026. Contact CEO Tor-André Skeie directly.', tag:'URGENT', tagClass:''},
        {icon:'<i data-lucide="zap" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i>', co:'Trumf (NorgesGruppen)', txt:'Rikke Etholm-Idsøe — new Commercial Director. 90-day honeymoon window open.', tag:'90-DAY WINDOW', tagClass:''},
        {icon:'<i data-lucide="plus-circle" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i>', co:'Vinmonopolet', txt:'Espen Terland new CDO (ex-XXL 15 yrs). Agenda not yet set.', tag:'NEW EXEC', tagClass:'amber'},
        {icon:'<i data-lucide="gem" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i>', co:'Siteimprove', txt:'Jen Jones — CMO day 5. Honeymoon window: 3-5x response rate vs 30 days later.', tag:'URGENT', tagClass:''},
      ];
      let html = hardcoded.map(s =>
        '<div class="signal"><div class="signal-icon">' + s.icon + '</div><div style="flex:1"><div class="signal-co">' + s.co + '</div><div class="signal-txt">' + s.txt + '</div></div><div class="signal-tag ' + s.tagClass + '">' + s.tag + '</div></div>'
      ).join('');
      if (d.live_signals && d.live_signals.length) {
        html += '<div style="font-size:9px;font-weight:800;color:var(--muted2);text-transform:uppercase;letter-spacing:2px;margin:12px 0 6px">From daily radar</div>';
        html += d.live_signals.slice(0,3).map(s =>
          '<div class="signal"><div class="signal-icon"><i data-lucide="radio" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i></div><div style="flex:1"><div class="signal-co">' + s.company + '</div><div class="signal-txt">' + (s.text || '').slice(0,100) + '</div></div><div class="signal-tag" style="background:var(--blue-dim2);color:var(--blue-light)">' + (s.date || 'RADAR') + '</div></div>'
        ).join('');
      }
      sl.innerHTML = html;
      if (typeof lucide !== 'undefined') lucide.createIcons();
    }

  } catch(e) {
    console.warn('Live dashboard load failed:', e);
  }
}

function animateDashboard() {
  // KPIs now loaded via loadLiveDashboard — fallback only
  ['kpi-accounts','strat-dru','strat-ai','strat-co','strat-xt'].forEach(id => {
    const el = document.getElementById(id);
    if (el) {
      const t = parseInt(el.getAttribute('data-target') || el.textContent) || 0;
      countUp(id, t, 1.2, 0);
    }
  });
  document.querySelectorAll('[data-w]').forEach(el => {
    setTimeout(() => { el.style.width = el.getAttribute('data-w'); }, 100);
  });
}

function countUp(id, target, duration, decimals) {
  const el = document.getElementById(id);
  if (!el) return;
  const start = performance.now();
  const ms = duration * 1000;
  function frame(now) {
    const p = Math.min((now - start) / ms, 1);
    const e = 1 - Math.pow(1 - p, 3);
    el.textContent = decimals > 0 ? (e * target).toFixed(decimals) : Math.round(e * target);
    if (p < 1) requestAnimationFrame(frame);
  }
  requestAnimationFrame(frame);
}

// ── Account sidebar ───────────────────────────────────────────────────────────
function renderAccountList(accounts) {
  const list = document.getElementById('account-list');
  list.innerHTML = accounts.map(a => {
    const icp  = a.icp  !== '\u2014' ? a.icp  : null;
    const high = icp && parseInt(icp) >= 8;
    const sel  = selectedAccount === a.slug;
    return '<div class="account-item' + (sel ? ' selected' : '') + '" data-slug="' + a.slug + '" data-name="' + a.name + '" onclick="selectAccount(this.dataset.slug,this.dataset.name)">' +
      '<span class="aname">' + a.name + '</span>' +
      '<span class="acc-badges">' +
        (icp ? '<span class="spill' + (high ? ' g' : '') + '">ICP ' + icp + '</span>' : '') +
        (a.country !== '\u2014' ? '<span class="spill c">' + a.country + '</span>' : '') +
      '</span>' +
    '</div>';
  }).join('');
}

function toggleSkillsPanel() {
  const panel = document.getElementById('skills-panel');
  const toggle = document.getElementById('skills-toggle');
  const chevron = document.getElementById('skills-chevron');
  const isOpen = panel.style.maxHeight !== '0px' && panel.style.maxHeight !== '';
  if (isOpen) {
    panel.style.maxHeight = '0';
    toggle.classList.remove('open');
    chevron.style.transform = 'rotate(0deg)';
  } else {
    panel.style.maxHeight = '300px';
    toggle.classList.add('open');
    chevron.style.transform = 'rotate(180deg)';
  }
}

function toggleAccountsPanel() {
  const panel = document.getElementById('accounts-panel');
  const toggle = document.getElementById('accounts-toggle');
  const chevron = document.getElementById('accounts-chevron');
  const isOpen = panel.style.maxHeight !== '0px' && panel.style.maxHeight !== '';
  if (isOpen) {
    panel.style.maxHeight = '0';
    panel.style.overflow = 'hidden';
    toggle.classList.remove('open');
    chevron.style.transform = 'rotate(0deg)';
  } else {
    panel.style.maxHeight = '560px';
    panel.style.overflow = 'hidden';
    toggle.classList.add('open');
    chevron.style.transform = 'rotate(180deg)';
  }
}

function filterAccounts() {
  const q = document.getElementById('account-search').value.toLowerCase();
  const filtered = q ? allAccounts.filter(a => a.name.toLowerCase().includes(q) || a.slug.includes(q)) : allAccounts;
  renderAccountList(filtered);
}

function selectAccount(slug, name) {
  selectedAccount = slug;
  renderAccountList(allAccounts);
  document.getElementById('selected-account-tag').style.display = 'flex';
  document.getElementById('selected-account-name').textContent = name;
  document.getElementById('deck-btn').style.display = 'block';
  showTab('chat');
  document.getElementById('welcome').style.display = 'none';
  addSystemNote('Account loaded: ' + name);
}

function clearAccount() {
  selectedAccount = null;
  renderAccountList(allAccounts);
  document.getElementById('selected-account-tag').style.display = 'none';
  document.getElementById('deck-btn').style.display = 'none';
}

async function generateDeck() {
  if (!selectedAccount) return;
  const btn = document.getElementById('deck-btn');
  btn.classList.add('loading');
  btn.textContent = 'Building deck…';
  try {
    const res = await fetch('/api/generate-deck/' + selectedAccount, { method: 'POST' });
    if (!res.ok) { showToast('Deck generation failed', true); return; }
    const blob = await res.blob();
    const url  = URL.createObjectURL(blob);
    const a    = document.createElement('a');
    a.href = url;
    a.download = 'JAKALA-' + selectedAccount + '-discovery.pptx';
    a.click();
    URL.revokeObjectURL(url);
    showToast('Deck downloaded \u2713');
  } catch(e) {
    showToast('Download failed', true);
  } finally {
    btn.classList.remove('loading');
    btn.textContent = '\u2b07 Generate Deck';
  }
}

// ── Accounts grid ─────────────────────────────────────────────────────────────
function renderAccountsGrid(accounts) {
  const grid = document.getElementById('accounts-grid');
  const list = gridFilter === 'all'    ? accounts :
               gridFilter === 'scored' ? accounts.filter(a => a.deal !== '\u2014') :
               accounts.filter(a => a.country === gridFilter);
  grid.innerHTML = list.map(a => {
    const icp  = a.icp  !== '\u2014' ? parseInt(a.icp)  : 0;
    const deal = a.deal !== '\u2014' ? parseInt(a.deal) : 0;
    const high = icp >= 8 || deal >= 8;
    return '<div class="account-card" data-slug="' + a.slug + '" data-name="' + a.name + '" onclick="selectAccount(this.dataset.slug,this.dataset.name)">' +
      '<div class="card-top">' +
        '<div class="card-name">' + a.name + '</div>' +
        '<div class="card-badges">' +
          (a.country !== '\u2014' ? '<span class="badge c">' + a.country + '</span>' : '') +
          (icp  ? '<span class="badge' + (high ? ' g' : '') + '">ICP ' + icp + '</span>' : '') +
          (deal ? '<span class="badge' + (high ? ' g' : '') + '">Deal ' + deal + '</span>' : '') +
        '</div>' +
      '</div>' +
      (icp || deal ? '<div class="score-bars">' +
        (icp  ? '<div class="sb-item"><div class="sb-label">ICP</div><div class="sb-track"><div class="sb-fill icp" style="width:' + (icp*10) + '%"></div></div></div>' : '') +
        (deal ? '<div class="sb-item"><div class="sb-label">Deal</div><div class="sb-track"><div class="sb-fill deal" style="width:' + (deal*10) + '%"></div></div></div>' : '') +
      '</div>' : '') +
    '</div>';
  }).join('');
}

function filterGrid(btn, filter) {
  document.querySelectorAll('.filter-btn').forEach(b => b.classList.remove('active'));
  btn.classList.add('active');
  gridFilter = filter;
  renderAccountsGrid(allAccounts);
}

function populateNotesSelect(accounts) {
  const sel = document.getElementById('notes-account');
  accounts.forEach(a => {
    const opt = document.createElement('option');
    opt.value = a.slug;
    opt.textContent = a.name;
    sel.appendChild(opt);
  });
}

// ── Chat ──────────────────────────────────────────────────────────────────────
function addSystemNote(text) {
  const msgs = document.getElementById('messages');
  const el = document.createElement('div');
  el.style.cssText = 'text-align:center;font-size:10.5px;color:var(--muted);padding:2px 0;';
  el.textContent = '\u00b7 ' + text + ' \u00b7';
  msgs.appendChild(el);
  msgs.scrollTop = msgs.scrollHeight;
}

function appendMessage(role, content) {
  const msgs = document.getElementById('messages');
  const el = document.createElement('div');
  el.className = 'msg ' + role;
  el.innerHTML = '<div class="msg-role">' + (role === 'user' ? 'You' : 'GTM OS') + '</div>' +
                 '<div class="msg-bubble">' + renderMarkdown(content) + '</div>';
  msgs.appendChild(el);
  msgs.scrollTop = msgs.scrollHeight;
  return el.querySelector('.msg-bubble');
}

function renderMarkdown(text) {
  return text
    .replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;')
    .replace(/```[\s\S]*?```/g, m => '<pre><code>' + m.slice(3,-3).replace(/^[a-z]+\\n/,'') + '</code></pre>')
    .replace(/`([^`]+)`/g,'<code>$1</code>')
    .replace(/^### (.+)$/gm,'<h3>$1</h3>')
    .replace(/^## (.+)$/gm,'<h2>$1</h2>')
    .replace(/^# (.+)$/gm,'<h1>$1</h1>')
    .replace(/\*\*(.+?)\*\*/g,'<strong>$1</strong>')
    .replace(/\*(.+?)\*/g,'<em>$1</em>')
    .replace(/^> (.+)$/gm,'<blockquote>$1</blockquote>')
    .replace(/^---$/gm,'<hr>')
    .replace(/^\| (.+) \|$/gm, m => {
      const cells = m.slice(2,-2).split(' | ');
      return '<tr>' + cells.map(c => '<td>' + c.replace(/\*\*(.+?)\*\*/g,'<strong>$1</strong>') + '</td>').join('') + '</tr>';
    })
    .replace(/(<tr>.*<\/tr>\\n?)+/gs, m => '<table>' + m + '</table>')
    .replace(/<table>(<tr><td>[-:| ]+<\/td><\/tr>)<\/table>/g,'')
    .replace(/^- (.+)$/gm,'<li>$1</li>')
    .replace(/(<li>.*<\/li>\\n?)+/gs, m => '<ul>' + m + '</ul>')
    .replace(/^\\d+\. (.+)$/gm,'<li>$1</li>')
    .replace(/\\n/g,'<br>');
}

async function sendMessage() {
  const input = document.getElementById('chat-input');
  const text = input.value.trim();
  if (!text) return;

  document.getElementById('welcome').style.display = 'none';
  input.value = '';
  input.style.height = 'auto';
  document.getElementById('send-btn').disabled = true;

  messages.push({ role: 'user', content: text });
  appendMessage('user', text);

  document.getElementById('typing-row').style.display = 'block';
  document.getElementById('typing-indicator').classList.add('visible');

  const bubble = appendMessage('assistant', '');
  let full = '';

  try {
    const res = await fetch('/api/chat', {
      method: 'POST',
      headers: { 'Content-Type': 'application/json' },
      body: JSON.stringify({ messages, account: selectedAccount })
    });

    document.getElementById('typing-row').style.display = 'none';
    document.getElementById('typing-indicator').classList.remove('visible');

    const reader = res.body.getReader();
    const decoder = new TextDecoder();

    while (true) {
      const { done, value } = await reader.read();
      if (done) break;
      const chunk = decoder.decode(value);
      for (const line of chunk.split('\\n')) {
        if (line.startsWith('data: ') && line !== 'data: [DONE]') {
          try {
            const { text: t } = JSON.parse(line.slice(6));
            full += t;
            bubble.innerHTML = renderMarkdown(full);
            document.getElementById('messages').scrollTop = 99999;
          } catch(e) {}
        }
      }
    }
    messages.push({ role: 'assistant', content: full });
  } catch(err) {
    document.getElementById('typing-row').style.display = 'none';
    bubble.innerHTML = '<em style="color:var(--red)">Connection error — check server and API key.</em>';
  }

  document.getElementById('send-btn').disabled = false;
  input.focus();
}

function handleKey(e) {
  if (e.key === 'Enter' && !e.shiftKey) { e.preventDefault(); sendMessage(); }
}

function autoResize(el) {
  el.style.height = 'auto';
  el.style.height = Math.min(el.scrollHeight, 140) + 'px';
}

// ── Quick skills ──────────────────────────────────────────────────────────────
function insertSkill(key) {
  showTab('chat');
  const acc = selectedAccount ? selectedAccount.replace(/-/g,' ') : null;
  const prompts = {
    contact:  'Run /contact-today \u2014 who should I contact today and with exactly what message?',
    blueprint: acc ? 'Run /deal-blueprint for ' + acc + ' \u2014 full commercial execution plan: narrative, offer architecture, buying journey, risk register, and 7-day sprint.' : 'Run /deal-blueprint \u2014 which account? Select one from the sidebar first, or tell me the account name.',
    prospect: 'Run /prospect-hunt \u2014 tell me: which industry, which country (DK/NO/SE or Nordic), and any keyword or signal to focus on?',
    morning:  'Give me the morning CCO briefing \u2014 top signals, pipeline health, and my 3 priorities for today.',
    warroom:  'Run the commercial war room \u2014 full situation assessment. Nordic pipeline only (DK/NO/SE).',
    forecast: 'Run the Q2 2026 commercial forecast \u2014 probability-weighted, all active Nordic accounts.',
    outreach: acc ? 'Write a LinkedIn outreach message for ' + acc + '. Language: English.' : 'Write a LinkedIn outreach message. Select an account from the sidebar first, or tell me the company name.',
    pitch:    acc ? 'Run the pitch partner brief for ' + acc + '. Who is the buyer and what is the best service match?' : 'Run the pitch partner brief. Select an account from the sidebar first.',
    brief:    acc ? 'Give me the pre-meeting brief for ' + acc + '. Meeting type: discovery.' : 'Give me the pre-meeting brief. Select an account from the sidebar first.',
    revenue:  acc ? 'Run the revenue simulation for ' + acc + '. Show all three levers and three scenarios.' : 'Run the revenue simulation. Select an account from the sidebar first.',
    signal:   'I have a new market signal to analyse. Here it is:',
    strategic: acc ? 'Run /strategic for ' + acc + ' \u2014 apply the DreamTRUE framework: identify the tension, define the ICP (commercial + emotional), sharpen our uniqueness, and find the emotional trigger that drives action.' : 'Run /strategic \u2014 apply the DreamTRUE framework. What brand, product or GTM challenge should we analyse? (or select an account from the sidebar first)',
    websiteopt: 'Run /website-optimizer \u2014 I need a website audit. Here is what I have: [paste URL, PageSpeed data, or HTML snippet]',
  };
  const input = document.getElementById('chat-input');
  input.value = prompts[key] || '';
  autoResize(input);
  input.focus();
  document.getElementById('welcome').style.display = 'none';
}

// ── Meeting notes ─────────────────────────────────────────────────────────────
async function processNotes() {
  const account = document.getElementById('notes-account').value;
  const notes   = document.getElementById('notes-text').value.trim();
  if (!account) { showToast('Select an account first', true); return; }
  if (!notes)   { showToast('Paste meeting notes first', true); return; }

  const btn = document.getElementById('process-btn');
  btn.disabled = true; btn.textContent = 'Processing\u2026';

  try {
    const res  = await fetch('/api/process-notes', {
      method: 'POST',
      headers: { 'Content-Type': 'application/json' },
      body: JSON.stringify({ account, notes })
    });
    const data = await res.json();
    if (data.error) { showToast('Error: ' + data.error, true); return; }

    processingResult = { account, ...data };
    document.getElementById('result-insight').textContent  = data.key_insight    || '\u2014';
    document.getElementById('result-summary').textContent  = data.summary        || '\u2014';
    document.getElementById('result-meeting').textContent  = data.meeting_entry  || '\u2014';
    document.getElementById('result-actions').textContent  = data.next_actions_updated || '\u2014';
    document.getElementById('notes-result').style.display = 'block';
    document.getElementById('notes-result').scrollIntoView({ behavior: 'smooth' });
  } catch(err) {
    showToast('Request failed \u2014 check connection', true);
  } finally {
    btn.disabled = false; btn.textContent = 'Process Notes';
  }
}

async function saveNotes() {
  if (!processingResult) return;
  const btn = document.getElementById('save-btn');
  btn.disabled = true; btn.textContent = 'Saving\u2026';

  try {
    const res  = await fetch('/api/save-notes', {
      method: 'POST',
      headers: { 'Content-Type': 'application/json' },
      body: JSON.stringify({
        account: processingResult.account,
        meeting_entry: processingResult.meeting_entry,
        next_actions_updated: processingResult.next_actions_updated
      })
    });
    const data = await res.json();
    if (data.ok) {
      showToast('Saved to account files \u2713');
      document.getElementById('notes-result').style.display = 'none';
      document.getElementById('notes-text').value = '';
    }
  } catch { showToast('Save failed', true); }
  finally { btn.disabled = false; btn.textContent = 'Save to Account Files'; }
}

function cancelNotes() {
  document.getElementById('notes-result').style.display = 'none';
  processingResult = null;
}

// ── Toast ──────────────────────────────────────────────────────────────────────
function showToast(msg, error = false) {
  const t = document.getElementById('toast');
  t.textContent = msg;
  t.className = 'show' + (error ? ' error' : '');
  setTimeout(() => t.className = '', 3000);
}

// ══════════════════════════════════════════════════════════════════════════════
// PARTNERSHIP VALIDATOR
// ══════════════════════════════════════════════════════════════════════════════
async function validatePartner() {
  const name    = document.getElementById('partner-name').value.trim();
  const market  = document.getElementById('partner-market').value;
  const context = document.getElementById('partner-context').value.trim();
  if (!name) { showToast('Angiv partnerens navn', true); return; }

  const btn = document.getElementById('validate-btn');
  btn.disabled = true; btn.textContent = 'Analysing\u2026';
  document.getElementById('partner-result').style.display = 'none';

  try {
    const res  = await fetch('/api/validate-partner', {
      method: 'POST', headers: { 'Content-Type': 'application/json' },
      body: JSON.stringify({ partner: name, market, context })
    });
    const data = await res.json();
    if (data.error) { showToast('Fejl: ' + data.error, true); return; }
    renderPartnerResult(data);
    loadPartnerHistory();
  } catch(e) {
    showToast('Request fejlede', true);
  } finally {
    btn.disabled = false; btn.textContent = 'Validate \u2197';
  }
}

function renderPartnerResult(d) {
  const verdict = (d.verdict || '').toLowerCase().replace(/\s+/g, '-').replace('not-recommended','not-rec');
  const icons   = {
    'strong-fit':   '<i data-lucide="check-circle" style="width:22px;height:22px;display:inline-block;vertical-align:middle;color:#00A67E;"></i>',
    'potential-fit':'<i data-lucide="lightbulb" style="width:22px;height:22px;display:inline-block;vertical-align:middle;color:#D97706;"></i>',
    'weak-fit':     '<i data-lucide="alert-triangle" style="width:22px;height:22px;display:inline-block;vertical-align:middle;color:#D97706;"></i>',
    'not-rec':      '<i data-lucide="x-circle" style="width:22px;height:22px;display:inline-block;vertical-align:middle;color:#DC2626;"></i>'
  };
  const colors  = { 'strong-fit': 'strong', 'potential-fit': 'potential', 'weak-fit': 'weak', 'not-rec': 'not-rec' };
  const cls     = colors[verdict] || 'potential';

  // Banner
  const banner = document.getElementById('verdict-banner');
  banner.className = 'verdict-banner ' + cls;
  document.getElementById('verdict-icon').innerHTML = icons[verdict] || icons['potential-fit'];
  if (typeof lucide !== 'undefined') lucide.createIcons();
  document.getElementById('verdict-label').textContent  = d.verdict || '\u2014';
  document.getElementById('verdict-label').className    = 'verdict-label ' + cls;
  document.getElementById('verdict-reason').textContent = d.verdict_reason || '\u2014';
  document.getElementById('verdict-score').textContent  = d.overall_score || '\u2014';

  // Dimension grid
  const dimColors = { 10:'#00D4A0', 9:'#00D4A0', 8:'#4B6EF7', 7:'#4B6EF7', 6:'#F5A623', 5:'#F5A623', 4:'#F6574A', 3:'#F6574A', 2:'#F6574A', 1:'#F6574A' };
  const dimGrid = document.getElementById('dim-grid');
  dimGrid.innerHTML = (d.dimensions || []).map(dim => {
    const pct = Math.round((dim.score / dim.max) * 100);
    const col = dimColors[dim.score] || '#8080B0';
    return '<div class="dim-card">' +
      '<div class="dim-name">' + dim.name + '</div>' +
      '<div class="dim-score-row">' +
        '<div class="dim-score-num">' + dim.score + '</div>' +
        '<div class="dim-score-max">/' + dim.max + '</div>' +
        '<div class="dim-bar-track"><div class="dim-bar-fill" style="width:0;background:' + col + '" data-w="' + pct + '%"></div></div>' +
      '</div>' +
      '<div class="dim-finding">' + (dim.finding || '') + '</div>' +
    '</div>';
  }).join('');

  // Animate bars
  setTimeout(() => {
    dimGrid.querySelectorAll('.dim-bar-fill').forEach(el => {
      el.style.width = el.getAttribute('data-w');
    });
  }, 80);

  // Meta fields
  document.getElementById('pm-gtm').innerHTML      = (d.gtm_match    || []).map(t => '<span class="pmeta-tag">' + t + '</span>').join('');
  document.getElementById('pm-buyers').innerHTML   = (d.buyer_overlap || []).map(t => '<span class="pmeta-tag">' + t + '</span>').join('');
  document.getElementById('pm-offer').textContent  = d.joint_offer   || '\u2014';
  document.getElementById('pm-accounts').innerHTML = (d.target_accounts || []).map(t => '<span class="pmeta-tag">' + t + '</span>').join('');
  document.getElementById('pm-step').textContent   = d.first_step    || '\u2014';
  document.getElementById('pm-risks').innerHTML    = (d.risks || []).map(r => '\u2022 ' + r).join('<br>');

  document.getElementById('partner-result').style.display = 'block';
  document.getElementById('partner-result').scrollIntoView({ behavior: 'smooth', block: 'start' });
}

async function loadPartnerHistory() {
  const container = document.getElementById('partner-history');
  if (!container) return;
  try {
    const res  = await fetch('/api/partner-history');
    const data = await res.json();
    if (!data.length) {
      container.innerHTML = '<div style="color:var(--muted);font-size:12px">No previous validations yet.</div>';
      return;
    }
    const vmap = { 'STRONG FIT':'strong', 'POTENTIAL FIT':'potential', 'WEAK FIT':'weak', 'NOT RECOMMENDED':'not-rec' };
    container.innerHTML = data.map(p => {
      const cls = vmap[p.verdict] || 'potential';
      return '<div class="ph-row">' +
        '<div class="ph-name">' + p.name + '</div>' +
        '<div class="ph-market">' + p.market + '</div>' +
        '<div class="ph-verdict ' + cls + '">' + p.verdict + '</div>' +
        '<div class="ph-score">' + p.score + '/10</div>' +
      '</div>';
    }).join('');
  } catch(e) {}
}

// ══════════════════════════════════════════════════════════════════════════════
// BOARD REPORT
// ══════════════════════════════════════════════════════════════════════════════
async function generateBoardReport() {
  const btn = document.getElementById('board-report-btn');
  btn.disabled = true; btn.textContent = 'Building report\u2026';
  try {
    const res = await fetch('/api/board-report', { method: 'POST' });
    if (!res.ok) { showToast('Board report failed', true); return; }
    const blob = await res.blob();
    const url = URL.createObjectURL(blob);
    const a = document.createElement('a');
    a.href = url;
    a.download = 'JAKALA-Board-Report-' + new Date().toISOString().slice(0,10) + '.pptx';
    a.click();
    URL.revokeObjectURL(url);
    showToast('Board Report downloaded \u2713');
  } catch(e) {
    showToast('Download failed', true);
  } finally {
    btn.disabled = false; btn.innerHTML = '<i data-lucide="bar-chart-2" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i> Board Report'; lucide.createIcons();
  }
}

// ══════════════════════════════════════════════════════════════════════════════
// PIPELINE RADAR
// ══════════════════════════════════════════════════════════════════════════════
const STRAT_MAP = {
  'hm':'dru','varner-group':'dru','trumf':'dru','norgesgruppen':'dru','oda':'dru',
  'europris':'dru','nille':'dru','vinmonopolet':'dru','coop-norge':'dru','naf':'dru',
  'matas':'ai','dustin-group':'ai','bestseller':'ai','dnb':'ai','saxo-bank':'ai',
  'lyko':'ai','apotea':'ai','lindex':'ai','komplett':'ai','imerco':'ai','pandora':'ai',
  'elkjop':'co','clas-ohlson':'co','jysk':'co','boozt':'co','xxl-sport':'co',
  'bohus':'co','skeidar':'co','halfords':'co','sport-outlet':'co','jernia':'co',
  'webhallen':'co','xxl-fraser-group':'co','salling-group':'co',
  'helly-hansen':'xt','loccitane':'xt','plantasjen':'xt','norrona':'xt',
  'kapphahl':'xt','gant-norway':'xt','follestad':'xt','polarn-o-pyret':'xt',
};
const RADAR_VALS = {
  'hm':'€900K','matas':'€700K','elkjop':'€700K','varner-group':'€450K',
  'trumf':'€450K','clas-ohlson':'€350K','boozt':'€300K','jysk':'€280K',
  'helly-hansen':'€250K','skeidar':'€220K','dnb':'€200K','komplett':'€200K',
  'vinmonopolet':'€180K','norgesgruppen':'€400K',
};

function slugToQuadrant(slug) {
  if (STRAT_MAP[slug]) return STRAT_MAP[slug];
  const dru = ['trumf','norges','oda','europris','nille','vino','coop','naf','xxl-fraser'];
  const ai  = ['matas','dustin','best','dnb','saxo','lyko','apotea','lindex','komplett'];
  const co  = ['elkjop','clas','jysk','boozt','sport','bohus','skeidar','halfords','jernia'];
  const xt  = ['helly','locc','plantasjen','norrona','kapphahl','gant','follestad'];
  if (dru.some(k => slug.includes(k))) return 'dru';
  if (ai.some(k  => slug.includes(k))) return 'ai';
  if (co.some(k  => slug.includes(k))) return 'co';
  if (xt.some(k  => slug.includes(k))) return 'xt';
  let h = 0; for (const c of slug) h = ((h * 31) + c.charCodeAt(0)) >>> 0;
  return ['dru','ai','co','xt'][h % 4];
}

function radarPos(slug, deal) {
  const baseAngles = { dru:45, ai:315, co:135, xt:225 };
  let hash = 0; for (const c of slug) hash = ((hash * 31) + c.charCodeAt(0)) >>> 0;
  const q = slugToQuadrant(slug);
  const spread = ((hash % 50000) / 50000 - 0.5) * 46;
  const deg = (baseAngles[q] || 45) + spread;
  const rad = (deg - 90) * Math.PI / 180;
  const ds = parseInt(deal) || 0;
  const r = ds >= 9 ? 32 + (hash % 40) :
            ds >= 7 ? 82 + (hash % 58) :
            ds >= 5 ? 148 + (hash % 56) :
                      198 + (hash % 36);
  return { x: 300 + r * Math.cos(rad), y: 300 + r * Math.sin(rad) };
}

function renderRadar() {
  const blipsG = document.getElementById('radar-blips');
  const hotList = document.getElementById('radar-hot-list');
  if (!blipsG || !allAccounts.length) return;

  const scored = allAccounts
    .filter(a => a.deal !== '\u2014' && parseInt(a.deal) >= 5)
    .sort((a,b) => parseInt(b.deal) - parseInt(a.deal));

  if (hotList) {
    hotList.innerHTML = scored.slice(0, 9).map(a => {
      const ds = parseInt(a.deal);
      const col = ds >= 8 ? '#00D4A0' : ds >= 6 ? '#4B6EF7' : '#F5A623';
      return '<div class="rhi" data-slug="' + a.slug + '" data-name="' + a.name + '" onclick="selectAccount(this.dataset.slug,this.dataset.name)">' +
        '<div class="rhi-dot" style="background:' + col + ';box-shadow:0 0 6px ' + col + '"></div>' +
        '<div class="rhi-name">' + a.name + '</div>' +
        '<div class="rhi-score">' + ds + '/10</div>' +
      '</div>';
    }).join('');
  }

  let html = '';
  for (const a of scored) {
    const ds = parseInt(a.deal) || 0;
    const pos = radarPos(a.slug, ds);
    const col = ds >= 8 ? '#00D4A0' : ds >= 6 ? '#4B6EF7' : ds >= 5 ? '#F5A623' : '#F6574A';
    const r   = Math.max(5, Math.min(13, 4 + ds * 1.0));
    const val = RADAR_VALS[a.slug] || '';
    html += '<g class="r-blip" style="cursor:pointer" data-slug="' + a.slug + '" data-name="' + a.name + '" data-ds="' + ds + '" data-val="' + val + '"' +
      ' onclick="selectAccount(this.dataset.slug,this.dataset.name)">' +
      '<circle cx="' + pos.x.toFixed(1) + '" cy="' + pos.y.toFixed(1) + '" r="' + (r+4) + '" fill="' + col + '" opacity="0.1"/>' +
      '<circle cx="' + pos.x.toFixed(1) + '" cy="' + pos.y.toFixed(1) + '" r="' + r + '" fill="' + col + '" opacity="0.88">' +
        '<animate attributeName="r" values="' + r + ';' + (r+1.5) + ';' + r + '" dur="2.2s" repeatCount="indefinite"/>' +
      '</circle>' +
    '</g>';
  }
  blipsG.innerHTML = html;

  blipsG.querySelectorAll('.r-blip').forEach(g => {
    g.addEventListener('mouseenter', function() {
      const tt = document.getElementById('radar-tt');
      const ttBg = document.getElementById('tt-bg');
      const ttName = document.getElementById('tt-name');
      const ttDeal = document.getElementById('tt-deal');
      const ttVal  = document.getElementById('tt-val');
      if (!tt) return;
      const cx = parseFloat(g.querySelector('circle').getAttribute('cx'));
      const cy = parseFloat(g.querySelector('circle').getAttribute('cy'));
      ttName.textContent = this.dataset.name;
      ttDeal.textContent = 'Deal ' + this.dataset.ds + '/10';
      ttVal.textContent  = this.dataset.val || '';
      const w = Math.max(ttName.textContent.length, 12) * 7.5 + 20;
      const tx = cx + 14, ty = cy - 50;
      ttBg.setAttribute('x', tx - 6); ttBg.setAttribute('y', ty - 14);
      ttBg.setAttribute('width', w); ttBg.setAttribute('height', 56);
      ttName.setAttribute('x', tx); ttName.setAttribute('y', ty);
      ttDeal.setAttribute('x', tx); ttDeal.setAttribute('y', ty + 17);
      ttVal.setAttribute('x',  tx); ttVal.setAttribute('y',  ty + 33);
      tt.setAttribute('display', 'block');
    });
    g.addEventListener('mouseleave', () => {
      const tt = document.getElementById('radar-tt');
      if (tt) tt.setAttribute('display', 'none');
    });
  });
}

// ══════════════════════════════════════════════════════════════════════════════
// PITCH SIMULATOR
// ══════════════════════════════════════════════════════════════════════════════
let simMessages = [];
let simAccount  = null;

function populateSimSelect() {
  const sel = document.getElementById('sim-account');
  if (!sel || sel.options.length > 1) return;
  allAccounts.forEach(a => {
    const opt = document.createElement('option');
    opt.value = a.slug; opt.textContent = a.name;
    sel.appendChild(opt);
  });
}

async function startPitchSession() {
  const slug = document.getElementById('sim-account').value;
  if (!slug) { showToast('Select an account first', true); return; }

  simAccount  = slug;
  simMessages = [];

  const name = allAccounts.find(a => a.slug === slug)?.name || slug.replace(/-/g,' ').replace(/\\b\\w/g, c => c.toUpperCase());
  const msgs = document.getElementById('sim-messages');
  msgs.innerHTML = '';
  addSimNote('Session started \u2014 ' + name + ' \u00b7 Start with your opening pitch');
  document.getElementById('sim-persona-name').textContent = name + ' Buyer';
  document.getElementById('sim-persona-sub').textContent  = 'Claude is playing the decision maker';
  document.getElementById('sim-persona-bar').classList.add('active');
  document.getElementById('sim-input-area').style.display = 'block';
  document.getElementById('sim-input').focus();

  // Trigger first buyer message
  simMessages.push({ role: 'user', content: 'Hello, I appreciate you taking the time to meet.' });
  const bubble = appendSimMsg('opponent', '');
  let full = '';
  try {
    const res = await fetch('/api/pitch', {
      method: 'POST', headers: { 'Content-Type': 'application/json' },
      body: JSON.stringify({ messages: simMessages, account: simAccount })
    });
    const reader = res.body.getReader(); const decoder = new TextDecoder();
    while (true) {
      const { done, value } = await reader.read(); if (done) break;
      for (const line of decoder.decode(value).split('\\n')) {
        if (line.startsWith('data: ') && line !== 'data: [DONE]') {
          try { const { text: t } = JSON.parse(line.slice(6)); full += t; bubble.innerHTML = renderMarkdown(full); msgs.scrollTop = 99999; } catch(e) {}
        }
      }
    }
    simMessages.push({ role: 'assistant', content: full });
    // Remove the fake opener from history, keep only the buyer reply
    simMessages = [{ role: 'assistant', content: full }];
  } catch(e) { bubble.innerHTML = '<em style="color:var(--red)">Connection error.</em>'; }
}

function addSimNote(text) {
  const msgs = document.getElementById('sim-messages');
  const el = document.createElement('div');
  el.className = 'sim-note';
  el.textContent = '\u00b7 ' + text + ' \u00b7';
  msgs.appendChild(el); msgs.scrollTop = msgs.scrollHeight;
}

function appendSimMsg(role, content) {
  const msgs = document.getElementById('sim-messages');
  const el = document.createElement('div');
  el.className = 'msg ' + (role === 'opponent' ? 'assistant sim-opp' : 'user');
  const label = role === 'opponent'
    ? (simAccount || '').replace(/-/g,' ').replace(/\\b\\w/g, c => c.toUpperCase()) + ' Buyer'
    : 'You';
  const roleStyle = role === 'opponent' ? ' style="color:var(--red)"' : '';
  el.innerHTML = '<div class="msg-role"' + roleStyle + '>' + label + '</div>' +
                 '<div class="msg-bubble">' + renderMarkdown(content) + '</div>';
  msgs.appendChild(el); msgs.scrollTop = msgs.scrollHeight;
  return el.querySelector('.msg-bubble');
}

async function sendSimMessage() {
  const input = document.getElementById('sim-input');
  const text = input.value.trim();
  if (!text || !simAccount) return;

  input.value = ''; input.style.height = 'auto';
  document.getElementById('sim-send-btn').disabled = true;

  simMessages.push({ role: 'user', content: text });
  appendSimMsg('user', text);
  const bubble = appendSimMsg('opponent', '');
  let full = '';

  try {
    const res = await fetch('/api/pitch', {
      method: 'POST', headers: { 'Content-Type': 'application/json' },
      body: JSON.stringify({ messages: simMessages, account: simAccount })
    });
    const reader = res.body.getReader(); const decoder = new TextDecoder();
    while (true) {
      const { done, value } = await reader.read(); if (done) break;
      const msgs = document.getElementById('sim-messages');
      for (const line of decoder.decode(value).split('\\n')) {
        if (line.startsWith('data: ') && line !== 'data: [DONE]') {
          try { const { text: t } = JSON.parse(line.slice(6)); full += t; bubble.innerHTML = renderMarkdown(full); msgs.scrollTop = 99999; } catch(e) {}
        }
      }
    }
    simMessages.push({ role: 'assistant', content: full });
  } catch(e) { bubble.innerHTML = '<em style="color:var(--red)">Error.</em>'; }
  document.getElementById('sim-send-btn').disabled = false;
  document.getElementById('sim-input').focus();
}

async function scorePitch() {
  if (!simMessages.length || !simAccount) { showToast('Start a pitch session first', true); return; }
  addSimNote('Scoring your pitch\u2026');
  const scorePrompt = [
    ...simMessages,
    { role: 'user', content: 'STOP the roleplay. You are now a senior B2B sales coach. Score this pitch conversation 1\u201310 across: Opening hook, Value proposition clarity, Objection handling, Buyer fit, Call to action. Give an Overall score. Then 1 key strength and 1 specific improvement to make it 20% more effective.' }
  ];
  const bubble = appendSimMsg('opponent', '');
  let full = '';
  try {
    const res = await fetch('/api/pitch', {
      method: 'POST', headers: { 'Content-Type': 'application/json' },
      body: JSON.stringify({ messages: scorePrompt, account: simAccount, scoring: true })
    });
    const reader = res.body.getReader(); const decoder = new TextDecoder();
    while (true) {
      const { done, value } = await reader.read(); if (done) break;
      const msgs = document.getElementById('sim-messages');
      for (const line of decoder.decode(value).split('\\n')) {
        if (line.startsWith('data: ') && line !== 'data: [DONE]') {
          try { const { text: t } = JSON.parse(line.slice(6)); full += t; bubble.innerHTML = renderMarkdown(full); msgs.scrollTop = 99999; } catch(e) {}
        }
      }
    }
  } catch(e) { bubble.innerHTML = '<em style="color:var(--red)">Scoring failed.</em>'; }
}

function handleSimKey(e) {
  if (e.key === 'Enter' && !e.shiftKey) { e.preventDefault(); sendSimMessage(); }
}

// ══════════════════════════════════════════════════════════════════════════════
// SIGNAL FEED
// ══════════════════════════════════════════════════════════════════════════════
async function renderSignalFeed() {
  const content = document.getElementById('signals-content');
  if (!content) return;
  content.innerHTML = '<div style="color:var(--muted);font-size:12px;padding:20px 0">Loading signals\u2026</div>';
  try {
    const res = await fetch('/api/signals');
    const data = await res.json();
    const sigs = data.signals || [];
    content.innerHTML = sigs.map(s => {
      const cardCls = s.tagColor === 'red' ? 'urg' : s.tagColor === 'amber' ? 'amb' : '';
      const icoCls  = s.tagColor === 'red' ? 'red' : s.tagColor === 'amber' ? 'amber' : 'blue';
      const hasAcc  = s.slug && allAccounts.some(a => a.slug === s.slug);
      return '<div class="signal-card ' + cardCls + '">' +
        '<div class="sc-ico ' + icoCls + '">' + (s.icon || '\u26a1') + '</div>' +
        '<div class="sc-body">' +
          '<div class="sc-co">' + s.company + '</div>' +
          '<div class="sc-txt">' + s.text + '</div>' +
        '</div>' +
        '<div class="sc-right">' +
          '<span class="sc-tag ' + s.tagColor + '">' + s.tag + '</span>' +
          (hasAcc ? '<button class="sc-act" data-slug="' + s.slug + '" data-name="' + s.company + '" onclick="signalOutreach(this.dataset.slug,this.dataset.name)">Write Outreach</button>' : '') +
        '</div>' +
      '</div>';
    }).join('');
    if (typeof lucide !== 'undefined') lucide.createIcons();
  } catch(e) {
    content.innerHTML = '<div style="color:var(--red);font-size:12px">Failed to load signals</div>';
  }
}

function signalOutreach(slug, name) {
  selectAccount(slug, name);
  document.getElementById('chat-input').value = 'Write a LinkedIn outreach message for ' + name + ' based on the active signal. Language: English. Match tone to the urgency.';
  autoResize(document.getElementById('chat-input'));
  showToast('Account loaded \u2014 outreach ready \u2191');
}

// ══════════════════════════════════════════════════════════════════
// ⌘K COMMAND PALETTE
// ══════════════════════════════════════════════════════════════════

const STATIC_COMMANDS = [
  { type:'skill', icon:'<i data-lucide="zap" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',           label:'Who to contact today?',      sub:'Who To Contact Today',              action: function(){ closePalette(); insertSkill('contact'); } },
  { type:'skill', icon:'<i data-lucide="building-2" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',     label:'Deal Blueprint',             sub:'Complete commercial execution plan', action: function(){ closePalette(); insertSkill('blueprint'); } },
  { type:'skill', icon:'<i data-lucide="search" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',         label:'Prospect Hunt',              sub:'Find new leads',                    action: function(){ closePalette(); insertSkill('prospect'); } },
  { type:'skill', icon:'<i data-lucide="sword" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',          label:'Commercial War Room',        sub:'Full situation assessment',          action: function(){ closePalette(); insertSkill('warroom'); } },
  { type:'skill', icon:'<i data-lucide="sun" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',            label:'Morning Briefing',           sub:'CCO daily briefing',                action: function(){ closePalette(); insertSkill('morning'); } },
  { type:'skill', icon:'<i data-lucide="bar-chart-2" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',    label:'Q2 Forecast',                sub:'Probability-weighted',              action: function(){ closePalette(); insertSkill('forecast'); } },
  { type:'skill', icon:'<i data-lucide="mail" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',           label:'Outreach Generator',         sub:'LinkedIn / email',                  action: function(){ closePalette(); insertSkill('outreach'); } },
  { type:'skill', icon:'<i data-lucide="target" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',         label:'Pitch Partner',              sub:'Prep for meeting',                  action: function(){ closePalette(); insertSkill('pitch'); } },
  { type:'skill', icon:'<i data-lucide="clipboard" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',      label:'Pre-Meeting Brief',          sub:'90-second battle card',             action: function(){ closePalette(); insertSkill('brief'); } },
  { type:'skill', icon:'<i data-lucide="dollar-sign" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',    label:'Revenue Simulation',         sub:'3 scenarios',                       action: function(){ closePalette(); insertSkill('revenue'); } },
  { type:'skill', icon:'<i data-lucide="radio" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',          label:'Signal to Action',           sub:'Convert signal to plan',            action: function(){ closePalette(); insertSkill('signal'); } },
  { type:'skill', icon:'<i data-lucide="layout-dashboard" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>', label:'Dashboard',                sub:'Command Center',                    action: function(){ closePalette(); showTab('dashboard'); } },
  { type:'skill', icon:'<i data-lucide="message-circle" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>', label:'GTM Assistant',              sub:'AI chat',                           action: function(){ closePalette(); showTab('chat'); } },
  { type:'skill', icon:'<i data-lucide="file-text" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',      label:'Meeting Notes',              sub:'Process meeting notes',             action: function(){ closePalette(); showTab('notes'); } },
  { type:'skill', icon:'<i data-lucide="building" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',       label:'Accounts',                   sub:'All accounts',                      action: function(){ closePalette(); showTab('accounts'); } },
  { type:'skill', icon:'<i data-lucide="gamepad-2" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',      label:'Pitch Simulator',            sub:'Practice your pitch',               action: function(){ closePalette(); showTab('simulator'); } },
];

let _cmdSelected = 0;
let _cmdVisible  = [];

function openPalette() {
  document.getElementById('cmd-overlay').classList.add('open');
  const inp = document.getElementById('cmd-input');
  inp.value = '';
  setTimeout(function(){ inp.focus(); }, 50);
  cmdFilter();
}

function closePalette() {
  document.getElementById('cmd-overlay').classList.remove('open');
}

function cmdFilter() {
  const q = (document.getElementById('cmd-input').value || '').toLowerCase().trim();
  const results = document.getElementById('cmd-results');

  // Build command list: static + accounts
  const accountCmds = (allAccounts || []).map(function(a) {
    return {
      type: 'account',
      icon: '<i data-lucide="building" style="width:14px;height:14px;display:inline-block;vertical-align:middle;"></i>',
      label: a.name,
      sub: (a.country !== '\u2014' ? a.country + ' · ' : '') + (a.icp !== '\u2014' ? 'ICP ' + a.icp : ''),
      action: (function(slug, name){ return function(){ closePalette(); selectAccount(slug, name); showTab('chat'); }; })(a.slug, a.name)
    };
  });

  const all = STATIC_COMMANDS.concat(accountCmds);
  _cmdVisible = q ? all.filter(function(c){
    return c.label.toLowerCase().indexOf(q) !== -1 || (c.sub || '').toLowerCase().indexOf(q) !== -1;
  }) : all;

  _cmdSelected = 0;

  if (_cmdVisible.length === 0) {
    results.innerHTML = '<div id="cmd-empty">No results for "' + q + '"</div>';
    return;
  }

  // Split into skills vs accounts
  const skillItems = _cmdVisible.filter(function(c){ return c.type !== 'account'; });
  const accItems   = _cmdVisible.filter(function(c){ return c.type === 'account'; });

  let html = '';
  if (skillItems.length) {
    if (!q) html += '<div class="cmd-section">Actions</div>';
    skillItems.forEach(function(c, i) {
      html += '<div class="cmd-item' + (i === 0 && !accItems.length - 1 ? ' selected' : '') + '" data-idx="' + i + '" onclick="cmdExecute(' + i + ')">' +
        '<span class="cmd-icon">' + c.icon + '</span>' +
        '<span class="cmd-label">' + c.label + '</span>' +
        (c.sub ? '<span class="cmd-sub">' + c.sub + '</span>' : '') +
        '<span class="cmd-arrow">↵</span>' +
      '</div>';
    });
  }
  if (accItems.length) {
    if (!q) html += '<div class="cmd-section">Accounts</div>';  // already English
    accItems.forEach(function(c, i) {
      const idx = skillItems.length + i;
      html += '<div class="cmd-item" data-idx="' + idx + '" onclick="cmdExecute(' + idx + ')">' +
        '<span class="cmd-icon">' + c.icon + '</span>' +
        '<span class="cmd-label">' + c.label + '</span>' +
        (c.sub ? '<span class="cmd-sub">' + c.sub + '</span>' : '') +
        '<span class="cmd-arrow">↵</span>' +
      '</div>';
    });
  }
  results.innerHTML = html;
  if (typeof lucide !== 'undefined') lucide.createIcons();
  cmdHighlight();
}

function cmdHighlight() {
  document.querySelectorAll('.cmd-item').forEach(function(el, i) {
    el.classList.toggle('selected', i === _cmdSelected);
    if (i === _cmdSelected) el.scrollIntoView({ block: 'nearest' });
  });
}

function cmdKey(e) {
  if (e.key === 'Escape') { closePalette(); return; }
  if (e.key === 'ArrowDown') { e.preventDefault(); _cmdSelected = Math.min(_cmdSelected + 1, _cmdVisible.length - 1); cmdHighlight(); return; }
  if (e.key === 'ArrowUp')   { e.preventDefault(); _cmdSelected = Math.max(_cmdSelected - 1, 0); cmdHighlight(); return; }
  if (e.key === 'Enter')     { e.preventDefault(); cmdExecute(_cmdSelected); return; }
}

function cmdExecute(idx) {
  if (_cmdVisible[idx]) _cmdVisible[idx].action();
}

// Global ⌘K / Ctrl+K listener
document.addEventListener('keydown', function(e) {
  if ((e.metaKey || e.ctrlKey) && e.key === 'k') {
    e.preventDefault();
    const isOpen = document.getElementById('cmd-overlay').classList.contains('open');
    if (isOpen) closePalette(); else openPalette();
  }
  if (e.key === 'Escape') closePalette();
});


// ══════════════════════════════════════════════════════════════════
// ONE-CLICK OUTREACH MODAL
// ══════════════════════════════════════════════════════════════════

let _outreachContext = { slug: '', name: '', buyer: '' };
let _outreachFull = '';

function openOutreachModal(slug, name, buyer) {
  _outreachContext = { slug: slug, name: name, buyer: buyer || '' };
  _outreachFull = '';

  document.getElementById('outreach-title').textContent = name || 'Outreach';
  document.getElementById('outreach-subtitle').textContent = (buyer ? buyer + ' · ' : '') + 'LinkedIn · ready to send';
  document.getElementById('outreach-loading').style.display = 'flex';
  document.getElementById('outreach-text').style.display = 'none';
  document.getElementById('outreach-text').textContent = '';
  document.getElementById('outreach-copy-btn').style.display = 'none';
  document.getElementById('outreach-chat-btn').style.display = 'none';
  document.getElementById('outreach-regen-btn').style.display = 'none';
  document.getElementById('outreach-overlay').classList.add('open');

  streamOutreach(slug, name, buyer);
}

function closeOutreachModal() {
  document.getElementById('outreach-overlay').classList.remove('open');
}

async function streamOutreach(slug, name, buyer) {
  const prompt = 'Write a ready-to-send LinkedIn outreach message for ' + name +
    (buyer ? ' — addressed to ' + buyer : '') +
    '. Rules: max 5 sentences, open with a specific public signal about this company, frame the problem (do not pitch), end with a soft ask for 20 minutes. Peer-to-peer tone. No price. No brackets or placeholders — this must be ready to send right now.';

  const textEl  = document.getElementById('outreach-text');
  const loadEl  = document.getElementById('outreach-loading');

  try {
    const res = await fetch('/api/chat', {
      method: 'POST',
      headers: { 'Content-Type': 'application/json' },
      body: JSON.stringify({ messages: [{ role: 'user', content: prompt }], account: slug })
    });

    loadEl.style.display = 'none';
    textEl.style.display = 'block';

    const reader  = res.body.getReader();
    const decoder = new TextDecoder();
    _outreachFull = '';

    while (true) {
      const { done, value } = await reader.read();
      if (done) break;
      const chunk = decoder.decode(value);
      for (const line of chunk.split('\\n')) {
        if (line.startsWith('data: ') && line !== 'data: [DONE]') {
          try {
            const tok = JSON.parse(line.slice(6));
            _outreachFull += tok.text || '';
            textEl.textContent = _outreachFull;
          } catch(err) {}
        }
      }
    }

    document.getElementById('outreach-copy-btn').style.display = 'inline-flex';
    document.getElementById('outreach-chat-btn').style.display = 'inline-flex';
    document.getElementById('outreach-regen-btn').style.display = 'inline-flex';

  } catch(err) {
    loadEl.style.display = 'none';
    textEl.style.display = 'block';
    textEl.textContent = 'Fejl ved generering — prøv igen.';
  }
}

function copyOutreach() {
  if (!_outreachFull) return;
  navigator.clipboard.writeText(_outreachFull).then(function() {
    const btn = document.getElementById('outreach-copy-btn');
    btn.textContent = '✓ Kopieret!';
    btn.style.background = '#00A880';
    setTimeout(function(){ btn.innerHTML = '<i data-lucide="clipboard" style="width:16px;height:16px;display:inline-block;vertical-align:middle;"></i> Copy message'; btn.style.background = '#153EED'; lucide.createIcons(); }, 2000);
  });
}

function openOutreachInChat() {
  closeOutreachModal();
  if (_outreachContext.slug) selectAccount(_outreachContext.slug, _outreachContext.name);
  const inp = document.getElementById('chat-input');
  inp.value = 'Her er den genererede outreach til ' + _outreachContext.name + '. Kan du justere tonen og gøre den mere personlig baseret på account-filerne?\\n\\n' + _outreachFull;
  autoResize(inp);
  showTab('chat');
  document.getElementById('welcome').style.display = 'none';
}

function regenOutreach() {
  openOutreachModal(_outreachContext.slug, _outreachContext.name, _outreachContext.buyer);
}
</script>
<script>lucide.createIcons();</script>
</body>
</html>"""


@app.route("/api/gtm/daily-plan", methods=["GET"])
def gtm_daily_plan():
    """Generate today's 3-account action plan using Claude."""
    try:
        today_str = datetime.date.today().strftime("%A %d %B %Y")

        # Build context from DB if available, else fall back to files
        accounts_context = ""
        signals_context  = ""
        actions_context  = ""

        if CC_DB_OK:
            try:
                db = SessionLocal()
                # Top 5 accounts scored by deal_score * win_probability
                top_accts = db.query(Account).filter(
                    Account.deal_score != None,
                    Account.win_probability != None
                ).all()
                top_accts = sorted(
                    top_accts,
                    key=lambda a: (a.deal_score or 0) * (a.win_probability or 0),
                    reverse=True
                )[:5]
                if top_accts:
                    lines = ["TOP ACCOUNTS (by deal_score × win_prob):"]
                    for a in top_accts:
                        pv = f"€{int((a.pipeline_value or 0)/1000)}K" if a.pipeline_value else "—"
                        lines.append(
                            f"- {a.name} ({a.country}) | ICP {a.icp_score} | Deal {a.deal_score} | "
                            f"Win {a.win_probability}% | Pipeline {pv} | "
                            f"Buyer: {a.named_buyer or 'TBD'} ({a.buyer_role or '—'}) | "
                            f"Stage: {a.deal_stage or 'identified'} | Notes: {(a.notes or '')[:120]}"
                        )
                    accounts_context = "\n".join(lines)

                # Latest 3 signals
                db_sigs = db.query(Signal).filter(Signal.is_active == True).order_by(
                    Signal.severity, Signal.date.desc()
                ).limit(3).all()
                if db_sigs:
                    sig_lines = ["LIVE SIGNALS:"]
                    for s in db_sigs:
                        sig_lines.append(f"- [{s.severity}] {s.title}: {(s.description or '')[:150]}")
                    signals_context = "\n".join(sig_lines)

                # Open actions
                open_actions = db.query(Action).filter(Action.status == "open").order_by(
                    Action.priority, Action.due_date
                ).limit(10).all()
                if open_actions:
                    act_lines = ["OPEN ACTIONS:"]
                    acct_map = {a.id: a.name for a in top_accts}
                    for ac in open_actions:
                        due = ac.due_date.strftime("%Y-%m-%d") if ac.due_date else "—"
                        acct_name = acct_map.get(ac.account_id, f"account_id:{ac.account_id}")
                        act_lines.append(f"- [{ac.priority}] {acct_name}: {ac.title} (due {due})")
                    actions_context = "\n".join(act_lines)

                db.close()
            except Exception as db_err:
                print(f"[GTM] daily-plan DB error: {db_err}")

        # Fallback to file context if DB gave nothing
        if not accounts_context:
            accounts_context = "TOP OPPORTUNITIES:\n" + (read_file("intelligence/top-opportunities.md") or "")[:2000]
        if not signals_context:
            signals_context = "LEAD LOG:\n" + (read_file("intelligence/lead-log.md") or "")[:2000]
        if not actions_context:
            actions_context = "PIPELINE:\n" + (read_file("intelligence/pipeline-dashboard.md") or "")[:1000]

        prompt = f"""You are the JAKALA GTM OS. Today is {today_str}.

Generate TODAY'S ACTION PLAN: exactly 3 accounts to contact today.

Select based on urgency signals: new executive in role (first 90 days), vacant leadership position, deal gone silent (14+ days), tech stack change, upcoming decision window.

Context:
{accounts_context}

{signals_context}

{actions_context}

Return ONLY valid JSON, no markdown, no explanation:
{{
  "date": "{today_str}",
  "plan": [
    {{
      "account": "Company Name",
      "buyer": "Full Name — Title",
      "urgency": "HOT|WARM|NOW",
      "why_today": "One sentence: specific reason to contact TODAY (e.g. new in role 47 days, CDO vacant since Feb, 21 days silent)",
      "opening_line": "First 2 sentences of the outreach message. Specific, relevant, not generic.",
      "gtm_strategy": "Commerce Optimization|Data Revenue Unlock|AI Readiness Accelerator|Experience Transformation",
      "pipeline_value": "€700K"
    }},
    ... (3 total)
  ]
}}"""

        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=1200,
            messages=[{"role": "user", "content": prompt}]
        )
        raw = response.content[0].text.strip()
        raw = re.sub(r'^```(?:json)?\s*', '', raw)
        raw = re.sub(r'\s*```$', '', raw)
        return jsonify(json.loads(raw))
    except Exception as e:
        return jsonify({"error": str(e), "plan": []}), 500


@app.route("/app")
def gtm_app():
    resp = app.make_response(render_template_string(HTML))
    resp.headers["Cache-Control"] = "no-store, no-cache, must-revalidate"
    resp.headers["Pragma"] = "no-cache"
    return resp


LANDING_HTML = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>JAKALA Nordic OS</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap" rel="stylesheet">
<style>
*,*::before,*::after{box-sizing:border-box;margin:0;padding:0}
html,body{height:100%}
body{
  font-family:'Inter',-apple-system,sans-serif;
  background:#060612;
  color:#fff;
  min-height:100vh;
  display:flex;
  flex-direction:column;
  align-items:center;
  justify-content:center;
  padding:40px 24px;
}
/* Subtle grid background */
body::before{
  content:'';
  position:fixed;inset:0;
  background-image:
    linear-gradient(rgba(21,62,237,.04) 1px,transparent 1px),
    linear-gradient(90deg,rgba(21,62,237,.04) 1px,transparent 1px);
  background-size:48px 48px;
  pointer-events:none;
}
/* Blue radial glow top */
body::after{
  content:'';
  position:fixed;top:-200px;left:50%;transform:translateX(-50%);
  width:800px;height:500px;
  background:radial-gradient(ellipse,rgba(21,62,237,.18) 0%,transparent 70%);
  pointer-events:none;
}
.wrapper{
  position:relative;z-index:1;
  width:100%;max-width:680px;
  display:flex;flex-direction:column;align-items:center;
  gap:0;
}
/* Logo */
.logo{
  font-size:11px;font-weight:800;letter-spacing:.22em;
  color:rgba(21,62,237,.9);text-transform:uppercase;
  margin-bottom:40px;
}
/* Headline */
h1{
  font-size:clamp(32px,5vw,52px);
  font-weight:800;
  letter-spacing:-.04em;
  line-height:1.05;
  text-align:center;
  margin-bottom:14px;
}
h1 span{color:#153EED}
.sub{
  font-size:15px;color:rgba(255,255,255,.45);
  text-align:center;margin-bottom:56px;
  line-height:1.6;
}
/* Platform cards */
.cards{
  display:grid;
  grid-template-columns:1fr 1fr;
  gap:16px;
  width:100%;
  margin-bottom:64px;
}
.card{
  background:rgba(255,255,255,.04);
  border:1px solid rgba(255,255,255,.08);
  border-radius:18px;
  padding:32px 28px;
  cursor:pointer;
  text-decoration:none;
  color:inherit;
  transition:all .25s cubic-bezier(.16,1,.3,1);
  display:flex;flex-direction:column;gap:0;
  position:relative;overflow:hidden;
}
.card::before{
  content:'';
  position:absolute;inset:0;
  background:var(--card-glow,transparent);
  opacity:0;
  transition:opacity .25s;
  border-radius:18px;
}
.card:hover{
  border-color:var(--card-border,rgba(255,255,255,.18));
  transform:translateY(-3px);
  box-shadow:0 20px 60px rgba(0,0,0,.4);
}
.card:hover::before{opacity:1}
.card-gtm{
  --card-glow:linear-gradient(135deg,rgba(21,62,237,.08) 0%,transparent 60%);
  --card-border:rgba(21,62,237,.4);
}
.card-cc{
  --card-glow:linear-gradient(135deg,rgba(0,212,160,.06) 0%,transparent 60%);
  --card-border:rgba(0,212,160,.3);
}
.card-icon{
  font-size:28px;margin-bottom:18px;
  width:52px;height:52px;
  border-radius:12px;
  background:rgba(255,255,255,.06);
  border:1px solid rgba(255,255,255,.08);
  display:flex;align-items:center;justify-content:center;
}
.card-tag{
  font-size:9.5px;font-weight:700;letter-spacing:.14em;
  text-transform:uppercase;color:var(--tag-color,rgba(255,255,255,.3));
  margin-bottom:8px;
}
.card-gtm .card-tag{color:rgba(21,62,237,.9);--tag-color:rgba(21,62,237,.9)}
.card-cc  .card-tag{color:rgba(0,212,160,.9);--tag-color:rgba(0,212,160,.9)}
.card-title{
  font-size:20px;font-weight:700;letter-spacing:-.02em;
  margin-bottom:10px;
}
.card-desc{
  font-size:13px;color:rgba(255,255,255,.45);
  line-height:1.65;flex:1;
}
.card-arrow{
  margin-top:24px;
  font-size:13px;font-weight:600;
  color:var(--arrow-color,rgba(255,255,255,.25));
  display:flex;align-items:center;gap:6px;
  transition:gap .2s,color .2s;
}
.card-gtm .card-arrow{color:rgba(21,62,237,.7)}
.card-cc  .card-arrow{color:rgba(0,212,160,.7)}
.card:hover .card-arrow{gap:10px}
/* Quote */
.quote-block{
  text-align:center;
  border-top:1px solid rgba(255,255,255,.06);
  padding-top:40px;
  width:100%;
}
.quote-text{
  font-size:17px;
  font-weight:500;
  font-style:italic;
  color:rgba(255,255,255,.55);
  line-height:1.6;
  letter-spacing:-.01em;
  margin-bottom:14px;
}
.quote-text em{
  color:rgba(255,255,255,.85);
  font-style:normal;
  font-weight:600;
}
.quote-attr{
  font-size:12px;
  font-weight:600;
  letter-spacing:.08em;
  text-transform:uppercase;
  color:rgba(255,255,255,.2);
}
@media(max-width:520px){
  .cards{grid-template-columns:1fr}
  h1{font-size:28px}
}
</style>
</head>
<body>
<div class="wrapper">

  <div class="logo">JAKALA Nordic</div>

  <h1>Commercial<br><span>Operating System</span></h1>
  <p class="sub">Two tools. One commercial engine.</p>

  <div class="cards">

    <a href="/login" class="card card-gtm">
      <div class="card-icon">◎</div>
      <div class="card-tag">AI Assistant</div>
      <div class="card-title">GTM Assistant</div>
      <div class="card-desc">Strategy, outreach, account analysis and pipeline intelligence — powered by Claude AI.</div>
      <div class="card-arrow">Open GTM Assistant →</div>
    </a>

    <a href="/cc" class="card card-cc">
      <div class="card-icon">◉</div>
      <div class="card-tag">Live Platform</div>
      <div class="card-title">Control Center</div>
      <div class="card-desc">Live pipeline, actions, meetings, outreach generator and AI intelligence in one dashboard.</div>
      <div class="card-arrow">Open Control Center →</div>
    </a>

  </div>

  <div class="quote-block">
    <div class="quote-text">"<em>Turn liquid gold into solid bricks</em>"</div>
    <div class="quote-attr">— Michael Drejer</div>
  </div>

</div>
</body>
</html>"""


@app.route("/")
def landing():
    landing_path = Path(__file__).parent / "landing.html"
    if landing_path.exists():
        return landing_path.read_text(encoding="utf-8")
    return render_template_string(LANDING_HTML)


# ══════════════════════════════════════════════════════════════════════════════
# CONTROL CENTER — routes, API, HTML
# ══════════════════════════════════════════════════════════════════════════════

COUNTRY_META = {
    "no": {"flag": "🇳🇴", "name": "Norway",  "color": "#153EED"},
    "dk": {"flag": "🇩🇰", "name": "Denmark", "color": "#C60C30"},
    "se": {"flag": "🇸🇪", "name": "Sweden",  "color": "#006AA7"},
    "uk": {"flag": "🇬🇧", "name": "UK",      "color": "#CF111C"},
    "fr": {"flag": "🇫🇷", "name": "France",  "color": "#0055A4"},
}

STAGE_ORDER = ["identified", "proposed", "negotiating", "active", "completed"]

# ── DB init on startup ────────────────────────────────────────────────────────
if CC_DB_OK:
    try:
        init_db()
    except Exception as e:
        print(f"[CC] DB init warning: {e}")

# ── Auth helpers ──────────────────────────────────────────────────────────────
def cc_current_user():
    uid = session.get("cc_uid")
    if not uid or not CC_DB_OK:
        return None
    db = SessionLocal()
    try:
        return db.query(User).get(uid)
    finally:
        db.close()

def cc_required(f):
    from functools import wraps
    @wraps(f)
    def decorated(*args, **kwargs):
        if not session.get("cc_uid"):
            return redirect("/cc/login")
        return f(*args, **kwargs)
    return decorated

# ── CC API routes ─────────────────────────────────────────────────────────────

@app.route("/cc/login", methods=["GET"])
def cc_login_page():
    if session.get("cc_uid"):
        return redirect("/cc")
    return render_template_string(CC_HTML)

@app.route("/api/cc/login", methods=["POST"])
def cc_api_login():
    if not CC_DB_OK:
        return jsonify({"error": "Database not available"}), 503
    import json as _json
    # Parse JSON from body; fall back to query params (Railway CDN may strip Content-Length)
    try:
        data = request.get_json(force=True, silent=True) or {}
    except Exception:
        data = {}
    email = (data.get("email") or request.args.get("email") or "").strip().lower()
    pw    = (data.get("password") or request.args.get("password") or "").encode()
    db    = SessionLocal()
    try:
        user = db.query(User).filter(User.email == email).first()
        if not user or not bcrypt.checkpw(pw, user.password_hash.encode()):
            return jsonify({"error": "Invalid email or password"}), 401
        session["cc_uid"] = user.id
        return jsonify({"ok": True, "role": user.role, "country": user.country, "name": user.name})
    finally:
        db.close()

@app.route("/api/cc/logout", methods=["POST"])
def cc_api_logout():
    session.pop("cc_uid", None)
    return jsonify({"ok": True})

@app.route("/api/cc/ping", methods=["GET", "POST"])
def cc_ping():
    import os as _os
    db_url = _os.getenv("DATABASE_URL","NOT SET")
    db_info = db_url[:30] + "..." if db_url != "NOT SET" else db_url
    user_count = 0
    try:
        if CC_DB_OK:
            _db = SessionLocal()
            user_count = _db.query(User).count()
            _db.close()
    except Exception as _e:
        db_info += f" | DB ERR: {_e}"
    return jsonify({"ok": True, "cc_db_ok": CC_DB_OK, "db_url": db_info, "users": user_count})

@app.route("/api/cc/me", methods=["GET"])
def cc_api_me():
    u = cc_current_user()
    if not u:
        return jsonify({"error": "Not logged in"}), 401
    return jsonify({"id": u.id, "name": u.name, "role": u.role,
                    "country": u.country, "initials": u.initials})

@app.route("/api/cc/country-data", methods=["GET"])
def cc_country_data():
    u = cc_current_user()
    if not u:
        return jsonify({"error": "Not logged in"}), 401
    country = request.args.get("country") or u.country
    if u.role == "country_head" and country != u.country:
        return jsonify({"error": "Forbidden"}), 403
    if not CC_DB_OK:
        return jsonify({"error": "DB unavailable"}), 503
    db = SessionLocal()
    try:
        accounts = db.query(Account).filter(Account.country == country).all()
        industries = db.query(Industry).all()
        services   = db.query(Service).all()
        signals    = db.query(Signal).filter(
            (Signal.country == country) | (Signal.country == None),
            Signal.is_active == True
        ).order_by(Signal.severity, Signal.date.desc()).limit(8).all()
        predictions = db.query(Prediction).filter(Prediction.country == country).order_by(Prediction.opportunity_score.desc()).limit(5).all()

        svc_map = {s.id: s for s in services}
        ind_map = {i.id: i for i in industries}

        def fmt_account(a):
            acts = [{"id": ac.id, "service_id": ac.service_id,
                     "service_name": svc_map[ac.service_id].short_name if ac.service_id in svc_map else "?",
                     "service_color": svc_map[ac.service_id].color if ac.service_id in svc_map else "#888",
                     "stage": ac.stage, "manager": ac.manager,
                     "cost": ac.cost_estimate, "weeks": ac.timeline_weeks, "roi": ac.roi_estimate}
                    for ac in a.activations]
            preds = [{"risk": p.risk_score, "opp": p.opportunity_score,
                      "service": svc_map[p.recommended_service_id].short_name if p.recommended_service_id in svc_map else "?",
                      "confidence": p.confidence, "trigger": p.trigger_summary,
                      "weeks": p.timeframe_weeks}
                     for p in a.predictions]
            return {"id": a.id, "name": a.name, "slug": a.slug,
                    "account_type": a.account_type,
                    "industry": ind_map[a.industry_id].name if a.industry_id in ind_map else "Other",
                    "industry_icon": ind_map[a.industry_id].icon if a.industry_id in ind_map else "building",
                    "industry_slug": ind_map[a.industry_id].slug if a.industry_id in ind_map else "",
                    "icp": a.icp_score, "deal": a.deal_score,
                    "pipeline": a.pipeline_value, "win_prob": a.win_probability,
                    "buyer": a.named_buyer, "buyer_role": a.buyer_role,
                    "revenue": a.revenue, "tech_stack": a.tech_stack,
                    "activations": acts, "predictions": preds}

        pipeline_total = sum(a.pipeline_value or 0 for a in accounts)
        named_buyers   = sum(1 for a in accounts if a.named_buyer and a.named_buyer != "TBD")
        active_acts    = sum(1 for a in accounts for ac in a.activations if ac.stage == "active")

        meta = COUNTRY_META.get(country, {"flag": "🌍", "name": country.upper(), "color": "#153EED"})

        return jsonify({
            "country": country,
            "meta": meta,
            "kpis": {"pipeline": pipeline_total, "accounts": len(accounts),
                     "buyers": named_buyers, "active_activations": active_acts},
            "accounts": [fmt_account(a) for a in sorted(accounts, key=lambda x: -(x.deal_score or 0))],
            "industries": [{"id": i.id, "name": i.name, "slug": i.slug, "icon": i.icon,
                            "count": sum(1 for a in accounts if a.industry_id == i.id)}
                           for i in industries if any(a.industry_id == i.id for a in accounts)],
            "signals": [{"type": sg.signal_type, "severity": sg.severity, "title": sg.title,
                         "description": sg.description, "action": sg.action_recommended,
                         "vertical": sg.vertical}
                        for sg in signals],
            "predictions": [{"account": svc_map.get(p.recommended_service_id, None) and
                              db.query(Account).get(p.account_id) and
                              db.query(Account).get(p.account_id).name,
                              "risk": p.risk_score, "opp": p.opportunity_score,
                              "trigger": p.trigger_summary, "confidence": p.confidence,
                              "service": svc_map[p.recommended_service_id].short_name if p.recommended_service_id in svc_map else "?",
                              "weeks": p.timeframe_weeks}
                             for p in predictions],
        })
    finally:
        db.close()

@app.route("/api/cc/global-data", methods=["GET"])
def cc_global_data():
    u = cc_current_user()
    if not u or u.role != "global":
        return jsonify({"error": "Forbidden"}), 403
    if not CC_DB_OK:
        return jsonify({"error": "DB unavailable"}), 503
    db = SessionLocal()
    try:
        countries_data = []
        for code, meta in COUNTRY_META.items():
            accs = db.query(Account).filter(Account.country == code).all()
            if not accs:
                countries_data.append({"code": code, "meta": meta,
                                        "pipeline": 0, "accounts": 0, "buyers": 0, "top_service": "—"})
                continue
            pipeline   = sum(a.pipeline_value or 0 for a in accs)
            buyers     = sum(1 for a in accs if a.named_buyer and a.named_buyer != "TBD")
            # most common service
            svc_counts: dict = {}
            for a in accs:
                for ac in a.activations:
                    svc_counts[ac.service_id] = svc_counts.get(ac.service_id, 0) + 1
            top_svc_id = max(svc_counts, key=svc_counts.get) if svc_counts else None
            top_svc    = db.query(Service).get(top_svc_id).short_name if top_svc_id else "—"
            countries_data.append({"code": code, "meta": meta,
                                    "pipeline": pipeline, "accounts": len(accs),
                                    "buyers": buyers, "top_service": top_svc})

        # service performance across all markets
        all_acts = db.query(Activation).all()
        svc_pipeline: dict = {}
        for ac in all_acts:
            acc = db.query(Account).get(ac.account_id)
            if acc:
                svc_pipeline[ac.service_id] = svc_pipeline.get(ac.service_id, 0) + (acc.pipeline_value or 0)
        services = db.query(Service).all()
        svc_perf = sorted([{"id": s.id, "name": s.short_name, "practice": s.practice,
                             "color": s.color, "pipeline": svc_pipeline.get(s.id, 0)}
                           for s in services], key=lambda x: -x["pipeline"])
        max_pipe = svc_perf[0]["pipeline"] if svc_perf else 1

        global_signals = db.query(Signal).filter(Signal.is_active == True).order_by(Signal.severity, Signal.date.desc()).limit(6).all()

        return jsonify({
            "countries": countries_data,
            "total_pipeline": sum(c["pipeline"] for c in countries_data),
            "total_accounts": sum(c["accounts"] for c in countries_data),
            "service_performance": [{"name": s["name"], "practice": s["practice"],
                                      "color": s["color"], "pipeline": s["pipeline"],
                                      "pct": round(s["pipeline"] / max_pipe * 100) if max_pipe else 0}
                                     for s in svc_perf],
            "signals": [{"type": sg.signal_type, "severity": sg.severity, "title": sg.title,
                          "description": sg.description, "vertical": sg.vertical,
                          "country": sg.country or "Global"}
                         for sg in global_signals],
        })
    finally:
        db.close()

@app.route("/api/cc/predict", methods=["POST"])
def cc_generate_prediction():
    u = cc_current_user()
    if not u:
        return jsonify({"error": "Not logged in"}), 401
    try:
        data = request.get_json(force=True, silent=True) or {}
    except Exception:
        data = {}
    account_id = data.get("account_id") or request.args.get("account_id")
    if not CC_DB_OK:
        return jsonify({"error": "DB unavailable"}), 503
    db = SessionLocal()
    try:
        acc     = db.query(Account).get(account_id)
        if not acc:
            return jsonify({"error": "Account not found"}), 404
        signals = db.query(Signal).filter(
            (Signal.country == acc.country) | (Signal.country == None),
            Signal.is_active == True
        ).all()
        ind     = db.query(Industry).get(acc.industry_id)
        acts    = acc.activations
        svcs    = db.query(Service).all()

        prompt = f"""You are a senior commercial strategist at JAKALA, a global data and digital experience company.

Analyze this account and generate an AI prediction for the next best commercial action.

ACCOUNT:
- Name: {acc.name}
- Country: {acc.country.upper()}
- Industry: {ind.name if ind else 'Unknown'}
- Pipeline value: €{acc.pipeline_value:,.0f}
- ICP score: {acc.icp_score}/10
- Named buyer: {acc.named_buyer or 'TBD'} ({acc.buyer_role or 'Unknown role'})
- Revenue: {acc.revenue or 'Unknown'}
- Tech stack: {acc.tech_stack or 'Unknown'}
- Current activations: {', '.join(a.service.name + ' (' + a.stage + ')' for a in acts) if acts else 'None'}

ACTIVE MARKET SIGNALS ({ind.name if ind else 'market'} + global):
{chr(10).join(f'- [{sg.severity.upper()}] {sg.title}: {sg.description[:200]}' for sg in signals[:4])}

JAKALA SERVICES AVAILABLE:
{chr(10).join(f'- {s.short_name} ({s.practice}): €{s.entry_price_min:,.0f}–{s.entry_price_max:,.0f} entry' for s in svcs)}

Generate a prediction with:
1. OPPORTUNITY SCORE (0-10): How strong is the commercial opportunity right now?
2. RISK SCORE (0-10): What is the risk of revenue loss / deal stalling?
3. RECOMMENDED SERVICE: Which JAKALA service is the best next action?
4. TRIGGER SUMMARY (2-3 sentences): Why now? What market forces, signals or account dynamics make this the moment to act?
5. CONFIDENCE (0-1): How confident are you in this prediction?
6. TIMEFRAME: How many weeks until the window closes or the opportunity peaks?

Respond in JSON format:
{{"opportunity_score": 8.5, "risk_score": 3.0, "recommended_service": "Data Revenue", "trigger_summary": "...", "confidence": 0.82, "timeframe_weeks": 6}}"""

        resp = client.messages.create(
            model=MODEL, max_tokens=512,
            messages=[{"role": "user", "content": prompt}]
        )
        raw = resp.content[0].text.strip()
        # extract JSON
        match = re.search(r'\{.*\}', raw, re.DOTALL)
        if not match:
            return jsonify({"error": "Could not parse prediction"}), 500
        pred_data = json.loads(match.group())

        rec_svc = next((s for s in svcs if s.short_name.lower() in pred_data.get("recommended_service", "").lower()), svcs[0])
        pred = Prediction(
            account_id=account_id, country=acc.country,
            vertical=ind.name if ind else "General",
            risk_score=pred_data.get("risk_score", 5),
            opportunity_score=pred_data.get("opportunity_score", 5),
            trigger_summary=pred_data.get("trigger_summary", ""),
            recommended_service_id=rec_svc.id,
            confidence=pred_data.get("confidence", 0.6),
            timeframe_weeks=pred_data.get("timeframe_weeks", 8),
        )
        db.add(pred); db.commit()
        pred_data["id"] = pred.id
        pred_data["recommended_service_name"] = rec_svc.short_name
        pred_data["recommended_service_color"] = rec_svc.color
        return jsonify(pred_data)
    finally:
        db.close()

# ── Fase 1–4: Actions / Meetings / Commits / Outreach ────────────────────────

@app.route("/api/cc/actions", methods=["GET", "POST"])
def cc_actions():
    u = cc_current_user()
    if not u:
        return jsonify({"error": "Not logged in"}), 401
    if not CC_DB_OK:
        return jsonify({"error": "DB unavailable"}), 503
    db = SessionLocal()
    try:
        if request.method == "GET":
            country  = request.args.get("country") or (u.country if u.role != "global" else None)
            status   = request.args.get("status", "open")
            account_id = request.args.get("account_id")
            q = db.query(Action).join(Account, Action.account_id == Account.id)
            if account_id:
                q = q.filter(Action.account_id == int(account_id))
            elif country:
                q = q.filter(Account.country == country)
            if status != "all":
                q = q.filter(Action.status == status)
            q = q.order_by(Action.due_date.asc().nullslast(), Action.created_at.desc())
            actions = q.all()
            return jsonify([{
                "id": a.id, "account_id": a.account_id,
                "account_name": a.account.name if a.account else "",
                "owner": a.owner, "title": a.title, "description": a.description,
                "due_date": a.due_date.isoformat() if a.due_date else None,
                "priority": a.priority, "status": a.status, "action_type": a.action_type,
                "created_at": a.created_at.isoformat() if a.created_at else None,
                "completed_at": a.completed_at.isoformat() if a.completed_at else None,
            } for a in actions])
        else:
            try:
                data = request.get_json(force=True, silent=True) or {}
            except Exception:
                data = {}
            account_id  = int(data.get("account_id") or request.args.get("account_id") or 0)
            title       = data.get("title") or request.args.get("title", "")
            description = data.get("description") or request.args.get("description", "")
            due_str     = data.get("due_date") or request.args.get("due_date")
            priority    = data.get("priority") or request.args.get("priority", "medium")
            action_type = data.get("action_type") or request.args.get("action_type", "follow-up")
            owner       = data.get("owner") or request.args.get("owner") or u.name
            due_date    = datetime.datetime.fromisoformat(due_str) if due_str else None
            if not title or not account_id:
                return jsonify({"error": "title and account_id required"}), 400
            action = Action(
                account_id=account_id, owner=owner, title=title,
                description=description, due_date=due_date,
                priority=priority, action_type=action_type, status="open"
            )
            db.add(action)
            # Update account last_activity
            acc = db.query(Account).get(account_id)
            if acc:
                acc.last_activity = datetime.datetime.utcnow()
            db.commit()
            return jsonify({"ok": True, "id": action.id})
    finally:
        db.close()

@app.route("/api/cc/actions/<int:action_id>", methods=["PATCH"])
def cc_action_update(action_id):
    u = cc_current_user()
    if not u:
        return jsonify({"error": "Not logged in"}), 401
    if not CC_DB_OK:
        return jsonify({"error": "DB unavailable"}), 503
    db = SessionLocal()
    try:
        try:
            data = request.get_json(force=True, silent=True) or {}
        except Exception:
            data = {}
        action = db.query(Action).get(action_id)
        if not action:
            return jsonify({"error": "Not found"}), 404
        new_status = data.get("status") or request.args.get("status")
        if new_status:
            action.status = new_status
            if new_status == "done":
                action.completed_at = datetime.datetime.utcnow()
        new_title = data.get("title") or request.args.get("title")
        if new_title:
            action.title = new_title
        db.commit()
        return jsonify({"ok": True})
    finally:
        db.close()

@app.route("/api/cc/accounts/<int:account_id>/stage", methods=["PATCH"])
def cc_account_stage(account_id):
    u = cc_current_user()
    if not u:
        return jsonify({"error": "Not logged in"}), 401
    if not CC_DB_OK:
        return jsonify({"error": "DB unavailable"}), 503
    db = SessionLocal()
    try:
        try:
            data = request.get_json(force=True, silent=True) or {}
        except Exception:
            data = {}
        stage = data.get("stage") or request.args.get("stage")
        if not stage:
            return jsonify({"error": "stage required"}), 400
        acc = db.query(Account).get(account_id)
        if not acc:
            return jsonify({"error": "Not found"}), 404
        acc.deal_stage = stage
        acc.deal_stage_updated = datetime.datetime.utcnow()
        acc.last_activity = datetime.datetime.utcnow()
        db.commit()
        return jsonify({"ok": True})
    finally:
        db.close()

@app.route("/api/cc/meetings", methods=["GET", "POST"])
def cc_meetings():
    u = cc_current_user()
    if not u:
        return jsonify({"error": "Not logged in"}), 401
    if not CC_DB_OK:
        return jsonify({"error": "DB unavailable"}), 503
    db = SessionLocal()
    try:
        if request.method == "GET":
            country    = request.args.get("country") or (u.country if u.role != "global" else None)
            account_id = request.args.get("account_id")
            q = db.query(Meeting)
            if account_id:
                q = q.filter(Meeting.account_id == int(account_id))
            elif country:
                q = q.filter(Meeting.country == country)
            meetings = q.order_by(Meeting.date.desc()).limit(50).all()
            return jsonify([{
                "id": m.id, "account_id": m.account_id,
                "account_name": m.account.name if m.account else "",
                "owner": m.owner, "date": m.date.isoformat() if m.date else None,
                "participants": m.participants, "summary": m.summary,
                "outcome": m.outcome, "next_step": m.next_step,
                "created_at": m.created_at.isoformat() if m.created_at else None,
            } for m in meetings])
        else:
            try:
                data = request.get_json(force=True, silent=True) or {}
            except Exception:
                data = {}
            account_id   = int(data.get("account_id") or request.args.get("account_id") or 0)
            date_str     = data.get("date") or request.args.get("date", "")
            participants = data.get("participants") or request.args.get("participants", "")
            summary      = data.get("summary") or request.args.get("summary", "")
            outcome      = data.get("outcome") or request.args.get("outcome", "neutral")
            next_step    = data.get("next_step") or request.args.get("next_step", "")
            owner        = data.get("owner") or request.args.get("owner") or u.name
            if not account_id or not date_str:
                return jsonify({"error": "account_id and date required"}), 400
            meet_date = datetime.datetime.fromisoformat(date_str)
            acc = db.query(Account).get(account_id)
            meeting = Meeting(
                account_id=account_id,
                country=acc.country if acc else u.country,
                owner=owner, date=meet_date, participants=participants,
                summary=summary, outcome=outcome, next_step=next_step
            )
            db.add(meeting)
            if acc:
                acc.last_activity = datetime.datetime.utcnow()
                if outcome == "positive" and acc.deal_stage in ("identified", "qualified"):
                    acc.deal_stage = "engaged"
                    acc.deal_stage_updated = datetime.datetime.utcnow()
            db.commit()
            return jsonify({"ok": True, "id": meeting.id})
    finally:
        db.close()

@app.route("/api/cc/weekly-commit", methods=["GET", "POST"])
def cc_weekly_commit():
    u = cc_current_user()
    if not u:
        return jsonify({"error": "Not logged in"}), 401
    if not CC_DB_OK:
        return jsonify({"error": "DB unavailable"}), 503
    db = SessionLocal()
    try:
        if request.method == "GET":
            # Return current week's commit for this user
            today = datetime.datetime.utcnow().date()
            monday = today - datetime.timedelta(days=today.weekday())
            week_start = datetime.datetime.combine(monday, datetime.time.min)
            commit = db.query(WeeklyCommit).filter(
                WeeklyCommit.user_id == u.id,
                WeeklyCommit.week_start >= week_start
            ).first()
            if not commit:
                return jsonify({"exists": False, "week_start": week_start.isoformat()})
            return jsonify({
                "exists": True, "id": commit.id,
                "week_start": commit.week_start.isoformat(),
                "commit_text": commit.commit_text,
                "target_value": commit.target_value,
                "accounts_committed": commit.accounts_committed,
                "status": commit.status,
            })
        else:
            try:
                data = request.get_json(force=True, silent=True) or {}
            except Exception:
                data = {}
            commit_text = data.get("commit_text") or request.args.get("commit_text", "")
            target_value = float(data.get("target_value") or request.args.get("target_value") or 0)
            accounts_committed = data.get("accounts_committed") or request.args.get("accounts_committed", "[]")
            today = datetime.datetime.utcnow().date()
            monday = today - datetime.timedelta(days=today.weekday())
            week_start = datetime.datetime.combine(monday, datetime.time.min)
            # Upsert
            existing = db.query(WeeklyCommit).filter(
                WeeklyCommit.user_id == u.id,
                WeeklyCommit.week_start >= week_start
            ).first()
            if existing:
                existing.commit_text = commit_text
                existing.target_value = target_value
                existing.accounts_committed = accounts_committed
            else:
                commit = WeeklyCommit(
                    user_id=u.id, country=u.country,
                    week_start=week_start, commit_text=commit_text,
                    target_value=target_value, accounts_committed=accounts_committed
                )
                db.add(commit)
            db.commit()
            return jsonify({"ok": True})
    finally:
        db.close()

@app.route("/api/cc/outreach", methods=["POST"])
def cc_generate_outreach():
    u = cc_current_user()
    if not u:
        return jsonify({"error": "Not logged in"}), 401
    if not CC_DB_OK:
        return jsonify({"error": "DB unavailable"}), 503
    try:
        data = request.get_json(force=True, silent=True) or {}
    except Exception:
        data = {}
    account_id = data.get("account_id") or request.args.get("account_id")
    channel    = data.get("channel") or request.args.get("channel", "linkedin")
    language   = data.get("language") or request.args.get("language", "en")
    if not account_id:
        return jsonify({"error": "account_id required"}), 400
    db = SessionLocal()
    try:
        acc = db.query(Account).get(account_id)
        if not acc:
            return jsonify({"error": "Account not found"}), 404
        ind  = db.query(Industry).get(acc.industry_id)
        acts = acc.activations
        preds = sorted(acc.predictions, key=lambda p: p.generated_at, reverse=True)[:1]
        pred_context = ""
        if preds:
            p = preds[0]
            svc = db.query(Service).get(p.recommended_service_id)
            pred_context = f"Latest AI prediction: {p.trigger_summary} (recommended: {svc.short_name if svc else 'N/A'}, confidence {p.confidence:.0%})"

        lang_map = {"en": "English", "no": "Norwegian (Bokmål)", "da": "Danish", "sv": "Swedish"}
        lang_name = lang_map.get(language, "English")

        prompt = f"""You are a senior commercial strategist at JAKALA writing a personalized {channel} outreach message.

ACCOUNT:
- Name: {acc.name}
- Country: {acc.country.upper()}
- Industry: {ind.name if ind else 'Unknown'}
- Revenue: {acc.revenue or 'Unknown'}
- Named buyer: {acc.named_buyer or 'TBD'} ({acc.buyer_role or 'Unknown role'})
- Tech stack: {acc.tech_stack or 'Unknown'}
- Pipeline value: €{acc.pipeline_value:,.0f}
- Current activations: {', '.join(a.service.name + ' (' + a.stage + ')' for a in acts) if acts else 'None'}
{pred_context}

JAKALA GTM strategies: Data Revenue Unlock, AI Readiness Accelerator, Commerce Optimization, Experience Transformation.

Write a personalized {channel} outreach message in {lang_name}:
- Max 200 words for LinkedIn, 300 for email
- No bullet points in the body
- Lead with a relevant insight or observation about their business
- One clear, soft question at the end
- No jargon, no hard sell
- Peer-to-peer tone — senior consultant to senior buyer
- Sign off as Jacob (JAKALA Nordic)

Return ONLY the message text, nothing else."""

        resp = client.messages.create(
            model=MODEL, max_tokens=600,
            messages=[{"role": "user", "content": prompt}]
        )
        return jsonify({"ok": True, "message": resp.content[0].text.strip(), "channel": channel, "language": language})
    finally:
        db.close()

@app.route("/api/cc/today", methods=["GET"])
def cc_today():
    """Today's priorities: overdue + due-today actions, cold accounts, this week's commit."""
    u = cc_current_user()
    if not u:
        return jsonify({"error": "Not logged in"}), 401
    if not CC_DB_OK:
        return jsonify({"error": "DB unavailable"}), 503
    db = SessionLocal()
    try:
        country = u.country if u.role != "global" else request.args.get("country")
        now = datetime.datetime.utcnow()
        today_end = now.replace(hour=23, minute=59, second=59)

        # Overdue + due today actions
        q = db.query(Action).join(Account, Action.account_id == Account.id).filter(
            Action.status == "open",
            Action.due_date <= today_end
        )
        if country:
            q = q.filter(Account.country == country)
        urgent_actions = q.order_by(Action.due_date.asc()).all()

        # Cold accounts: no activity in 14+ days
        cold_cutoff = now - datetime.timedelta(days=14)
        q2 = db.query(Account).filter(
            Account.account_type == "prospect",
            Account.deal_stage.notin_(["closed_won", "closed_lost"])
        )
        if country:
            q2 = q2.filter(Account.country == country)
        cold_accounts = [a for a in q2.all() if (not a.last_activity) or (a.last_activity < cold_cutoff)][:5]

        # This week's commit
        monday = (now.date() - datetime.timedelta(days=now.weekday()))
        week_start = datetime.datetime.combine(monday, datetime.time.min)
        commit = db.query(WeeklyCommit).filter(
            WeeklyCommit.user_id == u.id,
            WeeklyCommit.week_start >= week_start
        ).first()

        # Recent meetings (last 7 days)
        recent_cutoff = now - datetime.timedelta(days=7)
        q3 = db.query(Meeting)
        if country:
            q3 = q3.filter(Meeting.country == country)
        recent_meetings = q3.filter(Meeting.date >= recent_cutoff).order_by(Meeting.date.desc()).limit(5).all()

        return jsonify({
            "urgent_actions": [{
                "id": a.id, "account_id": a.account_id,
                "account_name": a.account.name if a.account else "",
                "title": a.title, "priority": a.priority, "action_type": a.action_type,
                "due_date": a.due_date.isoformat() if a.due_date else None,
                "is_overdue": bool(a.due_date and a.due_date < now),
            } for a in urgent_actions],
            "cold_accounts": [{
                "id": a.id, "name": a.name,
                "last_activity": a.last_activity.isoformat() if a.last_activity else None,
                "deal_stage": a.deal_stage, "pipeline_value": a.pipeline_value,
            } for a in cold_accounts],
            "weekly_commit": {
                "exists": bool(commit),
                "commit_text": commit.commit_text if commit else None,
                "target_value": commit.target_value if commit else None,
                "week_start": week_start.isoformat(),
            },
            "recent_meetings": [{
                "id": m.id, "account_name": m.account.name if m.account else "",
                "date": m.date.isoformat() if m.date else None,
                "outcome": m.outcome, "next_step": m.next_step,
            } for m in recent_meetings],
        })
    finally:
        db.close()

@app.route("/api/cc/weekly-brief", methods=["GET", "POST"])
def cc_weekly_brief():
    """Auto-generate the weekly strategic brief for leadership."""
    country = request.args.get("country", "no")
    try:
        # Load context
        top_opps = read_file("intelligence/top-opportunities.md") or ""
        pipeline  = read_file("intelligence/pipeline-dashboard.md") or ""
        lead_log  = read_file("intelligence/lead-log.md") or ""

        # Get open actions count from DB
        open_actions = 0
        if CC_DB_OK:
            try:
                db = SessionLocal()
                open_actions = db.query(Action).filter(
                    Action.status == "open",
                    Action.country == country
                ).count()
                db.close()
            except Exception:
                pass

        from datetime import date
        today_str = date.today().strftime("%d %B %Y")
        week_num  = date.today().isocalendar()[1]

        country_name = {"no": "Norway", "dk": "Denmark", "se": "Sweden"}.get(country, country.upper())

        prompt = f"""You are the strategic intelligence engine for JAKALA {country_name}. Today is {today_str}, Week {week_num}.

Generate THIS WEEK'S STRATEGIC BRIEF for the CCO / country lead. This is NOT a market analysis. It is a decision-forcing document.

Context:
TOP OPPORTUNITIES:
{top_opps[:2500]}

PIPELINE DASHBOARD:
{pipeline[:1500]}

RECENT SIGNALS:
{lead_log[:1500]}

Output a JSON object with exactly this structure:
{{
  "week": "Week {week_num}",
  "date": "{today_str}",
  "status": "GREEN|AMBER|RED",
  "status_reason": "One line explaining pipeline health",
  "pipeline_total": "€X.XM",
  "base_case": "€XXXK",
  "risk": {{
    "headline": "Short risk statement (e.g. €1.2M in deals gone silent)",
    "accounts": ["Account A — 21 days", "Account B — 18 days"],
    "implication": "One sentence: what happens if this continues"
  }},
  "top_opportunity": {{
    "account": "Company Name",
    "buyer": "Name — Title",
    "window": "Specific window (e.g. first 90 days ends ~April 15)",
    "why_now": "Why this week specifically",
    "recommended_action": "Exact action to take (e.g. Assign account owner and book discovery call by Friday)"
  }},
  "the_decision": {{
    "question": "The single most important decision to make this week — phrased as a question",
    "deadline": "By when",
    "options": ["Option A", "Option B"],
    "recommendation": "Our recommended choice and why in one sentence"
  }},
  "three_numbers": [
    {{"label": "Pipeline", "value": "€X.XM", "sub": "unweighted Nordic"}},
    {{"label": "Base Case", "value": "€XXXK", "sub": "Q2 2026"}},
    {{"label": "Open Actions", "value": "{open_actions}", "sub": "requiring follow-up"}}
  ]
}}

Be specific. Use real account names. Be direct. No fluff."""

        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=1500,
            messages=[{"role": "user", "content": prompt}]
        )
        import re as _re
        raw = response.content[0].text.strip()
        raw = _re.sub(r'^```(?:json)?\s*', '', raw)
        raw = _re.sub(r'\s*```$', '', raw)
        return jsonify(json.loads(raw))
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ── Fase 4: Intelligence ──────────────────────────────────────────────────────

@app.route("/api/cc/intelligence", methods=["GET"])
def cc_intelligence():
    """Returns cold accounts + churn risks + at-a-glance stats. No AI calls — pure data."""
    u = cc_current_user()
    if not u:
        return jsonify({"error": "Not logged in"}), 401
    if not CC_DB_OK:
        return jsonify({"error": "DB unavailable"}), 503
    db = SessionLocal()
    try:
        country = u.country if u.role != "global" else request.args.get("country")
        now = datetime.datetime.utcnow()

        # Cold prospects: no activity 14+ days, not closed
        cold_cutoff = now - datetime.timedelta(days=14)
        q = db.query(Account).filter(
            Account.account_type == "prospect",
            Account.deal_stage.notin_(["closed_won", "closed_lost"])
        )
        if country:
            q = q.filter(Account.country == country)
        cold_prospects = [a for a in q.all()
                          if (not a.last_activity) or (a.last_activity < cold_cutoff)]

        # Churn risk: existing clients — last activity > 30 days OR negative meeting outcome
        churn_cutoff = now - datetime.timedelta(days=30)
        q2 = db.query(Account).filter(Account.account_type == "existing")
        if country:
            q2 = q2.filter(Account.country == country)
        existing = q2.all()
        churn_risks = []
        for a in existing:
            risk_factors = []
            if not a.last_activity:
                risk_factors.append("No activity recorded")
            elif a.last_activity < churn_cutoff:
                days_cold = (now - a.last_activity).days
                risk_factors.append(f"No activity for {days_cold} days")
            neg_meetings = [m for m in a.meetings if m.outcome == "negative"]
            if neg_meetings:
                risk_factors.append(f"{len(neg_meetings)} negative meeting(s)")
            open_actions = [ac for ac in a.actions if ac.status == "open" and ac.due_date and ac.due_date < now]
            if open_actions:
                risk_factors.append(f"{len(open_actions)} overdue action(s)")
            if risk_factors:
                churn_risks.append({"account": a, "risk_factors": risk_factors})

        # Won accounts (deal_stage = closed_won) for win pattern analysis
        q3 = db.query(Account).filter(Account.deal_stage == "closed_won")
        if country:
            q3 = q3.filter(Account.country == country)
        won_accounts = q3.all()

        def _acc_dict(a, risk_factors=None):
            ind = db.query(Industry).get(a.industry_id) if a.industry_id else None
            days_cold = None
            if a.last_activity:
                days_cold = (now - a.last_activity).days
            return {
                "id": a.id, "name": a.name, "country": a.country,
                "industry": ind.name if ind else "—",
                "pipeline_value": a.pipeline_value,
                "deal_stage": a.deal_stage,
                "named_buyer": a.named_buyer,
                "buyer_role": a.buyer_role,
                "revenue": a.revenue,
                "tech_stack": a.tech_stack,
                "icp_score": a.icp_score,
                "last_activity": a.last_activity.isoformat() if a.last_activity else None,
                "days_cold": days_cold,
                "risk_factors": risk_factors or [],
                "activations": [{"service_name": ac.service.name if ac.service else "?", "stage": ac.stage} for ac in a.activations],
                "meeting_count": len(a.meetings),
                "action_count": len([ac for ac in a.actions if ac.status == "open"]),
            }

        return jsonify({
            "cold_prospects": [_acc_dict(a) for a in sorted(cold_prospects, key=lambda x: x.last_activity or datetime.datetime(2000,1,1))[:10]],
            "churn_risks": [{"account": _acc_dict(r["account"]), "risk_factors": r["risk_factors"]} for r in churn_risks[:5]],
            "won_accounts": [_acc_dict(a) for a in won_accounts[:5]],
            "stats": {
                "cold_count": len(cold_prospects),
                "churn_count": len(churn_risks),
                "won_count": len(won_accounts),
            }
        })
    finally:
        db.close()


@app.route("/api/cc/intelligence/diagnose", methods=["POST"])
def cc_intelligence_diagnose():
    """AI diagnosis for one account — what is the single best action to take right now?"""
    u = cc_current_user()
    if not u:
        return jsonify({"error": "Not logged in"}), 401
    if not CC_DB_OK:
        return jsonify({"error": "DB unavailable"}), 503
    try:
        data = request.get_json(force=True, silent=True) or {}
    except Exception:
        data = {}
    account_id   = data.get("account_id") or request.args.get("account_id")
    insight_type = data.get("insight_type") or request.args.get("insight_type", "cold_reactivation")
    if not account_id:
        return jsonify({"error": "account_id required"}), 400
    db = SessionLocal()
    try:
        acc  = db.query(Account).get(account_id)
        if not acc:
            return jsonify({"error": "Not found"}), 404
        ind  = db.query(Industry).get(acc.industry_id)
        now  = datetime.datetime.utcnow()
        days_cold = (now - acc.last_activity).days if acc.last_activity else None

        # Gather context
        meetings_summary = ""
        recent_meetings = sorted(acc.meetings, key=lambda m: m.date, reverse=True)[:3]
        if recent_meetings:
            meetings_summary = "\n".join(
                f"- {m.date.strftime('%Y-%m-%d')}: {m.outcome} — {(m.summary or '')[:150]}"
                for m in recent_meetings
            )

        open_actions = [a for a in acc.actions if a.status == "open"]
        actions_summary = ", ".join(a.title for a in open_actions[:3]) if open_actions else "None"

        signals = db.query(Signal).filter(
            (Signal.country == acc.country) | (Signal.country == None),
            Signal.is_active == True
        ).limit(3).all()
        signal_txt = "\n".join(f"- [{s.severity}] {s.title}" for s in signals) or "None"

        if insight_type == "churn_risk":
            task = """This is an EXISTING CLIENT. Identify:
1. CHURN RISK LEVEL: Low / Medium / High — and why
2. THE SINGLE most important action to protect and grow this relationship
3. EXPANSION OPPORTUNITY: What is the next logical service to propose?
4. A suggested message opening (1–2 sentences) to use in the next outreach"""
        else:
            task = f"""This account has gone cold ({days_cold or '?'} days without activity).
1. DIAGNOSIS: Why is this account likely cold? What has probably happened on their side?
2. REACTIVATION PLAY: The single best action to restart this relationship right now
3. HOOK: A specific, relevant insight or observation to open with (not generic)
4. A suggested LinkedIn message opening (2–3 sentences) — concrete, no jargon"""

        prompt = f"""You are a senior commercial director at JAKALA analyzing an account.

ACCOUNT: {acc.name} ({acc.country.upper()})
Industry: {ind.name if ind else '—'}
Revenue: {acc.revenue or '—'}
Named buyer: {acc.named_buyer or 'TBD'} ({acc.buyer_role or '—'})
Tech stack: {acc.tech_stack or '—'}
ICP score: {acc.icp_score}/10
Pipeline value: €{acc.pipeline_value:,.0f}
Deal stage: {acc.deal_stage or 'identified'}
Last activity: {acc.last_activity.strftime('%Y-%m-%d') if acc.last_activity else 'Never'}
Open actions: {actions_summary}

Recent meetings:
{meetings_summary or 'None logged'}

Active market signals:
{signal_txt}

Current activations: {', '.join(a.service.name + ' (' + a.stage + ')' for a in acc.activations) if acc.activations else 'None'}

TASK:
{task}

Be specific and direct. Reference actual facts about this company. Max 250 words total.
Format: use short bold headers for each section. No bullet soup."""

        resp = client.messages.create(
            model=MODEL, max_tokens=600,
            messages=[{"role": "user", "content": prompt}]
        )
        return jsonify({
            "ok": True,
            "account_id": acc.id,
            "account_name": acc.name,
            "insight_type": insight_type,
            "diagnosis": resp.content[0].text.strip(),
        })
    finally:
        db.close()


@app.route("/api/cc/intelligence/win-patterns", methods=["POST"])
def cc_intelligence_win_patterns():
    """AI analysis of what won deals and active clients have in common."""
    u = cc_current_user()
    if not u:
        return jsonify({"error": "Not logged in"}), 401
    if not CC_DB_OK:
        return jsonify({"error": "DB unavailable"}), 503
    db = SessionLocal()
    try:
        country = u.country if u.role != "global" else request.args.get("country")
        # Won deals + existing active clients
        q = db.query(Account).filter(
            (Account.deal_stage == "closed_won") | (Account.account_type == "existing")
        )
        if country:
            q = q.filter(Account.country == country)
        won = q.all()
        # Also high-ICP prospects in advanced stages
        q2 = db.query(Account).filter(
            Account.account_type == "prospect",
            Account.deal_stage.in_(["proposed", "negotiating"]),
            Account.icp_score >= 7
        )
        if country:
            q2 = q2.filter(Account.country == country)
        advanced = q2.all()

        if not won and not advanced:
            return jsonify({"ok": True, "analysis": "Not enough data yet. As you log meetings and close deals, patterns will emerge here."})

        def _acc_line(a):
            ind = db.query(Industry).get(a.industry_id) if a.industry_id else None
            svcs = ", ".join(ac.service.name for ac in a.activations) if a.activations else "none"
            return f"- {a.name} ({(ind.name if ind else '—')}, {a.country.upper()}, ICP {a.icp_score}, revenue {a.revenue or '?'}, services: {svcs}, buyer: {a.buyer_role or '?'}, tech: {(a.tech_stack or '—')[:80]}"

        won_lines   = "\n".join(_acc_line(a) for a in won[:8])
        adv_lines   = "\n".join(_acc_line(a) for a in advanced[:5])

        # High-ICP cold prospects for comparison
        q3 = db.query(Account).filter(Account.icp_score >= 8, Account.account_type == "prospect")
        if country:
            q3 = q3.filter(Account.country == country)
        cold_high = q3.order_by(Account.icp_score.desc()).limit(5).all()
        cold_lines = "\n".join(_acc_line(a) for a in cold_high)

        prompt = f"""You are a senior commercial director at JAKALA analyzing win patterns across the Nordic pipeline.

ACTIVE CLIENTS / WON DEALS:
{won_lines or 'None yet'}

ADVANCED PROSPECTS (proposed/negotiating):
{adv_lines or 'None yet'}

HIGH-ICP COLD PROSPECTS (for comparison):
{cold_lines or 'None yet'}

Analyze and answer:
1. WIN PATTERN: What do the won/active accounts have in common? (industry, tech stack, buyer role, revenue size, service entry point)
2. IDEAL ENTRY WEDGE: Which JAKALA service has the highest conversion rate based on this data?
3. NEXT BEST BET: Which 2–3 cold prospects most closely match the win pattern — and why?
4. BLIND SPOT: What type of account or buyer are we systematically NOT winning? What should we test?

Be specific. Reference actual account names. Max 300 words. Use bold section headers."""

        resp = client.messages.create(
            model=MODEL, max_tokens=700,
            messages=[{"role": "user", "content": prompt}]
        )
        return jsonify({"ok": True, "analysis": resp.content[0].text.strip()})
    finally:
        db.close()

# ── Control Center HTML ────────────────────────────────────────────────────────

@app.route('/api/cc/strategic-brief', methods=['POST'])
def cc_strategic_brief():
    if not CC_DB_OK:
        return jsonify({'error': 'DB unavailable'}), 503
    db = SessionLocal()
    try:
        country = request.args.get('country') or request.json.get('country', 'no')
        signals = db.query(Signal).filter(Signal.is_active == True).order_by(Signal.date.desc()).limit(10).all()
        accounts = db.query(Account).filter(Account.country == country).all()
        pipeline_value = sum(a.pipeline_value or 0 for a in accounts)
        high_icp = [a for a in accounts if (a.icp_score or 0) >= 8]
        engaged = [a for a in accounts if a.deal_stage in ('engaged','proposed','negotiating')]
        signals_text = '\n'.join([f'- [{s.severity.upper()}] {s.title}: {s.description}' for s in signals[:8]])
        top_accounts = sorted(high_icp, key=lambda a: (a.icp_score or 0) + (a.deal_score or 0), reverse=True)[:5]
        accounts_text = '\n'.join([f'- {a.name} (ICP {a.icp_score}, stage: {a.deal_stage}, value: €{int(a.pipeline_value or 0):,})' for a in top_accounts])
        prompt = f"""You are the Chief Commercial Officer of JAKALA Nordic. Generate a concise strategic market intelligence brief for this week.

Market: {country.upper()} | Pipeline: €{int(pipeline_value):,} | High-ICP accounts: {len(high_icp)} | Active engagements: {len(engaged)}

MARKET SIGNALS:
{signals_text}

TOP PRIORITY ACCOUNTS:
{accounts_text}

Write a 3-paragraph executive brief:
1. MARKET PULSE: What is happening in the market right now that demands attention?
2. STRATEGIC OPPORTUNITY: Where is the clearest path to new revenue in the next 90 days?
3. RECOMMENDED ACTION: What should the commercial leadership team do THIS WEEK to stay ahead?

Be direct, concrete, and commercially sharp. No fluff. Reference specific signals and accounts where relevant."""
        resp = client.messages.create(
            model=MODEL, max_tokens=600,
            messages=[{'role':'user','content':prompt}]
        )
        return jsonify({'brief': resp.content[0].text})
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        db.close()


@app.route('/api/cc/foresight', methods=['POST'])
def cc_foresight():
    if not CC_DB_OK:
        return jsonify({'error': 'DB unavailable'}), 503
    db = SessionLocal()
    try:
        body = request.get_json(silent=True) or {}
        country = request.args.get('country') or body.get('country', 'no')
        signals = db.query(Signal).filter(Signal.is_active == True).order_by(Signal.date.desc()).limit(15).all()
        accounts = db.query(Account).filter(Account.country == country).all()

        top_accounts = sorted(accounts, key=lambda a: (a.icp_score or 0) + (a.deal_score or 0), reverse=True)[:20]
        acc_lines = []
        for a in top_accounts:
            acc_lines.append(
                f"- {a.name} | ICP:{a.icp_score} Deal:{a.deal_score} | Stage:{a.deal_stage} | "
                f"Value:€{int(a.pipeline_value or 0):,} | Buyer:{a.named_buyer or 'TBD'} | "
                f"Tech:{(a.tech_stack or 'Unknown')[:80]} | Notes:{(a.notes or '')[:60]}"
            )
        sig_lines = [f"- [{s.severity.upper()}|{s.signal_type}] {s.title}: {(s.description or '')[:120]}" for s in signals]

        knowledge_files = ['knowledge/gtm-strategy.md', 'knowledge/offerings.md', 'knowledge/strategy-mapping.md']
        knowledge = ''
        for f in knowledge_files:
            p = BASE_DIR / f
            if p.exists():
                knowledge += p.read_text(encoding='utf-8', errors='replace')[:800] + '\n'

        prompt = f"""You are JAKALA Nordic's strategic intelligence engine. Generate a FORESIGHT REPORT with exactly three sections.

CONTEXT — {country.upper()} MARKET:
Accounts tracked: {len(accounts)} | Total pipeline: €{int(sum(a.pipeline_value or 0 for a in accounts)):,}

TOP ACCOUNTS:
{chr(10).join(acc_lines[:15])}

ACTIVE MARKET SIGNALS:
{chr(10).join(sig_lines[:12])}

JAKALA OFFERINGS & GTM STRATEGY (summary):
{knowledge[:600]}

Generate the following three sections. Use this exact format — no deviations:

---BUYING_WINDOWS---
List 5 accounts with imminent buying windows. For each use this exact format:
ACCOUNT: [company name]
TRIGGER: [specific signal or pattern that drives this prediction — be concrete]
WINDOW: [15-30 days / 30-60 days / 60-90 days]
ENTRY: [recommended GTM strategy and specific entry offer]
VALUE: [estimated deal value in €]
CONFIDENCE: [High / Medium / Low]
---

---PARTNERSHIPS---
Recommend 3 technology or agency partnerships for JAKALA Nordic. For each:
PARTNER: [company/platform name]
EVIDENCE: [why now — what market data supports this]
PROPOSITION: [specific joint offer or activation]
VALUE: [estimated annual revenue unlock in €]
ACTION: [one concrete step to initiate this week]
---

---NEW_SERVICES---
Identify 2 service/activation gaps the Nordic market needs but no agency offers yet. For each:
SERVICE: [name the service offering]
EVIDENCE: [what market signals / account patterns reveal this gap]
CLIENTS: [name 3-5 specific accounts that would buy this]
ENTRY_OFFER: [the first thing to sell — scope, price, duration]
REVENUE: [12-month potential in €]
BUILD_OR_PARTNER: [should JAKALA build this or partner?]
---

Be commercially sharp, specific, and grounded in the account and signal data. No generic advice. Reference named accounts and real signals where possible."""

        resp = client.messages.create(model=MODEL, max_tokens=2000,
            messages=[{'role': 'user', 'content': prompt}])
        raw = resp.content[0].text

        def parse_section(text, start_marker, end_marker):
            try:
                start = text.index(start_marker) + len(start_marker)
                end = text.index(end_marker, start)
                return text[start:end].strip()
            except ValueError:
                return text

        buying = parse_section(raw, '---BUYING_WINDOWS---', '---')
        partnerships = parse_section(raw, '---PARTNERSHIPS---', '---')
        services = parse_section(raw, '---NEW_SERVICES---', '---')

        def parse_blocks(text, keys):
            blocks = []
            current = {}
            for line in text.split('\n'):
                line = line.strip()
                if not line:
                    if current:
                        blocks.append(current)
                        current = {}
                    continue
                for k in keys:
                    if line.startswith(k + ':'):
                        current[k.lower()] = line[len(k)+1:].strip()
                        break
            if current:
                blocks.append(current)
            return blocks

        buying_blocks = parse_blocks(buying, ['ACCOUNT','TRIGGER','WINDOW','ENTRY','VALUE','CONFIDENCE'])
        partner_blocks = parse_blocks(partnerships, ['PARTNER','EVIDENCE','PROPOSITION','VALUE','ACTION'])
        service_blocks = parse_blocks(services, ['SERVICE','EVIDENCE','CLIENTS','ENTRY_OFFER','REVENUE','BUILD_OR_PARTNER'])

        return jsonify({
            'buying_windows': buying_blocks,
            'partnerships': partner_blocks,
            'new_services': service_blocks,
            'raw': raw
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        db.close()


CC_HTML = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>JAKALA Control Center</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap" rel="stylesheet">
<style>
*,*::before,*::after{box-sizing:border-box;margin:0;padding:0}
:root{
  --bg:#F4F6FC;--sb:#FFFFFF;--card:#FFFFFF;
  --border:#E4E7F0;--border2:#EEF0F8;
  --blue:#153EED;--blue-dim:rgba(21,62,237,.09);--blue-glow:rgba(21,62,237,.2);
  --green:#00A67E;--red:#DC2626;--amber:#D97706;--purple:#7C3AED;
  --w:#0F0F1A;--t:#1F2937;--m:#6B7280;--m2:#9CA3AF;
  --font:'Inter',-apple-system,sans-serif;--sb-w:258px;--radius:12px;
}
body{font-family:var(--font);background:var(--bg);color:var(--t);min-height:100vh;overflow:hidden;-webkit-font-smoothing:antialiased}
/* ── SCREENS ── */
.screen{position:fixed;inset:0;display:flex;align-items:center;justify-content:center;transition:opacity .3s}
.screen.hidden{opacity:0;pointer-events:none}
/* ── LOGIN ── */
#login-screen{background:radial-gradient(ellipse 80% 60% at 50% 0%,rgba(21,62,237,.07) 0%,transparent 70%),var(--bg)}
.login-card{width:380px;padding:48px 40px;background:#FFFFFF;border:1px solid var(--border);border-radius:20px;box-shadow:0 4px 32px rgba(21,62,237,.07)}
.login-logo{font-size:13px;font-weight:800;letter-spacing:.18em;color:var(--blue);margin-bottom:10px}
.login-title{font-size:28px;font-weight:700;letter-spacing:-.03em;margin-bottom:6px;color:var(--w)}
.login-sub{font-size:13px;color:var(--m);margin-bottom:32px}
.login-field{display:flex;flex-direction:column;gap:6px;margin-bottom:16px}
.login-field label{font-size:12px;font-weight:600;color:var(--m);letter-spacing:.06em;text-transform:uppercase}
.login-field input{background:#F8F9FE;border:1px solid var(--border);border-radius:8px;padding:12px 14px;font:500 14px var(--font);color:var(--w);outline:none;transition:border-color .2s}
.login-field input:focus{border-color:rgba(21,62,237,.5)}
.login-btn{width:100%;padding:13px;background:var(--blue);border:none;border-radius:8px;font:700 14px var(--font);color:#fff;cursor:pointer;transition:opacity .2s;margin-top:8px}
.login-btn:hover{opacity:.85}
.login-err{font-size:12px;color:var(--red);margin-top:10px;min-height:18px}
/* ── APP SHELL ── */
#app-shell{position:fixed;inset:0;display:flex;flex-direction:column}
#app-shell.hidden{display:none}
/* ── TOP BAR ── */
.topbar{height:54px;display:flex;align-items:center;padding:0 24px;border-bottom:1px solid var(--border);background:var(--sb);flex-shrink:0;gap:16px;box-shadow:0 1px 0 var(--border)}
.topbar-logo{font-size:11px;font-weight:800;letter-spacing:.18em;color:var(--blue)}
.topbar-sep{width:1px;height:18px;background:var(--border)}
.topbar-title{font-size:13px;font-weight:600;color:var(--m)}
.topbar-country{display:flex;align-items:center;gap:7px;font-size:13px;font-weight:600;color:var(--w)}
.topbar-right{margin-left:auto;display:flex;align-items:center;gap:12px}
.topbar-user{display:flex;align-items:center;gap:10px;font-size:13px;color:var(--m)}
.topbar-avatar{width:30px;height:30px;border-radius:50%;display:flex;align-items:center;justify-content:center;font-size:11px;font-weight:700;color:#fff;background:var(--blue);flex-shrink:0}
.topbar-link{font-size:12px;color:var(--m);text-decoration:none;padding:6px 10px;border-radius:6px;transition:background .15s;cursor:pointer}
.topbar-link:hover{background:rgba(21,62,237,.07);color:var(--blue)}
/* ── MAIN LAYOUT ── */
.main-layout{display:flex;flex:1;overflow:hidden}
/* ── SIDEBAR ── */
.sidebar{width:var(--sb-w);background:var(--sb);border-right:1px solid var(--border);display:flex;flex-direction:column;flex-shrink:0;overflow-y:auto}
.sb-section{padding:20px 14px 0}
.sb-label{font-size:9.5px;font-weight:700;letter-spacing:.1em;text-transform:uppercase;color:var(--m2);padding:0 8px;margin-bottom:8px}
.nav-item{display:flex;align-items:center;gap:10px;padding:9px 10px;border-radius:8px;font-size:13px;font-weight:500;color:var(--m);cursor:pointer;transition:all .15s;margin-bottom:2px}
.nav-item:hover{background:rgba(21,62,237,.06);color:var(--t)}
.nav-item.active{background:var(--blue-dim);color:var(--blue);font-weight:600}
.nav-item .nav-icon{width:16px;text-align:center;font-size:14px;flex-shrink:0}
.sb-divider{height:1px;background:var(--border);margin:14px 14px}
.sb-signals{padding:14px 14px 0}
.sig-row{display:flex;align-items:center;gap:8px;font-size:12px;padding:6px 8px;border-radius:6px}
.sig-dot{width:6px;height:6px;border-radius:50%;flex-shrink:0}
.sig-dot.critical{background:var(--red)}
.sig-dot.warning{background:var(--amber)}
.sig-dot.info{background:var(--blue)}
.sb-country-switcher{padding:14px}
.country-pill{display:flex;align-items:center;gap:8px;padding:8px 10px;border-radius:8px;background:#F8F9FE;border:1px solid var(--border);font-size:12px;font-weight:500;color:var(--m);cursor:pointer;margin-bottom:6px;transition:all .15s}
.country-pill:hover{border-color:rgba(21,62,237,.4);color:var(--blue)}
.country-pill.active{background:var(--blue-dim);border-color:rgba(21,62,237,.4);color:var(--blue);font-weight:600}
/* ── CONTENT ── */
.content{flex:1;overflow-y:auto;padding:28px 32px}
/* ── VIEWS ── */
.view{display:none}.view.active{display:block}
/* ── KPI ROW ── */
.kpi-row{display:grid;grid-template-columns:repeat(4,1fr);gap:14px;margin-bottom:28px}
.kpi-card{background:var(--card);border:1px solid var(--border);border-radius:var(--radius);padding:20px 22px}
.kpi-label{font-size:11px;font-weight:600;letter-spacing:.06em;text-transform:uppercase;color:var(--m);margin-bottom:10px}
.kpi-value{font-size:30px;font-weight:700;letter-spacing:-.04em;line-height:1}
.kpi-sub{font-size:11px;color:var(--m2);margin-top:6px}
/* ── SECTION HEADER ── */
.sec-header{display:flex;align-items:center;justify-content:space-between;margin-bottom:16px}
.sec-title{font-size:15px;font-weight:700;letter-spacing:-.02em}
/* ── INDUSTRY PILLS ── */
.industry-filter{display:flex;flex-wrap:wrap;gap:8px;margin-bottom:24px}
.ind-pill{padding:6px 14px;border-radius:20px;font-size:12px;font-weight:600;border:1px solid var(--border);color:var(--m);background:#fff;cursor:pointer;transition:all .2s;white-space:nowrap}
.ind-pill:hover{border-color:rgba(21,62,237,.4);color:var(--blue)}
.ind-pill.active{background:var(--blue-dim);border-color:rgba(21,62,237,.5);color:var(--blue);font-weight:700}
.ind-count{font-size:10px;background:rgba(0,0,0,.07);border-radius:10px;padding:1px 6px;margin-left:4px}
/* ── ACCOUNT GRID ── */
.account-grid{display:grid;grid-template-columns:repeat(auto-fill,minmax(310px,1fr));gap:14px}
.account-card{background:var(--card);border:1px solid var(--border);border-radius:var(--radius);padding:18px 20px;cursor:pointer;transition:all .2s;position:relative;overflow:hidden}
.account-card::before{content:'';position:absolute;left:0;top:0;bottom:0;width:3px;background:var(--svc-color,var(--blue));border-radius:3px 0 0 3px}
.account-card:hover{border-color:rgba(21,62,237,.3);transform:translateY(-1px);box-shadow:0 6px 24px rgba(21,62,237,.08)}
.ac-head{display:flex;justify-content:space-between;align-items:flex-start;margin-bottom:10px}
.ac-name{font-size:15px;font-weight:700;letter-spacing:-.02em}
.ac-value{font-size:13px;font-weight:700;color:var(--m);white-space:nowrap;margin-left:8px}
.ac-tags{display:flex;flex-wrap:wrap;gap:5px;margin-bottom:12px}
.ac-tag{padding:3px 8px;border-radius:20px;font-size:10px;font-weight:600;background:#F0F2F9;color:var(--m)}
.ac-tag.icp-high{background:rgba(0,166,126,.12);color:var(--green)}
.ac-tag.icp-mid{background:rgba(21,62,237,.09);color:var(--blue)}
.ac-tag.existing{background:rgba(0,166,126,.1);color:var(--green)}
/* ── ACTIVATION PILLS ── */
.act-section{margin-bottom:12px}
.act-label{font-size:9.5px;font-weight:700;letter-spacing:.07em;text-transform:uppercase;color:var(--m2);margin-bottom:6px}
.act-pills{display:flex;flex-wrap:wrap;gap:5px}
.act-pill{display:flex;align-items:center;gap:5px;padding:4px 9px;border-radius:6px;font-size:11px;font-weight:600;border:1px solid;cursor:default}
.act-stage-dot{width:5px;height:5px;border-radius:50%;background:currentColor;flex-shrink:0}
/* ── BUYER ROW ── */
.buyer-row{display:flex;align-items:center;justify-content:space-between;margin-top:12px;padding-top:12px;border-top:1px solid var(--border2)}
.buyer-left{display:flex;align-items:center;gap:8px}
.buyer-avatar{width:26px;height:26px;border-radius:50%;display:flex;align-items:center;justify-content:center;font-size:9px;font-weight:700;color:#fff;background:var(--blue);flex-shrink:0}
.buyer-name{font-size:12px;font-weight:600;color:var(--t)}
.buyer-role{font-size:10px;color:var(--m)}
/* ── WIN BAR ── */
.win-bar-wrap{display:flex;align-items:center;gap:8px}
.win-bar{width:60px;height:4px;background:rgba(0,0,0,.08);border-radius:2px;overflow:hidden}
.win-fill{height:100%;border-radius:2px;background:var(--fill-color,var(--green));transition:width .6s ease}
.win-pct{font-size:11px;font-weight:700;color:var(--m)}
/* ── SIGNALS VIEW ── */
.signal-list{display:flex;flex-direction:column;gap:12px}
.signal-card{background:var(--card);border:1px solid var(--border);border-radius:var(--radius);padding:18px 20px}
.signal-card.critical{border-left:3px solid var(--red)}
.signal-card.warning{border-left:3px solid var(--amber)}
.signal-card.info{border-left:3px solid var(--blue)}
.sig-head{display:flex;align-items:flex-start;gap:12px;margin-bottom:8px}
.sig-badge{padding:3px 8px;border-radius:4px;font-size:10px;font-weight:700;letter-spacing:.05em;text-transform:uppercase;white-space:nowrap;flex-shrink:0}
.sig-badge.critical{background:rgba(246,87,74,.15);color:var(--red)}
.sig-badge.warning{background:rgba(245,158,11,.15);color:var(--amber)}
.sig-badge.info{background:rgba(21,62,237,.09);color:var(--blue)}
.sig-title{font-size:14px;font-weight:700;line-height:1.3}
.sig-desc{font-size:12.5px;color:var(--m);line-height:1.6;margin-bottom:10px}
.sig-action{background:rgba(21,62,237,.06);border:1px solid rgba(21,62,237,.18);border-radius:8px;padding:10px 14px;font-size:12px;color:var(--blue);line-height:1.5}
.sig-action-label{font-size:9.5px;font-weight:700;letter-spacing:.08em;text-transform:uppercase;color:var(--blue);margin-bottom:4px}
/* ── PREDICTIONS VIEW ── */
.pred-grid{display:grid;grid-template-columns:repeat(auto-fill,minmax(300px,1fr));gap:14px}
.pred-card{background:var(--card);border:1px solid var(--border);border-radius:var(--radius);padding:20px}
.pred-scores{display:flex;gap:16px;margin:14px 0}
.pred-score{text-align:center;flex:1}
.pred-score-val{font-size:28px;font-weight:800;letter-spacing:-.04em}
.pred-score-label{font-size:9.5px;font-weight:700;letter-spacing:.06em;text-transform:uppercase;color:var(--m);margin-top:3px}
.pred-trigger{font-size:12px;color:var(--m);line-height:1.6;margin-top:10px;padding-top:10px;border-top:1px solid var(--border2)}
.pred-confidence{font-size:11px;font-weight:600;margin-top:8px}
.confidence-bar{height:3px;background:rgba(0,0,0,.08);border-radius:2px;margin-top:4px;overflow:hidden}
.confidence-fill{height:100%;border-radius:2px;background:var(--green)}
.pred-gen-btn{width:100%;margin-top:14px;padding:9px;background:rgba(21,62,237,.08);border:1px solid rgba(21,62,237,.2);border-radius:8px;color:var(--blue);font:600 12px var(--font);cursor:pointer;transition:all .2s}
.pred-gen-btn:hover{background:rgba(21,62,237,.15)}
/* ── GLOBAL VIEW ── */
.country-cards{display:grid;grid-template-columns:repeat(5,1fr);gap:12px;margin-bottom:28px}
.cc-card{background:var(--card);border:1px solid var(--border);border-radius:var(--radius);padding:18px;cursor:pointer;transition:all .2s;text-align:center}
.cc-card:hover{border-color:rgba(21,62,237,.3);transform:translateY(-2px);box-shadow:0 4px 16px rgba(21,62,237,.07)}
.cc-flag{font-size:28px;margin-bottom:8px}
.cc-cname{font-size:13px;font-weight:700;margin-bottom:12px}
.cc-kpi{margin-bottom:6px}
.cc-kpi-v{font-size:20px;font-weight:700;letter-spacing:-.03em}
.cc-kpi-l{font-size:10px;color:var(--m)}
/* ── SERVICE CHART ── */
.svc-chart{display:flex;flex-direction:column;gap:12px}
.svc-bar-row{display:flex;align-items:center;gap:14px}
.svc-bar-name{font-size:12px;font-weight:600;width:180px;flex-shrink:0}
.svc-bar-track{flex:1;height:8px;background:rgba(0,0,0,.06);border-radius:4px;overflow:hidden}
.svc-bar-fill{height:100%;border-radius:4px;transition:width .8s ease}
.svc-bar-val{font-size:12px;font-weight:600;color:var(--m);width:80px;text-align:right;flex-shrink:0}
/* ── DETAIL PANEL (slide-in) ── */
.detail-overlay{position:fixed;inset:0;background:rgba(0,0,0,.3);z-index:100;opacity:0;pointer-events:none;transition:opacity .25s}
.detail-overlay.open{opacity:1;pointer-events:all}
.detail-panel{position:fixed;right:0;top:0;bottom:0;width:480px;background:#FFFFFF;border-left:1px solid var(--border);z-index:101;overflow-y:auto;transform:translateX(100%);transition:transform .3s cubic-bezier(.16,1,.3,1);box-shadow:-8px 0 40px rgba(21,62,237,.06)}
.detail-panel.open{transform:none}
.dp-head{padding:24px;border-bottom:1px solid var(--border);position:sticky;top:0;background:#FFFFFF;z-index:1}
.dp-close{float:right;background:none;border:none;color:var(--m);font-size:20px;cursor:pointer;padding:2px 6px;border-radius:4px}
.dp-close:hover{color:var(--w)}
.dp-body{padding:24px}
.dp-section{margin-bottom:24px}
.dp-sec-title{font-size:11px;font-weight:700;letter-spacing:.08em;text-transform:uppercase;color:var(--m2);margin-bottom:12px}
.dp-row{display:flex;justify-content:space-between;padding:8px 0;border-bottom:1px solid var(--border2);font-size:13px}
.dp-row:last-child{border-bottom:none}
.dp-row-label{color:var(--m)}
.dp-row-val{font-weight:600;max-width:220px;text-align:right}
.activation-timeline{display:flex;flex-direction:column;gap:10px}
.actl-row{display:flex;align-items:flex-start;gap:12px;padding:12px;background:#F8F9FE;border-radius:8px;border:1px solid var(--border)}
.actl-dot{width:8px;height:8px;border-radius:50%;margin-top:4px;flex-shrink:0}
.actl-svc{font-size:13px;font-weight:600;margin-bottom:3px}
.actl-meta{font-size:11px;color:var(--m)}
/* ── EXISTING ACCOUNTS ── */
.ea-table{width:100%;border-collapse:collapse}
.ea-table th{font-size:10px;font-weight:700;letter-spacing:.07em;text-transform:uppercase;color:var(--m2);padding:8px 12px;text-align:left;border-bottom:1px solid var(--border)}
.ea-table td{padding:12px 12px;font-size:13px;border-bottom:1px solid var(--border2);vertical-align:top}
.ea-table tr:hover td{background:rgba(21,62,237,.03)}
.gap-badge{padding:3px 8px;border-radius:4px;font-size:10px;font-weight:700;background:rgba(245,159,11,.12);color:var(--amber)}
/* ── LOADING ── */
.loading-pulse{display:flex;align-items:center;gap:8px;color:var(--m);font-size:13px;padding:40px 0}
.pulse-dot{width:6px;height:6px;border-radius:50%;background:var(--blue);animation:pulse 1s ease-in-out infinite}
.pulse-dot:nth-child(2){animation-delay:.15s}
.pulse-dot:nth-child(3){animation-delay:.3s}
@keyframes pulse{0%,100%{opacity:.3;transform:scale(.8)}50%{opacity:1;transform:scale(1.2)}}
/* ── TOAST ── */
#cc-toast{position:fixed;bottom:24px;left:50%;transform:translateX(-50%) translateY(20px);background:#FFFFFF;backdrop-filter:blur(12px);border:1px solid var(--border);border-radius:8px;padding:10px 18px;font-size:13px;color:var(--w);opacity:0;transition:all .3s;pointer-events:none;z-index:200;box-shadow:0 4px 20px rgba(0,0,0,.12)}
#cc-toast.show{opacity:1;transform:translateX(-50%) translateY(0)}
/* ── SCROLLBAR ── */
/* ── DEAL STAGE ── */
.stage-pill{display:inline-flex;align-items:center;gap:5px;padding:3px 9px;border-radius:20px;font-size:10px;font-weight:700;letter-spacing:.04em;text-transform:uppercase;cursor:pointer;transition:opacity .15s}
.stage-pill:hover{opacity:.8}
/* ── ACTION CARDS ── */
.action-list{display:flex;flex-direction:column;gap:8px}
.action-card{background:var(--card);border:1px solid var(--border);border-radius:10px;padding:14px 16px;display:flex;align-items:flex-start;gap:12px;transition:all .2s}
.action-card:hover{border-color:rgba(21,62,237,.25);box-shadow:0 2px 12px rgba(21,62,237,.06)}
.action-card.overdue{border-left:3px solid var(--red)}
.action-card.high{border-left:3px solid var(--amber)}
.action-check{width:18px;height:18px;border-radius:4px;border:1.5px solid var(--border);background:none;cursor:pointer;flex-shrink:0;display:flex;align-items:center;justify-content:center;transition:all .2s;margin-top:1px}
.action-check:hover{border-color:var(--green);background:rgba(0,212,160,.1)}
.action-check.done{background:var(--green);border-color:var(--green)}
.action-title{font-size:13px;font-weight:600;margin-bottom:3px}
.action-title.done{text-decoration:line-through;color:var(--m2)}
.action-meta{font-size:11px;color:var(--m);display:flex;gap:10px;flex-wrap:wrap}
.action-meta .overdue-label{color:var(--red);font-weight:700}
.priority-dot{width:6px;height:6px;border-radius:50%;flex-shrink:0;margin-top:6px}
.priority-dot.critical{background:var(--red)}
.priority-dot.high{background:var(--amber)}
.priority-dot.medium{background:var(--blue)}
.priority-dot.low{background:var(--m2)}
.action-type-badge{padding:2px 7px;border-radius:4px;font-size:10px;font-weight:600;background:#F0F2F9;color:var(--m)}
/* ── MEETING CARDS ── */
.meeting-card{background:var(--card);border:1px solid var(--border);border-radius:10px;padding:16px;margin-bottom:10px}
.meeting-head{display:flex;justify-content:space-between;align-items:flex-start;margin-bottom:10px}
.meeting-account{font-size:14px;font-weight:700}
.meeting-date{font-size:11px;color:var(--m)}
.meeting-outcome{padding:3px 8px;border-radius:4px;font-size:10px;font-weight:700;text-transform:uppercase}
.meeting-outcome.positive{background:rgba(0,212,160,.15);color:var(--green)}
.meeting-outcome.neutral{background:rgba(21,62,237,.09);color:var(--blue)}
.meeting-outcome.negative{background:rgba(246,87,74,.12);color:var(--red)}
.meeting-outcome.no-show{background:rgba(128,128,128,.12);color:var(--m)}
.meeting-summary{font-size:12.5px;color:var(--m);line-height:1.6;margin-bottom:8px}
.meeting-next{background:rgba(21,62,237,.06);border:1px solid rgba(21,62,237,.18);border-radius:6px;padding:8px 12px;font-size:12px;color:var(--blue)}
/* ── TODAY VIEW ── */
.today-hero{background:linear-gradient(135deg,rgba(21,62,237,.12) 0%,rgba(21,62,237,.04) 100%);border:1px solid rgba(21,62,237,.25);border-radius:16px;padding:24px 28px;margin-bottom:24px}
.today-date{font-size:11px;font-weight:700;letter-spacing:.1em;text-transform:uppercase;color:var(--blue);margin-bottom:6px}
.today-headline{font-size:22px;font-weight:800;letter-spacing:-.03em;margin-bottom:4px}
.today-sub{font-size:13px;color:var(--m)}
.today-grid{display:grid;grid-template-columns:1.4fr 1fr;gap:20px}
.today-section{background:var(--card);border:1px solid var(--border);border-radius:12px;padding:20px}
.today-sec-title{font-size:11px;font-weight:700;letter-spacing:.08em;text-transform:uppercase;color:var(--m2);margin-bottom:14px;display:flex;justify-content:space-between;align-items:center}
.cold-account-row{display:flex;justify-content:space-between;align-items:center;padding:9px 0;border-bottom:1px solid var(--border2);font-size:13px}
.cold-account-row:last-child{border-bottom:none}
.cold-days{font-size:11px;font-weight:600;color:var(--amber)}
/* ── WEEKLY COMMIT WIDGET ── */
.commit-widget{background:linear-gradient(135deg,rgba(0,212,160,.08) 0%,rgba(0,212,160,.03) 100%);border:1px solid rgba(0,212,160,.2);border-radius:12px;padding:20px;margin-bottom:24px}
.commit-title{font-size:13px;font-weight:700;margin-bottom:4px;display:flex;align-items:center;gap:8px}
.commit-text{font-size:13px;color:var(--m);line-height:1.6;margin:10px 0}
.commit-cta{padding:8px 16px;background:rgba(0,212,160,.15);border:1px solid rgba(0,212,160,.3);border-radius:8px;color:var(--green);font:600 12px var(--font);cursor:pointer;transition:all .2s}
.commit-cta:hover{background:rgba(0,212,160,.25)}
/* ── DETAIL PANEL TABS ── */
.dp-tabs{display:flex;gap:2px;padding:0 24px;border-bottom:1px solid var(--border);background:#FFFFFF}
.dp-tab{padding:10px 14px;font-size:12px;font-weight:600;color:var(--m);cursor:pointer;border-bottom:2px solid transparent;transition:all .2s;margin-bottom:-1px}
.dp-tab:hover{color:var(--t)}
.dp-tab.active{color:var(--blue);border-bottom-color:var(--blue)}
.dp-tab-content{display:none}.dp-tab-content.active{display:block}
/* ── OUTREACH PANEL ── */
.outreach-controls{display:flex;gap:8px;margin-bottom:14px;flex-wrap:wrap}
.outreach-select{background:#F8F9FE;border:1px solid var(--border);border-radius:7px;padding:7px 10px;font:500 12px var(--font);color:var(--t);cursor:pointer}
.outreach-select:focus{outline:none;border-color:rgba(21,62,237,.5)}
.outreach-msg{background:#F8F9FE;border:1px solid var(--border);border-radius:8px;padding:14px;font-size:13px;color:var(--t);line-height:1.7;white-space:pre-wrap;min-height:120px;margin-bottom:10px}
.outreach-copy-btn{padding:8px 16px;background:var(--blue-dim);border:1px solid rgba(21,62,237,.3);border-radius:8px;color:var(--blue);font:600 12px var(--font);cursor:pointer;transition:all .2s}
.outreach-copy-btn:hover{background:rgba(21,62,237,.15)}
/* ── MODAL ── */
.modal-overlay{position:fixed;inset:0;background:rgba(0,0,0,.3);z-index:200;display:flex;align-items:center;justify-content:center;opacity:0;pointer-events:none;transition:opacity .25s}
.modal-overlay.open{opacity:1;pointer-events:all}
.modal{background:#FFFFFF;border:1px solid var(--border);border-radius:16px;padding:28px;width:480px;max-width:94vw;max-height:85vh;overflow-y:auto;box-shadow:0 20px 60px rgba(0,0,0,.12)}
.modal-title{font-size:16px;font-weight:700;margin-bottom:20px;color:var(--w)}
.modal-field{margin-bottom:14px}
.modal-field label{display:block;font-size:11px;font-weight:600;letter-spacing:.06em;text-transform:uppercase;color:var(--m);margin-bottom:6px}
.modal-field input,.modal-field select,.modal-field textarea{width:100%;background:#F8F9FE;border:1px solid var(--border);border-radius:8px;padding:10px 12px;font:400 13px var(--font);color:var(--t);outline:none;transition:border-color .2s}
.modal-field textarea{resize:vertical;min-height:80px}
.modal-field input:focus,.modal-field select:focus,.modal-field textarea:focus{border-color:rgba(21,62,237,.5)}
.modal-actions{display:flex;gap:10px;justify-content:flex-end;margin-top:20px}
.modal-btn{padding:9px 20px;border-radius:8px;font:600 13px var(--font);cursor:pointer;transition:all .2s}
.modal-btn.primary{background:var(--blue);border:none;color:#fff}
.modal-btn.primary:hover{opacity:.85}
.modal-btn.secondary{background:none;border:1px solid var(--border);color:var(--m)}
.modal-btn.secondary:hover{color:var(--t);border-color:rgba(0,0,0,.2)}
/* ── ADD BUTTON ── */
.add-btn{padding:8px 16px;background:var(--blue);border:none;border-radius:8px;color:#fff;font:600 12px var(--font);cursor:pointer;transition:opacity .2s;white-space:nowrap}
.add-btn:hover{opacity:.85}
/* ── INTELLIGENCE VIEW ── */
.intel-grid{display:grid;grid-template-columns:1fr 1fr;gap:16px;margin-bottom:24px}
.intel-card{background:var(--card);border:1px solid var(--border);border-radius:12px;padding:20px}
.intel-card.full{grid-column:1/-1}
.intel-card.risk{border-left:3px solid var(--red)}
.intel-card.cold{border-left:3px solid var(--amber)}
.intel-card.won{border-left:3px solid var(--green)}
.intel-sec-title{font-size:11px;font-weight:700;letter-spacing:.08em;text-transform:uppercase;margin-bottom:14px;display:flex;justify-content:space-between;align-items:center}
.intel-account-row{padding:12px 0;border-bottom:1px solid var(--border2);display:flex;align-items:flex-start;justify-content:space-between;gap:12px}
.intel-account-row:last-child{border-bottom:none}
.intel-account-name{font-size:13px;font-weight:700;margin-bottom:2px}
.intel-account-meta{font-size:11px;color:var(--m)}
.intel-risk-badges{display:flex;flex-wrap:wrap;gap:4px;margin-top:5px}
.intel-risk-badge{padding:2px 7px;border-radius:4px;font-size:10px;font-weight:600;background:rgba(246,87,74,.1);color:var(--red)}
.intel-cold-days{font-size:12px;font-weight:700;color:var(--amber);white-space:nowrap}
.diagnose-btn{padding:6px 12px;background:rgba(21,62,237,.08);border:1px solid rgba(21,62,237,.2);border-radius:7px;color:var(--blue);font:600 11px var(--font);cursor:pointer;transition:all .2s;white-space:nowrap;flex-shrink:0}
.diagnose-btn:hover{background:rgba(21,62,237,.15)}
.diagnose-btn.loading{opacity:.5;pointer-events:none}
.diagnosis-box{background:#F8F9FE;border:1px solid var(--border);border-radius:8px;padding:14px;margin-top:10px;font-size:12.5px;color:var(--t);line-height:1.7;white-space:pre-wrap}
.diagnosis-box strong{color:var(--w)}
.win-analysis-box{background:rgba(0,212,160,.04);border:1px solid rgba(0,212,160,.15);border-radius:10px;padding:18px;font-size:13px;color:var(--t);line-height:1.8;white-space:pre-wrap}
.win-analysis-box strong{color:var(--green)}
.intel-hero{background:linear-gradient(135deg,rgba(139,92,246,.1) 0%,rgba(21,62,237,.06) 100%);border:1px solid rgba(139,92,246,.2);border-radius:16px;padding:20px 24px;margin-bottom:20px;display:flex;align-items:center;gap:20px}
.intel-stat{text-align:center;padding:0 20px;border-right:1px solid var(--border)}
.intel-stat:last-child{border-right:none}
.intel-stat-val{font-size:28px;font-weight:800;letter-spacing:-.04em}
.intel-stat-label{font-size:10px;font-weight:700;letter-spacing:.07em;text-transform:uppercase;color:var(--m);margin-top:3px}
::-webkit-scrollbar{width:4px}::-webkit-scrollbar-track{background:transparent}::-webkit-scrollbar-thumb{background:rgba(0,0,0,.12);border-radius:2px}
/* ── FORESIGHT ── */
.fsight-shell{display:flex;flex-direction:column;gap:0;min-height:calc(100vh - 54px)}
.fsight-hero{background:linear-gradient(135deg,#0A0A1E 0%,#0D1A3A 100%);padding:32px 36px 28px;flex-shrink:0;position:relative;overflow:hidden}
.fsight-hero::before{content:'';position:absolute;inset:0;background:radial-gradient(ellipse 70% 80% at 80% 50%,rgba(21,62,237,.25) 0%,transparent 70%),radial-gradient(ellipse 40% 60% at 20% 80%,rgba(124,58,237,.15) 0%,transparent 60%);pointer-events:none}
.fsight-hero-inner{position:relative;z-index:1;display:flex;align-items:flex-start;justify-content:space-between;gap:20px;flex-wrap:wrap}
.fsight-eyebrow{font-size:10px;font-weight:700;letter-spacing:.18em;text-transform:uppercase;color:rgba(21,62,237,.9);margin-bottom:8px;display:flex;align-items:center;gap:8px}
.fsight-eyebrow-dot{width:6px;height:6px;border-radius:50%;background:#153EED;animation:livePulse 2s infinite}
@keyframes livePulse{0%,100%{opacity:1}50%{opacity:.3}}
.fsight-title{font-size:30px;font-weight:900;letter-spacing:-.04em;color:#fff;margin-bottom:6px;line-height:1}
.fsight-title span{color:#153EED}
.fsight-sub{font-size:13px;color:rgba(255,255,255,.55);line-height:1.5;max-width:480px}
.fsight-run-btn{padding:12px 24px;background:linear-gradient(135deg,#153EED,#2952D8);border:none;border-radius:12px;color:#fff;font:800 13px var(--font);cursor:pointer;transition:all .25s;white-space:nowrap;letter-spacing:.02em;box-shadow:0 6px 24px rgba(21,62,237,.4);flex-shrink:0}
.fsight-run-btn:hover{transform:translateY(-2px);box-shadow:0 10px 32px rgba(21,62,237,.5)}
.fsight-run-btn:disabled{opacity:.5;cursor:not-allowed;transform:none;box-shadow:none}
.fsight-stats{display:flex;gap:28px;margin-top:24px;position:relative;z-index:1}
.fsight-stat{text-align:left}
.fsight-stat-val{font-size:22px;font-weight:800;letter-spacing:-.03em;color:#fff}
.fsight-stat-label{font-size:10px;font-weight:600;letter-spacing:.06em;text-transform:uppercase;color:rgba(255,255,255,.4);margin-top:2px}
.fsight-last{font-size:11px;color:rgba(255,255,255,.3);margin-top:4px;position:relative;z-index:1}

.fsight-body{flex:1;padding:28px 36px;display:flex;flex-direction:column;gap:28px;background:var(--bg)}
.fsight-empty{display:flex;flex-direction:column;align-items:center;justify-content:center;min-height:320px;color:var(--m);text-align:center;gap:12px}
.fsight-empty-icon{font-size:48px;opacity:.4}
.fsight-empty-text{font-size:14px;font-weight:500}
.fsight-empty-sub{font-size:12px;color:var(--m2)}

.fsight-section{background:#fff;border:1px solid var(--border);border-radius:16px;overflow:hidden}
.fsight-sec-header{padding:18px 24px;border-bottom:1px solid var(--border);display:flex;align-items:center;justify-content:space-between}
.fsight-sec-icon{width:36px;height:36px;border-radius:10px;display:flex;align-items:center;justify-content:center;font-size:16px;flex-shrink:0}
.fsight-sec-icon.window{background:rgba(220,38,38,.08);border:1px solid rgba(220,38,38,.15)}
.fsight-sec-icon.partner{background:rgba(124,58,237,.08);border:1px solid rgba(124,58,237,.15)}
.fsight-sec-icon.service{background:rgba(0,166,126,.08);border:1px solid rgba(0,166,126,.15)}
.fsight-sec-title{font-size:15px;font-weight:800;color:var(--w);letter-spacing:-.02em}
.fsight-sec-sub{font-size:11px;color:var(--m);margin-top:2px}
.fsight-sec-count{font-size:12px;font-weight:700;padding:3px 10px;border-radius:20px}
.fsight-sec-count.window{background:rgba(220,38,38,.08);color:var(--red)}
.fsight-sec-count.partner{background:rgba(124,58,237,.08);color:var(--purple)}
.fsight-sec-count.service{background:rgba(0,166,126,.08);color:var(--green)}

.fsight-blocks{display:grid;gap:0}
.fsight-block{padding:20px 24px;border-bottom:1px solid var(--border);display:grid;gap:0}
.fsight-block:last-child{border-bottom:none}
.fsight-block:hover{background:#FAFBFF}

/* Buying window block */
.fw-head{display:flex;align-items:flex-start;justify-content:space-between;margin-bottom:12px;gap:12px}
.fw-name{font-size:16px;font-weight:800;letter-spacing:-.02em;color:var(--w)}
.fw-window{padding:4px 12px;border-radius:20px;font-size:10px;font-weight:800;letter-spacing:.05em;text-transform:uppercase;white-space:nowrap;flex-shrink:0}
.fw-window.hot{background:rgba(220,38,38,.1);color:var(--red);border:1px solid rgba(220,38,38,.2)}
.fw-window.warm{background:rgba(217,119,6,.1);color:var(--amber);border:1px solid rgba(217,119,6,.2)}
.fw-window.cool{background:rgba(21,62,237,.08);color:var(--blue);border:1px solid rgba(21,62,237,.18)}
.fw-grid{display:grid;grid-template-columns:1fr 1fr;gap:10px}
.fw-cell{background:var(--bg);border-radius:8px;padding:10px 12px}
.fw-cell-label{font-size:9px;font-weight:700;letter-spacing:.08em;text-transform:uppercase;color:var(--m2);margin-bottom:4px}
.fw-cell-val{font-size:12.5px;color:var(--t);line-height:1.45}
.fw-cell.full{grid-column:1/-1}
.fw-foot{display:flex;align-items:center;justify-content:space-between;margin-top:12px}
.fw-value{font-size:18px;font-weight:800;letter-spacing:-.02em;color:var(--blue)}
.fw-conf{font-size:10px;font-weight:700;letter-spacing:.06em;text-transform:uppercase;padding:3px 10px;border-radius:20px}
.fw-conf.High{background:rgba(0,166,126,.1);color:var(--green);border:1px solid rgba(0,166,126,.2)}
.fw-conf.Medium{background:rgba(217,119,6,.1);color:var(--amber);border:1px solid rgba(217,119,6,.2)}
.fw-conf.Low{background:rgba(220,38,38,.08);color:var(--red);border:1px solid rgba(220,38,38,.15)}

/* Partnership block */
.fp-head{display:flex;align-items:center;justify-content:space-between;margin-bottom:10px}
.fp-name{font-size:16px;font-weight:800;letter-spacing:-.02em;color:var(--w)}
.fp-value{font-size:15px;font-weight:800;color:var(--purple)}
.fp-body{display:grid;grid-template-columns:1fr 1fr;gap:8px;margin-bottom:10px}
.fp-cell{background:rgba(124,58,237,.04);border:1px solid rgba(124,58,237,.1);border-radius:8px;padding:10px 12px}
.fp-cell-label{font-size:9px;font-weight:700;letter-spacing:.08em;text-transform:uppercase;color:rgba(124,58,237,.6);margin-bottom:4px}
.fp-cell-val{font-size:12px;color:var(--t);line-height:1.45}
.fp-cell.full{grid-column:1/-1}
.fp-action{background:rgba(124,58,237,.07);border:1px solid rgba(124,58,237,.15);border-radius:8px;padding:10px 14px;font-size:12px;color:var(--purple);font-weight:600;line-height:1.5}
.fp-action-label{font-size:9px;font-weight:700;letter-spacing:.08em;text-transform:uppercase;color:var(--purple);opacity:.7;margin-bottom:4px}

/* New service block */
.fs-head{display:flex;align-items:flex-start;justify-content:space-between;margin-bottom:12px;gap:12px}
.fs-name{font-size:16px;font-weight:800;letter-spacing:-.02em;color:var(--w)}
.fs-revenue{font-size:15px;font-weight:800;color:var(--green)}
.fs-body{display:grid;grid-template-columns:1fr 1fr;gap:8px;margin-bottom:10px}
.fs-cell{background:rgba(0,166,126,.04);border:1px solid rgba(0,166,126,.1);border-radius:8px;padding:10px 12px}
.fs-cell-label{font-size:9px;font-weight:700;letter-spacing:.08em;text-transform:uppercase;color:rgba(0,166,126,.6);margin-bottom:4px}
.fs-cell-val{font-size:12px;color:var(--t);line-height:1.45}
.fs-cell.full{grid-column:1/-1}
.fs-bop{display:inline-flex;align-items:center;gap:6px;padding:5px 12px;background:rgba(0,166,126,.08);border:1px solid rgba(0,166,126,.2);border-radius:20px;font-size:11px;font-weight:700;color:var(--green)}

/* ── MARKET COMMAND ── */
.cmd-hero{background:linear-gradient(135deg,rgba(21,62,237,.07) 0%,rgba(21,62,237,.02) 100%);border:1px solid rgba(21,62,237,.15);border-radius:18px;padding:28px 32px;margin-bottom:28px}
.cmd-label{font-size:10px;font-weight:700;letter-spacing:.12em;text-transform:uppercase;color:var(--blue);margin-bottom:6px}
.cmd-title{font-size:26px;font-weight:800;letter-spacing:-.03em;color:var(--w);margin-bottom:4px}
.cmd-sub{font-size:13px;color:var(--m);line-height:1.5}
.cmd-brief-btn{padding:10px 20px;background:var(--blue);border:none;border-radius:10px;color:#fff;font:700 13px var(--font);cursor:pointer;transition:all .2s;white-space:nowrap;align-self:flex-start}
.cmd-brief-btn:hover{opacity:.88;transform:translateY(-1px);box-shadow:0 6px 20px rgba(21,62,237,.25)}
.cmd-brief-btn:disabled{opacity:.5;cursor:not-allowed;transform:none}
.cmd-pulse-row{display:flex;align-items:center;gap:0;margin-top:24px;background:#fff;border:1px solid var(--border);border-radius:12px;overflow:hidden}
.cmd-pulse-stat{flex:1;padding:16px 20px;text-align:center}
.cmd-pulse-val{font-size:24px;font-weight:800;letter-spacing:-.03em;color:var(--w)}
.cmd-pulse-label{font-size:10px;font-weight:600;letter-spacing:.05em;text-transform:uppercase;color:var(--m);margin-top:4px}
.cmd-pulse-sep{width:1px;height:48px;background:var(--border);flex-shrink:0}
.cmd-grid{display:grid;grid-template-columns:1fr 1fr;gap:24px}
.cmd-section-header{display:flex;align-items:center;justify-content:space-between;margin-bottom:14px}
.cmd-section-title{font-size:13px;font-weight:700;letter-spacing:-.01em;color:var(--w)}
.cmd-section-link{font-size:12px;font-weight:600;color:var(--blue);cursor:pointer;transition:opacity .15s}
.cmd-section-link:hover{opacity:.7}
.cmd-signal-card{background:#fff;border:1px solid var(--border);border-radius:10px;padding:14px 16px;margin-bottom:10px;border-left:3px solid var(--border);transition:box-shadow .2s}
.cmd-signal-card:hover{box-shadow:0 4px 16px rgba(21,62,237,.07)}
.cmd-signal-card.critical{border-left-color:var(--red)}
.cmd-signal-card.warning{border-left-color:var(--amber)}
.cmd-signal-card.info{border-left-color:var(--blue)}
.cmd-signal-sev{font-size:9px;font-weight:700;letter-spacing:.08em;text-transform:uppercase;margin-bottom:4px}
.cmd-signal-sev.critical{color:var(--red)}
.cmd-signal-sev.warning{color:var(--amber)}
.cmd-signal-sev.info{color:var(--blue)}
.cmd-signal-title{font-size:13px;font-weight:700;color:var(--w);margin-bottom:4px;line-height:1.3}
.cmd-signal-desc{font-size:11.5px;color:var(--m);line-height:1.5}
.cmd-target-card{background:#fff;border:1px solid var(--border);border-radius:10px;padding:14px 16px;margin-bottom:10px;display:flex;align-items:center;justify-content:space-between;gap:12px;transition:all .2s;cursor:pointer}
.cmd-target-card:hover{border-color:rgba(21,62,237,.3);box-shadow:0 4px 16px rgba(21,62,237,.07)}
.cmd-target-name{font-size:13px;font-weight:700;color:var(--w);margin-bottom:3px}
.cmd-target-meta{font-size:11px;color:var(--m)}
.cmd-target-score{text-align:right;flex-shrink:0}
.cmd-target-val{font-size:15px;font-weight:800;color:var(--blue)}
.cmd-target-icp{font-size:10px;color:var(--m);margin-top:2px}
.cmd-brief-card{background:#fff;border:1px solid rgba(21,62,237,.2);border-radius:14px;padding:24px;box-shadow:0 4px 20px rgba(21,62,237,.06)}
</style>
<script src="https://unpkg.com/lucide@latest/dist/umd/lucide.min.js"></script>
</head>
<body>

<!-- ══ LOGIN SCREEN ══════════════════════════════════════════════════════════ -->
<div class="screen" id="login-screen">
  <div class="login-card">
    <div class="login-logo">JAKALA</div>
    <div class="login-title">Control Center</div>
    <div class="login-sub">Strategic market command · Nordic</div>
    <div class="login-field">
      <label>Email</label>
      <input type="email" id="login-email" placeholder="you@jakala.com" autocomplete="username">
    </div>
    <div class="login-field">
      <label>Password</label>
      <input type="password" id="login-pw" placeholder="••••••••" autocomplete="current-password">
    </div>
    <button class="login-btn" onclick="doLogin()">Sign in →</button>
    <div class="login-err" id="login-err"></div>
  </div>
</div>

<!-- ══ APP SHELL ══════════════════════════════════════════════════════════════ -->
<div id="app-shell" class="hidden">

  <!-- Top bar -->
  <div class="topbar">
    <span class="topbar-logo">JAKALA</span>
    <span class="topbar-sep"></span>
    <span class="topbar-title">Market Command</span>
    <span class="topbar-sep"></span>
    <span class="topbar-country" id="tb-country"></span>
    <div class="topbar-right">
      <a class="topbar-link" href="/app" target="_blank">GTM Assistant ↗</a>
      <div class="topbar-user">
        <div class="topbar-avatar" id="tb-avatar"></div>
        <span id="tb-name"></span>
      </div>
      <span class="topbar-link" onclick="doLogout()">Sign out</span>
    </div>
  </div>

  <!-- Main layout -->
  <div class="main-layout">

    <!-- Sidebar -->
    <aside class="sidebar" id="sidebar">
      <div class="sb-section">
        <div class="sb-label">Foresight</div>
        <div class="nav-item active" data-view="foresight" onclick="switchView('foresight')" style="background:linear-gradient(135deg,rgba(21,62,237,.1),rgba(124,58,237,.08));border:1px solid rgba(21,62,237,.18);color:var(--blue);font-weight:700;margin-bottom:4px">
          <span class="nav-icon">✦</span> Market Foresight
        </div>
        <div class="nav-item" data-view="intelligence" onclick="switchView('intelligence')">
          <i data-lucide="eye" style="width:16px;height:16px;display:inline-block;vertical-align:middle;" class="nav-icon"></i> Strategic AI
        </div>
        <div class="nav-item" data-view="trends" onclick="switchView('trends')">
          <i data-lucide="radio" style="width:16px;height:16px;display:inline-block;vertical-align:middle;" class="nav-icon"></i> Signal Radar
        </div>
      </div>
      <div class="sb-divider"></div>
      <div class="sb-section">
        <div class="sb-label">Strategic Pipeline</div>
        <div class="nav-item" data-view="command" onclick="switchView('command')">
          <span class="nav-icon">◉</span> Market Command
        </div>
        <div class="nav-item" data-view="overview" onclick="switchView('overview')">
          <span class="nav-icon">◎</span> Pipeline Overview
        </div>
        <div class="nav-item" data-view="predictions" onclick="switchView('predictions')">
          <i data-lucide="zap" style="width:16px;height:16px;display:inline-block;vertical-align:middle;" class="nav-icon"></i> Deal Predictions
        </div>
        <div class="nav-item" data-view="new-biz" onclick="switchView('new-biz')">
          <span class="nav-icon">◈</span> Opportunity Map
        </div>
        <div class="nav-item" data-view="existing" onclick="switchView('existing')">
          <span class="nav-icon">⊙</span> Existing Accounts
        </div>
      </div>
      <div class="sb-divider"></div>
      <div class="sb-section">
        <div class="sb-label">Team Activity</div>
        <div class="nav-item" data-view="today" onclick="switchView('today')">
          <span class="nav-icon">☑</span> Open Actions <span id="sb-action-count" style="margin-left:auto;background:var(--red);border-radius:10px;padding:1px 6px;font-size:10px;font-weight:700;display:none"></span>
        </div>
        <div class="nav-item" data-view="meetings" onclick="switchView('meetings')">
          <span class="nav-icon">◷</span> Recent Meetings
        </div>
      </div>
      <div class="sb-divider"></div>
      <div class="sb-signals" id="sb-signals-summary"></div>
      <!-- Global: country switcher -->
      <div id="sb-country-switcher" style="display:none">
        <div class="sb-divider"></div>
        <div style="padding:0 14px 6px">
          <div class="sb-label">Markets</div>
          <div id="country-switcher-list"></div>
        </div>
      </div>
    </aside>

    <!-- Content -->
    <main class="content" id="main-content">

      <!-- ── FORESIGHT ── -->
      <div class="view" id="view-foresight" style="padding:0">
        <div class="fsight-shell">

          <!-- Hero bar -->
          <div class="fsight-hero">
            <div class="fsight-hero-inner">
              <div>
                <div class="fsight-eyebrow"><span class="fsight-eyebrow-dot"></span>Predictive Market Intelligence</div>
                <div class="fsight-title">JAKALA <span>Foresight</span></div>
                <div class="fsight-sub">Buying windows · Partnership opportunities · New service gaps. Generated from live market signals and account intelligence.</div>
                <div class="fsight-stats">
                  <div class="fsight-stat"><div class="fsight-stat-val" id="fs-windows-count">—</div><div class="fsight-stat-label">Buying windows</div></div>
                  <div class="fsight-stat"><div class="fsight-stat-val" id="fs-partner-count">—</div><div class="fsight-stat-label">Partnerships</div></div>
                  <div class="fsight-stat"><div class="fsight-stat-val" id="fs-service-count">—</div><div class="fsight-stat-label">Service gaps</div></div>
                </div>
                <div class="fsight-last" id="fs-last-run">Not yet generated this session</div>
              </div>
              <button class="fsight-run-btn" id="fsight-run-btn" onclick="runForesight()">✦ Run Foresight</button>
            </div>
          </div>

          <!-- ── WEEKLY STRATEGIC BRIEF ── -->
          <div id="weekly-brief-section" style="margin:0 32px 28px">
            <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:14px">
              <div>
                <div style="font-size:10px;font-weight:800;color:#153EED;text-transform:uppercase;letter-spacing:2px;margin-bottom:3px">◉ This week's strategic picture</div>
                <div style="font-size:12px;color:var(--m)" id="brief-meta">Generating…</div>
              </div>
              <button onclick="loadWeeklyBrief(true)" style="background:none;border:1px solid var(--border);color:var(--m);border-radius:6px;padding:5px 12px;font-size:11px;cursor:pointer;font-family:var(--font);transition:all .15s" onmouseenter="this.style.borderColor='#153EED';this.style.color='#153EED'" onmouseleave="this.style.borderColor='var(--border)';this.style.color='var(--m)'">↻ Refresh</button>
            </div>

            <!-- Loading -->
            <div id="brief-loading" style="background:var(--card);border:1px solid var(--border);border-radius:var(--radius);padding:32px;text-align:center;color:var(--m);font-size:13px">
              <div class="loading-pulse" style="justify-content:center"><div class="pulse-dot"></div><div class="pulse-dot"></div><div class="pulse-dot"></div><span style="margin-left:10px">Analysing pipeline and market signals…</span></div>
            </div>

            <!-- Content (hidden until loaded) -->
            <div id="brief-content" style="display:none">
              <!-- Status bar -->
              <div id="brief-status-bar" style="background:var(--card);border:1px solid var(--border);border-radius:var(--radius);padding:18px 22px;margin-bottom:12px;display:flex;align-items:center;gap:24px">
                <div id="brief-status-badge" style="padding:5px 12px;border-radius:6px;font-size:11px;font-weight:800;letter-spacing:.06em;text-transform:uppercase">—</div>
                <div style="flex:1">
                  <div id="brief-status-reason" style="font-size:13px;color:var(--t);font-weight:500">—</div>
                </div>
                <div id="brief-numbers" style="display:flex;gap:20px"></div>
              </div>

              <!-- Three columns: Risk · Opportunity · Decision -->
              <div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:12px">

                <!-- RISK -->
                <div style="background:rgba(229,62,62,.05);border:1px solid rgba(229,62,62,.2);border-radius:var(--radius);padding:18px 20px">
                  <div style="font-size:9.5px;font-weight:800;color:#E53E3E;text-transform:uppercase;letter-spacing:.1em;margin-bottom:10px">⚠ Pipeline Risiko</div>
                  <div id="brief-risk-headline" style="font-size:14px;font-weight:700;color:var(--t);margin-bottom:8px">—</div>
                  <div id="brief-risk-accounts" style="font-size:12px;color:var(--m);line-height:1.7;margin-bottom:10px"></div>
                  <div id="brief-risk-implication" style="font-size:12px;color:rgba(229,62,62,.8);font-style:italic;padding-top:8px;border-top:1px solid rgba(229,62,62,.15)">—</div>
                </div>

                <!-- OPPORTUNITY -->
                <div style="background:rgba(0,166,126,.05);border:1px solid rgba(0,166,126,.2);border-radius:var(--radius);padding:18px 20px">
                  <div style="font-size:9.5px;font-weight:800;color:#00A67E;text-transform:uppercase;letter-spacing:.1em;margin-bottom:10px">✦ Top Mulighed</div>
                  <div id="brief-opp-account" style="font-size:14px;font-weight:700;color:var(--t);margin-bottom:4px">—</div>
                  <div id="brief-opp-buyer" style="font-size:11.5px;color:var(--m);margin-bottom:10px">—</div>
                  <div id="brief-opp-window" style="font-size:12px;color:#00A67E;font-weight:600;margin-bottom:6px">—</div>
                  <div id="brief-opp-why" style="font-size:12px;color:var(--m);line-height:1.6;margin-bottom:10px">—</div>
                  <div id="brief-opp-action" style="font-size:12px;color:rgba(0,166,126,.9);font-weight:600;padding-top:8px;border-top:1px solid rgba(0,166,126,.15)">→ <span id="brief-opp-action-text">—</span></div>
                </div>

                <!-- THE DECISION -->
                <div style="background:rgba(21,62,237,.06);border:1px solid rgba(21,62,237,.25);border-radius:var(--radius);padding:18px 20px">
                  <div style="font-size:9.5px;font-weight:800;color:#153EED;text-transform:uppercase;letter-spacing:.1em;margin-bottom:10px">◉ Den ene beslutning</div>
                  <div id="brief-dec-question" style="font-size:14px;font-weight:700;color:var(--t);line-height:1.4;margin-bottom:10px">—</div>
                  <div id="brief-dec-options" style="font-size:12px;color:var(--m);line-height:1.8;margin-bottom:10px"></div>
                  <div style="padding-top:8px;border-top:1px solid rgba(21,62,237,.15)">
                    <div style="font-size:9px;font-weight:700;color:#153EED;text-transform:uppercase;letter-spacing:.08em;margin-bottom:4px">Anbefaling</div>
                    <div id="brief-dec-rec" style="font-size:12px;color:var(--t);line-height:1.6">—</div>
                    <div id="brief-dec-deadline" style="font-size:11px;color:var(--m);margin-top:6px;font-weight:600">—</div>
                  </div>
                </div>

              </div>
            </div>
          </div>
          <div style="margin:0 32px 16px;height:1px;background:var(--border)"></div>

          <!-- Body -->
          <div class="fsight-body" id="fsight-body">
            <div class="fsight-empty" id="fsight-empty">
              <div class="fsight-empty-icon">✦</div>
              <div class="fsight-empty-text">No foresight generated yet</div>
              <div class="fsight-empty-sub">Click "Run Foresight" to generate predictive market intelligence based on live signals and account data.</div>
            </div>
          </div>

        </div>
      </div>

      <!-- ── MARKET COMMAND ── -->
      <div class="view" id="view-command">
        <div class="cmd-hero">
          <div style="display:flex;align-items:flex-start;justify-content:space-between;flex-wrap:wrap;gap:16px">
            <div>
              <div class="cmd-label" id="cmd-week-label">MARKET INTELLIGENCE</div>
              <div class="cmd-title">Market Command</div>
              <div class="cmd-sub" id="cmd-sub">Loading market pulse…</div>
            </div>
            <button class="cmd-brief-btn" id="cmd-brief-btn" onclick="generateStrategicBrief()">◈ Generate strategic brief</button>
          </div>
          <div class="cmd-pulse-row" id="cmd-pulse-row">
            <div class="cmd-pulse-stat">
              <div class="cmd-pulse-val" id="cmd-signals-count">—</div>
              <div class="cmd-pulse-label">Critical signals</div>
            </div>
            <div class="cmd-pulse-sep"></div>
            <div class="cmd-pulse-stat">
              <div class="cmd-pulse-val" id="cmd-targets-count">—</div>
              <div class="cmd-pulse-label">High-ICP targets</div>
            </div>
            <div class="cmd-pulse-sep"></div>
            <div class="cmd-pulse-stat">
              <div class="cmd-pulse-val" id="cmd-pipeline-val">—</div>
              <div class="cmd-pulse-label">Total pipeline</div>
            </div>
            <div class="cmd-pulse-sep"></div>
            <div class="cmd-pulse-stat">
              <div class="cmd-pulse-val" id="cmd-engaged-count">—</div>
              <div class="cmd-pulse-label">Active engagements</div>
            </div>
          </div>
        </div>

        <!-- Strategic Brief -->
        <div id="cmd-brief-section" style="display:none;margin-bottom:24px">
          <div class="cmd-brief-card">
            <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:14px">
              <div class="cmd-section-title" style="color:var(--blue)">◈ Strategic Brief — AI Generated</div>
              <div id="cmd-brief-date" style="font-size:11px;color:var(--m)"></div>
            </div>
            <div id="cmd-brief-loading" class="loading-pulse" style="display:none"><div class="pulse-dot"></div><div class="pulse-dot"></div><div class="pulse-dot"></div><span>Generating market intelligence…</span></div>
            <div id="cmd-brief-text" style="font-size:13.5px;line-height:1.8;color:var(--t);white-space:pre-wrap"></div>
          </div>
        </div>

        <div class="cmd-grid">
          <!-- Critical Signals -->
          <div>
            <div class="cmd-section-header">
              <div class="cmd-section-title">Critical Signals</div>
              <span class="cmd-section-link" onclick="switchView('trends')">View all →</span>
            </div>
            <div id="cmd-signals-list"></div>
          </div>

          <!-- High-Value Targets -->
          <div>
            <div class="cmd-section-header">
              <div class="cmd-section-title">High-Value Targets</div>
              <span class="cmd-section-link" onclick="switchView('new-biz')">View all →</span>
            </div>
            <div id="cmd-targets-list"></div>
          </div>
        </div>
      </div>

      <!-- ── OVERVIEW ── -->
      <div class="view active" id="view-overview">
        <div id="overview-kpis" class="kpi-row"></div>
        <div class="sec-header">
          <div class="sec-title" id="industry-filter-title">All Industries</div>
        </div>
        <div class="industry-filter" id="industry-filter"></div>
        <div class="account-grid" id="account-grid"></div>
      </div>

      <!-- ── OPPORTUNITY MAP ── -->
      <div class="view" id="view-new-biz">
        <div class="sec-header">
          <div>
            <div class="sec-title">Opportunity Map</div>
            <div style="font-size:12px;color:var(--m);margin-top:3px">Unengaged accounts ranked by ICP score and commercial timing</div>
          </div>
        </div>
        <div class="industry-filter" id="nb-industry-filter"></div>
        <div class="account-grid" id="nb-account-grid"></div>
      </div>

      <!-- ── EXISTING ── -->
      <div class="view" id="view-existing">
        <div class="sec-header"><div class="sec-title">Existing Accounts — Activation Gaps</div></div>
        <div id="existing-content">
          <div class="loading-pulse"><div class="pulse-dot"></div><div class="pulse-dot"></div><div class="pulse-dot"></div><span>Loading...</span></div>
        </div>
      </div>

      <!-- ── TRENDS ── -->
      <div class="view" id="view-trends">
        <div class="sec-header" style="margin-bottom:20px">
          <div>
            <div class="sec-title">Signal Radar</div>
            <div style="font-size:12px;color:var(--m);margin-top:3px">Market intelligence feed — regulation, leadership transitions, tech shifts, and competitive moves</div>
          </div>
        </div>
        <div class="signal-list" id="signal-list"></div>
      </div>

      <!-- ── PREDICTIONS ── -->
      <div class="view" id="view-predictions">
        <div class="sec-header">
          <div>
            <div class="sec-title">Predictive Intelligence</div>
            <div style="font-size:12px;color:var(--m);margin-top:3px">AI-scored risk and opportunity signals per account, based on market dynamics</div>
          </div>
        </div>
        <div class="pred-grid" id="pred-grid"></div>
      </div>

      <!-- ── TODAY ── -->
      <div class="view" id="view-today">
        <div class="today-hero">
          <div class="today-date" id="today-date-label"></div>
          <div class="today-headline" id="today-headline">Good morning</div>
          <div class="today-sub" id="today-sub">Loading your priorities…</div>
        </div>
        <div id="commit-widget-area"></div>
        <div class="today-grid" id="today-grid">
          <div>
            <div class="sec-header"><div class="sec-title">Urgent Actions</div><button class="add-btn" onclick="openActionModal()">+ Add action</button></div>
            <div id="today-actions-list"><div class="loading-pulse"><div class="pulse-dot"></div><div class="pulse-dot"></div><div class="pulse-dot"></div><span>Loading…</span></div></div>
          </div>
          <div>
            <div class="sec-header"><div class="sec-title">Cold Accounts (14+ days)</div></div>
            <div id="today-cold-list"></div>
            <div style="margin-top:20px">
              <div class="sec-header"><div class="sec-title">Recent Meetings</div></div>
              <div id="today-meetings-list"></div>
            </div>
          </div>
        </div>
      </div>

      <!-- ── ACTIONS ── -->
      <div class="view" id="view-actions">
        <div class="sec-header">
          <div class="sec-title">All Actions</div>
          <div style="display:flex;gap:10px;align-items:center">
            <select id="action-status-filter" onchange="loadActions()" style="background:rgba(255,255,255,.06);border:1px solid var(--border);border-radius:7px;padding:6px 10px;font:500 12px var(--font);color:var(--w)">
              <option value="open">Open</option>
              <option value="all">All</option>
              <option value="done">Done</option>
            </select>
            <button class="add-btn" onclick="openActionModal()">+ New action</button>
          </div>
        </div>
        <div id="actions-list"><div class="loading-pulse"><div class="pulse-dot"></div><div class="pulse-dot"></div><div class="pulse-dot"></div><span>Loading…</span></div></div>
      </div>

      <!-- ── MEETINGS ── -->
      <div class="view" id="view-meetings">
        <div class="sec-header">
          <div class="sec-title">Meeting Log</div>
          <button class="add-btn" onclick="openMeetingModal()">+ Log meeting</button>
        </div>
        <div id="meetings-list"><div class="loading-pulse"><div class="pulse-dot"></div><div class="pulse-dot"></div><div class="pulse-dot"></div><span>Loading…</span></div></div>
      </div>

      <!-- ── STRATEGIC AI ── -->
      <div class="view" id="view-intelligence">
        <div class="sec-header">
          <div>
            <div class="sec-title">Strategic AI</div>
            <div style="font-size:12px;color:var(--m);margin-top:3px">Win pattern analysis, churn risk detection, and account-level diagnosis</div>
          </div>
          <div style="display:flex;gap:10px">
            <button class="add-btn" style="background:rgba(124,58,237,.15);border:1px solid rgba(124,58,237,.3);color:var(--purple)" onclick="loadWinPatterns()">◈ Analyse win patterns</button>
          </div>
        </div>
        <div id="intel-hero" class="intel-hero" style="display:none">
          <div class="intel-stat"><div class="intel-stat-val" id="intel-cold-count" style="color:var(--amber)">—</div><div class="intel-stat-label">Cold accounts</div></div>
          <div class="intel-stat"><div class="intel-stat-val" id="intel-churn-count" style="color:var(--red)">—</div><div class="intel-stat-label">Churn risks</div></div>
          <div class="intel-stat"><div class="intel-stat-val" id="intel-won-count" style="color:var(--green)">—</div><div class="intel-stat-label">Won / Active</div></div>
        </div>
        <div id="intel-win-box" style="display:none;margin-bottom:20px">
          <div class="sec-header"><div class="sec-title" style="color:var(--green)">Win Pattern Analysis</div></div>
          <div id="intel-win-loading" class="loading-pulse" style="display:none"><div class="pulse-dot"></div><div class="pulse-dot"></div><div class="pulse-dot"></div><span>Analysing patterns…</span></div>
          <div id="intel-win-content" class="win-analysis-box"></div>
        </div>
        <div class="intel-grid" id="intel-grid">
          <div class="loading-pulse" style="grid-column:1/-1"><div class="pulse-dot"></div><div class="pulse-dot"></div><div class="pulse-dot"></div><span>Loading intelligence…</span></div>
        </div>
      </div>

      <!-- ── GLOBAL ── -->
      <div class="view" id="view-global">
        <div id="global-kpis" class="kpi-row"></div>
        <div class="sec-header" style="margin-bottom:16px">
          <div>
            <div class="sec-title">Global Market Overview</div>
            <div style="font-size:12px;color:var(--m);margin-top:3px">Pipeline health, service mix, and strategic signals across all Nordic markets</div>
          </div>
        </div>
        <div class="country-cards" id="country-cards"></div>
        <div style="display:grid;grid-template-columns:1.4fr 1fr;gap:20px;margin-top:4px">
          <div>
            <div class="sec-header"><div class="sec-title">Activation Services — Pipeline by Service</div></div>
            <div class="svc-chart" id="svc-chart"></div>
          </div>
          <div>
            <div class="sec-header"><div class="sec-title">Global Signals</div></div>
            <div class="signal-list" id="global-signal-list"></div>
          </div>
        </div>
      </div>

    </main>
  </div>
</div>

<!-- Account detail panel -->
<div class="detail-overlay" id="detail-overlay" onclick="closeDetail()"></div>
<div class="detail-panel" id="detail-panel">
  <div class="dp-head">
    <button class="dp-close" onclick="closeDetail()">×</button>
    <div style="font-size:11px;font-weight:700;letter-spacing:.1em;text-transform:uppercase;color:var(--m2);margin-bottom:6px" id="dp-type-badge"></div>
    <div style="font-size:20px;font-weight:800;letter-spacing:-.03em" id="dp-name"></div>
    <div style="font-size:12px;color:var(--m);margin-top:4px" id="dp-meta"></div>
  </div>
  <div class="dp-tabs" id="dp-tabs">
    <div class="dp-tab active" data-tab="overview" onclick="switchDpTab('overview')">Overview</div>
    <div class="dp-tab" data-tab="actions" onclick="switchDpTab('actions')">Actions</div>
    <div class="dp-tab" data-tab="meetings" onclick="switchDpTab('meetings')">Meetings</div>
    <div class="dp-tab" data-tab="outreach" onclick="switchDpTab('outreach')">Outreach</div>
  </div>
  <div class="dp-body">
    <div class="dp-tab-content active" id="dp-tab-overview"></div>
    <div class="dp-tab-content" id="dp-tab-actions">
      <div style="padding:16px 0 10px;display:flex;justify-content:flex-end">
        <button class="add-btn" onclick="openActionModal(currentDetailAccount && currentDetailAccount.id)">+ Add action</button>
      </div>
      <div id="dp-actions-list"><div class="loading-pulse"><div class="pulse-dot"></div><div class="pulse-dot"></div><div class="pulse-dot"></div><span>Loading…</span></div></div>
    </div>
    <div class="dp-tab-content" id="dp-tab-meetings">
      <div style="padding:16px 0 10px;display:flex;justify-content:flex-end">
        <button class="add-btn" onclick="openMeetingModal(currentDetailAccount && currentDetailAccount.id)">+ Log meeting</button>
      </div>
      <div id="dp-meetings-list"><div class="loading-pulse"><div class="pulse-dot"></div><div class="pulse-dot"></div><div class="pulse-dot"></div><span>Loading…</span></div></div>
    </div>
    <div class="dp-tab-content" id="dp-tab-outreach">
      <div style="padding-top:16px">
        <div class="outreach-controls">
          <select class="outreach-select" id="dp-outreach-channel">
            <option value="linkedin">LinkedIn</option>
            <option value="email">Email</option>
          </select>
          <select class="outreach-select" id="dp-outreach-lang">
            <option value="en">English</option>
            <option value="no">Norwegian</option>
            <option value="da">Danish</option>
            <option value="sv">Swedish</option>
          </select>
          <button class="add-btn" onclick="generateOutreach()">Generate ✦</button>
        </div>
        <div id="dp-outreach-loading" style="display:none" class="loading-pulse"><div class="pulse-dot"></div><div class="pulse-dot"></div><div class="pulse-dot"></div><span>Generating…</span></div>
        <div id="dp-outreach-msg" class="outreach-msg" style="display:none"></div>
        <button id="dp-outreach-copy" class="outreach-copy-btn" style="display:none" onclick="copyOutreach()">Copy to clipboard</button>
      </div>
    </div>
  </div>
</div>

<!-- Action modal -->
<div class="modal-overlay" id="action-modal">
  <div class="modal">
    <div class="modal-title">New Action</div>
    <input type="hidden" id="action-account-id">
    <div class="modal-field">
      <label>Account</label>
      <select id="action-account-select" style="width:100%;background:rgba(255,255,255,.06);border:1px solid var(--border);border-radius:8px;padding:10px 12px;font:400 13px var(--font);color:var(--w)"></select>
    </div>
    <div class="modal-field">
      <label>Action Title</label>
      <input type="text" id="action-title" placeholder="e.g. Follow up with Morten Syversen">
    </div>
    <div style="display:grid;grid-template-columns:1fr 1fr;gap:10px">
      <div class="modal-field">
        <label>Type</label>
        <select id="action-type" style="width:100%;background:rgba(255,255,255,.06);border:1px solid var(--border);border-radius:8px;padding:10px 12px;font:400 13px var(--font);color:var(--w)">
          <option value="follow-up">Follow-up</option>
          <option value="call">Call</option>
          <option value="email">Email</option>
          <option value="linkedin">LinkedIn</option>
          <option value="meeting">Meeting</option>
          <option value="proposal">Proposal</option>
        </select>
      </div>
      <div class="modal-field">
        <label>Priority</label>
        <select id="action-priority" style="width:100%;background:rgba(255,255,255,.06);border:1px solid var(--border);border-radius:8px;padding:10px 12px;font:400 13px var(--font);color:var(--w)">
          <option value="medium">Medium</option>
          <option value="high">High</option>
          <option value="critical">Critical</option>
          <option value="low">Low</option>
        </select>
      </div>
    </div>
    <div class="modal-field">
      <label>Due Date</label>
      <input type="date" id="action-due">
    </div>
    <div class="modal-field">
      <label>Notes (optional)</label>
      <textarea id="action-desc" placeholder="Context or details…"></textarea>
    </div>
    <div class="modal-actions">
      <button class="modal-btn secondary" onclick="closeActionModal()">Cancel</button>
      <button class="modal-btn primary" onclick="saveAction()">Save Action</button>
    </div>
  </div>
</div>

<!-- Meeting modal -->
<div class="modal-overlay" id="meeting-modal">
  <div class="modal">
    <div class="modal-title">Log Meeting</div>
    <input type="hidden" id="meeting-account-id">
    <div class="modal-field">
      <label>Account</label>
      <select id="meeting-account-select" style="width:100%;background:rgba(255,255,255,.06);border:1px solid var(--border);border-radius:8px;padding:10px 12px;font:400 13px var(--font);color:var(--w)"></select>
    </div>
    <div style="display:grid;grid-template-columns:1fr 1fr;gap:10px">
      <div class="modal-field">
        <label>Date</label>
        <input type="date" id="meeting-date">
      </div>
      <div class="modal-field">
        <label>Outcome</label>
        <select id="meeting-outcome" style="width:100%;background:rgba(255,255,255,.06);border:1px solid var(--border);border-radius:8px;padding:10px 12px;font:400 13px var(--font);color:var(--w)">
          <option value="positive">Positive</option>
          <option value="neutral">Neutral</option>
          <option value="negative">Negative</option>
          <option value="no-show">No-show</option>
        </select>
      </div>
    </div>
    <div class="modal-field">
      <label>Participants</label>
      <input type="text" id="meeting-participants" placeholder="e.g. Morten Syversen, Jacob">
    </div>
    <div class="modal-field">
      <label>Summary</label>
      <textarea id="meeting-summary" placeholder="What was discussed? Key insights?"></textarea>
    </div>
    <div class="modal-field">
      <label>Next Step</label>
      <input type="text" id="meeting-next" placeholder="e.g. Send proposal by Friday">
    </div>
    <div class="modal-actions">
      <button class="modal-btn secondary" onclick="closeMeetingModal()">Cancel</button>
      <button class="modal-btn primary" onclick="saveMeeting()">Save Meeting</button>
    </div>
  </div>
</div>

<!-- Weekly commit modal -->
<div class="modal-overlay" id="commit-modal">
  <div class="modal">
    <div class="modal-title">Weekly Commit</div>
    <div class="modal-field">
      <label>What do you commit to closing this week?</label>
      <textarea id="commit-text" placeholder="e.g. Book first meeting with Elkjøp. Send proposal to Vinmonopolet. Follow up Trumf after intro call." style="min-height:100px"></textarea>
    </div>
    <div class="modal-field">
      <label>Target Value (€)</label>
      <input type="number" id="commit-value" placeholder="50000">
    </div>
    <div class="modal-actions">
      <button class="modal-btn secondary" onclick="closeCommitModal()">Cancel</button>
      <button class="modal-btn primary" onclick="saveCommit()">Commit →</button>
    </div>
  </div>
</div>

<div id="cc-toast"></div>

<script>
// ══ STATE ══════════════════════════════════════════════════════════════════════
let currentUser = null;
let countryData = null;
let globalData  = null;
let activeIndustry = 'all';
let activeView = 'overview';
let currentDetailAccount = null;
const STAGE_COLORS = {
  identified: {color:'#6B8EF7',bg:'rgba(21,62,237,.15)'},
  proposed:   {color:'#F59E0B',bg:'rgba(245,158,11,.15)'},
  negotiating:{color:'#F97316',bg:'rgba(249,115,22,.15)'},
  active:     {color:'#00D4A0',bg:'rgba(0,212,160,.15)'},
  completed:  {color:'#888',bg:'rgba(128,128,128,.15)'},
};
const DEAL_STAGES = ['identified','qualified','engaged','proposed','negotiating','closed_won'];
const DEAL_STAGE_LABELS = {identified:'Identified',qualified:'Qualified',engaged:'Engaged',proposed:'Proposed',negotiating:'Negotiating',closed_won:'Won',closed_lost:'Lost'};
const PRIORITY_COLORS = {critical:'var(--red)',high:'var(--amber)',medium:'var(--blue)',low:'var(--m2)'};

// ══ AUTH ═══════════════════════════════════════════════════════════════════════
async function doLogin() {
  const email = document.getElementById('login-email').value.trim();
  const pw    = document.getElementById('login-pw').value;
  const btn   = document.querySelector('.login-btn');
  const err   = document.getElementById('login-err');
  btn.disabled = true; btn.textContent = 'Signing in…';
  err.textContent = '';
  try {
    const params = new URLSearchParams({email, password: pw});
    const r = await fetch(`/api/cc/login?${params}`, {method:'POST'});
    const d = await r.json();
    if (!r.ok) { err.textContent = d.error || 'Login failed'; return; }
    await initApp(d);
  } catch(e) { err.textContent = 'Connection error'; }
  finally { btn.disabled = false; btn.textContent = 'Sign in →'; }
}

async function doLogout() {
  await fetch('/api/cc/logout',{method:'POST'});
  currentUser = null; countryData = null; globalData = null;
  document.getElementById('app-shell').classList.add('hidden');
  document.getElementById('login-screen').classList.remove('hidden');
}

document.getElementById('login-pw').addEventListener('keydown', e => { if(e.key==='Enter') doLogin(); });

// ══ INIT ═══════════════════════════════════════════════════════════════════════
async function initApp(user) {
  currentUser = user;
  document.getElementById('login-screen').classList.add('hidden');
  document.getElementById('app-shell').classList.remove('hidden');

  const r2 = await fetch('/api/cc/me');
  currentUser = await r2.json();

  document.getElementById('tb-name').textContent  = currentUser.name;
  document.getElementById('tb-avatar').textContent = currentUser.initials || currentUser.name.slice(0,2).toUpperCase();

  if (currentUser.role === 'global') {
    // Global view
    document.getElementById('tb-country').textContent = '🌍  Global Markets';
    document.querySelectorAll('.nav-item').forEach(n => n.style.display = 'none');
    const gi = document.querySelector('[data-view="global"]') ||
      (() => { const d = document.createElement('div'); d.className='nav-item active';d.dataset.view='global';d.onclick=()=>switchView('global');d.innerHTML='<span class="nav-icon">🌍</span> Global Overview'; document.querySelector('.sb-section').appendChild(d); return d; })();
    gi.style.display = '';
    document.getElementById('sb-country-switcher').style.display = 'block';
    buildCountrySwitcher();
    switchView('global');
    loadGlobalData();
  } else {
    // Country head view
    const meta = {'no':'🇳🇴 Norway','dk':'🇩🇰 Denmark','se':'🇸🇪 Sweden','uk':'🇬🇧 UK','fr':'🇫🇷 France'};
    document.getElementById('tb-country').textContent = meta[currentUser.country] || currentUser.country.toUpperCase();
    switchView('foresight');
    loadCountryData(currentUser.country);
  }
}

// On page load, check if already logged in
(async () => {
  try {
    const r = await fetch('/api/cc/me');
    if (r.ok) { const u = await r.json(); await initApp(u); }
  } catch(e) {}
})();

// ══ DATA LOADING ══════════════════════════════════════════════════════════════
async function loadCountryData(country) {
  showLoading();
  const r = await fetch('/api/cc/country-data?country=' + country);
  countryData = await r.json();
  renderCountryDashboard();
  if (activeView === 'command') loadCommandView();
}

async function loadGlobalData() {
  showLoading();
  const r = await fetch('/api/cc/global-data');
  globalData = await r.json();
  renderGlobalDashboard();
}

function showLoading() {
  ['account-grid','nb-account-grid','signal-list','pred-grid'].forEach(id => {
    const el = document.getElementById(id);
    if(el) el.innerHTML = '<div class="loading-pulse"><div class="pulse-dot"></div><div class="pulse-dot"></div><div class="pulse-dot"></div><span>Loading…</span></div>';
  });
}

// ══ RENDER: COUNTRY ═══════════════════════════════════════════════════════════
function renderCountryDashboard() {
  if (!countryData) return;
  const d = countryData;

  // KPIs
  const pFmt = d.kpis.pipeline >= 1000000 ? '€' + (d.kpis.pipeline/1000000).toFixed(1) + 'M' : '€' + (d.kpis.pipeline/1000).toFixed(0) + 'K';
  document.getElementById('overview-kpis').innerHTML = `
    <div class="kpi-card"><div class="kpi-label">Pipeline Value</div><div class="kpi-value">${pFmt}</div><div class="kpi-sub">${d.kpis.accounts} accounts</div></div>
    <div class="kpi-card"><div class="kpi-label">Named Buyers</div><div class="kpi-value">${d.kpis.buyers}</div><div class="kpi-sub">of ${d.kpis.accounts} accounts</div></div>
    <div class="kpi-card"><div class="kpi-label">Active Activations</div><div class="kpi-value">${d.kpis.active_activations}</div><div class="kpi-sub">services in delivery</div></div>
    <div class="kpi-card"><div class="kpi-label">Avg ICP Score</div><div class="kpi-value">${(d.accounts.reduce((s,a)=>s+(a.icp||0),0)/Math.max(d.accounts.length,1)).toFixed(1)}</div><div class="kpi-sub">out of 10</div></div>`;

  // Industry filter
  const inds = d.industries;
  const buildFilter = (filterId, gridId) => {
    const f = document.getElementById(filterId);
    if (!f) return;
    f.innerHTML = `<button class="ind-pill ${activeIndustry==='all'?'active':''}" onclick="filterIndustry('all','${filterId}','${gridId}')">All <span class="ind-count">${d.accounts.filter(a=>filterId.includes('nb')?a.account_type==='prospect':true).length}</span></button>` +
      inds.map(i => {
        const cnt = d.accounts.filter(a => a.industry_slug === i.slug && (filterId.includes('nb') ? a.account_type==='prospect' : true)).length;
        return cnt ? `<button class="ind-pill" onclick="filterIndustry('${i.slug}','${filterId}','${gridId}')">${i.icon} ${i.name} <span class="ind-count">${cnt}</span></button>` : '';
      }).join('');
  };
  buildFilter('industry-filter','account-grid');
  buildFilter('nb-industry-filter','nb-account-grid');

  // Account grids
  renderAccountGrid('account-grid', d.accounts);
  renderAccountGrid('nb-account-grid', d.accounts.filter(a => a.account_type === 'prospect'));

  // Signals
  renderSignals('signal-list', d.signals);

  // Predictions
  renderPredictions('pred-grid', d.predictions, d.accounts);

  // Existing accounts
  renderExisting(d.accounts.filter(a => a.account_type === 'existing'));

  // Sidebar signals summary
  const crit = d.signals.filter(s=>s.severity==='critical').length;
  const warn = d.signals.filter(s=>s.severity==='warning').length;
  document.getElementById('sb-signals-summary').innerHTML = `
    <div class="sb-label">Signals</div>
    ${crit ? `<div class="sig-row"><span class="sig-dot critical"></span><span style="font-size:12px;color:var(--m)">${crit} Critical</span></div>` : ''}
    ${warn ? `<div class="sig-row"><span class="sig-dot warning"></span><span style="font-size:12px;color:var(--m)">${warn} Warnings</span></div>` : ''}
    ${!crit && !warn ? `<div class="sig-row"><span class="sig-dot info"></span><span style="font-size:12px;color:var(--m)">No critical signals</span></div>` : ''}`;
}

function renderAccountGrid(gridId, accounts) {
  const grid = document.getElementById(gridId);
  if (!grid) return;
  const filtered = activeIndustry === 'all' ? accounts : accounts.filter(a => a.industry_slug === activeIndustry);
  if (!filtered.length) { grid.innerHTML = '<div style="color:var(--m);font-size:13px;padding:20px 0">No accounts in this industry.</div>'; return; }
  grid.innerHTML = filtered.map(a => renderAccountCard(a)).join('');
}

function renderAccountCard(a) {
  const topAct = a.activations[0];
  const svcColor = topAct ? topAct.service_color : '#153EED';
  const icpClass = (a.icp >= 8) ? 'icp-high' : (a.icp >= 6) ? 'icp-mid' : '';
  const pFmt = a.pipeline >= 1000000 ? '€'+(a.pipeline/1000000).toFixed(1)+'M' : '€'+(a.pipeline/1000).toFixed(0)+'K';
  const winPct = Math.round((a.win_prob || 0) * 100);
  const winColor = winPct >= 60 ? 'var(--green)' : winPct >= 35 ? 'var(--amber)' : 'var(--red)';
  const buyerInitials = (a.buyer || 'TBD').split(' ').filter(Boolean).map(w=>w[0]).slice(0,2).join('');
  const actPills = a.activations.map(ac => {
    const sc = STAGE_COLORS[ac.stage] || {color:'#888',bg:'rgba(128,128,128,.15)'};
    return `<span class="act-pill" style="color:${sc.color};border-color:${sc.color}30;background:${sc.bg}"><span class="act-stage-dot"></span>${ac.service_name}</span>`;
  }).join('');

  return `<div class="account-card" style="--svc-color:${svcColor}" onclick="openDetail(${JSON.stringify(a).replace(/"/g,'&quot;')})">
    <div class="ac-head">
      <div>
        <div class="ac-name">${a.name}</div>
      </div>
      <div class="ac-value">${pFmt}</div>
    </div>
    <div class="ac-tags">
      <span class="ac-tag">${a.industry_icon} ${a.industry}</span>
      ${a.icp ? `<span class="ac-tag ${icpClass}">ICP ${a.icp}</span>` : ''}
      ${a.account_type === 'existing' ? '<span class="ac-tag existing">✓ Active</span>' : ''}
    </div>
    ${actPills ? `<div class="act-section"><div class="act-label">Activation Services</div><div class="act-pills">${actPills}</div></div>` : '<div style="font-size:12px;color:var(--m2);margin-bottom:12px">No activations mapped yet</div>'}
    <div class="buyer-row">
      <div class="buyer-left">
        <div class="buyer-avatar">${buyerInitials || '?'}</div>
        <div>
          <div class="buyer-name">${a.buyer || 'Buyer TBD'}</div>
          <div class="buyer-role">${(a.buyer_role || '').slice(0,40)}${(a.buyer_role||'').length>40?'…':''}</div>
        </div>
      </div>
      <div class="win-bar-wrap">
        <div class="win-bar"><div class="win-fill" style="width:${winPct}%;--fill-color:${winColor}"></div></div>
        <span class="win-pct">${winPct}%</span>
      </div>
    </div>
  </div>`;
}

function filterIndustry(slug, filterId, gridId) {
  activeIndustry = slug;
  document.querySelectorAll(`#${filterId} .ind-pill`).forEach(p => {
    p.classList.toggle('active', p.onclick.toString().includes(`'${slug}'`));
  });
  const accounts = countryData ? (gridId.includes('nb') ? countryData.accounts.filter(a=>a.account_type==='prospect') : countryData.accounts) : [];
  renderAccountGrid(gridId, accounts);
}

function renderSignals(listId, signals) {
  const list = document.getElementById(listId);
  if (!list) return;
  if (!signals.length) { list.innerHTML = '<div style="color:var(--m);font-size:13px">No active signals.</div>'; return; }
  list.innerHTML = signals.map(s => `
    <div class="signal-card ${s.severity}">
      <div class="sig-head">
        <span class="sig-badge ${s.severity}">${s.severity.toUpperCase()}</span>
        <div>
          <div class="sig-title">${s.title}</div>
          <div style="font-size:11px;color:var(--m2);margin-top:3px">${(s.type||s.signal_type||'').charAt(0).toUpperCase()+(s.type||s.signal_type||'').slice(1)} · ${s.vertical}</div>
        </div>
      </div>
      <div class="sig-desc">${s.description}</div>
      ${s.action ? `<div class="sig-action"><div class="sig-action-label">→ Recommended Action</div>${s.action}</div>` : ''}
    </div>`).join('');
}

function renderPredictions(gridId, predictions, accounts) {
  const grid = document.getElementById(gridId);
  if (!grid) return;
  const accMap = {};
  if (accounts) accounts.forEach(a => accMap[a.id] = a.name);

  let html = predictions.map(p => {
    const oppColor = p.opp >= 7 ? 'var(--green)' : p.opp >= 4 ? 'var(--amber)' : 'var(--red)';
    const riskColor = p.risk >= 7 ? 'var(--red)' : p.risk >= 4 ? 'var(--amber)' : 'var(--green)';
    const conf = Math.round((p.confidence || 0) * 100);
    return `<div class="pred-card">
      <div style="display:flex;justify-content:space-between;align-items:flex-start">
        <div>
          <div style="font-size:14px;font-weight:700">${p.account || accMap[p.account_id] || 'Account'}</div>
          <div style="font-size:11px;color:var(--m);margin-top:2px">${p.service} · ${p.weeks}w window</div>
        </div>
        <span style="font-size:10px;font-weight:700;padding:3px 8px;border-radius:4px;background:rgba(21,62,237,.15);color:#6B8EF7">AI PREDICTION</span>
      </div>
      <div class="pred-scores">
        <div class="pred-score"><div class="pred-score-val" style="color:${oppColor}">${(p.opp||0).toFixed(1)}</div><div class="pred-score-label">Opportunity</div></div>
        <div class="pred-score"><div class="pred-score-val" style="color:${riskColor}">${(p.risk||0).toFixed(1)}</div><div class="pred-score-label">Risk</div></div>
      </div>
      <div class="pred-trigger">${p.trigger || p.trigger_summary || ''}</div>
      <div class="pred-confidence">Confidence: ${conf}%
        <div class="confidence-bar"><div class="confidence-fill" style="width:${conf}%"></div></div>
      </div>
    </div>`;
  }).join('');

  // Add "Generate new prediction" cards for accounts without predictions
  if (accounts) {
    const accountsWithPreds = new Set(predictions.map(p => p.account || ''));
    const without = accounts.filter(a => !accountsWithPreds.has(a.name) && a.account_type === 'prospect').slice(0,3);
    html += without.map(a => `
      <div class="pred-card" style="border-style:dashed;opacity:.7">
        <div style="font-size:14px;font-weight:700">${a.name}</div>
        <div style="font-size:12px;color:var(--m);margin:10px 0">No prediction generated yet.</div>
        <button class="pred-gen-btn" onclick="generatePrediction(${a.id},'${a.name}')">Generate AI Prediction →</button>
      </div>`).join('');
  }
  grid.innerHTML = html || '<div style="color:var(--m);font-size:13px">No predictions yet.</div>';
}

function renderExisting(accounts) {
  const el = document.getElementById('existing-content');
  if (!el) return;
  if (!accounts.length) {
    el.innerHTML = '<div style="color:var(--m);font-size:13px;padding:20px 0">No existing accounts yet.</div>'; return;
  }
  el.innerHTML = `<table class="ea-table">
    <thead><tr><th>Account</th><th>Industry</th><th>Active Services</th><th>Revenue</th><th>Activation Gap</th><th>Next Service</th></tr></thead>
    <tbody>${accounts.map(a => {
      const activeActs = a.activations.filter(ac => ac.stage === 'active');
      const gap = a.activations.length === 0 ? 'No activations mapped' : activeActs.length === 0 ? 'All services in pipeline' : '';
      const allSvcs = ['Data Revenue','AI Readiness','Commerce Optim.','Shopify Build'];
      const activeSvcNames = a.activations.map(ac => ac.service_name);
      const nextSvc = allSvcs.find(s => !activeSvcNames.some(n => n.includes(s.split(' ')[0]))) || '—';
      return `<tr onclick="openDetail(${JSON.stringify(a).replace(/"/g,'&quot;')})" style="cursor:pointer">
        <td><strong>${a.name}</strong></td>
        <td>${a.industry_icon} ${a.industry}</td>
        <td>${activeActs.map(ac=>`<span style="font-size:11px;font-weight:600;color:var(--green);background:rgba(0,212,160,.1);padding:2px 7px;border-radius:4px;margin-right:4px">${ac.service_name}</span>`).join('') || '<span style="color:var(--m2);font-size:12px">None active</span>'}</td>
        <td style="font-size:12px;color:var(--m)">${a.revenue || '—'}</td>
        <td>${gap ? `<span class="gap-badge">${gap}</span>` : ''}</td>
        <td style="font-size:12px;color:var(--blue)">${nextSvc}</td>
      </tr>`;
    }).join('')}</tbody></table>`;
}

// ══ RENDER: GLOBAL ════════════════════════════════════════════════════════════
function renderGlobalDashboard() {
  if (!globalData) return;
  const d = globalData;

  const pFmt = d.total_pipeline >= 1000000 ? '€'+(d.total_pipeline/1000000).toFixed(1)+'M' : '€'+(d.total_pipeline/1000).toFixed(0)+'K';
  document.getElementById('global-kpis').innerHTML = `
    <div class="kpi-card"><div class="kpi-label">Total Pipeline</div><div class="kpi-value">${pFmt}</div><div class="kpi-sub">All markets</div></div>
    <div class="kpi-card"><div class="kpi-label">Total Accounts</div><div class="kpi-value">${d.total_accounts}</div><div class="kpi-sub">5 markets</div></div>
    <div class="kpi-card"><div class="kpi-label">Active Markets</div><div class="kpi-value">${d.countries.filter(c=>c.accounts>0).length}</div><div class="kpi-sub">of 5</div></div>
    <div class="kpi-card"><div class="kpi-label">Services Mapped</div><div class="kpi-value">${d.service_performance.filter(s=>s.pipeline>0).length}</div><div class="kpi-sub">of 8 offerings</div></div>`;

  document.getElementById('country-cards').innerHTML = d.countries.map(c => {
    const pFmt = c.pipeline >= 1000000 ? '€'+(c.pipeline/1000000).toFixed(1)+'M' : c.pipeline > 0 ? '€'+(c.pipeline/1000).toFixed(0)+'K' : '—';
    return `<div class="cc-card" onclick="drillCountry('${c.code}')">
      <div class="cc-flag">${c.meta.flag}</div>
      <div class="cc-cname">${c.meta.name}</div>
      <div class="cc-kpi"><div class="cc-kpi-v">${pFmt}</div><div class="cc-kpi-l">Pipeline</div></div>
      <div class="cc-kpi" style="margin-top:8px"><div style="font-size:14px;font-weight:700">${c.accounts}</div><div class="cc-kpi-l">Accounts</div></div>
      <div style="margin-top:10px;font-size:10px;color:var(--m);padding-top:8px;border-top:1px solid var(--border2)">${c.top_service}</div>
    </div>`;
  }).join('');

  const maxP = Math.max(...d.service_performance.map(s=>s.pipeline), 1);
  document.getElementById('svc-chart').innerHTML = d.service_performance.filter(s=>s.pipeline>0).map(s => {
    const pFmt = s.pipeline >= 1000000 ? '€'+(s.pipeline/1000000).toFixed(1)+'M' : '€'+(s.pipeline/1000).toFixed(0)+'K';
    return `<div class="svc-bar-row">
      <div class="svc-bar-name">${s.name}</div>
      <div class="svc-bar-track"><div class="svc-bar-fill" style="width:${Math.round(s.pipeline/maxP*100)}%;background:${s.color}"></div></div>
      <div class="svc-bar-val">${pFmt}</div>
    </div>`;
  }).join('');

  renderSignals('global-signal-list', d.signals);
}

function drillCountry(code) {
  // Global head drills into a specific country
  loadCountryData(code);
  const meta = {'no':'🇳🇴 Norway','dk':'🇩🇰 Denmark','se':'🇸🇪 Sweden','uk':'🇬🇧 UK','fr':'🇫🇷 France'};
  document.getElementById('tb-country').textContent = meta[code];
  // Show country nav items
  document.querySelectorAll('.nav-item').forEach(n => { if(n.dataset.view && n.dataset.view !== 'global') n.style.display = ''; });
  switchView('overview');
}

function buildCountrySwitcher() {
  const codes = [{code:'no',flag:'🇳🇴',name:'Norway'},{code:'dk',flag:'🇩🇰',name:'Denmark'},{code:'se',flag:'🇸🇪',name:'Sweden'},{code:'uk',flag:'🇬🇧',name:'UK'},{code:'fr',flag:'🇫🇷',name:'France'}];
  document.getElementById('country-switcher-list').innerHTML = codes.map(c =>
    `<div class="country-pill" onclick="drillCountry('${c.code}')">${c.flag} ${c.name}</div>`
  ).join('') + `<div class="country-pill active" onclick="switchView('global');loadGlobalData()">🌍 All Markets</div>`;
}

// ══ DETAIL PANEL ══════════════════════════════════════════════════════════════
function openDetail(account) {
  if (typeof account === 'string') account = JSON.parse(account);
  currentDetailAccount = account;

  document.getElementById('dp-type-badge').textContent = account.account_type === 'existing' ? '✓ EXISTING CLIENT' : '◎ PROSPECT';
  document.getElementById('dp-name').textContent = account.name;
  document.getElementById('dp-meta').textContent = `${account.industry_icon} ${account.industry}  ·  ${countryData?.meta?.flag || ''} ${countryData?.meta?.name || ''}`;

  // Reset to overview tab
  switchDpTab('overview');

  const pFmt = account.pipeline >= 1000000 ? '€'+(account.pipeline/1000000).toFixed(1)+'M' : '€'+(account.pipeline/1000).toFixed(0)+'K';
  const winPct = Math.round((account.win_prob||0)*100);
  const stage = account.deal_stage || 'identified';
  const stageIdx = DEAL_STAGES.indexOf(stage);

  const overviewEl = document.getElementById('dp-tab-overview');
  overviewEl.innerHTML = `
    <div class="dp-section" style="margin-top:0;padding-top:16px">
      <div class="dp-sec-title" style="margin-bottom:8px">Deal Stage</div>
      <div style="display:flex;gap:6px;flex-wrap:wrap;margin-bottom:4px">
        ${DEAL_STAGES.map((s,i) => {
          const active = s === stage;
          const past   = i < stageIdx;
          const bg = active ? 'var(--blue)' : past ? 'rgba(21,62,237,.15)' : 'rgba(255,255,255,.04)';
          const col = active ? '#fff' : past ? '#6B8EF7' : 'var(--m)';
          return `<span class="stage-pill" style="background:${bg};color:${col};border:1px solid ${active?'var(--blue)':'var(--border)'}" onclick="updateStage(${account.id},'${s}')">${DEAL_STAGE_LABELS[s]}</span>`;
        }).join('')}
      </div>
    </div>
    <div class="dp-section">
      <div class="dp-sec-title">Account Overview</div>
      ${[['Pipeline Value', pFmt],['Win Probability', winPct + '%'],['ICP Score', (account.icp||'—') + '/10'],['Deal Score', (account.deal||'—') + '/10'],['Revenue', account.revenue||'—']].map(([l,v])=>`<div class="dp-row"><span class="dp-row-label">${l}</span><span class="dp-row-val">${v}</span></div>`).join('')}
    </div>
    <div class="dp-section">
      <div class="dp-sec-title">Named Buyer</div>
      ${[['Name', account.buyer||'TBD'],['Role', account.buyer_role||'—']].map(([l,v])=>`<div class="dp-row"><span class="dp-row-label">${l}</span><span class="dp-row-val">${v}</span></div>`).join('')}
    </div>
    <div class="dp-section">
      <div class="dp-sec-title">Tech Stack</div>
      <div style="font-size:12px;color:var(--m);line-height:1.7">${account.tech_stack||'No data'}</div>
    </div>
    <div class="dp-section">
      <div class="dp-sec-title">Activation Map</div>
      <div class="activation-timeline">
        ${account.activations.length ? account.activations.map(ac => {
          const sc = STAGE_COLORS[ac.stage]||{color:'#888',bg:'rgba(128,128,128,.15)'};
          const costFmt = ac.cost ? '€'+(ac.cost/1000).toFixed(0)+'K entry' : '';
          const roiFmt  = ac.roi  ? '· Est. ROI €'+(ac.roi/1000).toFixed(0)+'K' : '';
          return `<div class="actl-row">
            <div class="actl-dot" style="background:${sc.color}"></div>
            <div>
              <div class="actl-svc">${ac.service_name}</div>
              <div class="actl-meta" style="color:${sc.color}">${ac.stage.toUpperCase()}</div>
              <div class="actl-meta">${ac.manager||''} ${costFmt} ${roiFmt}</div>
              ${ac.notes ? `<div class="actl-meta" style="margin-top:4px;font-style:italic">${ac.notes}</div>` : ''}
            </div>
          </div>`;
        }).join('') : '<div style="color:var(--m2);font-size:12px">No activations mapped.</div>'}
      </div>
    </div>
    ${account.predictions.length ? `<div class="dp-section">
      <div class="dp-sec-title">AI Predictions</div>
      ${account.predictions.map(p=>`<div style="background:rgba(255,255,255,.03);border-radius:8px;padding:12px;margin-bottom:8px">
        <div style="font-size:12px;font-weight:600;color:var(--green);margin-bottom:4px">${p.service} — ${Math.round(p.confidence*100)}% confidence · ${p.weeks}w</div>
        <div style="font-size:12px;color:var(--m);line-height:1.6">${p.trigger}</div>
      </div>`).join('')}
    </div>` : ''}
    <button class="pred-gen-btn" style="margin-top:4px" onclick="generatePrediction(${account.id},'${account.name}')">
      ◈ Generate New AI Prediction
    </button>`;

  document.getElementById('detail-overlay').classList.add('open');
  document.getElementById('detail-panel').classList.add('open');
}

async function updateStage(accountId, stage) {
  const params = new URLSearchParams({stage});
  const r = await fetch('/api/cc/accounts/' + accountId + '/stage?' + params, {method:'PATCH'});
  if (r.ok) {
    showToast('Stage updated → ' + DEAL_STAGE_LABELS[stage]);
    if (currentDetailAccount) { currentDetailAccount.deal_stage = stage; openDetail(currentDetailAccount); }
  }
}

function closeDetail() {
  document.getElementById('detail-overlay').classList.remove('open');
  document.getElementById('detail-panel').classList.remove('open');
}

// ══ PREDICTION GENERATION ════════════════════════════════════════════════════
async function generatePrediction(accountId, accountName) {
  showToast('Generating AI prediction for ' + accountName + '…');
  const r = await fetch(`/api/cc/predict?account_id=${accountId}`, {method:'POST'});
  const d = await r.json();
  if (!r.ok) { showToast('Error: ' + (d.error||'Failed')); return; }
  showToast('Prediction generated ✓');
  // Reload data
  const country = currentUser.country || (countryData && countryData.country);
  if (country) loadCountryData(country);
}

// ══ TODAY VIEW ════════════════════════════════════════════════════════════════
async function loadTodayData() {
  const r = await fetch('/api/cc/today');
  if (!r.ok) return;
  const d = await r.json();

  // Date headline
  const days = ['Sunday','Monday','Tuesday','Wednesday','Thursday','Friday','Saturday'];
  const months = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec'];
  const now = new Date();
  document.getElementById('today-date-label').textContent = days[now.getDay()] + ' · ' + months[now.getMonth()] + ' ' + now.getDate();

  const urgentCount = d.urgent_actions.length;
  const coldCount = d.cold_accounts.length;
  document.getElementById('today-headline').textContent = urgentCount
    ? urgentCount + ' action' + (urgentCount>1?'s':'') + ' need your attention'
    : "You're clear — no overdue actions";
  document.getElementById('today-sub').textContent =
    coldCount ? coldCount + ' accounts have gone cold · ' + d.recent_meetings.length + ' meetings this week'
               : 'Pipeline is healthy · ' + d.recent_meetings.length + ' meetings this week';

  // Sidebar badge
  const badge = document.getElementById('sb-action-count');
  if (urgentCount > 0) { badge.textContent = urgentCount; badge.style.display = ''; }
  else { badge.style.display = 'none'; }

  // Weekly commit widget
  const commitArea = document.getElementById('commit-widget-area');
  if (d.weekly_commit.exists) {
    const val = d.weekly_commit.target_value ? ' · Target: €' + Number(d.weekly_commit.target_value).toLocaleString() : '';
    commitArea.innerHTML = '<div class="commit-widget"><div class="commit-title"><span style="color:var(--green)">✓</span> This weeks commit' + val + '</div><div class="commit-text">' + (d.weekly_commit.commit_text || '—') + '</div><button class="commit-cta" onclick="openCommitModal()">Update →</button></div>';
  } else {
    commitArea.innerHTML = '<div class="commit-widget"><div class="commit-title">📌 Set your weekly commit</div><div class="commit-text" style="color:var(--m2)">What are you committing to close this week? Hold yourself accountable.</div><button class="commit-cta" onclick="openCommitModal()">Commit now →</button></div>';
  }

  // Urgent actions
  const actList = document.getElementById('today-actions-list');
  if (!d.urgent_actions.length) {
    actList.innerHTML = '<div style="color:var(--m);font-size:13px;padding:20px 0">No overdue or due-today actions. 🎉</div>';
  } else {
    actList.innerHTML = d.urgent_actions.map(a => renderActionCard(a)).join('');
  }

  // Cold accounts
  const coldList = document.getElementById('today-cold-list');
  if (!d.cold_accounts.length) {
    coldList.innerHTML = '<div style="color:var(--m);font-size:13px">No cold accounts.</div>';
  } else {
    coldList.innerHTML = d.cold_accounts.map(a => {
      const days = a.last_activity ? Math.floor((Date.now() - new Date(a.last_activity).getTime()) / 86400000) : null;
      const daysLabel = days ? days + 'd ago' : 'Never contacted';
      const pFmt = a.pipeline_value >= 1e6 ? '€'+(a.pipeline_value/1e6).toFixed(1)+'M' : '€'+(a.pipeline_value/1e3).toFixed(0)+'K';
      return '<div class="cold-account-row"><div><div style="font-weight:600">' + a.name + '</div><div style="font-size:11px;color:var(--m)">' + DEAL_STAGE_LABELS[a.deal_stage] + ' · ' + pFmt + '</div></div><span class="cold-days">' + daysLabel + '</span></div>';
    }).join('');
  }

  // Recent meetings
  const meetList = document.getElementById('today-meetings-list');
  if (!d.recent_meetings.length) {
    meetList.innerHTML = '<div style="color:var(--m);font-size:13px">No meetings this week.</div>';
  } else {
    meetList.innerHTML = d.recent_meetings.map(m => {
      const dateStr = m.date ? new Date(m.date).toLocaleDateString('en-GB',{weekday:'short',day:'numeric',month:'short'}) : '?';
      return '<div class="cold-account-row"><div><div style="font-weight:600">' + m.account_name + '</div><div style="font-size:11px;color:var(--m)">' + dateStr + '</div></div><span class="meeting-outcome ' + m.outcome + '">' + (m.outcome||'?') + '</span></div>';
    }).join('');
  }
}

function renderActionCard(a) {
  const isOverdue = a.is_overdue || (a.due_date && new Date(a.due_date) < new Date());
  const cardClass = isOverdue ? 'overdue' : (a.priority === 'high' || a.priority === 'critical') ? 'high' : '';
  const dueStr = a.due_date ? new Date(a.due_date).toLocaleDateString('en-GB',{day:'numeric',month:'short'}) : '';
  const dueLabelClass = isOverdue ? 'overdue-label' : '';
  return '<div class="action-card ' + cardClass + '" id="action-card-' + a.id + '">' +
    '<div class="priority-dot ' + (a.priority||'medium') + '"></div>' +
    '<div onclick="markActionDone(' + a.id + ')" class="action-check" id="check-' + a.id + '" title="Mark done" style="cursor:pointer"></div>' +
    '<div style="flex:1">' +
      '<div class="action-title">' + a.title + '</div>' +
      '<div class="action-meta">' +
        (a.account_name ? '<span>' + a.account_name + '</span>' : '') +
        '<span class="action-type-badge">' + (a.action_type||'follow-up') + '</span>' +
        (dueStr ? '<span class="' + dueLabelClass + '">' + (isOverdue ? '⚠ Overdue · ' : '') + dueStr + '</span>' : '') +
      '</div>' +
    '</div>' +
  '</div>';
}

async function markActionDone(actionId) {
  const params = new URLSearchParams({status:'done'});
  const r = await fetch('/api/cc/actions/' + actionId + '?' + params, {method:'PATCH'});
  if (r.ok) {
    const card = document.getElementById('action-card-' + actionId);
    if (card) { card.style.opacity = '0.4'; card.style.pointerEvents = 'none'; }
    showToast('Action marked as done ✓');
    setTimeout(() => { loadTodayData(); loadActions(); }, 800);
  }
}

// ══ ACTIONS VIEW ══════════════════════════════════════════════════════════════
async function loadActions() {
  const statusFilter = document.getElementById('action-status-filter');
  const status = statusFilter ? statusFilter.value : 'open';
  const r = await fetch('/api/cc/actions?status=' + status);
  if (!r.ok) return;
  const actions = await r.json();
  const list = document.getElementById('actions-list');
  if (!list) return;
  if (!actions.length) { list.innerHTML = '<div style="color:var(--m);font-size:13px;padding:20px 0">No ' + status + ' actions.</div>'; return; }
  list.innerHTML = actions.map(a => renderActionCard(a)).join('');
}

// ══ MEETINGS VIEW ═════════════════════════════════════════════════════════════
async function loadMeetings(accountId) {
  let url = '/api/cc/meetings';
  if (accountId) url += '?account_id=' + accountId;
  const r = await fetch(url);
  if (!r.ok) return [];
  return await r.json();
}

async function renderMeetingsList(listId, accountId) {
  const el = document.getElementById(listId);
  if (!el) return;
  const meetings = await loadMeetings(accountId);
  if (!meetings.length) { el.innerHTML = '<div style="color:var(--m);font-size:13px;padding:20px 0">No meetings logged yet.</div>'; return; }
  el.innerHTML = meetings.map(m => {
    const dateStr = m.date ? new Date(m.date).toLocaleDateString('en-GB',{weekday:'short',day:'numeric',month:'short',year:'numeric'}) : '?';
    return '<div class="meeting-card">' +
      '<div class="meeting-head"><div><div class="meeting-account">' + (accountId ? '' : m.account_name + ' · ') + m.owner + '</div><div class="meeting-date">' + dateStr + (m.participants ? ' · ' + m.participants : '') + '</div></div>' +
      '<span class="meeting-outcome ' + (m.outcome||'neutral') + '">' + (m.outcome||'neutral') + '</span></div>' +
      (m.summary ? '<div class="meeting-summary">' + m.summary + '</div>' : '') +
      (m.next_step ? '<div class="meeting-next"><strong>Next step:</strong> ' + m.next_step + '</div>' : '') +
    '</div>';
  }).join('');
}

// ══ DETAIL PANEL TABS ═════════════════════════════════════════════════════════
function switchDpTab(tab) {
  document.querySelectorAll('.dp-tab').forEach(t => t.classList.toggle('active', t.dataset.tab === tab));
  document.querySelectorAll('.dp-tab-content').forEach(c => c.classList.remove('active'));
  const el = document.getElementById('dp-tab-' + tab);
  if (el) el.classList.add('active');

  if (tab === 'actions' && currentDetailAccount) {
    const listEl = document.getElementById('dp-actions-list');
    if (listEl) {
      fetch('/api/cc/actions?account_id=' + currentDetailAccount.id).then(r=>r.json()).then(actions => {
        if (!actions.length) { listEl.innerHTML = '<div style="color:var(--m);font-size:13px;padding:20px 0">No actions for this account yet.</div>'; return; }
        listEl.innerHTML = actions.map(a => renderActionCard(a)).join('');
      });
    }
  }
  if (tab === 'meetings' && currentDetailAccount) {
    renderMeetingsList('dp-meetings-list', currentDetailAccount.id);
  }
  if (tab === 'outreach') {
    // Reset outreach panel
    document.getElementById('dp-outreach-msg').style.display = 'none';
    document.getElementById('dp-outreach-copy').style.display = 'none';
  }
}

// ══ OUTREACH GENERATOR ════════════════════════════════════════════════════════
async function generateOutreach() {
  if (!currentDetailAccount) return;
  const channel  = document.getElementById('dp-outreach-channel').value;
  const language = document.getElementById('dp-outreach-lang').value;
  const msgEl    = document.getElementById('dp-outreach-msg');
  const copyBtn  = document.getElementById('dp-outreach-copy');
  const loading  = document.getElementById('dp-outreach-loading');
  loading.style.display = 'flex'; msgEl.style.display = 'none'; copyBtn.style.display = 'none';
  const params = new URLSearchParams({account_id: currentDetailAccount.id, channel, language});
  const r = await fetch('/api/cc/outreach?' + params, {method:'POST'});
  const d = await r.json();
  loading.style.display = 'none';
  if (!r.ok) { showToast('Error: ' + (d.error||'Failed')); return; }
  msgEl.textContent = d.message; msgEl.style.display = 'block';
  copyBtn.style.display = 'block';
}

function copyOutreach() {
  const msg = document.getElementById('dp-outreach-msg').textContent;
  navigator.clipboard.writeText(msg).then(() => showToast('Copied to clipboard ✓'));
}

// ══ MODALS ════════════════════════════════════════════════════════════════════
function openActionModal(accountId) {
  const modal = document.getElementById('action-modal');
  modal.classList.add('open');
  // Populate account select
  if (countryData) {
    const sel = document.getElementById('action-account-select');
    sel.innerHTML = countryData.accounts.map(a => '<option value="' + a.id + '"' + (a.id == accountId ? ' selected' : '') + '>' + a.name + '</option>').join('');
  }
  // Default due date = tomorrow
  const tomorrow = new Date(); tomorrow.setDate(tomorrow.getDate()+1);
  document.getElementById('action-due').value = tomorrow.toISOString().split('T')[0];
}
function closeActionModal() { document.getElementById('action-modal').classList.remove('open'); }

async function saveAction() {
  const accountId = document.getElementById('action-account-select').value;
  const title     = document.getElementById('action-title').value.trim();
  if (!title) { showToast('Please enter a title'); return; }
  const due     = document.getElementById('action-due').value;
  const type    = document.getElementById('action-type').value;
  const priority= document.getElementById('action-priority').value;
  const desc    = document.getElementById('action-desc').value;
  const params  = new URLSearchParams({account_id:accountId, title, due_date:due, action_type:type, priority, description:desc});
  const r = await fetch('/api/cc/actions?' + params, {method:'POST'});
  if (r.ok) {
    closeActionModal();
    showToast('Action saved ✓');
    document.getElementById('action-title').value = '';
    document.getElementById('action-desc').value = '';
    loadTodayData(); loadActions();
    if (currentDetailAccount && currentDetailAccount.id == accountId) switchDpTab('actions');
  } else {
    const d = await r.json(); showToast('Error: ' + (d.error||'Failed'));
  }
}

function openMeetingModal(accountId) {
  const modal = document.getElementById('meeting-modal');
  modal.classList.add('open');
  if (countryData) {
    const sel = document.getElementById('meeting-account-select');
    sel.innerHTML = countryData.accounts.map(a => '<option value="' + a.id + '"' + (a.id == accountId ? ' selected' : '') + '>' + a.name + '</option>').join('');
  }
  document.getElementById('meeting-date').value = new Date().toISOString().split('T')[0];
}
function closeMeetingModal() { document.getElementById('meeting-modal').classList.remove('open'); }

async function saveMeeting() {
  const accountId    = document.getElementById('meeting-account-select').value;
  const date         = document.getElementById('meeting-date').value;
  const participants = document.getElementById('meeting-participants').value;
  const summary      = document.getElementById('meeting-summary').value;
  const outcome      = document.getElementById('meeting-outcome').value;
  const nextStep     = document.getElementById('meeting-next').value;
  if (!date) { showToast('Please enter a date'); return; }
  const params = new URLSearchParams({account_id:accountId, date, participants, summary, outcome, next_step:nextStep});
  const r = await fetch('/api/cc/meetings?' + params, {method:'POST'});
  if (r.ok) {
    closeMeetingModal();
    showToast('Meeting logged ✓');
    ['meeting-participants','meeting-summary','meeting-next'].forEach(id => document.getElementById(id).value = '');
    loadTodayData();
    if (activeView === 'meetings') renderMeetingsList('meetings-list', null);
    if (currentDetailAccount && currentDetailAccount.id == accountId) switchDpTab('meetings');
  } else {
    const d = await r.json(); showToast('Error: ' + (d.error||'Failed'));
  }
}

function openCommitModal() {
  const modal = document.getElementById('commit-modal');
  modal.classList.add('open');
  // Pre-fill if exists
  fetch('/api/cc/weekly-commit').then(r=>r.json()).then(d => {
    if (d.exists) {
      document.getElementById('commit-text').value = d.commit_text || '';
      document.getElementById('commit-value').value = d.target_value || '';
    }
  });
}
function closeCommitModal() { document.getElementById('commit-modal').classList.remove('open'); }

async function saveCommit() {
  const text  = document.getElementById('commit-text').value.trim();
  const value = document.getElementById('commit-value').value;
  if (!text) { showToast('Please describe your commit'); return; }
  const params = new URLSearchParams({commit_text:text, target_value:value||0});
  const r = await fetch('/api/cc/weekly-commit?' + params, {method:'POST'});
  if (r.ok) {
    closeCommitModal();
    showToast('Weekly commit saved ✓');
    loadTodayData();
  }
}

// ══ INTELLIGENCE VIEW ════════════════════════════════════════════════════════
let intelData = null;

async function loadIntelligence() {
  const r = await fetch('/api/cc/intelligence');
  if (!r.ok) return;
  intelData = await r.json();
  renderIntelligence();
}

function renderIntelligence() {
  if (!intelData) return;
  const d = intelData;

  // Hero stats
  const heroEl = document.getElementById('intel-hero');
  heroEl.style.display = 'flex';
  document.getElementById('intel-cold-count').textContent = d.stats.cold_count;
  document.getElementById('intel-churn-count').textContent = d.stats.churn_count;
  document.getElementById('intel-won-count').textContent = d.stats.won_count;

  const grid = document.getElementById('intel-grid');
  let html = '';

  // Cold prospects card
  html += '<div class="intel-card cold">';
  html += '<div class="intel-sec-title"><span style="color:var(--amber)">⚠ Cold Accounts (' + d.cold_prospects.length + ')</span><span style="font-size:11px;color:var(--m)">No activity 14+ days</span></div>';
  if (!d.cold_prospects.length) {
    html += '<div style="color:var(--m);font-size:13px">No cold accounts.</div>';
  } else {
    d.cold_prospects.forEach(a => {
      const pFmt = a.pipeline_value >= 1e6 ? '€'+(a.pipeline_value/1e6).toFixed(1)+'M' : '€'+(a.pipeline_value/1e3|0)+'K';
      const daysLabel = a.days_cold != null ? a.days_cold + 'd cold' : 'Never contacted';
      html += '<div class="intel-account-row" id="intel-row-' + a.id + '">';
      html += '<div style="flex:1"><div class="intel-account-name">' + a.name + '</div>';
      html += '<div class="intel-account-meta">' + (a.named_buyer || 'Buyer TBD') + ' · ' + (DEAL_STAGE_LABELS[a.deal_stage]||a.deal_stage) + ' · ' + pFmt + '</div>';
      html += '<div id="diagnosis-' + a.id + '"></div></div>';
      html += '<div style="display:flex;flex-direction:column;align-items:flex-end;gap:6px">';
      html += '<span class="intel-cold-days">' + daysLabel + '</span>';
      html += '<button class="diagnose-btn" id="diag-btn-' + a.id + '" onclick="diagnoseAccount(' + a.id + ',&quot;cold_reactivation&quot;)">Diagnose →</button>';
      html += '</div></div>';
    });
  }
  html += '</div>';

  // Churn risk card
  html += '<div class="intel-card risk">';
  html += '<div class="intel-sec-title"><span style="color:var(--red)"><span style="width:8px;height:8px;background:#DC2626;border-radius:50%;display:inline-block;margin-right:4px;vertical-align:middle;"></span> Churn Risk (' + d.churn_risks.length + ')</span><span style="font-size:11px;color:var(--m)">Existing clients</span></div>';
  if (!d.churn_risks.length) {
    html += '<div style="color:var(--m);font-size:13px">No churn risks detected.</div>';
  } else {
    d.churn_risks.forEach(r => {
      const a = r.account;
      const pFmt = a.pipeline_value >= 1e6 ? '€'+(a.pipeline_value/1e6).toFixed(1)+'M' : '€'+(a.pipeline_value/1e3|0)+'K';
      html += '<div class="intel-account-row" id="intel-row-churn-' + a.id + '">';
      html += '<div style="flex:1"><div class="intel-account-name">' + a.name + '</div>';
      html += '<div class="intel-account-meta">' + (a.named_buyer || 'No buyer') + ' · ' + pFmt + '</div>';
      html += '<div class="intel-risk-badges">' + r.risk_factors.map(f => '<span class="intel-risk-badge">' + f + '</span>').join('') + '</div>';
      html += '<div id="diagnosis-churn-' + a.id + '"></div></div>';
      html += '<button class="diagnose-btn" style="align-self:flex-start;margin-top:2px" id="diag-btn-churn-' + a.id + '" onclick="diagnoseAccount(' + a.id + ',&quot;churn_risk&quot;,true)">Diagnose →</button>';
      html += '</div>';
    });
  }
  html += '</div>';

  grid.innerHTML = html;
}

async function diagnoseAccount(accountId, insightType, isChurn) {
  const btnId = 'diag-btn-' + (isChurn ? 'churn-' : '') + accountId;
  const diagId = 'diagnosis-' + (isChurn ? 'churn-' : '') + accountId;
  const btn = document.getElementById(btnId);
  const diagEl = document.getElementById(diagId);
  if (!btn || !diagEl) return;

  btn.classList.add('loading'); btn.textContent = 'Analysing…';
  diagEl.innerHTML = '<div class="loading-pulse" style="margin-top:8px"><div class="pulse-dot"></div><div class="pulse-dot"></div><div class="pulse-dot"></div><span>Thinking…</span></div>';

  const params = new URLSearchParams({account_id: accountId, insight_type: insightType});
  const r = await fetch('/api/cc/intelligence/diagnose?' + params, {method:'POST'});
  const d = await r.json();

  btn.classList.remove('loading'); btn.style.display = 'none';
  if (!r.ok) {
    diagEl.innerHTML = '<div style="color:var(--red);font-size:12px">' + (d.error||'Failed') + '</div>';
    return;
  }
  // Format markdown-ish bold
  const formatted = d.diagnosis.replace(/\*\*([^*]+)\*\*/g, '<strong>$1</strong>');
  diagEl.innerHTML = '<div class="diagnosis-box">' + formatted + '</div>';
}

async function loadWinPatterns() {
  const box = document.getElementById('intel-win-box');
  const loading = document.getElementById('intel-win-loading');
  const content = document.getElementById('intel-win-content');
  box.style.display = 'block';
  loading.style.display = 'flex'; content.textContent = '';
  const r = await fetch('/api/cc/intelligence/win-patterns', {method:'POST'});
  const d = await r.json();
  loading.style.display = 'none';
  if (!r.ok) { content.textContent = d.error || 'Failed'; return; }
  const formatted = d.analysis.replace(/\*\*([^*]+)\*\*/g, '<strong>$1</strong>');
  content.innerHTML = formatted;
}

// ══ NAVIGATION ════════════════════════════════════════════════════════════════
function switchView(view) {
  activeView = view;
  document.querySelectorAll('.view').forEach(v => v.classList.remove('active'));
  const target = document.getElementById('view-' + view);
  if (target) target.classList.add('active');
  // Sidebar active state — foresight item doesn't use data-view matching due to inline style
  document.querySelectorAll('.nav-item[data-view]').forEach(n => {
    n.classList.toggle('active', n.dataset.view === view);
    if (n.dataset.view === 'foresight') {
      // keep foresight item's gradient even when active
      if (n.dataset.view === view) {
        n.style.background = 'linear-gradient(135deg,rgba(21,62,237,.18),rgba(124,58,237,.14))';
      } else {
        n.style.background = 'linear-gradient(135deg,rgba(21,62,237,.1),rgba(124,58,237,.08))';
      }
    }
  });
  if (view === 'foresight') { /* user clicks Run Foresight manually */ }
  if (view === 'command') loadCommandView();
  if (view === 'actions') loadActions();
  if (view === 'today') loadTodayData();
  if (view === 'meetings') renderMeetingsList('meetings-list', null);
  if (view === 'foresight') loadWeeklyBrief();
  if (view === 'intelligence') loadIntelligence();
}

// ══ WEEKLY STRATEGIC BRIEF ═══════════════════════════════════════════════════
let _weeklyBriefLoaded = false;

async function loadWeeklyBrief(force) {
  if (_weeklyBriefLoaded && !force) return;
  const loading = document.getElementById('brief-loading');
  const content = document.getElementById('brief-content');
  const meta    = document.getElementById('brief-meta');
  if (!loading) return;

  loading.style.display = 'block';
  content.style.display = 'none';

  const statusColors = {
    GREEN: { bg: 'rgba(0,166,126,.15)', color: '#00A67E', border: 'rgba(0,166,126,.3)' },
    AMBER: { bg: 'rgba(217,119,6,.15)', color: '#D97706', border: 'rgba(217,119,6,.3)' },
    RED:   { bg: 'rgba(229,62,62,.15)', color: '#E53E3E', border: 'rgba(229,62,62,.3)' }
  };

  try {
    const country = (currentUser && currentUser.country) ? currentUser.country : 'no';
    const r = await fetch('/api/cc/weekly-brief?country=' + country, { method: 'POST' });
    const d = await r.json();
    if (d.error) throw new Error(d.error);

    const sc = statusColors[d.status] || statusColors.AMBER;
    const badge = document.getElementById('brief-status-badge');
    if (badge) {
      badge.textContent = '● ' + d.status;
      badge.style.background = sc.bg;
      badge.style.color = sc.color;
      badge.style.border = '1px solid ' + sc.border;
    }

    const sr = document.getElementById('brief-status-reason');
    if (sr) sr.textContent = d.status_reason || '';
    if (meta) meta.textContent = (d.week || '') + ' · ' + (d.date || '');

    const nums = document.getElementById('brief-numbers');
    if (nums && d.three_numbers) {
      nums.innerHTML = d.three_numbers.map(function(n) {
        return '<div style="text-align:right">'
          + '<div style="font-size:18px;font-weight:800;color:var(--t);letter-spacing:-.02em">' + n.value + '</div>'
          + '<div style="font-size:9.5px;color:var(--m);font-weight:600">' + n.label + '</div>'
          + '<div style="font-size:9px;color:var(--m2)">' + (n.sub||'') + '</div>'
          + '</div>';
      }).join('<div style="width:1px;background:var(--border);margin:0 4px"></div>');
    }

    if (d.risk) {
      const rh = document.getElementById('brief-risk-headline');
      const ra = document.getElementById('brief-risk-accounts');
      const ri = document.getElementById('brief-risk-implication');
      if (rh) rh.textContent = d.risk.headline || '';
      if (ra) ra.innerHTML = (d.risk.accounts || []).map(function(a){ return '<div>· ' + a + '</div>'; }).join('');
      if (ri) ri.textContent = d.risk.implication || '';
    }

    if (d.top_opportunity) {
      const oa = document.getElementById('brief-opp-account');
      const ob = document.getElementById('brief-opp-buyer');
      const ow = document.getElementById('brief-opp-window');
      const oy = document.getElementById('brief-opp-why');
      const oat = document.getElementById('brief-opp-action-text');
      if (oa) oa.textContent = d.top_opportunity.account || '';
      if (ob) ob.textContent = d.top_opportunity.buyer || '';
      if (ow) ow.textContent = d.top_opportunity.window || '';
      if (oy) oy.textContent = d.top_opportunity.why_now || '';
      if (oat) oat.textContent = d.top_opportunity.recommended_action || '';
    }

    if (d.the_decision) {
      const dq = document.getElementById('brief-dec-question');
      const dopt = document.getElementById('brief-dec-options');
      const dr = document.getElementById('brief-dec-rec');
      const dd = document.getElementById('brief-dec-deadline');
      if (dq) dq.textContent = d.the_decision.question || '';
      if (dopt) dopt.innerHTML = (d.the_decision.options || []).map(function(o,i){ return '<div>' + String.fromCharCode(65+i) + '. ' + o + '</div>'; }).join('');
      if (dr) dr.textContent = d.the_decision.recommendation || '';
      if (dd) dd.textContent = d.the_decision.deadline ? 'Deadline: ' + d.the_decision.deadline : '';
    }

    loading.style.display = 'none';
    content.style.display = 'block';
    _weeklyBriefLoaded = true;
  } catch(e) {
    if (loading) loading.innerHTML = '<div style="color:var(--red);font-size:13px">Fejl: ' + e.message + '</div>';
  }
}

// ══ FORESIGHT ════════════════════════════════════════════════════════════════
async function runForesight() {
  const btn = document.getElementById('fsight-run-btn');
  const body = document.getElementById('fsight-body');
  const empty = document.getElementById('fsight-empty');
  btn.disabled = true;
  btn.textContent = '✦ Generating…';
  empty.style.display = 'none';
  body.innerHTML = '<div class="loading-pulse" style="padding:60px 0;justify-content:center"><div class="pulse-dot"></div><div class="pulse-dot"></div><div class="pulse-dot"></div><span style="margin-left:10px;font-size:14px">Scanning market intelligence — this takes ~15 seconds…</span></div>';

  try {
    const country = currentUser.country || 'no';
    const r = await fetch('/api/cc/foresight?country=' + country, {
      method: 'POST', headers: {'Content-Type':'application/json'}, body: JSON.stringify({country})
    });
    const d = await r.json();
    if (d.error) throw new Error(d.error);

    const windows = d.buying_windows || [];
    const partners = d.partnerships || [];
    const services = d.new_services || [];

    document.getElementById('fs-windows-count').textContent = windows.length;
    document.getElementById('fs-partner-count').textContent = partners.length;
    document.getElementById('fs-service-count').textContent = services.length;
    document.getElementById('fs-last-run').textContent = 'Generated ' + new Date().toLocaleTimeString('en-GB', {hour:'2-digit',minute:'2-digit'}) + ' · ' + new Date().toLocaleDateString('en-GB', {day:'numeric',month:'short',year:'numeric'});

    body.innerHTML = renderForesightSections(windows, partners, services);
    if (typeof lucide !== 'undefined') lucide.createIcons();
  } catch(e) {
    body.innerHTML = '<div style="padding:40px;color:var(--red);font-size:14px">Error: ' + e.message + '</div>';
  }
  btn.disabled = false;
  btn.textContent = '↺ Refresh Foresight';
}

function windowClass(w) {
  if (!w) return 'cool';
  if (w.includes('15') || w.includes('30')) return 'hot';
  if (w.includes('60')) return 'warm';
  return 'cool';
}

function renderForesightSections(windows, partners, services) {
  const windowsHtml = windows.length ? windows.map(w => `
    <div class="fsight-block">
      <div class="fw-head">
        <div class="fw-name">${w.account || 'Account'}</div>
        <div class="fw-window ${windowClass(w.window)}">${w.window || '30-60 days'}</div>
      </div>
      <div class="fw-grid">
        <div class="fw-cell full"><div class="fw-cell-label">Trigger</div><div class="fw-cell-val">${w.trigger || '—'}</div></div>
        <div class="fw-cell full"><div class="fw-cell-label">Recommended entry</div><div class="fw-cell-val">${w.entry || '—'}</div></div>
      </div>
      <div class="fw-foot">
        <div class="fw-value">${w.value || '—'}</div>
        <div class="fw-conf ${w.confidence || 'Medium'}">${w.confidence || 'Medium'} confidence</div>
      </div>
    </div>`).join('') : '<div style="padding:20px 24px;color:var(--m);font-size:13px">No buying windows identified.</div>';

  const partnersHtml = partners.length ? partners.map(p => `
    <div class="fsight-block">
      <div class="fp-head">
        <div class="fp-name">${p.partner || 'Partner'}</div>
        <div class="fp-value">${p.value || '—'}</div>
      </div>
      <div class="fp-body">
        <div class="fp-cell full"><div class="fp-cell-label">Market evidence</div><div class="fp-cell-val">${p.evidence || '—'}</div></div>
        <div class="fp-cell full"><div class="fp-cell-label">Joint proposition</div><div class="fp-cell-val">${p.proposition || '—'}</div></div>
      </div>
      <div class="fp-action"><div class="fp-action-label">Action this week</div>${p.action || '—'}</div>
    </div>`).join('') : '<div style="padding:20px 24px;color:var(--m);font-size:13px">No partnership opportunities identified.</div>';

  const servicesHtml = services.length ? services.map(s => `
    <div class="fsight-block">
      <div class="fs-head">
        <div class="fs-name">${s.service || 'Service'}</div>
        <div class="fs-revenue">${s.revenue || '—'}</div>
      </div>
      <div class="fs-body">
        <div class="fs-cell full"><div class="fs-cell-label">Market evidence</div><div class="fs-cell-val">${s.evidence || '—'}</div></div>
        <div class="fs-cell"><div class="fs-cell-label">Target clients</div><div class="fs-cell-val">${s.clients || '—'}</div></div>
        <div class="fs-cell"><div class="fs-cell-label">Entry offer</div><div class="fs-cell-val">${s.entry_offer || '—'}</div></div>
      </div>
      <div style="margin-top:8px"><span class="fs-bop">⊕ ${s.build_or_partner || 'TBD'}</span></div>
    </div>`).join('') : '<div style="padding:20px 24px;color:var(--m);font-size:13px">No service gaps identified.</div>';

  return `
    <div class="fsight-section">
      <div class="fsight-sec-header">
        <div style="display:flex;align-items:center;gap:12px">
          <div class="fsight-sec-icon window"><i data-lucide="target" style="width:20px;height:20px;display:inline-block;vertical-align:middle;"></i></div>
          <div>
            <div class="fsight-sec-title">Buying Windows</div>
            <div class="fsight-sec-sub">Accounts predicted to enter a buying decision in the next 90 days</div>
          </div>
        </div>
        <div class="fsight-sec-count window">${windows.length} accounts</div>
      </div>
      <div class="fsight-blocks">${windowsHtml}</div>
    </div>

    <div class="fsight-section">
      <div class="fsight-sec-header">
        <div style="display:flex;align-items:center;gap:12px">
          <div class="fsight-sec-icon partner"><i data-lucide="handshake" style="width:20px;height:20px;display:inline-block;vertical-align:middle;"></i></div>
          <div>
            <div class="fsight-sec-title">Partnership Opportunities</div>
            <div class="fsight-sec-sub">Technology and agency alliances that unlock new revenue for JAKALA Nordic</div>
          </div>
        </div>
        <div class="fsight-sec-count partner">${partners.length} partners</div>
      </div>
      <div class="fsight-blocks">${partnersHtml}</div>
    </div>

    <div class="fsight-section">
      <div class="fsight-sec-header">
        <div style="display:flex;align-items:center;gap:12px">
          <div class="fsight-sec-icon service"><i data-lucide="zap" style="width:20px;height:20px;display:inline-block;vertical-align:middle;"></i></div>
          <div>
            <div class="fsight-sec-title">New Service Opportunities</div>
            <div class="fsight-sec-sub">Service gaps the Nordic market needs — that no competitor offers yet</div>
          </div>
        </div>
        <div class="fsight-sec-count service">${services.length} opportunities</div>
      </div>
      <div class="fsight-blocks">${servicesHtml}</div>
    </div>`;
}

async function loadCommandView() {
  if (!countryData) return;
  const accounts = countryData.accounts || [];
  const signals = countryData.signals || [];
  const critical = signals.filter(s => s.severity === 'critical');
  const highIcp = accounts.filter(a => (a.icp_score || 0) >= 8 && a.account_type !== 'existing');
  const engaged = accounts.filter(a => ['engaged','proposed','negotiating'].includes(a.deal_stage));
  const pipeline = accounts.reduce((s, a) => s + (a.pipeline_value || 0), 0);

  // Pulse stats
  document.getElementById('cmd-signals-count').textContent = critical.length || signals.length;
  document.getElementById('cmd-targets-count').textContent = highIcp.length;
  document.getElementById('cmd-pipeline-val').textContent = '€' + (pipeline >= 1000000 ? (pipeline/1000000).toFixed(1)+'M' : Math.round(pipeline/1000)+'K');
  document.getElementById('cmd-engaged-count').textContent = engaged.length;

  // Sub text
  const now = new Date();
  document.getElementById('cmd-week-label').textContent = 'MARKET INTELLIGENCE · W' + getWeekNumber(now) + ' ' + now.getFullYear();
  document.getElementById('cmd-sub').textContent = `${critical.length} critical signals · ${highIcp.length} high-ICP targets unengaged · pipeline trending`;

  // Critical signals (top 4)
  const sigList = document.getElementById('cmd-signals-list');
  const topSigs = [...signals].sort((a,b) => {
    const sev = {critical:0,warning:1,info:2};
    return (sev[a.severity]||2) - (sev[b.severity]||2);
  }).slice(0, 4);
  sigList.innerHTML = topSigs.length ? topSigs.map(s => `
    <div class="cmd-signal-card ${s.severity}">
      <div class="cmd-signal-sev ${s.severity}">${s.signal_type || s.severity}</div>
      <div class="cmd-signal-title">${s.title}</div>
      <div class="cmd-signal-desc">${(s.description||'').substring(0,120)}${(s.description||'').length>120?'…':''}</div>
    </div>`).join('') : '<div style="color:var(--m);font-size:13px;padding:16px 0">No active signals. Add signals via the admin API.</div>';

  // High-value targets (top 5 by ICP+deal score)
  const tgtList = document.getElementById('cmd-targets-list');
  const topTargets = [...highIcp].sort((a,b) => ((b.icp_score||0)+(b.deal_score||0)) - ((a.icp_score||0)+(a.deal_score||0))).slice(0,5);
  tgtList.innerHTML = topTargets.length ? topTargets.map(a => `
    <div class="cmd-target-card" onclick="openDetail(${JSON.stringify(a).replace(/"/g,'&quot;')})">
      <div>
        <div class="cmd-target-name">${a.name}</div>
        <div class="cmd-target-meta">${a.deal_stage||'identified'} · ${a.named_buyer || 'Buyer TBD'}</div>
      </div>
      <div class="cmd-target-score">
        <div class="cmd-target-val">€${a.pipeline_value >= 1000000 ? (a.pipeline_value/1000000).toFixed(1)+'M' : Math.round((a.pipeline_value||0)/1000)+'K'}</div>
        <div class="cmd-target-icp">ICP ${a.icp_score || '—'}</div>
      </div>
    </div>`).join('') : '<div style="color:var(--m);font-size:13px;padding:16px 0">No high-ICP unengaged accounts found.</div>';
}

function getWeekNumber(d) {
  const date = new Date(Date.UTC(d.getFullYear(), d.getMonth(), d.getDate()));
  date.setUTCDate(date.getUTCDate() + 4 - (date.getUTCDay()||7));
  const yearStart = new Date(Date.UTC(date.getUTCFullYear(),0,1));
  return Math.ceil((((date - yearStart) / 86400000) + 1)/7);
}

async function generateStrategicBrief() {
  const btn = document.getElementById('cmd-brief-btn');
  const section = document.getElementById('cmd-brief-section');
  const loading = document.getElementById('cmd-brief-loading');
  const textEl = document.getElementById('cmd-brief-text');
  btn.disabled = true;
  btn.textContent = 'Generating…';
  section.style.display = 'block';
  loading.style.display = 'flex';
  textEl.textContent = '';
  document.getElementById('cmd-brief-date').textContent = new Date().toLocaleDateString('en-GB', {day:'numeric',month:'long',year:'numeric'});
  try {
    const country = currentUser.country || 'no';
    const r = await fetch('/api/cc/strategic-brief?country=' + country, {
      method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({country})
    });
    const d = await r.json();
    loading.style.display = 'none';
    textEl.textContent = d.brief || d.error || 'Could not generate brief.';
  } catch(e) {
    loading.style.display = 'none';
    textEl.textContent = 'Error generating brief.';
  }
  btn.disabled = false;
  btn.textContent = '◈ Regenerate brief';
}

// ══ TOAST ════════════════════════════════════════════════════════════════════
function showToast(msg) {
  const t = document.getElementById('cc-toast');
  t.textContent = msg; t.classList.add('show');
  setTimeout(() => t.classList.remove('show'), 3500);
}
</script>
<script>lucide.createIcons();</script>
</body>
</html>"""

@app.route("/cc")
def cc_index():
    return render_template_string(CC_HTML)


if __name__ == "__main__":
    port = int(os.getenv("PORT", 5050))
    print(f"\n  JAKALA GTM OS running at http://localhost:{port}\n")
    app.run(debug=False, port=port, threaded=True)
